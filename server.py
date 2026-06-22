"""DeepAgent Server - AI Agentic Server with Web UI

Features: Web search, Deep research, MCP server support, Custom tools, Streaming responses, Headless browser control
Supported Models: Claude (Anthropic), GPT-4 (OpenAI), Gemini (Google), Ollama (Local)

CONTEXT WINDOW MANAGEMENT:
- Tool Response Summarization: Large MCP responses are summarized via GPT-4o-mini before entering message history
- Conversation History Summarization: When message token count exceeds threshold, older messages are summarized
- Full data preservation: Raw responses are always saved to disk for reference
- Uses ONLY ChatGPT (OpenAI) for all summarization tasks
"""

import asyncio
import contextvars
import json
import os
import logging
import re
import time
import uuid
from typing import List, Dict, Any, Optional, Literal
from datetime import datetime

from fastapi import FastAPI, WebSocket, HTTPException, WebSocketDisconnect, Request
from fastapi.staticfiles import StaticFiles
from fastapi.responses import HTMLResponse, StreamingResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

from deepagents import create_deep_agent
from deepagents_patches import disable_write_todos

# Remove the built-in deepagents `write_todos` tool + its prompt injection from
# every agent/subagent. Must run before any create_deep_agent call.
disable_write_todos()
from langchain_core.tools import tool
from langchain_community.tools import DuckDuckGoSearchRun
from langgraph.checkpoint.memory import InMemorySaver

_AGENT_CHECKPOINTER = InMemorySaver()
import httpx
import dotenv

dotenv.load_dotenv()

def _parse_recursion_limit() -> int:
    raw = os.getenv("RECURSION_LIMIT", "1000000")
    try:
        v = int(raw)
        if v < 1:
            raise ValueError(f"must be >= 1, got {v}")
        return v
    except (ValueError, TypeError) as e:
        print(f"⚠️  RECURSION_LIMIT env var invalid ({raw!r}: {e}) — falling back to 1000000")
        return 1000000

_RECURSION_LIMIT = _parse_recursion_limit()
print(f"LangGraph recursion_limit: {_RECURSION_LIMIT} (override with RECURSION_LIMIT env var)")

if not os.getenv("LANGCHAIN_API_KEY"):
    os.environ["LANGCHAIN_TRACING_V2"] = "false"
    os.environ["LANGCHAIN_API_KEY"] = ""

from fastapi.responses import RedirectResponse
try:
    from google_sheets_auth import GoogleSheetsAuth
    sheets_auth = GoogleSheetsAuth()
    GOOGLE_SHEETS_ENABLED = sheets_auth.is_authenticated() or os.path.exists(
        "client_secrets.json")
    print(
        f"Google Sheets integration: {'READY' if sheets_auth.is_authenticated() else 'NOT AUTHENTICATED'}"
    )
except Exception as e:
    sheets_auth = None
    GOOGLE_SHEETS_ENABLED = False
    print(f"Google Sheets integration disabled: {e}")

logging.getLogger("fastmcp").setLevel(logging.ERROR)
logging.getLogger("mcp").setLevel(logging.ERROR)

import lake as _lake

try:
    from supabase import create_client, Client
except ImportError:
    print(
        "Warning: supabase package not found. Install with: pip install supabase"
    )
    Client = None


class Config:
    """Server configuration"""
    HOST = os.getenv("HOST", "0.0.0.0")
    PORT = int(os.getenv("PORT", "5000"))
    MODEL = os.getenv("MODEL", "anthropic:claude-sonnet-4-6-20260901")
    ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY", "")
    OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
    GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "")
    XAI_API_KEY = os.getenv("XAI_API_KEY", "")
    # Fireworks AI (OpenAI-compatible). Used when a model id is prefixed
    # "fireworks:" — currently a super-admin-only sandbox surfaced from VIBE.
    FIREWORKS_API_KEY = os.getenv("FIREWORKS_API_KEY", "")
    # Dedicated output budget for Fireworks. gpt-oss are reasoning models whose
    # reasoning tokens count toward the output cap, so the 8192 Anthropic-sized
    # MAX_OUTPUT_TOKENS truncates long turns. Mirror DEAL_SWEEP_MAX_TOKENS.
    FIREWORKS_MAX_TOKENS = int(os.getenv("FIREWORKS_MAX_TOKENS", "32000"))

    MCP_CONFIG_FILE = os.getenv("MCP_CONFIG_FILE", "mcp_config.json")
    CUSTOM_TOOLS_DIR = os.getenv("CUSTOM_TOOLS_DIR", "custom_tools")

    MCP_MAX_RESPONSE_SIZE = int(os.getenv("MCP_MAX_RESPONSE_SIZE", "500000"))
    MCP_MAX_STRING_LENGTH = int(os.getenv("MCP_MAX_STRING_LENGTH", "50000"))
    MCP_MAX_LIST_ITEMS = int(os.getenv("MCP_MAX_LIST_ITEMS", "500"))

    SUMMARIZER_MODEL = os.getenv("SUMMARIZER_MODEL", "gpt-4o-mini")
    TOOL_RESPONSE_SUMMARIZE_THRESHOLD = int(
        os.getenv("TOOL_RESPONSE_SUMMARIZE_THRESHOLD", "5000"))
    # Lowered from 100000 → 60000 on 2026-05-22 (cost task #40). 100k meant
    # the pre-astream summarizer almost never fired in practice; the bulk of
    # the cost on long ABM runs (e.g. chat 15dbbb61 = $17.40) came from
    # intra-run tool_result accumulation, not the restored user history. The
    # new ContextTrimMiddleware addresses the intra-run case; this knob still
    # controls the pre-astream cross-turn summarizer.
    CONVERSATION_SUMMARIZE_TOKEN_THRESHOLD = int(
        os.getenv("CONVERSATION_SUMMARIZE_TOKEN_THRESHOLD", "60000"))
    CONVERSATION_KEEP_RECENT_MESSAGES = int(
        os.getenv("CONVERSATION_KEEP_RECENT_MESSAGES", "20"))
    SUMMARIZER_INPUT_LIMIT = int(os.getenv("SUMMARIZER_INPUT_LIMIT", "50000"))

    # Intra-run context-trim middleware (cost task #40). Fires before EVERY
    # LLM call inside a single astream() and shrinks long ToolMessage content
    # to a one-line placeholder when the in-flight message budget exceeds
    # CONTEXT_TRIM_THRESHOLD_TOKENS. Safer than cross-turn summarization
    # because it never drops or reorders messages, so tool_call_id pairing
    # always holds. See `agent_checklist/context_trim_middleware.py`.
    CONTEXT_TRIM_ENABLED = os.getenv("CONTEXT_TRIM_ENABLED", "true").lower() in ("1", "true", "yes")
    CONTEXT_TRIM_THRESHOLD_TOKENS = int(
        os.getenv("CONTEXT_TRIM_THRESHOLD_TOKENS", "60000"))
    CONTEXT_TRIM_KEEP_RECENT_MESSAGES = int(
        os.getenv("CONTEXT_TRIM_KEEP_RECENT_MESSAGES", "10"))
    CONTEXT_TRIM_PLACEHOLDER_MAX_CHARS = int(
        os.getenv("CONTEXT_TRIM_PLACEHOLDER_MAX_CHARS", "400"))

    # --- Token-budget knobs (rate-limit mitigation) ---
    # Compress every MCP tool description to roughly this many characters when
    # registering tools with the LLM. The full docstring stays on the underlying
    # MCP function — this only trims what is sent to the model on every turn.
    # Set to 0 to disable compression and send full descriptions.
    TOOL_DESCRIPTION_MAX_CHARS = int(os.getenv("TOOL_DESCRIPTION_MAX_CHARS", "400"))
    # Comma-separated MCP server names to forcibly disable at startup, even if
    # mcp_config.json marks them enabled. Useful for pruning the tool catalog
    # for a deployment that won't use heavy servers (e.g. "mailchimp,eloqua").
    MCP_SERVER_DENYLIST = [
        s.strip() for s in os.getenv("MCP_SERVER_DENYLIST", "").split(",") if s.strip()
    ]
    # If non-empty, ONLY these MCP server names are loaded (overrides denylist
    # and the mcp_config.json `enabled` flag, except the flag still has to be
    # true). Use to whitelist a phase, e.g. MCP_SERVER_ALLOWLIST="lemlist,salesforce".
    MCP_SERVER_ALLOWLIST = [
        s.strip() for s in os.getenv("MCP_SERVER_ALLOWLIST", "").split(",") if s.strip()
    ]
    # Tool-level allow/denylist (cost task #40, 2026-05-22). Finer-grained
    # than the server-level lists above: filters individual MCP tools after
    # they've been loaded + wrapped. Patterns are fnmatch globs against the
    # tool name, comma-separated. Examples:
    #   MCP_TOOL_DENYLIST="mailchimp_*,showpad_*,linkedin_*,eloqua_*,gojiberry_*"
    #     drops 5 server's tools without disabling the servers (still
    #     reachable via the /mcp endpoint).
    #   MCP_TOOL_ALLOWLIST="salesforce_*,lemlist_*,zi_*,apollo_*,zerobounce_*"
    #     when non-empty, ONLY tools matching one of these globs are
    #     registered with the agent. Denylist is then ignored.
    # Custom tools (web_search, search_knowledge, send_to_clay, web_scrape,
    # web_search_with_urls, get_current_time, browser_*) are NEVER filtered
    # by these lists — they're considered core capability.
    MCP_TOOL_ALLOWLIST = [
        s.strip() for s in os.getenv("MCP_TOOL_ALLOWLIST", "").split(",") if s.strip()
    ]
    MCP_TOOL_DENYLIST = [
        s.strip() for s in os.getenv("MCP_TOOL_DENYLIST", "").split(",") if s.strip()
    ]
    # How many times the Anthropic SDK should retry on 429 / 5xx responses.
    # The SDK already honours the Retry-After header. Default lowered from 8
    # to 2 on 2026-05-20 after chat f4c06387 hung for ~16min: a single slow
    # LLM call timing out at 120s × 8 retries = 16min of silent retrying with
    # the UI frozen on "Thinking…". 2 retries (~6min worst case) is enough
    # for transient blips; the watchdog catches longer hangs.
    ANTHROPIC_MAX_RETRIES = int(os.getenv("ANTHROPIC_MAX_RETRIES", "2"))
    # Per-LLM-call timeout. LLM_REQUEST_TIMEOUT_S is the canonical env var
    # (per Batch-2 follow-up); ANTHROPIC_TIMEOUT_SECONDS is kept as a legacy
    # alias so existing deploys don't regress.
    ANTHROPIC_TIMEOUT_SECONDS = int(
        os.getenv("LLM_REQUEST_TIMEOUT_S",
                  os.getenv("ANTHROPIC_TIMEOUT_SECONDS", "180")))
    # Watchdog: if no agent.astream chunk arrives for this many seconds the
    # run is treated as stalled (hung LLM call, dropped stream, deadlock).
    # The run is cancelled and a loud `error` row is written so the UI moves
    # off "Thinking…" instead of sitting silent forever.
    #
    # CRITICAL: with stream_mode="values" no chunk is emitted *during* a single
    # LLM/tool superstep, so the gap between chunks equals one heavy LLM call
    # (large context — hundreds of tools, 8192 output tokens) PLUS the SDK's own
    # retry budget. The watchdog must therefore outlast the LLM SDK self-recovery
    # window — per-call timeout × (retries + 1) — otherwise it kills a
    # legitimately slow-but-progressing call before its retry can finish. That
    # 180/180 collision is what surfaced as "no agent chunk for 180s" on real
    # runs that make many API calls. The default below derives from the LLM
    # budget + slack and stays well under the pipeline PHASE_TIMEOUT_SECONDS
    # (1800s) so the phase cap remains the outer guard. Override explicitly via
    # the WATCHDOG_STALL_SECONDS env var if needed.
    _LLM_SELF_RECOVERY_S = ANTHROPIC_TIMEOUT_SECONDS * (ANTHROPIC_MAX_RETRIES + 1)
    WATCHDOG_STALL_SECONDS = int(
        os.getenv("WATCHDOG_STALL_SECONDS", str(_LLM_SELF_RECOVERY_S + 120))
    )
    # Anthropic per-response max output tokens. Default 8192 — keeps a single
    # heavy Phase-5 draft turn from being truncated to ~1024 (the langchain
    # default). Sonnet 4.x supports up to 8192; raise via env if a single
    # draft still truncates (and add the output-128k beta header).
    MAX_OUTPUT_TOKENS = int(os.getenv("MAX_OUTPUT_TOKENS", "8192"))
    # Auto-continuation circuit-breakers for run_agent_and_save / streaming
    # handler. These are runaway safeties, NOT cost targets — set high enough
    # that a legitimate long run always completes. Worst observed legitimate
    # run was ~$27 / ~26 min uncheckpointed; with the checkpointer a clean
    # single-turn complete run should land well under 30 min / $20.
    MAX_RUN_SECONDS = int(os.getenv("MAX_RUN_SECONDS", "1800"))
    MAX_RUN_COST_USD = float(os.getenv("MAX_RUN_COST_USD", "20.0"))
    MAX_AUTO_CONTINUATIONS = int(os.getenv("MAX_AUTO_CONTINUATIONS", "25"))
    RUN_COST_WARN_USD = float(os.getenv("RUN_COST_WARN_USD", "10.0"))

    MAX_CONCURRENT_SESSIONS = int(os.getenv("MAX_CONCURRENT_SESSIONS", "50"))
    SESSION_TIMEOUT_MINUTES = int(os.getenv("SESSION_TIMEOUT_MINUTES", "60"))

    SUPABASE_URL = os.getenv("SUPABASE_URL", "")
    # SUPABASE_SERVICE_ROLE_KEY is the canonical Supabase env var name (matches
    # their docs + SDK). We check it FIRST so it always wins. SUPABASE_SERVICE_KEY
    # is kept as a legacy fallback only — historically some deploys had an anon
    # key stored under that name, and reading it first would silently use the
    # anon role, making every RLS-protected read return zero rows.
    SUPABASE_SERVICE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY", "") or os.getenv("SUPABASE_SERVICE_KEY", "")


config = Config()
print(
    f"[STARTUP] Watchdog: stall={config.WATCHDOG_STALL_SECONDS}s "
    f"(derived from LLM timeout={config.ANTHROPIC_TIMEOUT_SECONDS}s "
    f"× (retries={config.ANTHROPIC_MAX_RETRIES}+1) + 120s slack; "
    f"override with WATCHDOG_STALL_SECONDS) — must outlast a single slow LLM "
    f"call + its retries so legitimate heavy runs aren't killed",
    flush=True,
)

supabase: Optional[Client] = None
if config.SUPABASE_URL and config.SUPABASE_SERVICE_KEY and Client:
    try:
        supabase = create_client(config.SUPABASE_URL,
                                 config.SUPABASE_SERVICE_KEY)
        print(f"Supabase client initialized: {config.SUPABASE_URL}")
    except Exception as e:
        print(f"Failed to initialize Supabase: {e}")


class ContextWindowManager:
    """Manages context window using ChatGPT for summarization."""

    def __init__(self):
        self._summarizer = None

    def _get_summarizer(self):
        if self._summarizer is None:
            try:
                from langchain_openai import ChatOpenAI
                self._summarizer = ChatOpenAI(
                    model=config.SUMMARIZER_MODEL,
                    temperature=0,
                    max_tokens=4096,
                    api_key=config.OPENAI_API_KEY
                    or os.getenv("OPENAI_API_KEY"),
                )
                print(
                    f"Context summarizer initialized: {config.SUMMARIZER_MODEL}"
                )
            except Exception as e:
                print(f"Failed to initialize summarizer: {e}")
                self._summarizer = None
        return self._summarizer

    def estimate_tokens(self, text: str) -> int:
        return len(text) // 4

    def estimate_messages_tokens(self, messages: list) -> int:
        total = 0
        for msg in messages:
            if hasattr(msg, 'content'):
                content = str(msg.content)
            elif isinstance(msg, dict):
                content = str(msg.get('content', ''))
            else:
                content = str(msg)
            total += self.estimate_tokens(content)
        return total

    async def summarize_tool_response(self, tool_name: str,
                                      result_str: str) -> str:
        summarizer = self._get_summarizer()
        if summarizer is None:
            return self._truncate_with_context(result_str)
        import time as _time
        _sum_start = _time.monotonic()
        try:
            from langchain_core.messages import HumanMessage, SystemMessage
            input_text = result_str[:config.SUMMARIZER_INPUT_LIMIT]
            print(
                f"[SUMMARIZER] ⏳ start | tool={tool_name}"
                f" | input={len(result_str):,} chars (capped at {len(input_text):,})",
                flush=True,
            )
            messages = [
                SystemMessage(
                    content=
                    """You are a data extraction assistant. Your job is to summarize large tool/API responses 
into a compact format that preserves ALL actionable information.

RULES:
- Preserve ALL: record IDs, names, amounts, dates, stages, statuses, owners, types
- Preserve ALL: relationships, lookup fields, reference IDs, counts, aggregates
- Preserve ALL: error messages, warnings, validation failures
- Remove: redundant metadata fields (attributes, urls, api types), duplicate nested references
- Remove: null/empty fields, system timestamps that aren't business-relevant
- Format: Use structured text, not JSON. Group related records logically.
- If data contains records/rows, present them as a concise numbered list with key fields
- Always state the total count of records at the top
- Keep your summary under 5000 words"""),
                HumanMessage(
                    content=
                    f"Summarize this {tool_name} response ({len(result_str):,} chars):\n\n{input_text}"
                )
            ]
            SUMMARIZER_TIMEOUT = int(os.getenv("SUMMARIZER_TIMEOUT", "45"))
            try:
                response = await asyncio.wait_for(
                    summarizer.ainvoke(messages),
                    timeout=SUMMARIZER_TIMEOUT,
                )
            except asyncio.TimeoutError:
                elapsed_sum = _time.monotonic() - _sum_start
                print(
                    f"[SUMMARIZER] ⏰ TIMEOUT after {elapsed_sum:.1f}s | tool={tool_name}"
                    f" — falling back to truncation",
                    flush=True,
                )
                return self._truncate_with_context(result_str)
            summary = response.content
            summary += f"\n\n[Summarized from {len(result_str):,} chars. Full data saved to disk.]"
            elapsed_sum = _time.monotonic() - _sum_start
            print(
                f"[SUMMARIZER] ✅ done | tool={tool_name}"
                f" | {len(result_str):,} → {len(summary):,} chars | took={elapsed_sum:.1f}s",
                flush=True,
            )
            return summary
        except Exception as e:
            elapsed_sum = _time.monotonic() - _sum_start
            print(f"[SUMMARIZER] ❌ failed | tool={tool_name} | took={elapsed_sum:.1f}s | err={e}", flush=True)
            return self._truncate_with_context(result_str)

    async def summarize_conversation_history(self, messages: list) -> list:
        total_tokens = self.estimate_messages_tokens(messages)
        if total_tokens < config.CONVERSATION_SUMMARIZE_TOKEN_THRESHOLD:
            return messages
        keep_count = config.CONVERSATION_KEEP_RECENT_MESSAGES
        if len(messages) <= keep_count:
            return messages
        older_messages = messages[:-keep_count]
        recent_messages = messages[-keep_count:]
        print(
            f"  Summarizing conversation: {len(messages)} messages ({total_tokens:,} tokens)"
        )
        summarizer = self._get_summarizer()
        if summarizer is None:
            return recent_messages
        try:
            from langchain_core.messages import HumanMessage, SystemMessage
            history_text = ""
            for msg in older_messages:
                if hasattr(msg, 'content'):
                    role = getattr(msg, 'type', 'unknown')
                    content = str(msg.content)
                elif isinstance(msg, dict):
                    role = msg.get('role', 'unknown')
                    content = str(msg.get('content', ''))
                else:
                    role = 'unknown'
                    content = str(msg)
                if len(content) > 10000:
                    content = content[:10000] + "...[truncated]"
                history_text += f"[{role}]: {content}\n\n"
            history_text = history_text[:config.SUMMARIZER_INPUT_LIMIT]
            summary_messages = [
                SystemMessage(
                    content=
                    """You are a conversation summarizer for an AI agent system.
Summarize the conversation history preserving ALL:
- What the user asked for and the agent's conclusions/answers
- Key data points retrieved (record IDs, names, amounts, statuses, dates)
- Decisions made, actions taken, tools called and their outcomes
- Any errors encountered and how they were resolved
- Current task context and what the user is working towards
Format as a structured summary with sections. Be thorough but concise.
Keep under 3000 words."""),
                HumanMessage(
                    content=
                    f"Summarize this conversation history ({len(older_messages)} messages):\n\n{history_text}"
                )
            ]
            response = await summarizer.ainvoke(summary_messages)
            summary_content = response.content
            summary_msg = {
                "role":
                "system",
                "content":
                f"[CONVERSATION HISTORY SUMMARY]\n{summary_content}\n[END SUMMARY - Recent messages follow below]"
            }
            compressed = [summary_msg] + [{
                "role":
                msg.get("role", "user") if isinstance(msg, dict) else getattr(
                    msg, 'type', 'user').replace('human', 'user').replace(
                        'ai', 'assistant'),
                "content":
                msg.get("content", "") if isinstance(msg, dict) else str(
                    getattr(msg, 'content', ''))
            } for msg in recent_messages]
            new_tokens = self.estimate_messages_tokens(compressed)
            print(
                f"  Conversation compressed: {total_tokens:,} -> {new_tokens:,} tokens"
            )
            return compressed
        except Exception as e:
            print(f"  Conversation summarization failed: {e}")
            return recent_messages

    def _truncate_with_context(self, text: str, max_length: int = None) -> str:
        max_len = max_length or config.MCP_MAX_STRING_LENGTH
        if len(text) <= max_len:
            return text
        head_size = int(max_len * 0.7)
        tail_size = int(max_len * 0.25)
        return (
            text[:head_size] +
            f"\n\n...[TRUNCATED {len(text) - head_size - tail_size:,} chars]...\n\n"
            + text[-tail_size:])


context_manager = ContextWindowManager()


_raw_duckduckgo_search = DuckDuckGoSearchRun()


@tool
def web_search(query: str) -> str:
    """Search the web using DuckDuckGo. Returns search results as text. Use this to find current information, news, company data, and facts from the internet."""
    try:
        return _raw_duckduckgo_search.run(query)
    except Exception as e:
        error_msg = str(e)
        if "429" in error_msg or "Too Many Requests" in error_msg or "RESOURCE_EXHAUSTED" in error_msg:
            return f"[Web search rate limited - too many requests. Try again in a moment or proceed without this search result. Query was: {query}]"
        return f"[Web search error: {error_msg}. The agent should continue with available information. Query was: {query}]"


@tool
def get_current_time() -> str:
    """Get the current date and time."""
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


class CustomToolsLoader:

    @staticmethod
    def load_tools_from_directory(directory: str) -> List[Any]:
        tools = []
        if not os.path.exists(directory):
            os.makedirs(directory)
            example_tool = '''"""Example custom tool"""
from langchain_core.tools import tool

@tool
def example_calculator(a: float, b: float, operation: str = "add") -> float:
    """Perform basic arithmetic operations.
    
    Args:
        a: First number
        b: Second number
        operation: Operation to perform (add, subtract, multiply, divide)
    
    Returns:
        Result of the operation
    """
    if operation == "add":
        return a + b
    elif operation == "subtract":
        return a - b
    elif operation == "multiply":
        return a * b
    elif operation == "divide":
        if b == 0:
            return "Error: Division by zero"
        return a / b
    else:
        return "Error: Unknown operation"
'''
            with open(os.path.join(directory, "example_tools.py"), "w") as f:
                f.write(example_tool)

        import sys
        import importlib.util

        for filename in os.listdir(directory):
            if filename.endswith(".py") and not filename.startswith("__"):
                filepath = os.path.join(directory, filename)
                try:
                    spec = importlib.util.spec_from_file_location(
                        filename[:-3], filepath)
                    module = importlib.util.module_from_spec(spec)
                    sys.modules[filename[:-3]] = module
                    spec.loader.exec_module(module)
                    for attr_name in dir(module):
                        attr = getattr(module, attr_name)
                        if hasattr(attr, "name") and hasattr(
                                attr, "description"):
                            tools.append(attr)
                            print(f"Loaded custom tool: {attr.name}")
                except Exception as e:
                    print(f"Error loading tools from {filename}: {e}")
        return tools


class MCPConfigManager:

    def __init__(self, config_file: str):
        self.config_file = config_file
        self.config = self._load_config()

    def _load_config(self) -> Dict[str, Any]:
        if os.path.exists(self.config_file):
            with open(self.config_file, 'r') as f:
                return json.load(f)
        else:
            default_config = {
                "mcp_servers": {
                    "example_filesystem": {
                        "command":
                        "npx",
                        "args": [
                            "-y", "@modelcontextprotocol/server-filesystem",
                            "./workspace"
                        ],
                        "transport":
                        "stdio",
                        "enabled":
                        False
                    }
                }
            }
            self.save_config(default_config)
            return default_config

    def save_config(self, config: Dict[str, Any]):
        with open(self.config_file, 'w') as f:
            json.dump(config, f, indent=2)
        self.config = config

    @staticmethod
    def _resolve_env_vars(value):
        """Resolve ${VAR} placeholders in config values from environment."""
        import re
        if isinstance(value, str):

            def replacer(match):
                var_name = match.group(1)
                return os.environ.get(var_name, match.group(0))

            return re.sub(r'\$\{(\w+)\}', replacer, value)
        elif isinstance(value, dict):
            return {
                k: MCPConfigManager._resolve_env_vars(v)
                for k, v in value.items()
            }
        elif isinstance(value, list):
            return [MCPConfigManager._resolve_env_vars(v) for v in value]
        return value

    def get_enabled_servers(self) -> Dict[str, Any]:
        mcp_servers = self.config.get("mcp_servers", {})
        denylist = set(config.MCP_SERVER_DENYLIST)
        allowlist = set(config.MCP_SERVER_ALLOWLIST)
        enabled_servers = {}
        skipped: list = []
        for name, cfg in mcp_servers.items():
            if not cfg.get("enabled", False):
                continue
            if allowlist and name not in allowlist:
                skipped.append(f"{name} (not in MCP_SERVER_ALLOWLIST)")
                continue
            if name in denylist:
                skipped.append(f"{name} (in MCP_SERVER_DENYLIST)")
                continue
            server_config = {
                k: v
                for k, v in cfg.items() if k != "enabled"
            }
            server_config = self._resolve_env_vars(server_config)
            enabled_servers[name] = server_config
        if skipped:
            print(f"[MCP-CONFIG] Skipping {len(skipped)} server(s): {', '.join(skipped)}")
        return enabled_servers


# ──────────────────────────────────────────────────────────────────────────────
# INTAKE GUARD — prevents the agent from passing a hallucinated Lemlist
# campaign_id to write-side tools. Built after chat c01fabb3-7f10-44e7-8d55-
# 08004e05d9b9 (2026-05-21) where the model invented `cam_aB7xno3KiW53FjMm`
# mid-run after 25 auto-continuations and then blamed the user for it.
# ──────────────────────────────────────────────────────────────────────────────
_CAMPAIGN_ID_RX = re.compile(r"cam_[A-Za-z0-9]{15,}")
_approved_campaigns_cache: Dict[str, set] = {}
_approved_campaigns_lock = asyncio.Lock()

# ──────────────────────────────────────────────────────────────────────────────
# TOOL-CALL DEDUPE — catches duplicate parallel tool-calls within one LLM turn.
# Built after chat c01fabb3 audit showed every lemlist push fired twice (10
# calls for 5 leads) and ~50% of search_knowledge/web_search calls were
# bag-of-words duplicates within ≤2s of each other. Generic across all tools.
# Two layers: (1) in-flight — second concurrent call awaits the first's future;
# (2) recently-completed — second call within TTL returns cached result.
# ──────────────────────────────────────────────────────────────────────────────
import hashlib  # stdlib, safe to import here
TOOL_DEDUPE_TTL_SECONDS = float(os.getenv("TOOL_DEDUPE_TTL_SECONDS", "10"))
_dedupe_inflight: Dict[str, Dict[tuple, asyncio.Future]] = {}
_dedupe_completed: Dict[str, Dict[tuple, tuple]] = {}  # (ts, result)


def _dedupe_key(tool_name: str, args, kwargs) -> tuple:
    """Stable hash of (tool_name, args, kwargs). Pydantic models are dumped to
    dicts so identical models hash equal."""
    def _norm(v):
        if hasattr(v, "model_dump"):
            try: return v.model_dump()
            except Exception: pass
        return v
    try:
        canon = json.dumps(
            [[_norm(a) for a in (args or ())], {k: _norm(v) for k, v in (kwargs or {}).items()}],
            sort_keys=True, default=str,
        )
    except Exception:
        canon = repr((args, kwargs))
    return (tool_name, hashlib.sha256(canon.encode()).hexdigest()[:16])


def _dedupe_evict_expired(chat_id: str, now: float) -> None:
    """Opportunistic cleanup — caller already holds no lock on these dicts.
    Safe because we only touch the chat's own dict, and worst case is a stale
    entry survives one extra call."""
    d = _dedupe_completed.get(chat_id)
    if not d:
        return
    expired = [k for k, (ts, _) in d.items() if now - ts > TOOL_DEDUPE_TTL_SECONDS]
    for k in expired:
        d.pop(k, None)
    if not d:
        _dedupe_completed.pop(chat_id, None)


async def _fetch_approved_campaign_ids(chat_id: str) -> set:
    """Scan every user-role message in this chat for `cam_xxx` tokens.
    The resulting set is the only campaign IDs the agent is allowed to push to.
    """
    if not supabase or not chat_id:
        return set()
    loop = asyncio.get_event_loop()
    def _q():
        return supabase.table("chat_messages").select("content").eq(
            "chat_id", chat_id).eq("role", "user").execute()
    try:
        r = await loop.run_in_executor(None, _q)
        ids = set()
        for row in (r.data or []):
            ids.update(_CAMPAIGN_ID_RX.findall(row.get("content") or ""))
        return ids
    except Exception as e:
        print(f"[INTAKE-GUARD] fetch failed for chat={chat_id}: {e}")
        return set()


def _extract_campaign_ids(args, kwargs) -> List[str]:
    """Pull every `campaign_id` / `campaignId` / `campaign_ids` value out of
    flat kwargs, a nested `params` dict, or positional pydantic models.
    Covers FastMCP flat-arg tools, pydantic-wrapped tools, and bulk variants.
    """
    sources: List[dict] = []
    if isinstance(kwargs, dict):
        sources.append(kwargs)
        if isinstance(kwargs.get("params"), dict):
            sources.append(kwargs["params"])
    # Positional pydantic model (e.g. wrapped(params_obj))
    for a in args or ():
        if hasattr(a, "model_dump"):
            try:
                d = a.model_dump()
                if isinstance(d, dict):
                    sources.append(d)
                    if isinstance(d.get("params"), dict):
                        sources.append(d["params"])
            except Exception:
                pass
        elif isinstance(a, dict):
            sources.append(a)

    found: List[str] = []
    for s in sources:
        for k in ("campaign_id", "campaignId"):
            v = s.get(k)
            if isinstance(v, str) and v.startswith("cam_"):
                found.append(v)
        for k in ("campaign_ids", "campaignIds"):
            v = s.get(k)
            if isinstance(v, (list, tuple)):
                for item in v:
                    if isinstance(item, str) and item.startswith("cam_"):
                        found.append(item)
    return found


class AgentManager:

    MCP_HEALTH_CHECK_INTERVAL = 90
    MCP_RELOAD_MAX_RETRIES = 3
    MCP_RELOAD_RETRY_DELAY = 5
    MCP_SERVER_LOAD_TIMEOUT = 60

    def __init__(self):
        self.mcp_config_manager = MCPConfigManager(config.MCP_CONFIG_FILE)
        self.custom_tools_loader = CustomToolsLoader()
        self.agent = None
        self.mcp_client = None
        self.mcp_loading_status = {}
        self.mcp_tools_loaded = False
        self._mcp_load_task = None
        self._cached_mcp_tools = []
        self._cached_mcp_tools_by_server = {}
        self._cached_custom_tools = []
        self._expected_tool_counts = {}
        self._mcp_health_task = None
        self._initial_load_complete = False
        self._init_lock = asyncio.Lock()
        self._cache_lock = asyncio.Lock()

    # Tools that must never be truncated or summarized — the agent needs
    # the raw output verbatim. Two groups:
    #   1. Avoma meeting tools (transcripts, notes, action items).
    #   2. Reporting "board" tools that return CSV / table rows the agent
    #      renders directly to the user; the gpt-4o-mini summariser would
    #      flatten headers, merge rows, and drop precision.
    _AVOMA_NO_TRUNCATE_TOOLS = {
        # Avoma — meeting data
        "list_meetings",
        "get_all_meetings_for_account",
        "get_all_meetings_for_opportunity",
        "get_all_meetings_for_attendee",
        "get_meeting",
        "get_meeting_transcript",
        "get_meeting_notes",
        "get_meeting_insights",
        "get_meeting_segments",
        "get_meeting_recording_url",
        "get_meeting_action_items",
        "get_meetings_summary_for_account",
        "get_meetings_summary_for_opportunity",
        # Reporting / board tools — preserve full tabular output
        "mailchimp_campaign_performance_board",
        "mailchimp_reports_performance_board",
        "mailchimp_full_campaign_report",
        "mailchimp_get_performance_report_by_date",
        "eloqua_campaign_performance_board",
        "eloqua_get_campaign_email_report",
        "eloqua_get_email_performance",
        "lemlist_get_push_receipts",
    }

    @staticmethod
    def _compress_description(desc: str, max_chars: int) -> str:
        """
        Compress a tool description for the LLM tool catalog.

        The full docstring stays on the underlying MCP function — this only
        trims the copy that gets serialized into the prompt on every turn.
        Strategy:
          1. If max_chars is 0, return as-is (compression disabled).
          2. Drop everything from the first "Args:", "Returns:", "Example",
             "Examples:", "Note:", or "Notes:" header onwards — those are
             redundant with the JSON schema the LLM already gets.
          3. Collapse runs of whitespace into single spaces.
          4. Hard-cap at max_chars, ending at the last sentence boundary.
        """
        if not desc or max_chars <= 0:
            return desc or ""
        text = desc
        for marker in ("\nArgs:", "\nReturns:", "\nExample:", "\nExamples:",
                       "\nNote:", "\nNotes:", "\n    Args:", "\n    Returns:",
                       "\n    Example:", "\n    Examples:"):
            idx = text.find(marker)
            if idx > 0:
                text = text[:idx]
        text = " ".join(text.split())
        if len(text) <= max_chars:
            return text
        cut = text[:max_chars]
        last_period = cut.rfind(". ")
        if last_period > max_chars * 0.6:
            return cut[:last_period + 1]
        return cut.rstrip() + "…"

    def _wrap_mcp_tool(self, original_tool):
        from langchain_core.tools import StructuredTool
        # hasattr() is always True on StructuredTool (Pydantic always declares
        # the `coroutine` field, even when it's None for sync-only tools).
        # Use truthy fallback so sync tools (e.g. custom_tools/web_search_with_urls)
        # don't end up with original_func = None and blow up with
        # 'NoneType is not callable' on first invoke.
        original_func = (
            getattr(original_tool, 'coroutine', None)
            or getattr(original_tool, 'func', None)
        )
        if original_func is None:
            raise ValueError(
                f"Tool {original_tool.name!r} has neither .coroutine nor .func"
            )

        async def wrapped_func(*args, **kwargs):
            # ── INTAKE GUARD: lemlist campaign_id hallucination check ──
            # Hard-fails if the agent passes a campaign_id that the user never
            # typed in any of their messages on this chat. Lazy-refetches on
            # miss to allow the user to introduce a new ID mid-conversation.
            if original_tool.name.startswith("lemlist_"):
                cids = _extract_campaign_ids(args, kwargs)
                chat_id = _current_chat_id.get(None) if cids else None
                if cids and chat_id:
                    # Lock around the cache read+refetch+write so concurrent
                    # tool calls on the same chat can't race a non-monotonic
                    # overwrite.
                    async with _approved_campaigns_lock:
                        approved = _approved_campaigns_cache.get(chat_id)
                        if approved is None or any(c not in approved for c in cids):
                            fresh = await _fetch_approved_campaign_ids(chat_id)
                            # Union, never shrink — protects against refetch
                            # failures or transient empty responses silently
                            # disabling the guard.
                            approved = fresh if approved is None else (approved | fresh)
                            _approved_campaigns_cache[chat_id] = approved
                    if approved:
                        bad_ids = [c for c in cids if c not in approved]
                        if bad_ids:
                            print(
                                f"  [INTAKE-GUARD] BLOCKED hallucinated "
                                f"campaign_id(s)={bad_ids} tool={original_tool.name} "
                                f"chat={chat_id} approved={sorted(approved)}"
                            )
                            return {
                                "error": "HALLUCINATED_CAMPAIGN_ID",
                                "tool": original_tool.name,
                                "rejected_campaign_ids": bad_ids,
                                "approved_campaign_ids": sorted(approved),
                                "message": (
                                    f"Campaign ID(s) {bad_ids} were never provided by the user "
                                    f"in this chat. Approved IDs from user intake messages: "
                                    f"{sorted(approved)}. Re-read the user's intake message "
                                    "verbatim and use the campaign_id they actually provided. "
                                    "Do NOT invent, substitute, or recall a campaign_id from "
                                    "memory — only use one of the approved IDs above."
                                ),
                            }

            # ── CHAT-ID INJECTION: lemlist receipt tools ──
            # `lemlist_validated_push` and `lemlist_get_push_receipts` both
            # take `chat_id` as their first arg and use it as the key in
            # `public.lemlist_push_receipts`. The LLM sometimes fabricates a
            # value (e.g. "current_session_aaf_001") which orphans receipts
            # under a key nobody queries. Override with the real chat UUID
            # from the per-agent ContextVar whenever it's set.
            real_chat_id_for_override = None
            if original_tool.name in (
                "lemlist_validated_push",
                "lemlist_get_push_receipts",
            ):
                _real_cid = _current_chat_id.get(None)
                if _real_cid:
                    if "chat_id" in kwargs:
                        _supplied = kwargs.get("chat_id")
                    elif len(args) >= 1:
                        _supplied = args[0]
                    else:
                        _supplied = None
                    if _supplied != _real_cid:
                        print(
                            f"  [CHAT-ID-INJECT] tool={original_tool.name} "
                            f"supplied={_supplied!r} -> real={_real_cid!r}"
                        )
                        if "chat_id" in kwargs or not args:
                            kwargs["chat_id"] = _real_cid
                        else:
                            args = (_real_cid,) + tuple(args[1:])
                        real_chat_id_for_override = _real_cid
            # Bound every async MCP tool call so one hung stdio subprocess (e.g. an Avoma
            # transcript read stuck under contention) can't pin this agent run for the
            # full ~660s watchdog window and hold a session slot. On timeout we return the
            # same {error,status:failed} shape so the agent's error-recovery path handles
            # it instead of hanging. The default (300s) sits ABOVE a worst-case legit Avoma
            # call (AVOMA_HTTP_TIMEOUT 60s x (retries 2 +1) ~= 180s) so it won't cut valid
            # slow reads, and BELOW the watchdog so it still frees the slot early. The
            # sweep worker raises this (MCP_TOOL_TIMEOUT_S=600 in deploy.ps1) for its longer
            # timeout posture. NOTE: only the async path is time-bounded; a purely
            # synchronous tool func relies on its own internal timeout (all MCP-adapter
            # tools here are async coroutines, so this covers them).
            _tool_timeout = float(os.getenv("MCP_TOOL_TIMEOUT_S", "300"))
            try:
                if asyncio.iscoroutinefunction(original_func):
                    result = await asyncio.wait_for(original_func(*args, **kwargs), timeout=_tool_timeout)
                else:
                    result = original_func(*args, **kwargs)
                    if asyncio.iscoroutine(result):
                        result = await asyncio.wait_for(result, timeout=_tool_timeout)
            except asyncio.TimeoutError:
                error_msg = f"MCP tool '{original_tool.name}' timed out after {int(_tool_timeout)}s"
                print(f"  [MCP TIMEOUT] {error_msg}")
                return {
                    "error": error_msg,
                    "tool": original_tool.name,
                    "status": "failed"
                }
            except Exception as tool_err:
                error_msg = f"MCP tool '{original_tool.name}' error: {str(tool_err)}"
                print(f"  [MCP ERROR] {error_msg}")
                return {
                    "error": error_msg,
                    "tool": original_tool.name,
                    "status": "failed"
                }
            if isinstance(result, tuple):
                result = result[0] if len(result) > 0 else result
            if isinstance(result, str):
                try:
                    result = json.loads(result)
                except:
                    pass

            # Mirror the chat_id override into the tool's returned payload so
            # the agent sees its receipts under the real chat UUID in context
            # replays. Tolerant of three shapes:
            #   (a) dict with top-level "chat_id"
            #   (b) JSON string
            #   (c) MCP content-block list: [{"type":"text","text":"<json>"}]
            if real_chat_id_for_override is not None:
                try:
                    if isinstance(result, dict) and "chat_id" in result:
                        result["chat_id"] = real_chat_id_for_override
                    elif isinstance(result, list):
                        for _blk in result:
                            if not isinstance(_blk, dict):
                                continue
                            _txt = _blk.get("text")
                            if not isinstance(_txt, str):
                                continue
                            try:
                                _inner = json.loads(_txt)
                            except Exception:
                                continue
                            if isinstance(_inner, dict) and "chat_id" in _inner:
                                _inner["chat_id"] = real_chat_id_for_override
                                _blk["text"] = json.dumps(_inner, indent=2, default=str)
                except Exception as _rewrite_err:
                    print(f"  [CHAT-ID-INJECT] payload rewrite skipped: {_rewrite_err}")

            result_str = json.dumps(result, default=str) if isinstance(
                result, (dict, list)) else str(result)

            # Avoma tools: always return full data, no summarization or truncation.
            if original_tool.name in AgentManager._AVOMA_NO_TRUNCATE_TOOLS:
                return result

            SUMMARIZE_THRESHOLD = config.TOOL_RESPONSE_SUMMARIZE_THRESHOLD
            MAX_RESPONSE_SIZE = config.MCP_MAX_RESPONSE_SIZE
            MAX_STRING_LENGTH = config.MCP_MAX_STRING_LENGTH
            MAX_LIST_ITEMS = config.MCP_MAX_LIST_ITEMS

            if len(result_str) <= SUMMARIZE_THRESHOLD:
                return result

            os.makedirs("mcp_output", exist_ok=True)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            unique_id = uuid.uuid4().hex[:8]
            filename = f"mcp_output/{original_tool.name}_{timestamp}_{unique_id}.json"
            with open(filename, 'w') as f:
                try:
                    if isinstance(result, (dict, list)):
                        json.dump(result, f, indent=2, default=str)
                    else:
                        f.write(result_str)
                except:
                    f.write(result_str)
            print(
                f"  Saved full response to {filename} ({len(result_str):,} chars)"
            )

            if len(result_str) <= MAX_RESPONSE_SIZE and not _skip_llm_summarizer.get():
                try:
                    summarized = await context_manager.summarize_tool_response(
                        original_tool.name, result_str)
                    return summarized
                except Exception as e:
                    print(
                        f"  Summarization failed, falling back to truncation: {e}"
                    )

            if isinstance(result, dict):
                truncated_result = {}
                for key, value in result.items():
                    if isinstance(value,
                                  str) and len(value) > MAX_STRING_LENGTH:
                        truncated_result[
                            key] = value[:MAX_STRING_LENGTH] + "...[truncated]"
                    elif isinstance(value,
                                    list) and len(value) > MAX_LIST_ITEMS:
                        truncated_result[key] = value[:MAX_LIST_ITEMS]
                        truncated_result[key].append(
                            f"... and {len(value) - MAX_LIST_ITEMS} more items"
                        )
                    else:
                        truncated_result[key] = value
                truncated_str = json.dumps(truncated_result, default=str)
            elif isinstance(result, list):
                truncated_list = []
                for item in result[:MAX_LIST_ITEMS]:
                    if isinstance(item, dict):
                        truncated_item = {}
                        for key, value in item.items():
                            if isinstance(
                                    value,
                                    str) and len(value) > MAX_STRING_LENGTH:
                                truncated_item[
                                    key] = value[:MAX_STRING_LENGTH] + "...[truncated]"
                            else:
                                truncated_item[key] = value
                        truncated_list.append(truncated_item)
                    else:
                        truncated_list.append(item)
                if len(result) > MAX_LIST_ITEMS:
                    truncated_list.append(
                        f"... and {len(result) - MAX_LIST_ITEMS} more items")
                truncated_str = json.dumps(truncated_list, default=str)
            else:
                truncated_str = result_str[:MAX_STRING_LENGTH]

            if _skip_llm_summarizer.get():
                # Analyzer path: skip the LLM summariser entirely (it times out
                # at 45s on big payloads). Return the deterministically-truncated
                # data directly — same result the summariser falls back to anyway.
                return truncated_str[:MAX_STRING_LENGTH] + f"\n\n[Response truncated. Full data: {filename}]"
            try:
                summarized = await context_manager.summarize_tool_response(
                    original_tool.name, truncated_str)
                return summarized
            except Exception as e:
                print(f"  Summarization failed on truncated data: {e}")
                return truncated_str[:MAX_STRING_LENGTH] + f"\n\n[Response truncated. Full data: {filename}]"

        HARD_CAP = int(os.getenv("MCP_HARD_CAP", "100000"))

        async def capped_func(*args, **kwargs):
            result = await wrapped_func(*args, **kwargs)
            # Avoma tools are exempt from the hard cap — full data must be preserved.
            if original_tool.name in AgentManager._AVOMA_NO_TRUNCATE_TOOLS:
                return result
            result_str = json.dumps(result, default=str) if isinstance(result, (dict, list)) else str(result)
            if len(result_str) > HARD_CAP:
                print(f"  [HARD CAP] {original_tool.name} result {len(result_str):,} chars -> {HARD_CAP:,} chars")
                return result_str[:HARD_CAP] + f"\n...[truncated from {len(result_str):,} chars]"
            return result

        async def deduped_func(*args, **kwargs):
            """Outermost wrapper: dedupes identical tool calls within
            TOOL_DEDUPE_TTL_SECONDS per-chat. Fails open if no chat context."""
            chat_id = _current_chat_id.get(None)
            if not chat_id:
                return await capped_func(*args, **kwargs)
            key = _dedupe_key(original_tool.name, args, kwargs)
            now = time.time()
            _dedupe_evict_expired(chat_id, now)
            completed = _dedupe_completed.get(chat_id, {})
            hit = completed.get(key)
            if hit is not None:
                ts, cached_result = hit
                if now - ts <= TOOL_DEDUPE_TTL_SECONDS:
                    print(f"  [DEDUPE] cache-hit {original_tool.name} chat={chat_id} age={now-ts:.1f}s")
                    return cached_result
            inflight = _dedupe_inflight.setdefault(chat_id, {})
            existing = inflight.get(key)
            if existing is not None and not existing.done():
                print(f"  [DEDUPE] in-flight await {original_tool.name} chat={chat_id}")
                return await existing
            fut: asyncio.Future = asyncio.get_event_loop().create_future()
            inflight[key] = fut
            try:
                result = await capped_func(*args, **kwargs)
                if not fut.done():
                    fut.set_result(result)
                _dedupe_completed.setdefault(chat_id, {})[key] = (time.time(), result)
                return result
            except Exception as e:
                if not fut.done():
                    fut.set_exception(e)
                raise
            finally:
                inflight.pop(key, None)
                if not inflight:
                    _dedupe_inflight.pop(chat_id, None)

        compressed_desc = self._compress_description(
            original_tool.description or "",
            config.TOOL_DESCRIPTION_MAX_CHARS,
        )
        return StructuredTool(name=original_tool.name,
                              description=compressed_desc,
                              coroutine=deduped_func,
                              args_schema=original_tool.args_schema)

    @staticmethod
    def _log_tool_catalog_size(tools) -> None:
        """
        Log a rough estimate of how many input tokens the tool catalog adds to
        every LLM call. ~4 chars/token is the standard rule of thumb. The
        catalog is re-sent on every turn, so this is the per-turn fixed cost.
        """
        if not tools:
            return
        total_chars = 0
        for t in tools:
            total_chars += len(t.name or "")
            total_chars += len(getattr(t, "description", "") or "")
            schema = getattr(t, "args_schema", None)
            if schema is not None:
                try:
                    schema_dict = schema.model_json_schema() if hasattr(schema, "model_json_schema") else {}
                    total_chars += len(json.dumps(schema_dict))
                except Exception:
                    pass
        approx_tokens = total_chars // 4
        print(
            f"[TOOL-CATALOG] {len(tools)} tools, ~{total_chars:,} chars, "
            f"~{approx_tokens:,} input tokens per LLM turn "
            f"(description cap: {config.TOOL_DESCRIPTION_MAX_CHARS} chars)"
        )

    async def _load_single_server_with_retry(self, name, cfg, max_retries=None):
        from langchain_mcp_adapters.client import MultiServerMCPClient
        if max_retries is None:
            max_retries = self.MCP_RELOAD_MAX_RETRIES
        for attempt in range(1, max_retries + 1):
            try:
                label = f"(attempt {attempt}/{max_retries}) " if max_retries > 1 else ""
                print(f"Loading MCP server: {name}... {label}")
                self.mcp_loading_status[name] = f"loading (attempt {attempt})"
                client = MultiServerMCPClient({name: cfg})
                server_tools = await asyncio.wait_for(
                    client.get_tools(),
                    timeout=self.MCP_SERVER_LOAD_TIMEOUT
                )
                if server_tools:
                    print(f"  Loaded {len(server_tools)} tools from {name}")
                    self.mcp_loading_status[name] = f"ready ({len(server_tools)} tools)"
                    wrapped = [self._wrap_mcp_tool(t) for t in server_tools]
                    async with self._cache_lock:
                        self._cached_mcp_tools_by_server[name] = wrapped
                        self._expected_tool_counts[name] = len(server_tools)
                    return wrapped
                else:
                    print(f"  {name} returned 0 tools (attempt {attempt})")
                    if attempt < max_retries:
                        await asyncio.sleep(self.MCP_RELOAD_RETRY_DELAY * attempt)
                    continue
            except asyncio.TimeoutError:
                print(f"  TIMEOUT loading {name} after {self.MCP_SERVER_LOAD_TIMEOUT}s (attempt {attempt})")
                if attempt < max_retries:
                    await asyncio.sleep(self.MCP_RELOAD_RETRY_DELAY * attempt)
            except Exception as e:
                print(f"  FAILED to load {name} (attempt {attempt}): {e}")
                if attempt < max_retries:
                    await asyncio.sleep(self.MCP_RELOAD_RETRY_DELAY * attempt)
        self.mcp_loading_status[name] = f"failed after {max_retries} attempts"
        print(f"  {name}: all {max_retries} attempts exhausted")
        return []

    async def _rebuild_flat_cache(self):
        async with self._cache_lock:
            all_tools = []
            for server_name, tools in self._cached_mcp_tools_by_server.items():
                all_tools.extend(tools)
            # Apply tool-level allow/denylist (cost task #40, 2026-05-22).
            # MCP_TOOL_ALLOWLIST wins if non-empty; otherwise MCP_TOOL_DENYLIST
            # applies. fnmatch glob support: e.g. "lemlist_*", "mailchimp_*".
            allow = config.MCP_TOOL_ALLOWLIST
            deny = config.MCP_TOOL_DENYLIST
            if allow or deny:
                import fnmatch
                before = len(all_tools)
                if allow:
                    all_tools = [
                        t for t in all_tools
                        if any(fnmatch.fnmatchcase(t.name, p) for p in allow)
                    ]
                    print(
                        f"[TOOL-FILTER] MCP_TOOL_ALLOWLIST={allow} → "
                        f"kept {len(all_tools)}/{before} MCP tools"
                    )
                elif deny:
                    dropped = [
                        t.name for t in all_tools
                        if any(fnmatch.fnmatchcase(t.name, p) for p in deny)
                    ]
                    all_tools = [
                        t for t in all_tools
                        if not any(fnmatch.fnmatchcase(t.name, p) for p in deny)
                    ]
                    print(
                        f"[TOOL-FILTER] MCP_TOOL_DENYLIST={deny} → "
                        f"dropped {len(dropped)}/{before} MCP tools "
                        f"(remaining: {len(all_tools)}); first 10 dropped: {dropped[:10]}"
                    )
            self._cached_mcp_tools = all_tools
            self.mcp_tools_loaded = bool(all_tools)
        self._log_tool_catalog_size(all_tools)
        return all_tools

    async def _load_mcp_tools_sync(self):
        enabled_mcp_servers = self.mcp_config_manager.get_enabled_servers()
        print(f"Enabled MCP servers: {list(enabled_mcp_servers.keys())}")
        if not enabled_mcp_servers:
            print("No enabled MCP servers found in config")
            return []

        try:
            failed_servers = []

            results = await asyncio.gather(*[
                self._load_single_server_with_retry(name, cfg)
                for name, cfg in enabled_mcp_servers.items()
            ])

            for name, tool_list in zip(enabled_mcp_servers.keys(), results):
                if not tool_list:
                    failed_servers.append(name)

            if failed_servers:
                print(
                    f"Warning: {len(failed_servers)} MCP server(s) failed to load: {failed_servers}"
                )

            all_tools = await self._rebuild_flat_cache()
            print(f"Successfully loaded {len(all_tools)} MCP tools total")
            print(
                f"Loaded {len(all_tools)} MCP tools from {len(enabled_mcp_servers)} servers (with ChatGPT summarization)"
            )
            print(f"  Tool names: {[t.name for t in all_tools]}")
            per_server = {k: len(v) for k, v in self._cached_mcp_tools_by_server.items()}
            print(f"  Per-server tool counts: {per_server}")
            return all_tools

        except ImportError:
            print(
                "Warning: langchain-mcp-adapters not installed. MCP support disabled."
            )
            return []
        except Exception as e:
            print(f"Error loading MCP tools: {e}")
            import traceback
            traceback.print_exc()
            return []

    async def _clean_npx_cache(self):
        try:
            import subprocess
            npx_cache = os.path.expanduser("~/.npm/_npx")
            if os.path.exists(npx_cache):
                print(f"[BACKGROUND] Cleaning npx cache at {npx_cache}...")
                result = subprocess.run(
                    ["rm", "-rf", npx_cache],
                    capture_output=True, text=True, timeout=30
                )
                if result.returncode == 0:
                    print("[BACKGROUND] npx cache cleaned successfully")
                else:
                    print(f"[BACKGROUND] npx cache clean warning: {result.stderr}")
            else:
                print("[BACKGROUND] No npx cache to clean")
        except Exception as e:
            print(f"[BACKGROUND] npx cache clean error (non-fatal): {e}")

    async def _reinitialize_agent_with_current_tools(self):
        all_tools = await self._rebuild_flat_cache()
        if all_tools:
            async with self._init_lock:
                await self.initialize_agent()
            print(f"[BACKGROUND] Agent reinitialized with {len(all_tools)} MCP tools (old agent replaced atomically)")

    async def _load_mcp_background(self):
        try:
            print("[BACKGROUND] Starting MCP server loading...")
            start_time = asyncio.get_event_loop().time()
            await self._clean_npx_cache()

            enabled_mcp_servers = self.mcp_config_manager.get_enabled_servers()
            print(f"Enabled MCP servers: {list(enabled_mcp_servers.keys())}")

            tasks = {
                name: asyncio.create_task(self._load_single_server_with_retry(name, cfg))
                for name, cfg in enabled_mcp_servers.items()
            }

            agent_initialized = False
            pending = set(tasks.values())

            while pending:
                done, pending = await asyncio.wait(pending, return_when=asyncio.FIRST_COMPLETED)

                new_tools_loaded = False
                for task in done:
                    result = task.result()
                    if result:
                        new_tools_loaded = True

                if new_tools_loaded and not agent_initialized:
                    agent_initialized = True
                    await self._reinitialize_agent_with_current_tools()
                elif new_tools_loaded and agent_initialized:
                    await self._reinitialize_agent_with_current_tools()

            elapsed = asyncio.get_event_loop().time() - start_time
            print(
                f"[BACKGROUND] MCP loading complete in {elapsed:.1f}s ({len(self._cached_mcp_tools)} tools)"
            )

            if not agent_initialized and self._cached_mcp_tools:
                await self._reinitialize_agent_with_current_tools()

            self._initial_load_complete = True
            failed = [n for n, s in self.mcp_loading_status.items() if "failed" in str(s)]
            if failed:
                print(f"[BACKGROUND] {len(failed)} server(s) failed initial load: {failed}")
                print(f"[BACKGROUND] Will retry failed servers in health check loop")
        except Exception as e:
            self._initial_load_complete = True
            print(f"[BACKGROUND] MCP loading error: {e}")
            import traceback
            traceback.print_exc()

    async def _reload_failed_servers(self):
        enabled_mcp_servers = self.mcp_config_manager.get_enabled_servers()
        failed_servers = {}
        for name in enabled_mcp_servers:
            status = str(self.mcp_loading_status.get(name, ""))
            if "loading" in status:
                continue
            if "failed" in status or name not in self._cached_mcp_tools_by_server or not self._cached_mcp_tools_by_server.get(name):
                failed_servers[name] = enabled_mcp_servers[name]

        if not failed_servers:
            return False

        print(f"[MCP-HEALTH] Reloading {len(failed_servers)} failed server(s): {list(failed_servers.keys())}")
        reloaded_any = False
        results = await asyncio.gather(*[
            self._load_single_server_with_retry(name, cfg)
            for name, cfg in failed_servers.items()
        ])
        for name, tool_list in zip(failed_servers.keys(), results):
            if tool_list:
                print(f"[MCP-HEALTH] Successfully reloaded {name} ({len(tool_list)} tools)")
                reloaded_any = True
            else:
                print(f"[MCP-HEALTH] {name} still failed to load")

        if reloaded_any:
            await self._rebuild_flat_cache()
            print(f"[MCP-HEALTH] Cache rebuilt: {len(self._cached_mcp_tools)} total tools")
            per_server = {k: len(v) for k, v in self._cached_mcp_tools_by_server.items()}
            print(f"[MCP-HEALTH] Per-server: {per_server}")
            async with self._init_lock:
                await self.initialize_agent()
                print(f"[MCP-HEALTH] Agent reinitialized with {len(self._cached_mcp_tools)} MCP tools")

        return reloaded_any

    async def _mcp_health_check_loop(self):
        while not self._initial_load_complete:
            await asyncio.sleep(5)
        print("[MCP-HEALTH] Initial load complete, starting health check loop")
        while True:
            try:
                await asyncio.sleep(self.MCP_HEALTH_CHECK_INTERVAL)
                enabled_mcp_servers = self.mcp_config_manager.get_enabled_servers()
                missing = []
                for name in enabled_mcp_servers:
                    status = str(self.mcp_loading_status.get(name, ""))
                    if "loading" in status:
                        continue
                    if name not in self._cached_mcp_tools_by_server or not self._cached_mcp_tools_by_server.get(name):
                        missing.append(name)
                    elif "failed" in status:
                        missing.append(name)

                if missing:
                    print(f"[MCP-HEALTH] Detected {len(missing)} missing/failed server(s): {missing}")
                    await self._reload_failed_servers()
                else:
                    current_total = len(self._cached_mcp_tools)
                    per_server = {k: len(v) for k, v in self._cached_mcp_tools_by_server.items()}
                    print(f"[MCP-HEALTH] All servers healthy. {current_total} tools cached. {per_server}")

            except Exception as e:
                print(f"[MCP-HEALTH] Health check error: {e}")

    def start_mcp_background_loading(self):
        enabled_servers = self.mcp_config_manager.get_enabled_servers()
        for name in enabled_servers:
            self.mcp_loading_status[name] = "pending"
        self._mcp_load_task = asyncio.create_task(self._load_mcp_background())
        self._mcp_health_task = asyncio.create_task(self._mcp_health_check_loop())
        print(
            f"[BACKGROUND] MCP loading task started for {len(enabled_servers)} servers"
        )
        print(f"[BACKGROUND] MCP health check loop started (interval: {self.MCP_HEALTH_CHECK_INTERVAL}s)")

    async def initialize_agent(self,
                               instructions: Optional[str] = None,
                               model: Optional[str] = None,
                               headless: bool = True,
                               skip_mcp: bool = False):
        tools = [web_search, get_current_time]
        custom_tools = self.custom_tools_loader.load_tools_from_directory(
            config.CUSTOM_TOOLS_DIR)
        wrapped_custom_tools = []
        for t in custom_tools:
            if t.name in [
                    'browser_research', 'browser_research_multiple',
                    'browser_interactive_research'
            ]:
                original_func = t.coroutine if hasattr(t,
                                                       'coroutine') else t.func

                def create_wrapped_browser_tool(original, headless_val):

                    async def wrapped_func(*args, **kwargs):
                        kwargs['headless'] = headless_val
                        if asyncio.iscoroutinefunction(original):
                            return await original(*args, **kwargs)
                        else:
                            return original(*args, **kwargs)

                    return wrapped_func

                wrapped_func = create_wrapped_browser_tool(
                    original_func, headless)
                from langchain_core.tools import StructuredTool
                wrapped_tool = StructuredTool(
                    name=t.name,
                    description=t.description +
                    f" (Browser mode: {'headless/invisible' if headless else 'visible/slow'})",
                    coroutine=wrapped_func,
                    args_schema=t.args_schema)
                wrapped_custom_tools.append(wrapped_tool)
                print(
                    f"Wrapped browser tool: {t.name} with headless={headless}")
            else:
                wrapped_custom_tools.append(t)
        # Pipe custom tools through the same wrapper as MCP tools so they get
        # dedupe + summarisation + hard-cap + (lemlist) intake-guard. Without
        # this, custom tools like web_search_with_urls / search_knowledge /
        # web_scrape / send_to_clay would bypass the dedupe layer entirely —
        # which is exactly what live-test chat 1e1fe9cf exposed on 2026-05-21
        # (agent fired 3 identical parallel web_search_with_urls calls and
        # all 3 hit the real backend instead of sharing one result).
        wrapped_custom_tools = [self._wrap_mcp_tool(t) for t in wrapped_custom_tools]
        self._cached_custom_tools = wrapped_custom_tools
        tools.extend(wrapped_custom_tools)

        if not skip_mcp:
            if self._cached_mcp_tools:
                tools.extend(self._cached_mcp_tools)
                print(f"Using {len(self._cached_mcp_tools)} cached MCP tools")
            else:
                mcp_tools = await self._load_mcp_tools_sync()
                tools.extend(mcp_tools)

        if instructions is None:
            browser_mode = "headless=True (invisible, fast)" if headless else "headless=False (visible, slow)"
            instructions = f"""You are an expert AI assistant specialized in B2B sales intelligence, CRM operations, and business research. You have access to 100+ tools across 8 integrated data platforms and web search.

TODAY'S DATE: Use get_current_time to get the current date/time when needed for date ranges, lookback windows, or time-sensitive queries.

THINKING OUT LOUD RULE: You MUST always explain your reasoning between tool calls. Before calling any tool, briefly state:
- What information you have so far
- What you need next and why
- Which tool you will call and what you expect to learn from it
After receiving a tool result, briefly summarize what you learned before proceeding. This helps the user follow your thought process in real-time.

CRITICAL ERROR HANDLING RULE: If ANY tool returns an error (rate limit, timeout, 403, 404, or any other error), you MUST NEVER stop or give up. Instead:
1. Acknowledge the error briefly
2. Try an alternative approach (different tool, different query, different data source)
3. If no alternative works, continue with the information you already have
4. Complete the task to the best of your ability with available data
5. NEVER let a single tool failure stop your entire workflow

CONTEXT MANAGEMENT (cost-saving placeholders):
On long runs the system automatically compresses older tool results to save cost. If you see a message in your earlier context that looks like:
   `[Earlier <tool_name> result summarised: N chars omitted; first 200: <snippet>...]`
that means the full result was replaced with a short placeholder. Rules:
1. Do NOT re-call the tool to "recover" the result. Re-calling wastes a tool call and the next compression pass will just shrink it again.
2. Use the first-200-char snippet in the placeholder, or whatever is still visible in your most recent ~10 messages.
3. If the placeholder is for a result you genuinely still need in full, say so to the user and ask them to narrow the scope — don't loop on the same tool call.
4. This only happens on heavy multi-tool runs; for short tasks you'll never see it.

WEB SEARCH DISCIPLINE (MANDATORY):
1. MAXIMUM 15 web searches per task phase. If you haven't found what you need in 15, work with what you have. Stop searching and move on.
2. ONE search per topic. Never rephrase and retry the same topic. If the first search didn't return useful results, accept that and move on.
3. Before every web search, ask yourself: "Have I already searched for this or something very similar?" If yes, DO NOT search again.
4. Combine related queries into a single broad search instead of multiple narrow ones. Example: search "Jumeirah Group CEO leadership changes 2024 2025" ONCE, not "Jumeirah CEO", then "Jumeirah leadership", then "Jumeirah management changes" separately.
5. NEVER use web search to validate individual contacts (checking if someone still works at a company, finding their LinkedIn, verifying their title). That is the job of Lemlist, ZoomInfo, Apollo, or Seamless.ai contact enrichment tools. Web-searching a person's name + company + LinkedIn is unreliable and wasteful.
6. NEVER search for topics irrelevant to the user's specific request (generic industry trends, brand identity changes, sustainability initiatives — unless explicitly asked).

WEB SEARCH AND SCRAPING TOOLS:
- web_search(query) — Basic text search. Returns search result snippets as plain text. Fast but does NOT return URLs.
- web_search_with_urls(query, max_results=10) — Structured search. Returns results WITH URLs, titles, and snippets as JSON. Use this when you need the actual website links (e.g., to share URLs with the user, to scrape specific pages, or to reference sources).
- web_scrape(url) — Fetch and extract readable text content from a web page. Strips out navigation, ads, scripts, and non-content elements. Use this after web_search_with_urls to get full page content from specific URLs. Works with any http/https URL.
WORKFLOW: When you need both search results AND page content, use web_search_with_urls first to get URLs, then web_scrape on the most relevant URLs to get full content.

TOOL SELECTION RULES (MANDATORY):
1. Contact lookup/validation → Use ZoomInfo (zi_search_contacts, zi_enrich_contact), Apollo (apollo_search_people, apollo_enrich_person), or Seamless.ai (seamless_search_contacts). NEVER use web_search for this.
   APOLLO PHONE NUMBERS: NEVER pass reveal_phone_number=true to apollo_enrich_person or apollo_bulk_enrich_people. Apollo's phone reveal is asynchronous (requires a webhook callback) and will always fail with an error in this environment. For phone/mobile numbers use Lusha (lusha_enrich_person), ZoomInfo (zi_enrich_contact), Wiza (wiza_reveal_contact), or Lemlist (lemlist_enrich_lead) instead. Apollo is excellent for email, LinkedIn URL, title, and company data — just not for phone reveal.
2. Company enrichment → Use ZoomInfo (zi_search_companies, zi_enrich_company), Apollo (apollo_enrich_company, apollo_search_companies). Web search is a LAST RESORT only.
3. Email verification/validation → Use ZeroBounce tools (zerobounce_validate_email, zerobounce_batch_validate, zerobounce_find_email, zerobounce_email_scoring). ZeroBounce can validate emails, find emails, check deliverability, detect disposable/abuse emails, and score email quality. NEVER google someone's email.
4. CRM data → Use Salesforce tools (soql, get_record, search). NEVER approximate CRM data from web results.
5. Meeting notes/transcripts → Use Avoma tools. NEVER try to reconstruct meeting content from other sources.
6. Campaign management → Use Lemlist tools (lemlist_list_campaigns, lemlist_create_campaign, lemlist_add_lead_to_campaign, lemlist_add_leads_batch, lemlist_update_lead, lemlist_get_lead_in_campaign, lemlist_get_lead, lemlist_get_campaign_stats, lemlist_export_campaign, lemlist_get_activities, lemlist_send_email, lemlist_search_people, lemlist_enrich_lead, lemlist_get_enrichment_result, lemlist_create_webhook, etc.). NEVER manually compose what Lemlist can automate.
   LEAD CAMPAIGN MEMBERSHIP (CRITICAL — READ BEFORE USING LEMLIST):
   - To check which campaigns a lead is enrolled in → call lemlist_get_lead(email). ONE call returns ALL campaign memberships across the entire account.
   - NEVER call lemlist_get_lead_in_campaign in a loop across multiple campaigns to discover enrollment. This makes N API calls (one per campaign) and causes severe context window bloat — a 50-campaign check costs 50x the tokens of one lemlist_get_lead call.
   - lemlist_get_lead_in_campaign is ONLY for fetching the _id (lead ID) from ONE KNOWN campaign when you need to call lemlist_update_lead. It is a targeted lookup, not a discovery tool.
   UPDATING EXISTING LEADS: When you need to update ANY field on an existing lead (LinkedIn URL, phone, name, custom variables, contactOwner, etc.), use this two-step process:
   Step 1: Call lemlist_get_lead_in_campaign(campaign_id, email) → returns the lead details including _id (lead ID). Use this ONLY when you already know the campaign_id.
   Step 2: Call lemlist_update_lead(campaign_id, lead_id, ...) → PATCH only the fields you want to change. All other fields remain untouched.
   NEVER tell the user to "manually edit in Lemlist UI." You have the tools to update leads directly. NEVER delete and re-add a lead just to change one field — use lemlist_update_lead instead.
   ENRICHMENT: lemlist_enrich_lead can find emails, phone numbers, and LinkedIn profile data. Provide at least one of: email, linkedin_url, or (first_name + last_name + company_name). The tool polls automatically for up to 30 seconds. If results are not ready, use lemlist_get_enrichment_result with the returned enrichment_id to poll later.
7. Sales content/collateral → Use Showpad tools (list_assets, search_assets, list_shared_spaces). NEVER try to find sales content via web search.
8. LinkedIn Ads / advertising analytics → Use LinkedIn Ads tools (linkedin_list_ad_accounts, linkedin_get_ad_account, linkedin_get_campaign_groups, linkedin_get_campaigns, linkedin_get_campaign_details, linkedin_get_campaign_analytics, linkedin_get_account_analytics). Use these for ad account overviews, campaign performance, and campaign analysis. NEVER try to get LinkedIn Ads data from web search.
9. Mailchimp campaign performance reports:
   - SINGLE DATE → mailchimp_full_campaign_report(target_date="YYYY-MM-DD"). Do NOT call mailchimp_list_campaigns or mailchimp_get_report first — this single tool handles everything.
   - MULTI-DAY RANGE (2+ days, e.g. "last 10 days", "Apr 21-30") → mailchimp_reports_performance_board(since_send_time="...", before_send_time="..."). This tool hits the /reports endpoint directly — one paginated call returns all campaign metrics with no per-campaign fetches. Do NOT loop day-by-day, do NOT call mailchimp_full_campaign_report repeatedly — that exhausts the context window.
   - NEVER call mailchimp_get_all_sent_campaigns + individual report calls — this is the slow path that causes timeouts.

SHOWPAD SALES ENABLEMENT TOOLS (14 tools available):
Showpad provides sales enablement and content management tools. Key tools include:
- search_assets(keyword="BFSI") - PRIMARY SEARCH: Find assets by keyword using ShowQL name~ operator. Supports optional filters: asset_type ("document","video","image"), shareable_only. USE THIS for finding case studies, presentations, demos by topic/industry/customer.
- query_assets(showql_query="...") - RAW ShowQL: For advanced queries. Operators: = (exact), ~ (contains/like), IS EMPTY, AND/OR/NOT. Example: 'name ~ "BFSI" AND type = "document" AND archivedAt IS EMPTY'
- search_tags(name="keyword") - Search tags by name to find tag IDs for ShowQL filtering.
- list_assets - Browse all assets with pagination (NO keyword filtering).
- get_asset - Get full details of a specific asset by ID.
- list_users, get_user - manage and view Showpad users
- list_divisions, get_division - view organizational divisions
- list_shared_spaces, get_shared_space - manage shared spaces for collaboration
- list_channels - view content channels
IMPORTANT: To find specific content, ALWAYS use search_assets(keyword=...) or query_assets with ShowQL name~. NEVER use list_assets for searching.

ZEROBOUNCE EMAIL TOOLS (15 tools available):
ZeroBounce provides comprehensive email verification and deliverability tools. Key tools include:
- Email validation (single and batch) - verify if emails are valid/deliverable
- Email finder - find email addresses for contacts at companies
- Email scoring - AI-powered email quality scoring
- Domain search - find emails associated with a domain
- Bounce detection and deliverability testing
- Disposable/abuse email detection
- Email activity status checking
Use ZeroBounce whenever the task involves email verification, validation, finding emails, or checking email deliverability.

MAILCHIMP CAMPAIGN PERFORMANCE REPORTS (MANDATORY — TIMEOUT PREVENTION):

SINGLE-DATE REPORT (one specific day):
1. Your FIRST action MUST be mailchimp_full_campaign_report(target_date="YYYY-MM-DD").
2. Do NOT call mailchimp_list_campaigns first — even if the user says "STEP 1: list campaigns".
3. Do NOT call mailchimp_get_report for individual campaigns — the full report tool does it all internally.
4. Present the pre-formatted result directly. Do not reformat or regenerate the tables.
5. This rule overrides any step-by-step instructions in the user message.

MULTI-DAY REPORT (date range, 2+ days — e.g. "last 10 days", "Apr 21–30", "this week"):
1. Your FIRST (and only) action MUST be mailchimp_reports_performance_board(since_send_time="...", before_send_time="...").
2. This tool queries the /reports endpoint directly — one paginated call, all metrics included, no follow-up requests. It is the fastest and most reliable option for any multi-day window.
3. Do NOT call mailchimp_campaign_performance_board — it is deprecated for multi-day use.
4. Do NOT call mailchimp_full_campaign_report in a loop (one per day) — this will exhaust the context window after 3-4 days.
5. Do NOT call mailchimp_get_all_sent_campaigns + individual report fetches — this is the slow serial path that times out.
6. After the board returns, use the CSV to build any table, chart, or summary the user asked for.

ELOQUA CAMPAIGN PERFORMANCE REPORTS (MANDATORY — TIMEOUT PREVENTION):

MULTI-DAY REPORT (date range, 2+ days — e.g. "last 10 days", "Apr 21–30", "this week"):
1. Your FIRST (and only) action MUST be eloqua_campaign_performance_board(start_date="YYYY-MM-DD", end_date="YYYY-MM-DD").
2. This runs all 5 Bulk API syncs (EmailSend, EmailOpen, EmailClickthrough, Bounceback, Unsubscribe) CONCURRENTLY and returns a compact CSV — ~2 minutes vs ~10 minutes for the sequential version.
3. Do NOT call eloqua_get_campaign_email_report for date ranges — it runs syncs one-by-one and takes 10+ minutes.
4. Do NOT call eloqua_get_email_performance for campaign-based sends — it only covers quick sends, not campaign emails.
5. Use campaign_name_filter if the user specifies a region or team (e.g. 'US_ENT', 'APAC', 'EU', 'MM').
6. After the board returns, use the CSV to build any table, chart, or summary the user asked for.

SINGLE-CAMPAIGN DRILL-DOWN (one specific campaign): Use eloqua_get_campaign_email_report with a tight date range + campaign_name_filter.

OPPORTUNITY OBSERVATORY (pre-computed deal dossiers — 3 dedicated tools):
The Observatory holds pre-written, long-form intelligence dossiers for a curated set of opportunities, stored in a single Supabase table. Each dossier has header fields (name, opportunity_owner, close_date, amount, stage, account_name) plus 8 markdown sections: sf_90day_evidence, avoma_evidence, outbound_campaign_intelligence, bundle_a_deal_progress, bundle_b_competition_fit, bundle_c_stakeholder_map, bundle_d_vulnerabilities, diagnosis_sheet.
WHEN TO USE: Whenever the user asks about an opportunity's deal health, risks/vulnerabilities, competition, stakeholders, deal progress, or wants a briefing/dossier on a specific deal or account, CHECK THE OBSERVATORY FIRST. It is far cheaper and faster than re-deriving the same picture from live Salesforce + Avoma tool calls. Only fall back to live soql/Avoma tools if the opportunity is NOT in the Observatory or the user explicitly wants fresh real-time data.
THE 3 TOOLS (read-only, scoped to this one table — there is no write path):
1. list_opportunity_dossiers(limit, stage, account_name_contains, name_contains) — Lightweight discovery. Returns ONLY header rows (no heavy markdown). Use this to see what's available or to filter by stage/account/name. Start here when you don't already know the opportunity_id.
2. get_opportunity_dossier(opportunity_id, sections=None) — The full dossier for ONE opportunity. `sections` is an OPTIONAL comma-separated subset of the 8 section names above — pass it to pull only what you need (e.g. sections="bundle_d_vulnerabilities,diagnosis_sheet" for a risk question) and save tokens. Omit it to get the whole dossier. Unknown section names are rejected.
3. search_opportunity_dossiers(query, limit) — Fuzzy substring search over opportunity name + account name. Use when the user gives a company/deal name but not an ID. Returns header rows; follow up with get_opportunity_dossier using the returned opportunity_id.
TYPICAL FLOW: search_opportunity_dossiers OR list_opportunity_dossiers (find the opportunity_id) → get_opportunity_dossier(opportunity_id, sections=...) (pull the relevant sections). Do NOT use the generic supabase_query MCP tool for this data — these 3 tools are the correct, scoped path.

OUTPUT AND FILE RULES (MANDATORY):
1. ALWAYS deliver your final output directly in the chat response. NEVER write output to a file.
2. Present findings in clear, structured text within the conversation.
3. If the output is very long, break it into sections but still deliver it in chat — not as a file.
4. NEVER use write_file or edit_file tools. You do NOT have permission to create or modify files. You may use read_file ONLY to read existing files when needed.
5. NEVER delegate work to subagents or the task tool. Do ALL work yourself directly.

EFFICIENCY RULES (MANDATORY):
1. Plan your research before starting. List the 5-10 specific questions you need answered, then execute ONE search per question.
2. After completing research, STOP searching and synthesize what you have. Do not keep searching for "one more thing."
3. If a tool gives you partial information, USE IT. Do not discard partial results and re-search hoping for better data.
4. When multiple tools can answer the same question, pick the BEST one and use it. Do not query 3 different tools for the same data point.
5. Track what you've already learned. Before each tool call, review your existing findings to avoid redundant queries.

"""

        selected_model = model or config.MODEL
        print(f"Creating agent with model: {selected_model}")
        print(f"Number of tools: {len(tools)}")
        print(f"Tool names: {[t.name for t in tools]}")

        def make_error_safe(tool_obj):
            original_func = tool_obj.coroutine if hasattr(
                tool_obj,
                'coroutine') and tool_obj.coroutine else tool_obj.func
            if original_func is None:
                return tool_obj

            tool_name = tool_obj.name

            LEMLIST_TOOLS = {
                "lemlist_add_lead_to_campaign", "lemlist_add_leads_batch",
                "lemlist_create_campaign", "lemlist_update_campaign",
                "lemlist_pause_campaign", "lemlist_start_campaign",
                "lemlist_send_email", "lemlist_update_lead",
                "lemlist_create_webhook", "lemlist_delete_webhook",
                "lemlist_add_unsubscribe", "lemlist_remove_unsubscribe",
                "lemlist_unsubscribe_lead_from_campaign",
                "lemlist_mark_lead_interested", "lemlist_mark_lead_not_interested",
                "lemlist_delete_lead_from_campaign", "lemlist_complete_task",
                "lemlist_pause_lead_in_campaign", "lemlist_resume_lead_in_campaign",
            }
            is_lemlist_write = tool_name in LEMLIST_TOOLS

            def _extract_real_error(e):
                if isinstance(e, BaseExceptionGroup):
                    msgs = []
                    for sub in e.exceptions:
                        msgs.append(_extract_real_error(sub))
                    return "; ".join(msgs)
                return str(e)

            is_lemlist_any = tool_name.startswith("lemlist_")

            LEMLIST_RETRY_MAX = 2
            LEMLIST_RETRY_DELAY = 3

            async def _invoke_tool(*args, **kwargs):
                if asyncio.iscoroutinefunction(original_func):
                    result = await original_func(*args, **kwargs)
                else:
                    result = original_func(*args, **kwargs)
                if asyncio.iscoroutine(result):
                    result = await result
                return result

            async def safe_func(*args, **kwargs):
                chat_id = _current_chat_id.get(None)

                if chat_id:
                    args_summary = json.dumps(
                        kwargs, default=str)[:500] if kwargs else "{}"
                    await save_to_supabase(
                        chat_id, "tool_call", f"Calling {tool_name}...", {
                            "tool": tool_name,
                            "args": kwargs,
                            "source": "tool_wrapper",
                        })

                try:
                    if is_lemlist_any:
                        async with _lemlist_semaphore:
                            print(f"  [QUEUE] {tool_name} acquired semaphore")
                            last_err = None
                            for attempt in range(LEMLIST_RETRY_MAX + 1):
                                try:
                                    result = await _invoke_tool(*args, **kwargs)
                                    if attempt > 0:
                                        print(f"  [RETRY] {tool_name} succeeded on attempt {attempt + 1}")
                                    break
                                except (BaseExceptionGroup, ExceptionGroup) as eg:
                                    real_error = _extract_real_error(eg)
                                    last_err = eg
                                    print(f"  [RETRY] {tool_name} TaskGroup error on attempt {attempt + 1}/{LEMLIST_RETRY_MAX + 1}: {real_error}")
                                    if attempt < LEMLIST_RETRY_MAX:
                                        await asyncio.sleep(LEMLIST_RETRY_DELAY * (attempt + 1))
                                    else:
                                        raise
                                except Exception as e:
                                    if "TaskGroup" in str(e) and attempt < LEMLIST_RETRY_MAX:
                                        print(f"  [RETRY] {tool_name} error on attempt {attempt + 1}: {str(e)[:200]}")
                                        await asyncio.sleep(LEMLIST_RETRY_DELAY * (attempt + 1))
                                        last_err = e
                                    else:
                                        raise
                            else:
                                raise last_err
                            await asyncio.sleep(0.5)
                    else:
                        result = await _invoke_tool(*args, **kwargs)

                    if chat_id:
                        result_str = json.dumps(
                            result, default=str) if isinstance(
                                result, (dict, list)) else str(result)
                        display_result = result_str[:1000] + (
                            "..." if len(result_str) > 1000 else "")
                        await save_to_supabase(
                            chat_id, "tool_result", display_result, {
                                "tool": tool_name,
                                "full_length": len(result_str),
                                "source": "tool_wrapper"
                            })

                    return result
                except (BaseExceptionGroup, ExceptionGroup) as eg:
                    error_msg = _extract_real_error(eg)
                    print(f"  [TOOL ERROR] {tool_name} (TaskGroup unwrapped): {error_msg}")
                    import traceback
                    traceback.print_exc()

                    if chat_id:
                        await save_to_supabase(chat_id, "tool_error",
                                               error_msg[:500], {
                                                   "tool": tool_name,
                                                   "source": "tool_wrapper",
                                                   "error_type": "TaskGroup"
                                               })

                    return f"[Tool '{tool_name}' MCP transport error (retries exhausted): {error_msg}. The agent should retry or proceed with available information.]"
                except Exception as e:
                    error_msg = str(e)
                    print(f"  [TOOL ERROR] {tool_name}: {error_msg}")

                    if chat_id:
                        await save_to_supabase(chat_id, "tool_error",
                                               error_msg[:500], {
                                                   "tool": tool_name,
                                                   "source": "tool_wrapper"
                                               })

                    if "429" in error_msg or "Too Many Requests" in error_msg or "RESOURCE_EXHAUSTED" in error_msg or "rate" in error_msg.lower(
                    ):
                        return f"[Tool '{tool_name}' rate limited. Try again shortly or proceed without this result. Error: {error_msg}]"
                    if "404" in error_msg or "NOT_FOUND" in error_msg:
                        return f"[Tool '{tool_name}' resource not found. Error: {error_msg}]"
                    if "403" in error_msg or "FORBIDDEN" in error_msg or "Forbidden" in error_msg:
                        return f"[Tool '{tool_name}' access denied (no entitlement). Error: {error_msg}]"
                    if "timeout" in error_msg.lower(
                    ) or "timed out" in error_msg.lower():
                        return f"[Tool '{tool_name}' timed out. Try again or proceed without this result. Error: {error_msg}]"
                    return f"[Tool '{tool_name}' encountered an error: {error_msg}. The agent should continue with available information.]"

            from langchain_core.tools import StructuredTool
            return StructuredTool(name=tool_obj.name,
                                  description=tool_obj.description,
                                  coroutine=safe_func,
                                  args_schema=tool_obj.args_schema)

        tools = [make_error_safe(t) for t in tools]
        print(f"All {len(tools)} tools wrapped with error-safe handling")

        agent_model = selected_model
        if selected_model.startswith("anthropic:"):
            model_name = selected_model.split(":", 1)[1]
            from anthropic_cache import CachedChatAnthropic
            agent_model = CachedChatAnthropic(
                model_name=model_name,
                api_key=config.ANTHROPIC_API_KEY or None,
                max_retries=config.ANTHROPIC_MAX_RETRIES,
                timeout=config.ANTHROPIC_TIMEOUT_SECONDS,
                max_tokens=config.MAX_OUTPUT_TOKENS,
                stop=None,
            )
            print(
                f"Created Anthropic model: {model_name} "
                f"(max_retries={config.ANTHROPIC_MAX_RETRIES}, "
                f"timeout={config.ANTHROPIC_TIMEOUT_SECONDS}s, "
                f"max_tokens={config.MAX_OUTPUT_TOKENS} — Retry-After honoured on 429, "
                f"prompt caching ENABLED on system + tools)"
            )
        elif selected_model.startswith(
                "google_genai:") or selected_model.startswith("google:"):
            model_name = selected_model.split(":", 1)[1]
            from langchain_google_genai import ChatGoogleGenerativeAI
            agent_model = ChatGoogleGenerativeAI(
                model=model_name,
                google_api_key=config.GOOGLE_API_KEY,
            )
            print(f"Created Google Gemini model: {model_name}")
        elif selected_model.startswith("grok:") or selected_model.startswith("xai:"):
            model_name = selected_model.split(":", 1)[1]
            from langchain_openai import ChatOpenAI
            agent_model = ChatOpenAI(
                model=model_name,
                api_key=config.XAI_API_KEY,
                base_url="https://api.x.ai/v1",
            )
            print(f"Created xAI Grok model: {model_name}")
        elif selected_model.startswith("fireworks:"):
            # Fireworks AI is OpenAI-compatible — route through ChatOpenAI with
            # the Fireworks base URL. model_name is the full Fireworks model path,
            # e.g. "accounts/fireworks/models/gpt-oss-120b".
            model_name = selected_model.split(":", 1)[1]
            from langchain_openai import ChatOpenAI
            agent_model = ChatOpenAI(
                model=model_name,
                api_key=config.FIREWORKS_API_KEY or None,
                base_url="https://api.fireworks.ai/inference/v1",
                max_tokens=config.FIREWORKS_MAX_TOKENS,
            )
            print(f"Created Fireworks model: {model_name} (max_tokens={config.FIREWORKS_MAX_TOKENS})")

        # Intra-run context-trim middleware (cost task #40, 2026-05-22).
        # Shrinks old ToolMessage content before every LLM call once the
        # in-flight messages exceed CONTEXT_TRIM_THRESHOLD_TOKENS. Disabled
        # via CONTEXT_TRIM_ENABLED=false. See agent_checklist/context_trim_middleware.py.
        middlewares = []
        if config.CONTEXT_TRIM_ENABLED:
            from agent_checklist.context_trim_middleware import ContextTrimMiddleware
            middlewares.append(
                ContextTrimMiddleware(
                    threshold_tokens=config.CONTEXT_TRIM_THRESHOLD_TOKENS,
                    keep_recent_messages=config.CONTEXT_TRIM_KEEP_RECENT_MESSAGES,
                    placeholder_max_chars=config.CONTEXT_TRIM_PLACEHOLDER_MAX_CHARS,
                )
            )
            print(
                f"[CONTEXT-TRIM] enabled "
                f"(threshold={config.CONTEXT_TRIM_THRESHOLD_TOKENS:,} tokens, "
                f"keep_recent={config.CONTEXT_TRIM_KEEP_RECENT_MESSAGES}, "
                f"placeholder_cap={config.CONTEXT_TRIM_PLACEHOLDER_MAX_CHARS})"
            )
        # Fireworks' OpenAI-compat endpoint rejects LangChain content-block `id`
        # fields ("Extra inputs are not permitted"), which 400s every multi-step
        # tool-using run. Normalise message content right before each model call.
        if selected_model.startswith("fireworks:"):
            from agent_checklist.fireworks_compat_middleware import FireworksCompatMiddleware
            middlewares.append(FireworksCompatMiddleware())
            print("[FIREWORKS-COMPAT] message-normalisation middleware enabled (OpenAI-compat id strip)")
        self.agent = create_deep_agent(tools=tools,
                                       system_prompt=instructions,
                                       subagents=[],
                                       model=agent_model,
                                       middleware=middlewares,
                                       checkpointer=_AGENT_CHECKPOINTER,
                                       debug=False)
        print(
            f"Agent initialized with {len(tools)} tools (no subagents)"
        )
        return self.agent

    async def get_agent(self):
        if self.agent is None:
            async with self._init_lock:
                if self.agent is None:
                    await self.initialize_agent(skip_mcp=True)
        return self.agent

    def get_all_tools(self):
        """Flat catalog (custom tools + MCP tools) for reuse by sub-runners such
        as the per-cell AI-column agents. Reuses the SAME wrapped objects given
        to the main agent, so any MCP_TOOL_DENYLIST / allowlist already applied
        to ``_cached_mcp_tools`` is inherited (e.g. Salesforce write tools stay
        dropped). Returns whatever is loaded so far; safe to call any time."""
        tools = list(self._cached_custom_tools or [])
        tools.extend(self._cached_mcp_tools or [])
        return tools

    async def reinitialize_agent(self,
                                 instructions: Optional[str] = None,
                                 model: Optional[str] = None,
                                 headless: bool = True):
        async with self._init_lock:
            if self.mcp_client:
                self.mcp_client = None
            return await self.initialize_agent(instructions,
                                               model=model,
                                               headless=headless)


app = FastAPI(
    title="DeepAgent Server",
    description="AI Agentic Server with Context Window Management via ChatGPT",
    version="2.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    # allow_credentials must stay False while allow_origins is "*": the API now
    # authenticates via Bearer token (cross-origin clients) or a same-origin
    # cookie (the HTML UIs). Same-origin requests are not subject to CORS, so
    # credentialed cross-origin access is intentionally NOT enabled — the
    # invalid "*" + credentials=True combo would otherwise reflect any origin.
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

agent_manager = AgentManager()

# Let the per-cell AI-column agents (analysis_engine) reuse the SAME tool catalog
# the main chat agent has — custom tools + MCP tools, with any denylist already
# applied. The provider is a getter so it always reflects the live catalog
# (MCP servers load/reload in the background after startup).
try:
    import analysis_engine as _analysis_engine
    _analysis_engine.set_tool_provider(agent_manager.get_all_tools)
except Exception as _e:  # noqa: BLE001
    print(f"[ANALYSIS] could not register tool provider: {_e}")


class ChatMessage(BaseModel):
    role: str
    content: str
    images: Optional[List[str]] = None


class GoogleSheetConfig(BaseModel):
    spreadsheet_id: str
    sheet_name: Optional[str] = None


class ChatRequest(BaseModel):
    messages: List[ChatMessage]
    stream: bool = True
    system_prompt: Optional[str] = None
    model: Optional[str] = None
    headless: bool = True
    google_sheets: Optional[List[GoogleSheetConfig]] = None
    chat_id: Optional[str] = None
    project_id: Optional[str] = None


class StructuredChatRequest(BaseModel):
    messages: List[ChatMessage]
    structured_output_format: Dict[str, Any]
    system_prompt: Optional[str] = None
    model: Optional[str] = None
    headless: bool = True
    google_sheets: Optional[List[GoogleSheetConfig]] = None
    project_id: Optional[str] = None


class ConfigRequest(BaseModel):
    instructions: Optional[str] = None
    headless: bool = True


class JarvisSettingsRequest(BaseModel):
    # Both optional so the settings tab can save the analysis toggles and the
    # system prompt independently. Omit a field to leave it unchanged; send
    # system_prompt="" to reset it to the backend default.
    enabled_analysis_ids: Optional[List[str]] = None
    system_prompt: Optional[str] = None


class MCPServerConfig(BaseModel):
    command: str
    args: List[str]
    env: Optional[Dict[str, str]] = None
    enabled: bool = True


def _build_message_content(msg: ChatMessage) -> dict:
    if msg.images and len(msg.images) > 0:
        content_blocks = [{"type": "text", "text": msg.content}]
        for image_url in msg.images:
            content_blocks.append({
                "type": "image_url",
                "image_url": {"url": image_url}
            })
        return {"role": msg.role, "content": content_blocks}
    return {"role": msg.role, "content": msg.content}


_supabase_seq_counters = {}
_running_tasks: Dict[str, asyncio.Task] = {}
_cancelled_chats: set = set()
# chat_ids reserved in the window between the duplicate-run guard check and the
# moment the asyncio task is actually registered in `_running_tasks`. The task
# registration is often DEFERRED (it happens lazily inside a StreamingResponse
# generator, after the endpoint has already returned), so a plain
# "is there a live task?" check is not enough to stop a near-simultaneous
# double-submit. This set bridges that gap.
_starting_chats: set = set()


def _reserve_run_slot(chat_id: str) -> None:
    """Reject a duplicate concurrent run for the same chat and reserve the slot.

    Exactly ONE agent run may be in flight per chat_id at a time. Two runs on one
    chat corrupt the session: every tool fires twice, the dedupe/cleanup state
    races, and the shared agent gets reinitialised mid-run — which is what leaves
    the UI stuck on "Thinking…". The frontend should not double-submit, but this
    is the backend backstop that holds regardless of client behaviour.

    Single-threaded asyncio: this function performs NO `await`, so the
    check-and-reserve is atomic against other coroutines. Callers MUST call
    `_release_run_slot(chat_id)` once the task is registered in `_running_tasks`,
    or if setup fails before the task is created.

    Raises HTTPException(409) if a run is already active or starting.
    """
    if not chat_id:
        return
    existing = _running_tasks.get(chat_id)
    if (existing and not existing.done()) or chat_id in _starting_chats:
        raise HTTPException(
            status_code=409,
            detail=("A run is already in progress for this chat. Wait for it to "
                    "finish, or stop it, before sending another message."),
        )
    _starting_chats.add(chat_id)


def _release_run_slot(chat_id: str) -> None:
    """Release the start-reservation taken by `_reserve_run_slot`. Idempotent —
    safe to call multiple times and when no reservation is held."""
    if chat_id:
        _starting_chats.discard(chat_id)


# Captured server event loop, used to route cross-thread cancellations.
# LangChain sync tools run inside worker threads (run_in_executor), so a
# direct `task.cancel()` from there violates the asyncio thread-safety
# contract. We stash the loop here at FastAPI startup and use
# `loop.call_soon_threadsafe(task.cancel)` from non-loop threads. See the
# 2026-05-22 architect review for chat 8359d7a6 RAG-loop fixes.
_server_event_loop: Optional[asyncio.AbstractEventLoop] = None


def _register_server_event_loop() -> None:
    """Capture the running event loop. Called once at FastAPI startup so
    later `cancel_running_chat()` calls from worker threads can schedule
    cancellation onto the correct loop."""
    global _server_event_loop
    try:
        _server_event_loop = asyncio.get_running_loop()
    except RuntimeError:
        # Called outside a running loop — defer; later calls will fall
        # back to the loop the task itself is bound to.
        _server_event_loop = None


def cancel_running_chat(chat_id: str) -> bool:
    """Programmatic equivalent of POST /api/chat/stop, callable from other
    modules (e.g. pipeline_runner on phase timeout, search_knowledge on
    cap-escalation).

    Marks the chat as cancelled and, if an asyncio task is registered for it
    under `_running_tasks`, schedules `.cancel()` on it so the agent loop
    unwinds instead of continuing to call tools after the orchestrator has
    given up on it. Returns True iff a live task was actually cancelled
    (or successfully scheduled for cancellation from another thread).

    Thread-safe: if invoked from a non-loop thread (e.g. a LangChain sync
    tool running in `run_in_executor`), the cancel is routed through
    `loop.call_soon_threadsafe` so the asyncio contract is respected.

    Safe to call when no task is registered (returns False). Does not block,
    does not raise — caller may follow up with `await asyncio.sleep(...)` to
    give the cancellation a chance to propagate before continuing.
    """
    if not chat_id:
        return False
    _cancelled_chats.add(chat_id)
    task = _running_tasks.get(chat_id)
    if not (task and not task.done()):
        return False
    # Pick the right loop. Prefer the loop the task is bound to (most
    # accurate), fall back to the server loop captured at startup.
    target_loop: Optional[asyncio.AbstractEventLoop] = None
    try:
        target_loop = task.get_loop()
    except Exception:
        target_loop = _server_event_loop
    # Are we currently running inside that loop? If so, plain .cancel()
    # is fine and avoids an unnecessary callback hop.
    try:
        current_loop = asyncio.get_running_loop()
    except RuntimeError:
        current_loop = None
    try:
        if target_loop is None or current_loop is target_loop:
            task.cancel()
        else:
            target_loop.call_soon_threadsafe(task.cancel)
        print(f"[STOP] cancel_running_chat({chat_id}) — task cancelled "
              f"(cross_thread={current_loop is not target_loop})")
        return True
    except Exception as exc:  # noqa: BLE001
        print(f"[STOP] cancel_running_chat({chat_id}) — cancel failed: {exc}")
        return False
_current_chat_id = contextvars.ContextVar('current_chat_id', default=None)
_current_project_id = contextvars.ContextVar('current_project_id', default=None)
# When True, _wrap_mcp_tool skips the gpt-4o-mini LLM summariser for oversized
# tool outputs and goes straight to deterministic truncation. Set by the
# opportunity-analysis agent (opportunity_analyzer.py): the LLM summariser was
# timing out at 45s per large Salesforce get_record and grinding the analyzer
# into its 600s timeout. Truncation is what the summariser fell back to anyway,
# so this gives the same data far faster. Scoped to the analyzer only — the main
# chat agent still summarises.
_skip_llm_summarizer = contextvars.ContextVar('skip_llm_summarizer', default=False)
_lemlist_semaphore = asyncio.Semaphore(1)

# ---------------------------------------------------------------------------
# LLM Credit / Token Tracking
# Pricing: USD per million tokens (input, output)
# ---------------------------------------------------------------------------
_LLM_PRICING: Dict[str, tuple] = {
    # Anthropic — prices updated April 2026 (per million tokens)
    "claude-opus-4-7":     (5.0,   25.0),   # Opus 4.7: $5/$25 (released Apr 16 2026)
    "claude-opus-4":       (5.0,   25.0),   # Opus 4.6: $5/$25
    "claude-3-opus":       (15.0,  75.0),   # Claude 3 Opus (legacy)
    "claude-sonnet-4":     (3.0,   15.0),   # Sonnet 4.6: $3/$15
    "claude-3-5-sonnet":   (3.0,   15.0),
    "claude-3-sonnet":     (3.0,   15.0),
    "claude-haiku-4":      (1.0,   5.0),    # Haiku 4.5: $1/$5
    "claude-3-5-haiku":    (0.80,  4.0),
    "claude-3-haiku":      (0.25,  1.25),
    # OpenAI — most specific keys first to avoid substring collisions
    "gpt-5.4-pro":         (30.0,  180.0),  # GPT-5.4 Pro
    "gpt-5.4-nano":        (0.20,  1.25),   # GPT-5.4 Nano
    "gpt-5.4-mini":        (0.75,  4.50),   # GPT-5.4 Mini
    "gpt-5.4":             (2.50,  15.0),   # GPT-5.4
    "gpt-5.2":             (1.75,  14.0),   # GPT-5.2
    "gpt-5-mini":          (0.15,  0.60),   # GPT-5 Mini (must precede gpt-5 bare)
    "gpt-5":               (2.50,  15.0),   # GPT-5 (bare, last to avoid matching gpt-5.x)
    "gpt-4o-mini":         (0.15,  0.60),
    "gpt-4o":              (2.50,  10.0),
    "gpt-4-turbo":         (10.0,  30.0),
    "gpt-4":               (30.0,  60.0),
    "o1-mini":             (3.0,   12.0),
    "o1":                  (15.0,  60.0),
    "o3-mini":             (1.10,  4.40),
    "o3":                  (10.0,  40.0),
    # Google Gemini — most specific first
    "gemini-2.5-flash":    (0.15,  1.0),    # Gemini 2.5 Flash
    "gemini-2.5-pro":      (1.25,  10.0),   # Gemini 2.5 Pro
    "gemini-2.0-flash":    (0.10,  0.40),   # Gemini 2.0 Flash
    "gemini-1.5-flash":    (0.075, 0.30),
    "gemini-1.5-pro":      (3.50,  10.50),
    # xAI Grok — most specific first to avoid grok-4 matching grok-4.20
    "grok-4.20":           (2.0,   6.0),    # Grok 4.20
    "grok-4-1-fast":       (0.20,  0.50),   # Grok 4.1 Fast
    "grok-3-mini":         (0.30,  0.50),
    "grok-3":              (3.0,   15.0),
    "grok-4":              (2.0,   6.0),    # Grok 4 (base)
}

_chat_token_accumulators: Dict[str, dict] = {}
_chat_seen_ai_msg_ids: Dict[str, set] = {}


def _calculate_llm_cost(
    model: str,
    input_tokens: int,
    output_tokens: int,
    cache_creation_tokens: int = 0,
    cache_read_tokens: int = 0,
) -> float:
    """Compute USD cost.

    Anthropic prompt-caching pricing:
      * uncached input  -> 1.00 x input rate
      * cache CREATION  -> 1.25 x input rate (5-min ephemeral cache)
      * cache READ      -> 0.10 x input rate
      * output          -> 1.00 x output rate
    """
    model_lower = (model or "").lower()
    for key, (inp_cost, out_cost) in _LLM_PRICING.items():
        if key in model_lower:
            return (
                input_tokens * inp_cost
                + cache_creation_tokens * inp_cost * 1.25
                + cache_read_tokens * inp_cost * 0.10
                + output_tokens * out_cost
            ) / 1_000_000
    return 0.0


def _accumulate_tokens_from_msg(chat_id: str, model: str, ai_message) -> None:
    usage = getattr(ai_message, "usage_metadata", None)
    if not usage:
        return
    msg_id = getattr(ai_message, "id", None) or id(ai_message)
    if chat_id not in _chat_seen_ai_msg_ids:
        _chat_seen_ai_msg_ids[chat_id] = set()
    if msg_id in _chat_seen_ai_msg_ids[chat_id]:
        return
    _chat_seen_ai_msg_ids[chat_id].add(msg_id)
    if chat_id not in _chat_token_accumulators:
        _chat_token_accumulators[chat_id] = {
            "model": model or "unknown",
            "input_tokens": 0,
            "output_tokens": 0,
            "total_tokens": 0,
            "cache_creation_tokens": 0,
            "cache_read_tokens": 0,
            "per_model": {},
        }
    acc = _chat_token_accumulators[chat_id]
    model_key = model or acc["model"]
    acc["model"] = model_key
    inp = usage.get("input_tokens", 0) or 0
    out = usage.get("output_tokens", 0) or 0
    details = usage.get("input_token_details") or {}
    cc = (details.get("cache_creation") or 0) if isinstance(details, dict) else 0
    cr = (details.get("cache_read") or 0) if isinstance(details, dict) else 0
    # Anthropic returns `input_tokens` as the NEW uncached input only; the cache
    # creation/read tokens are reported separately. We store the *grand total*
    # input for visibility so historical SUM(input_tokens) comparisons stay
    # meaningful, and keep the cache split for accurate cost computation.
    total_input = inp + cc + cr
    tot = usage.get("total_tokens", 0) or (total_input + out)
    acc["input_tokens"] += total_input
    acc["output_tokens"] += out
    acc["total_tokens"] += tot
    acc["cache_creation_tokens"] += cc
    acc["cache_read_tokens"] += cr
    per = acc.setdefault("per_model", {})
    if model_key not in per:
        per[model_key] = {"input": 0, "output": 0, "cache_creation": 0, "cache_read": 0}
    per[model_key]["input"] += inp
    per[model_key]["output"] += out
    per[model_key]["cache_creation"] += cc
    per[model_key]["cache_read"] += cr


async def _save_chat_usage(chat_id: str) -> None:
    if not supabase or chat_id not in _chat_token_accumulators:
        return
    acc = _chat_token_accumulators.get(chat_id, {})
    if not acc or not acc.get("total_tokens"):
        return
    per_model = acc.get("per_model", {})
    if per_model:
        cost = sum(
            _calculate_llm_cost(
                m,
                v.get("input", 0),
                v.get("output", 0),
                v.get("cache_creation", 0),
                v.get("cache_read", 0),
            )
            for m, v in per_model.items()
        )
        model_display = ", ".join(sorted(per_model.keys()))
    else:
        cost = _calculate_llm_cost(
            acc["model"],
            acc["input_tokens"] - acc.get("cache_creation_tokens", 0) - acc.get("cache_read_tokens", 0),
            acc["output_tokens"],
            acc.get("cache_creation_tokens", 0),
            acc.get("cache_read_tokens", 0),
        )
        model_display = acc["model"]
    loop = asyncio.get_event_loop()
    try:
        await loop.run_in_executor(
            None,
            lambda: supabase.table("chat_usage").upsert(
                {
                    "chat_id": chat_id,
                    "model": model_display,
                    "input_tokens": acc["input_tokens"],
                    "output_tokens": acc["output_tokens"],
                    "total_tokens": acc["total_tokens"],
                    "cost_usd": round(cost, 6),
                },
                on_conflict="chat_id",
            ).execute(),
        )
        cc_tok = acc.get("cache_creation_tokens", 0)
        cr_tok = acc.get("cache_read_tokens", 0)
        cache_hit_pct = (100.0 * cr_tok / acc["input_tokens"]) if acc["input_tokens"] else 0.0
        print(
            f"[USAGE] chat_id={chat_id} models={model_display} "
            f"in={acc['input_tokens']} (cache_write={cc_tok} cache_read={cr_tok} "
            f"hit={cache_hit_pct:.1f}%) out={acc['output_tokens']} cost=${cost:.6f}"
        )
    except Exception as e:
        print(f"[USAGE ERROR] Failed to save usage for {chat_id}: {e}")

def _set_rag_context(chat_id: str, project_id: str = None):
    """Set RAG context vars so search_knowledge tool can scope its queries."""
    try:
        import rag_context
        rag_context.current_project_id.set(project_id)
        rag_context.current_chat_id.set(chat_id)
    except Exception:
        pass

_session_start_times: Dict[str, float] = {}
_sessions_rejected = 0
_server_start_time = None


async def _cleanup_expired_sessions():
    """Periodically clean up abandoned/timed-out sessions"""
    while True:
        try:
            await asyncio.sleep(300)  # Check every 5 minutes
            now = asyncio.get_event_loop().time()
            timeout_seconds = config.SESSION_TIMEOUT_MINUTES * 60
            expired = []
            for cid, start_time in list(_session_start_times.items()):
                if now - start_time > timeout_seconds:
                    expired.append(cid)
            for cid in expired:
                task = _running_tasks.get(cid)
                if task and not task.done():
                    print(f"[CLEANUP] Cancelling expired session {cid} (ran >{config.SESSION_TIMEOUT_MINUTES} min)")
                    _cancelled_chats.add(cid)
                    task.cancel()
                _running_tasks.pop(cid, None)
                _session_start_times.pop(cid, None)
                _supabase_seq_counters.pop(cid, None)
                _dedupe_completed.pop(cid, None)
                _dedupe_inflight.pop(cid, None)
                _approved_campaigns_cache.pop(cid, None)
            if expired:
                print(f"[CLEANUP] Cleaned up {len(expired)} expired sessions. Active: {len(_running_tasks)}")
        except Exception as e:
            print(f"[CLEANUP ERROR] {e}")


_chats_created: set = set()

async def ensure_chat_row(chat_id: str):
    if not supabase or not chat_id or chat_id in _chats_created:
        return
    try:
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, lambda: supabase.table("chats").insert({"id": chat_id}).execute())
        _chats_created.add(chat_id)
    except Exception:
        _chats_created.add(chat_id)

def _query_max_sequence(chat_id: str) -> int:
    """Return the current max `chat_messages.sequence` for `chat_id`, or 0.

    Used to seed `_supabase_seq_counters[chat_id]` on lazy-init so that
    second-turn writes (remediation re-run, post-restart, post-session-pop)
    continue past the existing rows instead of restarting at 0 — which
    the UI realtime feed drops as "older than current_max".
    """
    if not supabase:
        return 0
    try:
        res = (
            supabase.table("chat_messages")
            .select("sequence")
            .eq("chat_id", chat_id)
            .order("sequence", desc=True)
            .limit(1)
            .execute()
        )
        rows = res.data or []
        if rows and rows[0].get("sequence") is not None:
            return int(rows[0]["sequence"])
        return 0
    except Exception as e:
        print(f"[SUPABASE] _query_max_sequence failed for {chat_id} (defaulting to 0): {e}")
        return 0


# Bounded-retry insert: every chat_messages write goes through here so a
# transient Supabase blip (network hiccup, brief 5xx) doesn't silently drop a
# row the UI is waiting on. Failures that survive all attempts are logged
# loudly with chat_id + row type so dropped writes are visible in the console.
_INSERT_RETRY_ATTEMPTS = 3
_INSERT_RETRY_BASE_DELAY = 0.25  # seconds; exponential backoff: 0.25, 0.5, 1.0


async def _insert_chat_message_with_retry(row: dict, chat_id: str,
                                          msg_type: str):
    """Await a chat_messages insert with a short bounded retry. Returns the
    Supabase result on success, or None if every attempt failed (logged)."""
    loop = asyncio.get_event_loop()
    last_exc = None
    for attempt in range(1, _INSERT_RETRY_ATTEMPTS + 1):
        try:
            result = await loop.run_in_executor(
                None,
                lambda: supabase.table("chat_messages").insert(row).execute())
            # postgrest raises on hard errors; this guards the rare case of a
            # success-shaped response that carries no inserted row.
            if getattr(result, "data", None) is None:
                raise RuntimeError(
                    f"insert returned no data (error="
                    f"{getattr(result, 'error', None)})")
            return result
        except Exception as e:
            last_exc = e
            if attempt < _INSERT_RETRY_ATTEMPTS:
                delay = _INSERT_RETRY_BASE_DELAY * (2 ** (attempt - 1))
                print(f"[SUPABASE RETRY] insert type={msg_type} "
                      f"chat_id={chat_id} attempt {attempt}/"
                      f"{_INSERT_RETRY_ATTEMPTS} failed: {e} — retrying in "
                      f"{delay:.2f}s")
                await asyncio.sleep(delay)
            else:
                print(f"[SUPABASE DROP] insert type={msg_type} "
                      f"chat_id={chat_id} FAILED after "
                      f"{_INSERT_RETRY_ATTEMPTS} attempts: {e}")
    return None


# The parent `chats.updated_at` bump gives the client a cheap freshness signal
# that new activity happened even if a Realtime event was dropped. Best-effort:
# if the column is absent in this Supabase schema, we disable further bumps
# after one quiet log rather than spamming the console on every write.
_chats_updated_at_ok = True
_bump_tasks: set = set()


async def _bump_chat_updated_at(chat_id: str):
    global _chats_updated_at_ok
    if not supabase or not chat_id or not _chats_updated_at_ok:
        return
    try:
        from datetime import timezone as _tz
        now = datetime.now(_tz.utc).isoformat()
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(
            None,
            lambda: supabase.table("chats").update({"updated_at": now}).eq(
                "id", chat_id).execute())
    except Exception as e:
        # Disable on schema mismatch (column missing); otherwise just note it.
        msg = str(e).lower()
        if "updated_at" in msg or "column" in msg or "schema" in msg:
            _chats_updated_at_ok = False
            print(f"[SUPABASE] chats.updated_at bump disabled "
                  f"(column unavailable): {e}")
        else:
            print(f"[SUPABASE] chats.updated_at bump failed "
                  f"chat_id={chat_id} (non-fatal): {e}")


def _schedule_chat_bump(chat_id: str):
    """Fire-and-forget freshness bump so it never adds latency to the
    streaming/save path. Keeps a reference so the task isn't GC'd early."""
    if not _chats_updated_at_ok:
        return
    try:
        task = asyncio.create_task(_bump_chat_updated_at(chat_id))
        _bump_tasks.add(task)
        task.add_done_callback(_bump_tasks.discard)
    except RuntimeError:
        pass  # no running loop (best-effort)


async def save_to_supabase(chat_id: str,
                           msg_type: str,
                           content: str,
                           metadata: dict = None):
    if not supabase or not chat_id:
        return
    await ensure_chat_row(chat_id)
    try:
        if chat_id not in _supabase_seq_counters:
            loop = asyncio.get_event_loop()
            seeded = await loop.run_in_executor(None, _query_max_sequence, chat_id)
            _supabase_seq_counters[chat_id] = seeded
        _supabase_seq_counters[chat_id] += 1
        seq = _supabase_seq_counters[chat_id]

        row = {
            "chat_id": chat_id,
            "role": "assistant",
            "type": msg_type,
            "content": content,
            "sequence": seq,
        }
        if metadata:
            row["metadata"] = json.dumps(metadata)
        result = await _insert_chat_message_with_retry(row, chat_id, msg_type)
        if result is not None:
            print(
                f"[SUPABASE] Saved {msg_type} seq={seq} (chat_id={chat_id}, len={len(content)})"
            )
            _schedule_chat_bump(chat_id)
        return result
    except Exception as e:
        print(f"[SUPABASE ERROR] Failed to save {msg_type}: {e}")
        import traceback
        print(f"[SUPABASE ERROR] {traceback.format_exc()}")


async def _verify_and_maybe_remediate(chat_id: str, project_id: str | None) -> None:
    """Run the verifier and persist its verdict. ADVISORY ONLY.

    The auto-remediation re-prompt loop was removed (user request, 2026-05-20)
    because it doubled token cost on long ABM runs and rarely fixed gaps that
    the agent had not already considered. The verifier still runs and writes a
    `verifier_report` row so the gap data is visible in the UI / `/api/chat/
    {chat_id}/verifier_report`, but it no longer re-invokes the agent.

    All errors are swallowed — this never blocks or breaks the agent loop.
    """
    try:
        from verifier import run_verifier_for_chat

        verdict = await run_verifier_for_chat(chat_id, project_id=project_id)
        if verdict is None or verdict.passed:
            return
        print(
            f"[VERIFIER] chat={chat_id} verdict=dirty "
            f"missing={len(verdict.missed_ids)} (advisory; no remediation re-prompt)"
        )
    except Exception as e:
        print(f"[VERIFIER] post-run check failed (non-fatal): {e}")


def _trigger_verifier(chat_id: str, project_id: str | None) -> None:
    """Fire-and-forget post-run verifier (advisory only, no remediation).

    Reads tool_call rows from chat_messages, evaluates the registered flow
    for the chat's project, and writes a verifier_report row. The auto
    re-prompt loop was removed (2026-05-20). See `verifier/` package.
    """
    try:
        asyncio.create_task(
            _verify_and_maybe_remediate(chat_id, project_id)
        )
    except Exception as e:
        print(f"[VERIFIER] trigger failed (non-fatal): {e}")


def _get_running_cost(chat_id: str) -> float:
    """Compute current USD cost for a chat from accumulated token totals.
    Mirrors the math in `_save_chat_usage` but doesn't touch Supabase.
    """
    acc = _chat_token_accumulators.get(chat_id)
    if not acc:
        return 0.0
    per = acc.get("per_model") or {}
    if per:
        return sum(
            _calculate_llm_cost(
                m,
                v.get("input", 0),
                v.get("output", 0),
                v.get("cache_creation", 0),
                v.get("cache_read", 0),
            )
            for m, v in per.items()
        )
    return _calculate_llm_cost(
        acc.get("model", "unknown"),
        max(0, acc.get("input_tokens", 0)
                - acc.get("cache_creation_tokens", 0)
                - acc.get("cache_read_tokens", 0)),
        acc.get("output_tokens", 0),
        acc.get("cache_creation_tokens", 0),
        acc.get("cache_read_tokens", 0),
    )


async def _evict_checkpoint(chat_id: str) -> None:
    """Drop the InMemorySaver state for a finished chat thread. Called from
    the `finally` block of every run_agent path so accumulating chat threads
    don't grow the process heap without bound.
    """
    try:
        await _AGENT_CHECKPOINTER.adelete_thread(chat_id)
    except Exception as e:
        print(f"[CHECKPOINT] evict failed for {chat_id}: {e}")


async def _agent_astream_autocontinue(agent, messages, chat_id, terminal_meta_out):
    """Async generator wrapping `agent.astream` with auto-continuation on
    Anthropic `stop_reason='max_tokens'` truncation.

    - First iteration sends `{"messages": messages}`.
    - Subsequent iterations send `None` to resume from the InMemorySaver
      checkpoint keyed by `thread_id=chat_id` (no full re-send of history).
    - Stops when the model finishes cleanly (any stop_reason != max_tokens)
      OR a circuit-breaker fires: continuation count, wall-clock, or cost.
    - Terminal classification is written into `terminal_meta_out` (mutated
      in place so the caller can read it after the generator finishes).
    """
    cfg = {
        "recursion_limit": _RECURSION_LIMIT,
        "configurable": {"thread_id": chat_id},
    }
    start_ts = time.time()
    auto_continue_count = 0
    warned_cost = False
    last_ai = None
    first = True

    try:
        from custom_tools.search_knowledge import reset_search_knowledge_step
    except Exception:
        reset_search_knowledge_step = None

    while True:
        # Per-step RAG budget: each auto-continue model step gets a fresh
        # search_knowledge cap. Without this, a legitimately RAG-heavy run that
        # auto-continues across many steps accumulates past the per-step cap and
        # trips the runaway-loop cancel, killing the whole run with no output.
        if reset_search_knowledge_step is not None:
            reset_search_knowledge_step(chat_id)
        stream_input = {"messages": messages} if first else None
        first = False
        # Watchdog-wrapped iteration: if no chunk arrives within
        # WATCHDOG_STALL_SECONDS we raise — the caller's try/except writes a
        # loud `error` row so the UI moves off "Thinking…". This is the only
        # reliable catch for a hung Anthropic call (the LLM SDK may be
        # silently retrying inside its own retry loop, with no chunks being
        # emitted on the stream).
        _stream = agent.astream(stream_input, stream_mode="values", config=cfg)
        while True:
            try:
                chunk = await asyncio.wait_for(
                    _stream.__anext__(),
                    timeout=config.WATCHDOG_STALL_SECONDS,
                )
            except StopAsyncIteration:
                break
            except asyncio.TimeoutError:
                elapsed = time.time() - start_ts
                # Close the abandoned async generator so the underlying
                # httpx connection / langgraph resources are released
                # immediately rather than waiting for GC.
                try:
                    await _stream.aclose()
                except Exception:
                    pass
                raise RuntimeError(
                    f"WATCHDOG: no agent chunk for "
                    f"{config.WATCHDOG_STALL_SECONDS}s — run stalled "
                    f"(likely hung LLM call or dropped stream). "
                    f"auto_continue={auto_continue_count}, "
                    f"elapsed={elapsed:.0f}s, "
                    f"cost=${_get_running_cost(chat_id):.2f}"
                )
            if isinstance(chunk, dict) and chunk.get("messages"):
                m = chunk["messages"][-1]
                if type(m).__name__ == "AIMessage":
                    last_ai = m
            yield chunk

        stop_reason = (
            (getattr(last_ai, "response_metadata", None) or {}).get("stop_reason")
            if last_ai else None
        )
        cost = _get_running_cost(chat_id)
        elapsed = time.time() - start_ts

        if not warned_cost and cost >= config.RUN_COST_WARN_USD:
            print(f"[AUTOCONT] ⚠️  chat={chat_id} cost=${cost:.2f} crossed warn "
                  f"threshold ${config.RUN_COST_WARN_USD}")
            warned_cost = True

        def _set(reason):
            terminal_meta_out["reason"] = reason
            terminal_meta_out["stop_reason"] = stop_reason
            terminal_meta_out["auto_continue_count"] = auto_continue_count
            terminal_meta_out["cost_usd"] = round(cost, 4)
            terminal_meta_out["elapsed_s"] = round(elapsed, 1)

        if stop_reason != "max_tokens":
            _set("model_finished")
            if auto_continue_count > 0:
                print(f"[AUTOCONT] ✅ chat={chat_id} finished after "
                      f"{auto_continue_count} auto-continuation(s) | "
                      f"stop_reason={stop_reason} | cost=${cost:.2f} | "
                      f"elapsed={elapsed:.0f}s")
            return
        if auto_continue_count >= config.MAX_AUTO_CONTINUATIONS:
            _set("budget_continuations")
            print(f"[AUTOCONT] 🛑 chat={chat_id} hit MAX_AUTO_CONTINUATIONS="
                  f"{config.MAX_AUTO_CONTINUATIONS} | cost=${cost:.2f} | "
                  f"elapsed={elapsed:.0f}s")
            return
        if cost >= config.MAX_RUN_COST_USD:
            _set("budget_cost")
            print(f"[AUTOCONT] 🛑 chat={chat_id} hit MAX_RUN_COST_USD="
                  f"${config.MAX_RUN_COST_USD} (actual ${cost:.2f})")
            return
        if elapsed >= config.MAX_RUN_SECONDS:
            _set("budget_time")
            print(f"[AUTOCONT] 🛑 chat={chat_id} hit MAX_RUN_SECONDS="
                  f"{config.MAX_RUN_SECONDS} (actual {elapsed:.0f}s)")
            return

        auto_continue_count += 1
        print(f"[AUTOCONT] 🔄 chat={chat_id} truncated (max_tokens); "
              f"auto-continue #{auto_continue_count} | cost=${cost:.2f} | "
              f"elapsed={elapsed:.0f}s")


_BUDGET_BREAKER_REASONS = ("budget_cost", "budget_time", "budget_continuations")


def _budget_breaker_reason(terminal_meta) -> str:
    """Return the auto-continue circuit-breaker reason (cost / wall-clock /
    continuation cap) if the run ended on one, else ''.

    Pure read — no DB write. Budget exhaustion is folded into the single
    terminal row by the caller (a loud note appended to the `final` answer,
    or a standalone `error` row when there is no model text) so a turn never
    emits two terminal rows. Defensive: never raises.
    """
    try:
        reason = (terminal_meta or {}).get("reason")
        return reason if reason in _BUDGET_BREAKER_REASONS else ""
    except Exception:
        return ""


def _budget_breaker_note(terminal_meta, reason: str) -> str:
    return (
        f"Run hit circuit-breaker '{reason}' after "
        f"{(terminal_meta or {}).get('auto_continue_count', 0)} "
        f"auto-continuation(s), "
        f"cost=${(terminal_meta or {}).get('cost_usd', 0):.2f}, "
        f"elapsed={(terminal_meta or {}).get('elapsed_s', 0):.0f}s."
    )


async def run_agent_and_save(chat_id: str,
                             messages: list,
                             agent,
                             request_model: str = None,
                             project_id: str = None):
    print(
        f"\n[AGENT TASK] Starting background agent run for chat_id={chat_id}, project_id={project_id}")
    _current_chat_id.set(chat_id)
    _current_project_id.set(project_id)
    _set_rag_context(chat_id, project_id)
    try:
        from custom_tools.search_knowledge import reset_search_knowledge_counter
        reset_search_knowledge_counter(chat_id)
        await save_to_supabase(chat_id, "status", "processing",
                               {"status": "started"})

        final_response = ""
        step_count = 0
        seen_tool_calls = set()
        seen_tool_results = set()
        thinking_logs = []
        pending_thinking = ""
        terminal_meta: dict = {}
        terminal_written = False

        async for chunk in _agent_astream_autocontinue(
                agent, messages, chat_id, terminal_meta):
            if chat_id in _cancelled_chats:
                print(f"[AGENT TASK] CANCELLED chat_id={chat_id}")
                await save_to_supabase(chat_id, "status",
                                       "Agent stopped by user.", {
                                           "status": "cancelled",
                                           "tool_calls": step_count
                                       })
                # The UI only treats type='final'/'error' as terminal, so a
                # 'status' row alone leaves the spinner spinning. Write a
                # terminal 'error' row so user-stop reliably ends the turn.
                _res = await save_to_supabase(chat_id, "error",
                                       "Agent stopped by user.", {
                                           "status": "cancelled",
                                           "kind": "user_stopped",
                                           "tool_calls": step_count
                                       })
                terminal_written = _res is not None
                _cancelled_chats.discard(chat_id)
                return
            if "messages" not in chunk or not chunk["messages"]:
                continue
            last_message = chunk["messages"][-1]
            msg_type = type(last_message).__name__

            if msg_type == "AIMessage":
                _accumulate_tokens_from_msg(chat_id, request_model, last_message)
                has_tool_calls = hasattr(
                    last_message, 'tool_calls') and last_message.tool_calls

                ai_content = ""
                if hasattr(last_message, 'content') and last_message.content:
                    raw = last_message.content
                    if isinstance(raw, list):
                        text_parts = [
                            block.get("text", "") for block in raw
                            if isinstance(block, dict)
                            and block.get("type") == "text"
                        ]
                        ai_content = "\n".join(text_parts).strip()
                    else:
                        ai_content = str(raw).strip()

                if ai_content and ai_content != final_response:
                    if has_tool_calls:
                        print(
                            f"[AGENT TASK] AI thinking ({len(ai_content)} chars): {ai_content[:200]}"
                        )
                        thinking_logs.append(ai_content)
                        await save_to_supabase(chat_id, "thinking", ai_content,
                                               {"source": "agent_reasoning"})
                    else:
                        if pending_thinking:
                            print(
                                f"[AGENT TASK] AI thinking ({len(pending_thinking)} chars): {pending_thinking[:200]}"
                            )
                            thinking_logs.append(pending_thinking)
                            await save_to_supabase(
                                chat_id, "thinking", pending_thinking,
                                {"source": "agent_reasoning"})
                        pending_thinking = ai_content
                    final_response = ai_content

                if has_tool_calls:
                    for tool_call in last_message.tool_calls:
                        tool_call_id = f"{tool_call.get('name', 'unknown')}_{tool_call.get('id', '')}"
                        if tool_call_id in seen_tool_calls:
                            continue
                        seen_tool_calls.add(tool_call_id)
                        step_count += 1
                        tool_name = tool_call.get('name', 'unknown')
                        tool_args = tool_call.get('args', {})
                        log_entry = f"Calling **{tool_name}** with args: `{json.dumps(tool_args, indent=2)}`"
                        thinking_logs.append(log_entry)
                        print(
                            f"[AGENT TASK] Step {step_count}: Calling {tool_name}"
                        )
                        await save_to_supabase(
                            chat_id, "tool_call", f"Calling {tool_name}...", {
                                "tool": tool_name,
                                "args": tool_args,
                                "step": step_count,
                                "tool_call_id": tool_call.get("id"),
                                "source": "sync_handler",
                            })

            elif msg_type == "ToolMessage":
                tool_result_id = getattr(last_message, 'tool_call_id',
                                         'unknown')
                if tool_result_id not in seen_tool_results:
                    seen_tool_results.add(tool_result_id)
                    tool_content = str(last_message.content)
                    tool_name = getattr(last_message, 'name', 'unknown')
                    display_content = tool_content[:500] + (
                        "..." if len(tool_content) > 500 else "")
                    log_entry = f"Result from **{tool_name}**: \n> {display_content}"
                    thinking_logs.append(log_entry)
                    print(
                        f"[AGENT TASK] Tool result from {tool_name} ({len(tool_content)} chars)"
                    )
                    await save_to_supabase(chat_id, "tool_result",
                                           tool_content, {"tool": tool_name})

        # Exactly one terminal row per turn. If the run ended on a budget
        # circuit-breaker, fold that into the single terminal row instead of
        # writing a separate `error` + `final` pair. terminal_written only
        # flips on a *confirmed* insert (save_to_supabase returns None when an
        # insert is dropped after retries) so the finally safety-net can still
        # fire if the terminal write is lost.
        _breaker = _budget_breaker_reason(terminal_meta)
        if _breaker:
            print(f"[AGENT TASK] 🛑 chat={chat_id} circuit-breaker '{_breaker}' "
                  f"cost=${terminal_meta.get('cost_usd', 0):.2f} "
                  f"elapsed={terminal_meta.get('elapsed_s', 0):.0f}s")
        if final_response:
            # The `final` row is the CLEAN answer only. The thinking + each
            # tool call/result are already persisted as their own `thinking`/
            # `tool_call`/`tool_result` rows, so the frontend renders those as
            # collapsible activity — we must NOT inline that trace here.
            full_content = final_response
            if _breaker:
                full_content += ("\n\n---\n⚠️ " +
                                 _budget_breaker_note(terminal_meta, _breaker))
            res = await save_to_supabase(chat_id, "final", full_content, {
                "tool_calls": step_count,
                "status": "budget_exhausted" if _breaker else "completed",
                "auto_continue_count": terminal_meta.get("auto_continue_count", 0),
                "terminal_reason": terminal_meta.get("reason"),
            })
            terminal_written = res is not None
            print(
                f"[AGENT TASK] COMPLETED chat_id={chat_id}, tool_calls={step_count}, "
                f"response_len={len(full_content)}, "
                f"auto_continue={terminal_meta.get('auto_continue_count', 0)}, "
                f"terminal={terminal_meta.get('reason')}"
            )
            _trigger_verifier(chat_id, project_id)
            if project_id in _lake.OD_PROJECT_IDS:
                asyncio.create_task(_lake.write_lake_diagnosis(
                    chat_id=chat_id,
                    project_id=project_id,
                    final_response=final_response,
                    supabase_client=supabase,
                    openai_api_key=config.OPENAI_API_KEY or os.getenv("OPENAI_API_KEY", ""),
                ))
        elif _breaker:
            # No model text — surface budget exhaustion as the single terminal.
            res = await save_to_supabase(
                chat_id, "error", _budget_breaker_note(terminal_meta, _breaker),
                {"status": "budget_exhausted", "terminal": terminal_meta})
            terminal_written = res is not None
            print(f"[AGENT TASK] 🛑 chat={chat_id} {_breaker} → error terminal")
        else:
            res = await save_to_supabase(chat_id, "final", "Task completed.", {
                "tool_calls": step_count,
                "status": "completed",
                "auto_continue_count": terminal_meta.get("auto_continue_count", 0),
                "terminal_reason": terminal_meta.get("reason"),
            })
            terminal_written = res is not None
            print(
                f"[AGENT TASK] COMPLETED (no text response) chat_id={chat_id}, tool_calls={step_count}"
            )
            _trigger_verifier(chat_id, project_id)

    except Exception as e:
        import traceback
        error_msg = str(e)
        print(f"[AGENT TASK] FAILED chat_id={chat_id}: {error_msg}")
        print(f"[AGENT TASK] {traceback.format_exc()}")
        res = await save_to_supabase(chat_id, "error", error_msg,
                                     {"status": "failed"})
        terminal_written = res is not None
    finally:
        # Terminal-row safety net: guarantee every turn ends with exactly one
        # final/error row. If no terminal row was written on any path above
        # (e.g. an unexpected early return or a swallowed watchdog/runtime
        # condition), write a loud error row so the UI never hangs on
        # "Thinking…". Guarded by terminal_written to avoid double-writes.
        if not terminal_written:
            try:
                await save_to_supabase(
                    chat_id, "error",
                    "Run ended without a terminal status (internal).",
                    {"status": "failed", "kind": "missing_terminal"})
                print(f"[AGENT TASK] safety-net terminal error row written "
                      f"chat_id={chat_id}")
            except Exception as _e:
                print(f"[AGENT TASK] safety-net terminal write failed "
                      f"chat_id={chat_id}: {_e}")
        await _save_chat_usage(chat_id)
        await _evict_checkpoint(chat_id)
        try:
            from custom_tools.search_knowledge import reset_search_knowledge_counter
            reset_search_knowledge_counter(chat_id)
        except Exception:
            pass


@app.post("/api/chat/async")
async def chat_async(request: ChatRequest):
    """Async chat endpoint - runs agent inside a streaming response to keep Autoscale alive,
    saves all results to Supabase in real-time. Client reads results via Supabase realtime.
    The HTTP stream sends keepalive pings + the chat_id, but the client doesn't need to read it."""
    chat_id = request.chat_id or str(uuid.uuid4())
    # Backend backstop: reject a second concurrent run for this chat (e.g. a
    # frontend double-submit) instead of spawning a duplicate run that corrupts
    # the session. Done before the try so the 409 propagates cleanly and the
    # error handlers below never release a slot owned by another request.
    _reserve_run_slot(chat_id)
    try:
        print(f"\n{'='*60}")
        print(f"[API REQUEST] /api/chat/async")
        print(f"[ASYNC] chat_id={chat_id}, model={request.model}")
        print(f"[ASYNC] messages count={len(request.messages)}")
        print(
            f"[SUPABASE STATUS] client={'INITIALIZED' if supabase else 'NOT INITIALIZED'}"
        )

        if not supabase:
            raise HTTPException(
                status_code=500,
                detail="Supabase not configured - async mode requires Supabase"
            )

        active_count = sum(1 for t in _running_tasks.values() if t and not t.done())
        if active_count >= config.MAX_CONCURRENT_SESSIONS:
            global _sessions_rejected
            _sessions_rejected += 1
            raise HTTPException(
                status_code=503,
                detail=f"Server at capacity ({config.MAX_CONCURRENT_SESSIONS} concurrent sessions). Please try again shortly."
            )

        # Admin instruction override (see /api/chat): per-request prompt wins, else
        # the saved override, else the built-in default. Never blocks the chat path.
        _base_prompt = request.system_prompt
        if not (_base_prompt or "").strip():
            try:
                import agent_prompt_store as _aps
                _ov = await _aw(_aps.get_prompt)
                if (_ov or "").strip():
                    _base_prompt = _ov
            except Exception:  # noqa: BLE001
                pass
        _async_system_prompt = await _lake.inject_lake_context(
            request_messages=request.messages,
            project_id=request.project_id,
            system_prompt=_base_prompt,
            supabase_client=supabase,
        )
        _current_chat_id.set(chat_id)
        if _async_system_prompt or request.model or request.headless != True:
            await agent_manager.reinitialize_agent(
                instructions=_async_system_prompt or None,
                model=request.model,
                headless=request.headless)

        agent = await agent_manager.get_agent()

        messages = [_build_message_content(msg) for msg in request.messages]

        total_chars = sum(len(str(m.get("content", ""))) for m in messages)
        estimated_tokens = total_chars // 4
        if estimated_tokens > config.CONVERSATION_SUMMARIZE_TOKEN_THRESHOLD:
            messages = await context_manager.summarize_conversation_history(
                messages)
        elif len(messages) > 10:
            messages = messages[-10:]

        chunk_queue = asyncio.Queue()
        stream_done = asyncio.Event()

        async def consume_stream():
            try:
                await run_agent_and_save(chat_id, messages, agent,
                                         request.model,
                                         project_id=request.project_id)
            except Exception as e:
                await chunk_queue.put(e)
            finally:
                stream_done.set()

        # Create + register the background task in the endpoint body (NOT inside
        # the generator) so the start-reservation is converted into a live
        # tracked task before we return, regardless of whether/when the stream is
        # iterated. The `finally` below then releases the reservation; concurrency
        # is enforced from here on by the live-task check on _running_tasks.
        consumer_task = asyncio.create_task(consume_stream())
        _running_tasks[chat_id] = consumer_task
        _session_start_times[chat_id] = asyncio.get_event_loop().time()
        def _cleanup_session(t, cid=chat_id):
            _running_tasks.pop(cid, None)
            _session_start_times.pop(cid, None)
            _supabase_seq_counters.pop(cid, None)
            _dedupe_completed.pop(cid, None)
            _dedupe_inflight.pop(cid, None)
            _approved_campaigns_cache.pop(cid, None)
            _starting_chats.discard(cid)
        consumer_task.add_done_callback(_cleanup_session)

        async def run_and_keepalive():
            KEEPALIVE_INTERVAL = 5
            yield f"data: {json.dumps({'type': 'chat_id', 'chat_id': chat_id})}\n\n"
            print(f"[ASYNC] Starting agent run for chat_id={chat_id}")

            while not stream_done.is_set():
                try:
                    msg = await asyncio.wait_for(chunk_queue.get(),
                                                 timeout=KEEPALIVE_INTERVAL)
                    if isinstance(msg, Exception):
                        yield f"data: {json.dumps({'type': 'error', 'content': str(msg)})}\n\n"
                        break
                except asyncio.TimeoutError:
                    yield f"data: {json.dumps({'type': 'ping'})}\n\n"

            await consumer_task
            yield f"data: {json.dumps({'type': 'done', 'chat_id': chat_id})}\n\n"
            print(f"[ASYNC] Completed for chat_id={chat_id}")

        return StreamingResponse(run_and_keepalive(),
                                 media_type="text/event-stream",
                                 headers={
                                     "X-Accel-Buffering": "no",
                                     "Cache-Control": "no-cache",
                                     "Connection": "keep-alive",
                                 })

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"[ASYNC ERROR] {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        # Idempotent release on EVERY exit path (success return, HTTP error,
        # unexpected error, or cancellation/client-disconnect during setup). If a
        # task was registered above it already owns the slot via _running_tasks,
        # so this discard is harmless; if setup failed before registration, this
        # frees the chat so it isn't permanently 409-blocked. The duplicate-run
        # 409 is raised before the try, so it never reaches this finally.
        _release_run_slot(chat_id)


@app.post("/api/chat")
async def chat(request: ChatRequest):
    """Chat endpoint with streaming support, browser control, and context management"""
    chat_id = request.chat_id or str(uuid.uuid4())
    # Backend backstop: one run per chat. Reserved before the try so the
    # duplicate-run 409 propagates cleanly without the error handler releasing a
    # slot owned by another request.
    _reserve_run_slot(chat_id)
    try:
        active_count = sum(1 for t in _running_tasks.values() if t and not t.done())
        if active_count >= config.MAX_CONCURRENT_SESSIONS:
            global _sessions_rejected
            _sessions_rejected += 1
            raise HTTPException(
                status_code=503,
                detail=f"Server at capacity ({config.MAX_CONCURRENT_SESSIONS} concurrent sessions). Please try again shortly."
            )

        print(f"\n{'='*60}")
        print(f"[API REQUEST] /api/chat")
        print(
            f"[REQUEST] chat_id={chat_id}, stream={request.stream}, model={request.model}, headless={request.headless}"
        )
        print(f"[REQUEST] messages count={len(request.messages)}")
        for i, msg in enumerate(request.messages):
            img_count = len(msg.images) if msg.images else 0
            print(
                f"[REQUEST] msg[{i}] role={msg.role}, content_len={len(msg.content)}, images={img_count}, content_preview={msg.content[:200]}"
            )
        print(
            f"[SUPABASE STATUS] client={'INITIALIZED' if supabase else 'NOT INITIALIZED'}"
        )

        # Admin instruction override (Admin -> Agent Control 'Instructions'). A
        # per-request prompt always wins; otherwise fall back to the saved override
        # (then the built-in default). Never blocks the chat path.
        _base_prompt = request.system_prompt
        if not (_base_prompt or "").strip():
            try:
                import agent_prompt_store as _aps
                _ov = await _aw(_aps.get_prompt)
                if (_ov or "").strip():
                    _base_prompt = _ov
            except Exception:  # noqa: BLE001
                pass
        _chat_system_prompt = await _lake.inject_lake_context(
            request_messages=request.messages,
            project_id=getattr(request, 'project_id', None),
            system_prompt=_base_prompt,
            supabase_client=supabase,
        )
        _current_chat_id.set(chat_id)
        if _chat_system_prompt or request.model or request.headless != True:
            await agent_manager.reinitialize_agent(
                instructions=_chat_system_prompt or None,
                model=request.model,
                headless=request.headless)

        agent = await agent_manager.get_agent()

        messages = [_build_message_content(msg) for msg in request.messages]

        total_chars = sum(len(str(m.get("content", ""))) for m in messages)
        estimated_tokens = total_chars // 4
        if estimated_tokens > config.CONVERSATION_SUMMARIZE_TOKEN_THRESHOLD:
            messages = await context_manager.summarize_conversation_history(
                messages)
        elif len(messages) > 10:
            messages = messages[-10:]

        if request.stream:
            event_queue = asyncio.Queue()
            agent_done = asyncio.Event()

            async def run_agent_background():
                try:
                    final_response = ""
                    step_count = 0
                    seen_tool_calls = set()
                    seen_tool_results = set()
                    thinking_logs = []
                    pending_thinking = ""
                    terminal_meta: dict = {}
                    terminal_written = False

                    _current_chat_id.set(chat_id)
                    _set_rag_context(chat_id, getattr(request, 'project_id', None))
                    from custom_tools.search_knowledge import reset_search_knowledge_counter
                    reset_search_knowledge_counter(chat_id)
                    await save_to_supabase(chat_id, "status", "processing",
                                           {"status": "started"})
                    await event_queue.put({
                        "type": "chat_id",
                        "chat_id": chat_id
                    })

                    async for chunk in _agent_astream_autocontinue(
                            agent, messages, chat_id, terminal_meta):
                        if chat_id in _cancelled_chats:
                            print(f"[STREAM] CANCELLED chat_id={chat_id}")
                            await save_to_supabase(chat_id, "status",
                                                   "Agent stopped by user.", {
                                                       "status": "cancelled",
                                                       "tool_calls": step_count
                                                   })
                            # The UI only treats type='final'/'error' as
                            # terminal — a 'status' row leaves the spinner
                            # spinning. Write a terminal 'error' row so a
                            # user-stop reliably ends the turn.
                            _res = await save_to_supabase(chat_id, "error",
                                                   "Agent stopped by user.", {
                                                       "status": "cancelled",
                                                       "kind": "user_stopped",
                                                       "tool_calls": step_count
                                                   })
                            terminal_written = _res is not None
                            await event_queue.put({
                                "type": "status",
                                "content": "Agent stopped by user.",
                                "status": "cancelled"
                            })
                            _cancelled_chats.discard(chat_id)
                            return
                        if "messages" not in chunk or not chunk["messages"]:
                            continue
                        last_message = chunk["messages"][-1]
                        msg_type = type(last_message).__name__

                        if msg_type == "AIMessage":
                            _accumulate_tokens_from_msg(chat_id, request.model, last_message)
                            has_tool_calls = hasattr(
                                last_message,
                                'tool_calls') and last_message.tool_calls

                            ai_content = ""
                            if hasattr(last_message,
                                       'content') and last_message.content:
                                raw = last_message.content
                                if isinstance(raw, list):
                                    text_parts = [
                                        block.get("text", "") for block in raw
                                        if isinstance(block, dict)
                                        and block.get("type") == "text"
                                    ]
                                    ai_content = "\n".join(text_parts).strip()
                                else:
                                    ai_content = str(raw).strip()

                            if ai_content and ai_content != final_response:
                                if has_tool_calls:
                                    print(
                                        f"[STREAM] AI thinking ({len(ai_content)} chars): {ai_content[:200]}"
                                    )
                                    thinking_logs.append(ai_content)
                                    await event_queue.put({
                                        "type": "thinking",
                                        "content": ai_content
                                    })
                                    await save_to_supabase(
                                        chat_id, "thinking", ai_content,
                                        {"source": "agent_reasoning"})
                                else:
                                    if pending_thinking:
                                        print(
                                            f"[STREAM] AI thinking ({len(pending_thinking)} chars): {pending_thinking[:200]}"
                                        )
                                        thinking_logs.append(pending_thinking)
                                        await event_queue.put({
                                            "type":
                                            "thinking",
                                            "content":
                                            pending_thinking
                                        })
                                        await save_to_supabase(
                                            chat_id, "thinking",
                                            pending_thinking,
                                            {"source": "agent_reasoning"})
                                    pending_thinking = ai_content
                                final_response = ai_content

                            if has_tool_calls:
                                for tool_call in last_message.tool_calls:
                                    tool_call_id = f"{tool_call.get('name', 'unknown')}_{tool_call.get('id', '')}"
                                    if tool_call_id in seen_tool_calls:
                                        continue
                                    seen_tool_calls.add(tool_call_id)
                                    step_count += 1
                                    tool_name = tool_call.get(
                                        'name', 'unknown')
                                    tool_args = tool_call.get('args', {})
                                    log_entry = f"Calling **{tool_name}** with args: `{json.dumps(tool_args, indent=2)}`"
                                    thinking_logs.append(log_entry)
                                    print(
                                        f"[STREAM] Step {step_count}: Calling {tool_name}"
                                    )
                                    await event_queue.put({
                                        "type": "tool_call",
                                        "tool": tool_name,
                                        "args": tool_args
                                    })
                                    await save_to_supabase(
                                        chat_id, "tool_call",
                                        f"Calling {tool_name}...", {
                                            "tool": tool_name,
                                            "args": tool_args,
                                            "step": step_count,
                                            "tool_call_id": tool_call.get("id"),
                                            "source": "stream_handler",
                                        })

                        elif msg_type == "ToolMessage":
                            tool_result_id = getattr(last_message,
                                                     'tool_call_id', 'unknown')
                            if tool_result_id not in seen_tool_results:
                                seen_tool_results.add(tool_result_id)
                                tool_content = str(last_message.content)
                                tool_name = getattr(last_message, 'name',
                                                    'unknown')
                                display_content = tool_content[:500] + (
                                    "..." if len(tool_content) > 500 else "")
                                log_entry = f"Result from **{tool_name}**: \n> {display_content}"
                                thinking_logs.append(log_entry)
                                print(
                                    f"[STREAM] Tool result from {tool_name} ({len(tool_content)} chars)"
                                )
                                await event_queue.put({
                                    "type": "tool_result",
                                    "tool": tool_name,
                                    "result": tool_content
                                })
                                await save_to_supabase(chat_id, "tool_result",
                                                       tool_content,
                                                       {"tool": tool_name})

                    # Exactly one terminal row per turn. A budget circuit-breaker
                    # is folded into the single terminal row (loud note appended
                    # to the answer, or a standalone error row when there is no
                    # model text) rather than written as a separate error + final
                    # pair. terminal_written flips only on a confirmed insert.
                    _breaker = _budget_breaker_reason(terminal_meta)
                    if _breaker:
                        print(f"[STREAM] 🛑 chat={chat_id} circuit-breaker "
                              f"'{_breaker}' "
                              f"cost=${terminal_meta.get('cost_usd', 0):.2f} "
                              f"elapsed={terminal_meta.get('elapsed_s', 0):.0f}s")
                    if final_response:
                        # `final` row = CLEAN answer only. Thinking + tool
                        # calls/results are already saved as their own rows for
                        # the frontend's collapsible activity view; don't inline
                        # that trace into the answer.
                        full_content = final_response
                        if _breaker:
                            full_content += ("\n\n---\n⚠️ " +
                                _budget_breaker_note(terminal_meta, _breaker))
                        print(
                            f"[STREAM] COMPLETED chat_id={chat_id}, tool_calls={step_count}, "
                            f"response_len={len(full_content)}, "
                            f"auto_continue={terminal_meta.get('auto_continue_count', 0)}, "
                            f"terminal={terminal_meta.get('reason')}"
                        )
                        await event_queue.put({
                            "type": "final",
                            "content": final_response
                        })
                        res = await save_to_supabase(chat_id, "final", full_content,
                                               {
                                                   "tool_calls": step_count,
                                                   "status": "budget_exhausted" if _breaker else "completed",
                                                   "auto_continue_count": terminal_meta.get("auto_continue_count", 0),
                                                   "terminal_reason": terminal_meta.get("reason"),
                                               })
                        terminal_written = res is not None
                        _stream_project_id = getattr(request, 'project_id', None)
                        _trigger_verifier(chat_id, _stream_project_id)
                        if _stream_project_id in _lake.OD_PROJECT_IDS:
                            asyncio.create_task(_lake.write_lake_diagnosis(
                                chat_id=chat_id,
                                project_id=_stream_project_id,
                                final_response=final_response,
                                supabase_client=supabase,
                                openai_api_key=config.OPENAI_API_KEY or os.getenv("OPENAI_API_KEY", ""),
                            ))
                    elif _breaker:
                        # No model text — surface budget exhaustion as terminal.
                        _bnote = _budget_breaker_note(terminal_meta, _breaker)
                        await event_queue.put({"type": "error", "content": _bnote})
                        res = await save_to_supabase(chat_id, "error", _bnote,
                                               {"status": "budget_exhausted",
                                                "terminal": terminal_meta})
                        terminal_written = res is not None
                        print(f"[STREAM] 🛑 chat={chat_id} {_breaker} → error terminal")
                    else:
                        await event_queue.put({
                            "type": "final",
                            "content": "Task completed."
                        })
                        res = await save_to_supabase(chat_id, "final",
                                               "Task completed.", {
                                                   "tool_calls": step_count,
                                                   "status": "completed"
                                               })
                        terminal_written = res is not None
                        _trigger_verifier(chat_id, getattr(request, 'project_id', None))

                except Exception as e:
                    import traceback
                    print(f"[STREAM ERROR] {traceback.format_exc()}")
                    await event_queue.put({"type": "error", "content": str(e)})
                    res = await save_to_supabase(chat_id, "error", str(e),
                                           {"status": "failed"})
                    terminal_written = res is not None
                finally:
                    # Terminal-row safety net: guarantee exactly one
                    # final/error row per turn even on an unexpected early
                    # return or swallowed condition, so the UI never hangs on
                    # "Thinking…". Guarded to avoid double-writes.
                    if not terminal_written:
                        try:
                            await save_to_supabase(
                                chat_id, "error",
                                "Run ended without a terminal status (internal).",
                                {"status": "failed", "kind": "missing_terminal"})
                            print(f"[STREAM] safety-net terminal error row "
                                  f"written chat_id={chat_id}")
                        except Exception as _e:
                            print(f"[STREAM] safety-net terminal write failed "
                                  f"chat_id={chat_id}: {_e}")
                    await _save_chat_usage(chat_id)
                    await _evict_checkpoint(chat_id)
                    try:
                        from custom_tools.search_knowledge import reset_search_knowledge_counter
                        reset_search_knowledge_counter(chat_id)
                    except Exception:
                        pass
                    agent_done.set()

            agent_task = asyncio.create_task(run_agent_background())
            _running_tasks[chat_id] = agent_task
            _session_start_times[chat_id] = asyncio.get_event_loop().time()
            def _cleanup_session(t, cid=chat_id):
                _running_tasks.pop(cid, None)
                _session_start_times.pop(cid, None)
                _supabase_seq_counters.pop(cid, None)
                _dedupe_completed.pop(cid, None)
                _dedupe_inflight.pop(cid, None)
                _approved_campaigns_cache.pop(cid, None)
                _starting_chats.discard(cid)
            agent_task.add_done_callback(_cleanup_session)

            async def generate():
                KEEPALIVE_INTERVAL = 5
                try:
                    while not agent_done.is_set() or not event_queue.empty():
                        try:
                            event = await asyncio.wait_for(
                                event_queue.get(), timeout=KEEPALIVE_INTERVAL)
                            yield f"data: {json.dumps(event)}\n\n"
                            if event.get("type") in ("final", "error"):
                                break
                        except asyncio.TimeoutError:
                            yield f"data: {json.dumps({'type': 'ping'})}\n\n"
                except Exception:
                    print(
                        f"[STREAM] Client disconnected for chat_id={chat_id}, agent continues in background"
                    )

            return StreamingResponse(generate(),
                                     media_type="text/event-stream",
                                     headers={
                                         "X-Accel-Buffering": "no",
                                         "Cache-Control": "no-cache",
                                         "Connection": "keep-alive",
                                     })

        else:
            print(
                f"\n[NON-STREAMING REQUEST] chat_id={chat_id}, Messages: {len(messages)}\n"
            )
            _current_chat_id.set(chat_id)
            _set_rag_context(chat_id, getattr(request, 'project_id', None))
            await save_to_supabase(chat_id, "status", "processing",
                                   {"status": "started"})
            result = await agent.ainvoke({"messages": messages},
                                         config={"recursion_limit": _RECURSION_LIMIT})
            tool_call_count = 0
            for msg in result.get("messages", []):
                if hasattr(msg, 'tool_calls') and msg.tool_calls:
                    tool_call_count += len(msg.tool_calls)
                if type(msg).__name__ == "AIMessage":
                    _accumulate_tokens_from_msg(chat_id, request.model, msg)
            response_content = result["messages"][-1].content
            print(
                f"\n[NON-STREAMING COMPLETE] chat_id={chat_id}, Tool calls: {tool_call_count}\n"
            )
            res = await save_to_supabase(chat_id, "final", response_content, {
                "tool_calls": tool_call_count,
                "status": "completed"
            })
            if res is None:
                # Terminal-row safety net: the success `final` insert was dropped
                # after retries. Write a loud error terminal so the UI ends the
                # turn instead of hanging on "Thinking…" (the response is still
                # returned in the HTTP body below).
                await save_to_supabase(
                    chat_id, "error",
                    "Run completed but the result could not be persisted "
                    "(internal).",
                    {"status": "failed", "kind": "missing_terminal"})
            _trigger_verifier(chat_id, getattr(request, 'project_id', None))
            await _save_chat_usage(chat_id)
            return {
                "chat_id": chat_id,
                "response": response_content,
                "done": True
            }

    except Exception as e:
        import traceback
        error_detail = f"{str(e)}\n{traceback.format_exc()}"
        print(f"\n[ERROR] /api/chat failed:\n{error_detail}\n")
        # Write a terminal error row so the UI ends the turn instead of
        # hanging on "Thinking…". The non-streaming branch previously raised
        # without persisting any terminal row. (For the streaming branch this
        # except only fires on synchronous setup errors before the background
        # task starts — the background task owns its own terminal row.)
        try:
            await save_to_supabase(chat_id, "error", str(e),
                                   {"status": "failed",
                                    "kind": "chat_endpoint_error"})
        except Exception as _e:
            print(f"[ERROR] /api/chat terminal-row write failed "
                  f"chat_id={chat_id}: {_e}")
        await _save_chat_usage(chat_id)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        # Idempotent release on EVERY exit path. The streaming branch registers
        # a background task in _running_tasks (which then owns the slot), so this
        # discard is harmless there; the non-streaming branch runs inline while
        # the reservation blocks duplicates, and this frees it on completion or
        # failure/cancellation. The duplicate-run 409 is raised before the try.
        _release_run_slot(chat_id)


@app.post("/api/chat/stop")
async def stop_chat(chat_id: str = None):
    """Stop a running agent task by chat_id"""
    if not chat_id:
        raise HTTPException(status_code=400, detail="chat_id is required")

    print(f"[STOP] Received stop request for chat_id={chat_id}")

    if chat_id in _running_tasks:
        _cancelled_chats.add(chat_id)
        task = _running_tasks.get(chat_id)
        if task and not task.done():
            task.cancel()
            print(f"[STOP] Cancelled chat_id={chat_id}")
            try:
                await save_to_supabase(
                    chat_id, "status", "Agent stopped by user.",
                    {"status": "cancelled", "source": "stop_endpoint"},
                )
            except Exception as e:
                print(f"[STOP] supabase save failed (non-fatal): {e}")
            return {
                "chat_id": chat_id,
                "status": "stopped",
                "message": "Agent stopped.",
            }
        else:
            _running_tasks.pop(chat_id, None)
            return {
                "chat_id": chat_id,
                "status": "already_done",
                "message": "Agent task has already completed."
            }
    else:
        return {
            "chat_id":
            chat_id,
            "status":
            "not_found",
            "message":
            "No running agent task found for this chat_id. It may have already completed."
        }


@app.get("/api/chat/active")
async def list_active_chats():
    """List all currently running agent tasks"""
    active = []
    for cid, task in list(_running_tasks.items()):
        if task and not task.done():
            active.append({"chat_id": cid, "status": "running"})
        else:
            _running_tasks.pop(cid, None)
    return {"active_chats": active, "count": len(active)}


@app.get("/api/metrics")
async def get_metrics():
    """Server metrics for monitoring concurrent usage"""
    import psutil
    process = psutil.Process()
    mem_info = process.memory_info()

    active = sum(1 for t in _running_tasks.values() if t and not t.done())
    uptime = asyncio.get_event_loop().time() - _server_start_time if _server_start_time else 0

    return {
        "concurrent_sessions": {
            "active": active,
            "max_allowed": config.MAX_CONCURRENT_SESSIONS,
            "utilization_pct": round(active / config.MAX_CONCURRENT_SESSIONS * 100, 1),
            "total_rejected": _sessions_rejected,
        },
        "memory": {
            "rss_mb": round(mem_info.rss / 1024 / 1024, 1),
            "vms_mb": round(mem_info.vms / 1024 / 1024, 1),
        },
        "sessions": {
            "running": [{"chat_id": cid, "running_seconds": round(asyncio.get_event_loop().time() - _session_start_times.get(cid, 0), 0)} for cid, t in _running_tasks.items() if t and not t.done()],
            "timeout_minutes": config.SESSION_TIMEOUT_MINUTES,
        },
        "server": {
            "uptime_seconds": round(uptime, 0),
            "uptime_human": f"{int(uptime // 3600)}h {int((uptime % 3600) // 60)}m",
            "mcp_tools_loaded": agent_manager.mcp_tools_loaded,
        }
    }


LINKEDIN_CLIENT_ID = os.environ.get("LINKEDIN_CLIENT_ID", "")
LINKEDIN_CLIENT_SECRET = os.environ.get("LINKEDIN_CLIENT_SECRET", "")


def _get_linkedin_redirect_uri():
    domain = os.environ.get("REPLIT_DOMAINS", os.environ.get("REPLIT_DEV_DOMAIN", ""))
    if domain:
        return f"https://{domain}/auth/linkedin/callback"
    return "https://localhost:5000/auth/linkedin/callback"


@app.get("/auth/linkedin")
async def linkedin_auth_start():
    """Start LinkedIn OAuth flow — share this URL with your LinkedIn person."""
    redirect_uri = _get_linkedin_redirect_uri()
    scopes = "r_ads r_ads_reporting r_ads_leadgen_automation"
    auth_url = (
        f"https://www.linkedin.com/oauth/v2/authorization"
        f"?response_type=code"
        f"&client_id={LINKEDIN_CLIENT_ID}"
        f"&redirect_uri={redirect_uri}"
        f"&scope={scopes.replace(' ', '%20')}"
        f"&state=deepagent_linkedin_auth"
    )
    html = f"""
    <html>
    <head><title>LinkedIn Ads Authorization</title>
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, sans-serif; max-width: 600px; margin: 80px auto; padding: 20px; }}
        h1 {{ color: #0077B5; }}
        .btn {{ display: inline-block; background: #0077B5; color: white; padding: 14px 28px;
                text-decoration: none; border-radius: 6px; font-size: 16px; margin-top: 20px; }}
        .btn:hover {{ background: #005f8d; }}
        .info {{ background: #f0f7ff; border: 1px solid #cce0ff; padding: 16px; border-radius: 8px; margin: 20px 0; }}
        code {{ background: #e8e8e8; padding: 2px 6px; border-radius: 3px; }}
    </style>
    </head>
    <body>
        <h1>LinkedIn Ads Authorization</h1>
        <div class="info">
            <p><strong>What this does:</strong> Connects your LinkedIn Ads account to DeepAgent so the AI can read your campaign data, analytics, and ad performance.</p>
            <p><strong>Permissions requested:</strong></p>
            <ul>
                <li><code>r_ads</code> — Read your ad accounts and campaigns</li>
                <li><code>r_ads_reporting</code> — Read campaign analytics and reports</li>
                <li><code>r_ads_leadgen_automation</code> — Read Lead Gen Form submissions (Lead Sync API)</li>
            </ul>
            <p><strong>Note:</strong> This is read-only access. DeepAgent cannot modify your campaigns or leads.</p>
        </div>
        <a class="btn" href="{auth_url}">Authorize with LinkedIn</a>
        <p style="margin-top: 30px; color: #666; font-size: 13px;">
            Redirect URI: <code>{redirect_uri}</code><br>
            Make sure this URL is added to your LinkedIn app's authorized redirect URLs.
        </p>
    </body>
    </html>
    """
    return HTMLResponse(content=html)


@app.get("/auth/linkedin/callback")
async def linkedin_auth_callback(code: str = None, error: str = None, error_description: str = None, state: str = None):
    """LinkedIn OAuth callback — exchanges auth code for access token."""
    logger = logging.getLogger("linkedin_oauth")
    logger.info(f"LinkedIn callback hit: code={'YES' if code else 'NO'}, error={error}, state={state}")
    if error:
        html = f"""
        <html><head><title>LinkedIn Auth Error</title>
        <style>body {{ font-family: sans-serif; max-width: 600px; margin: 80px auto; padding: 20px; }}
        .error {{ background: #fff0f0; border: 1px solid #ffcccc; padding: 16px; border-radius: 8px; }}</style></head>
        <body><h1>Authorization Failed</h1>
        <div class="error"><p><strong>Error:</strong> {error}</p><p>{error_description or ''}</p></div>
        <p>Go back and try again, or check that your LinkedIn app has the Advertising API product enabled.</p></body></html>
        """
        return HTMLResponse(content=html, status_code=400)

    if not code:
        return HTMLResponse(content="<h1>Missing authorization code</h1>", status_code=400)

    redirect_uri = _get_linkedin_redirect_uri()

    try:
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.post(
                "https://www.linkedin.com/oauth/v2/accessToken",
                data={
                    "grant_type": "authorization_code",
                    "code": code,
                    "client_id": LINKEDIN_CLIENT_ID,
                    "client_secret": LINKEDIN_CLIENT_SECRET,
                    "redirect_uri": redirect_uri,
                },
                headers={"Content-Type": "application/x-www-form-urlencoded"},
            )

        if resp.status_code != 200:
            error_detail = resp.text
            html = f"""
            <html><head><title>Token Exchange Failed</title>
            <style>body {{ font-family: sans-serif; max-width: 600px; margin: 80px auto; padding: 20px; }}
            .error {{ background: #fff0f0; border: 1px solid #ffcccc; padding: 16px; border-radius: 8px; word-break: break-all; }}</style></head>
            <body><h1>Token Exchange Failed</h1>
            <div class="error"><p><strong>HTTP {resp.status_code}</strong></p><p>{error_detail}</p></div>
            <p>This usually means the authorization code expired (they only last 30 minutes) or the redirect URI doesn't match. Try authorizing again.</p></body></html>
            """
            return HTMLResponse(content=html, status_code=400)

        token_data = resp.json()
        access_token = token_data.get("access_token", "")
        expires_in = token_data.get("expires_in", 0)
        expires_days = round(expires_in / 86400)
        scope = token_data.get("scope", "")

        masked_token = access_token[:8] + "..." + access_token[-4:] if len(access_token) > 12 else "***"

        os.environ["LINKEDIN_ACCESS_TOKEN"] = access_token.strip()
        logger.info(f"LinkedIn token obtained! Preview: {masked_token}, expires_in: {expires_in}s, scope: {scope}")

        html = f"""
        <html><head><title>LinkedIn Authorization Successful</title>
        <style>
            body {{ font-family: -apple-system, BlinkMacSystemFont, sans-serif; max-width: 600px; margin: 80px auto; padding: 20px; }}
            h1 {{ color: #28a745; }}
            .success {{ background: #f0fff4; border: 1px solid #b7e4c7; padding: 16px; border-radius: 8px; }}
            .token-box {{ background: #f8f9fa; border: 1px solid #dee2e6; padding: 12px; border-radius: 6px;
                          font-family: monospace; word-break: break-all; margin: 12px 0; font-size: 13px; }}
            .warn {{ background: #fff8e1; border: 1px solid #ffe082; padding: 12px; border-radius: 6px; margin: 12px 0; }}
        </style></head>
        <body>
            <h1>Authorization Successful!</h1>
            <div class="success">
                <p>LinkedIn Ads access token has been obtained and loaded into the server.</p>
                <p><strong>Scopes granted:</strong> {scope}</p>
                <p><strong>Expires in:</strong> {expires_days} days ({expires_in} seconds)</p>
                <p><strong>Token preview:</strong> <code>{masked_token}</code></p>
            </div>
            <div class="warn">
                <p><strong>Important:</strong> The token is active in memory now, but to make it permanent,
                copy the full token below and add it as the <code>LINKEDIN_ACCESS_TOKEN</code> secret in Replit.</p>
            </div>
            <p><strong>Full access token (copy this):</strong></p>
            <div class="token-box" id="token">{access_token}</div>
            <button onclick="navigator.clipboard.writeText(document.getElementById('token').innerText).then(()=>this.innerText='Copied!')"
                    style="padding: 8px 16px; cursor: pointer; border: 1px solid #0077B5; background: #0077B5; color: white; border-radius: 4px;">
                Copy Token
            </button>
            <p style="margin-top: 30px; color: #666; font-size: 13px;">
                You can close this page now. The DeepAgent LinkedIn Ads tools are ready to use.
            </p>
        </body></html>
        """
        return HTMLResponse(content=html)

    except Exception as e:
        html = f"""
        <html><head><title>Error</title>
        <style>body {{ font-family: sans-serif; max-width: 600px; margin: 80px auto; padding: 20px; }}
        .error {{ background: #fff0f0; border: 1px solid #ffcccc; padding: 16px; border-radius: 8px; }}</style></head>
        <body><h1>Something went wrong</h1>
        <div class="error"><p>{str(e)}</p></div></body></html>
        """
        return HTMLResponse(content=html, status_code=500)


@app.post("/api/chat/structured")
async def structured_chat(request: StructuredChatRequest):
    """Chat endpoint with structured output support. Uses keep-alive pings to survive Autoscale timeout."""
    import re as _re

    chat_id = str(uuid.uuid4())
    agent_done = asyncio.Event()
    agent_result: dict = {}

    system_prompt = request.system_prompt
    if request.google_sheets:
        sheets_context = "\n\n## AVAILABLE GOOGLE SHEETS\n\nYou have access to the following Google Sheets:\n\n"
        for idx, sheet in enumerate(request.google_sheets, 1):
            sheets_context += f"{idx}. Spreadsheet ID: `{sheet.spreadsheet_id}`"
            if sheet.sheet_name:
                sheets_context += f" (Sheet: {sheet.sheet_name})"
            sheets_context += "\n"
        if system_prompt:
            system_prompt = system_prompt + sheets_context
        else:
            system_prompt = sheets_context

    async def run_agent():
        try:
            _current_chat_id.set(chat_id)
            _set_rag_context(chat_id, getattr(request, 'project_id', None))

            if system_prompt or request.model or request.headless != True:
                await agent_manager.reinitialize_agent(instructions=system_prompt,
                                                       model=request.model,
                                                       headless=request.headless)

            agent = await agent_manager.get_agent()
            messages = [_build_message_content(msg) for msg in request.messages]

            total_chars = sum(len(str(m.get("content", ""))) for m in messages)
            estimated_tokens = total_chars // 4
            if estimated_tokens > config.CONVERSATION_SUMMARIZE_TOKEN_THRESHOLD:
                messages = await context_manager.summarize_conversation_history(messages)
            elif len(messages) > 10:
                messages = messages[-10:]

            schema_str = json.dumps(request.structured_output_format, indent=2)
            structured_instruction = f"\n\nIMPORTANT: You MUST respond with valid JSON matching this exact schema:\n{schema_str}\n\nDo not include any text outside the JSON object."
            if messages and messages[-1]["role"] == "user":
                last_content = messages[-1]["content"]
                if isinstance(last_content, list):
                    messages[-1]["content"].append({"type": "text", "text": structured_instruction})
                else:
                    messages[-1]["content"] += structured_instruction

            result = await agent.ainvoke({"messages": messages},
                                         config={"recursion_limit": _RECURSION_LIMIT})
            response_content = result["messages"][-1].content

            if isinstance(response_content, list):
                text_parts = []
                for block in response_content:
                    if isinstance(block, dict) and block.get("type") == "text":
                        text_parts.append(block.get("text", ""))
                    elif isinstance(block, str):
                        text_parts.append(block)
                response_content = "\n".join(text_parts)

            try:
                json_match = _re.search(r'\{.*\}', str(response_content), _re.DOTALL)
                if json_match:
                    structured_data = json.loads(json_match.group(0))
                else:
                    structured_data = json.loads(str(response_content))
                agent_result["response"] = {
                    "data": structured_data,
                    "raw_response": str(response_content),
                    "success": True
                }
            except json.JSONDecodeError as e:
                agent_result["response"] = {
                    "data": None,
                    "raw_response": str(response_content),
                    "success": False,
                    "error": f"Failed to parse JSON: {str(e)}"
                }
        except Exception as e:
            import traceback
            print(f"Error in structured_chat agent: {traceback.format_exc()}")
            agent_result["response"] = {
                "data": None,
                "raw_response": None,
                "success": False,
                "error": str(e)
            }
        finally:
            agent_done.set()

    async def response_stream():
        task = asyncio.create_task(run_agent())
        _running_tasks[chat_id] = task
        try:
            while not agent_done.is_set():
                try:
                    await asyncio.wait_for(agent_done.wait(), timeout=15)
                except asyncio.TimeoutError:
                    yield b"\n"

            if task.done() and not task.cancelled() and task.exception():
                err = str(task.exception())
                print(f"[STRUCTURED] Agent task exception: {err}")
                yield json.dumps({"data": None, "success": False, "error": err, "chat_id": chat_id}, default=str).encode()
                return

            final_response = agent_result.get("response")
            if final_response is None:
                print(f"[STRUCTURED] No response in agent_result for chat_id={chat_id}")
                final_response = {"data": None, "success": False, "error": "Agent produced no response"}
            final_response["chat_id"] = chat_id
            yield json.dumps(final_response, default=str).encode()
        except GeneratorExit:
            print(f"[STRUCTURED] Client disconnected for chat_id={chat_id}")
            if not task.done():
                task.cancel()
        except Exception as e:
            import traceback
            print(f"[STRUCTURED] Stream error: {traceback.format_exc()}")
            try:
                yield json.dumps({"data": None, "success": False, "error": str(e), "chat_id": chat_id}, default=str).encode()
            except Exception:
                pass
        finally:
            _running_tasks.pop(chat_id, None)

    return StreamingResponse(
        response_stream(),
        media_type="text/plain"
    )


@app.post("/api/chat/structured/async")
async def structured_chat_async(request: StructuredChatRequest):
    """Async structured chat - returns chat_id immediately, runs agent in background, saves result to Supabase."""
    try:
        chat_id = str(uuid.uuid4())

        system_prompt = request.system_prompt
        if request.google_sheets:
            sheets_context = "\n\n## AVAILABLE GOOGLE SHEETS\n\nYou have access to the following Google Sheets:\n\n"
            for idx, sheet in enumerate(request.google_sheets, 1):
                sheets_context += f"{idx}. Spreadsheet ID: `{sheet.spreadsheet_id}`"
                if sheet.sheet_name:
                    sheets_context += f" (Sheet: {sheet.sheet_name})"
                sheets_context += "\n"
            if system_prompt:
                system_prompt = system_prompt + sheets_context
            else:
                system_prompt = sheets_context

        async def run_structured_agent():
            terminal_written = False
            try:
                _current_chat_id.set(chat_id)
                _set_rag_context(chat_id, getattr(request, 'project_id', None))
                seq = [0]

                await save_to_supabase(chat_id, "status", "processing", {"source": "structured_async"})
                seq[0] += 1

                if system_prompt or request.model or request.headless != True:
                    await agent_manager.reinitialize_agent(instructions=system_prompt,
                                                           model=request.model,
                                                           headless=request.headless)

                agent = await agent_manager.get_agent()
                messages = [_build_message_content(msg) for msg in request.messages]

                total_chars = sum(len(str(m.get("content", ""))) for m in messages)
                estimated_tokens = total_chars // 4
                if estimated_tokens > config.CONVERSATION_SUMMARIZE_TOKEN_THRESHOLD:
                    messages = await context_manager.summarize_conversation_history(messages)
                elif len(messages) > 10:
                    messages = messages[-10:]

                schema_str = json.dumps(request.structured_output_format, indent=2)
                structured_instruction = f"\n\nIMPORTANT: You MUST respond with valid JSON matching this exact schema:\n{schema_str}\n\nDo not include any text outside the JSON object."
                if messages and messages[-1]["role"] == "user":
                    last_content = messages[-1]["content"]
                    if isinstance(last_content, list):
                        messages[-1]["content"].append({"type": "text", "text": structured_instruction})
                    else:
                        messages[-1]["content"] += structured_instruction

                result = await agent.ainvoke({"messages": messages},
                                             config={"recursion_limit": _RECURSION_LIMIT})
                response_content = result["messages"][-1].content

                if isinstance(response_content, list):
                    text_parts = []
                    for block in response_content:
                        if isinstance(block, dict) and block.get("type") == "text":
                            text_parts.append(block.get("text", ""))
                        elif isinstance(block, str):
                            text_parts.append(block)
                    response_content = "\n".join(text_parts)

                try:
                    import re
                    json_match = re.search(r'\{.*\}', str(response_content), re.DOTALL)
                    if json_match:
                        structured_data = json.loads(json_match.group(0))
                    else:
                        structured_data = json.loads(str(response_content))

                    final_result = {
                        "data": structured_data,
                        "raw_response": str(response_content),
                        "success": True
                    }
                except json.JSONDecodeError as e:
                    final_result = {
                        "data": None,
                        "raw_response": str(response_content),
                        "success": False,
                        "error": f"Failed to parse JSON: {str(e)}"
                    }

                res = await save_to_supabase(chat_id, "final", json.dumps(final_result, default=str), {
                    "source": "structured_async",
                    "success": final_result.get("success", False)
                })
                terminal_written = res is not None
                print(f"[STRUCTURED ASYNC] Completed chat_id={chat_id}, success={final_result.get('success')}")

            except asyncio.CancelledError:
                res = await save_to_supabase(chat_id, "error", "Task was cancelled", {"source": "structured_async", "kind": "user_stopped"})
                terminal_written = res is not None
                print(f"[STRUCTURED ASYNC] Cancelled chat_id={chat_id}")
            except Exception as e:
                error_msg = str(e)
                print(f"[STRUCTURED ASYNC] Error chat_id={chat_id}: {error_msg}")
                import traceback
                traceback.print_exc()
                res = await save_to_supabase(chat_id, "error", error_msg[:2000], {"source": "structured_async"})
                terminal_written = res is not None
            finally:
                # Terminal-row safety net: guarantee exactly one final/error row
                # per turn even if a terminal insert was dropped after retries,
                # so the UI never hangs on "Thinking…".
                if not terminal_written:
                    try:
                        await save_to_supabase(
                            chat_id, "error",
                            "Run ended without a terminal status (internal).",
                            {"source": "structured_async", "kind": "missing_terminal"})
                        print(f"[STRUCTURED ASYNC] safety-net terminal error row "
                              f"written chat_id={chat_id}")
                    except Exception as _e:
                        print(f"[STRUCTURED ASYNC] safety-net terminal write "
                              f"failed chat_id={chat_id}: {_e}")

        task = asyncio.create_task(run_structured_agent())
        _running_tasks[chat_id] = task

        return {
            "chat_id": chat_id,
            "status": "processing",
            "message": "Structured agent task started. Results will be saved to Supabase chat_messages table. Listen for type='final' with your chat_id."
        }

    except Exception as e:
        import traceback
        print(f"Error in structured_chat_async: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/chat/{chat_id}/verifier_report")
async def get_chat_verifier_report(chat_id: str):
    """Return the most recent verifier verdict for a chat (advisory only).

    Returns the structured report (per-check pass/miss) if the verifier ran
    on this chat. 404 if no report exists (chat out of scope, verifier not
    yet run, or run still in progress).
    """
    if not supabase:
        raise HTTPException(status_code=500, detail="Supabase not configured")
    from verifier.runner import latest_report
    try:
        loop = asyncio.get_event_loop()
        report = await loop.run_in_executor(
            None, lambda: latest_report(chat_id, supabase)
        )
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Failed to load verifier report: {e}"
        )
    if report is None:
        raise HTTPException(
            status_code=404,
            detail=f"No verifier report for chat_id={chat_id}",
        )
    return JSONResponse(content=report)


@app.get("/api/chat/result/{chat_id}")
async def get_chat_result(chat_id: str):
    """Poll for chat result by chat_id. Returns the final result from Supabase if ready."""
    if not supabase:
        raise HTTPException(status_code=500, detail="Supabase not configured")
    try:
        loop = asyncio.get_event_loop()
        result = await loop.run_in_executor(
            None,
            lambda: supabase.table("chat_messages")
                .select("*")
                .eq("chat_id", chat_id)
                .eq("type", "final")
                .limit(1)
                .execute()
        )
        if result.data and len(result.data) > 0:
            row = result.data[0]
            try:
                content = json.loads(row["content"])
            except (json.JSONDecodeError, TypeError):
                content = row["content"]
            return JSONResponse(content={
                "status": "completed",
                "chat_id": chat_id,
                "result": content
            })

        error_result = await loop.run_in_executor(
            None,
            lambda: supabase.table("chat_messages")
                .select("*")
                .eq("chat_id", chat_id)
                .eq("type", "error")
                .limit(1)
                .execute()
        )
        if error_result.data and len(error_result.data) > 0:
            return JSONResponse(content={
                "status": "error",
                "chat_id": chat_id,
                "error": error_result.data[0].get("content", "Unknown error")
            })

        if chat_id in _running_tasks and not _running_tasks[chat_id].done():
            return JSONResponse(content={
                "status": "processing",
                "chat_id": chat_id
            })

        return JSONResponse(content={
            "status": "not_found",
            "chat_id": chat_id,
            "error": "No result found for this chat_id"
        })
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# Cached probe: does the documents table have a doc_type column yet? (The column
# is added by a Supabase migration the Next.js team applies; until then we fall
# back to encoding doc_type into the document name so nothing breaks pre-migration.)
_DOCS_DOC_TYPE_COL: Optional[bool] = None


# Upload extraction caps (defend against memory amplification + adversarial PDFs).
# Large files now arrive via S3 (browser → presigned PUT → backend pulls), so these
# are sized for real sales decks, not the old ~4.5 MB Vercel-proxy ceiling. All are
# env-overridable so we can tune without a redeploy.
_MAX_FILE_B64 = int(os.getenv("MASE_MAX_FILE_B64", "280000000"))   # ~210 MB decoded (legacy inline path)
_MAX_FILE_BYTES = int(os.getenv("MASE_MAX_UPLOAD_BYTES", "200000000"))  # 200 MB decoded
_MAX_PDF_PAGES = int(os.getenv("MASE_MAX_PDF_PAGES", "5000"))
_MAX_EXTRACT_CHARS = int(os.getenv("MASE_MAX_EXTRACT_CHARS", "4000000"))


def _extract_text_from_file(file_b64: str, name: str) -> str:
    """Extract text from a base64-encoded uploaded file (legacy inline path: small
    pastes / direct posts). Decodes then delegates to _extract_text_from_bytes."""
    import base64
    if len(file_b64 or "") > _MAX_FILE_B64:
        raise ValueError("file too large")
    try:
        raw = base64.b64decode(file_b64, validate=True)
    except Exception:
        raise ValueError("invalid base64 file payload")
    return _extract_text_from_bytes(raw, name)


def _extract_text_from_bytes(raw: bytes, name: str) -> str:
    """Extract plain text from raw uploaded file bytes. PDF via pypdf, DOCX via
    python-docx, XLSX via openpyxl, PPTX via python-pptx; anything else is decoded
    as UTF-8 text. CPU-bound + can be expensive on adversarial input, so the caller
    runs this in a thread pool with a wall-clock timeout. Bounded on every axis
    (decoded size, page count, output length); raises ValueError on a
    too-large/corrupt/invalid file so the caller 400s. Shared by the inline base64
    path and the S3 path (browser → presigned PUT → backend downloads → here)."""
    import io
    if len(raw) > _MAX_FILE_BYTES:
        raise ValueError(f"file too large (max ~{_MAX_FILE_BYTES // 1_000_000} MB)")
    low = (name or "").lower()
    if low.endswith(".pdf"):
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(raw))
        parts = []
        for i, p in enumerate(reader.pages):
            if i >= _MAX_PDF_PAGES:
                break
            parts.append(p.extract_text() or "")
        text = "\n\n".join(parts).strip()
    elif low.endswith(".docx"):
        import docx  # python-docx
        d = docx.Document(io.BytesIO(raw))
        text = "\n".join(p.text for p in d.paragraphs).strip()
    elif low.endswith((".xlsx", ".xlsm")):
        # Excel: flatten every sheet to tab-separated rows. read_only + data_only keeps
        # it memory-light and resolves formulas to values. Bounded by output length.
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        parts, total = [], 0
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if not cells:
                    continue
                line = "\t".join(cells)
                parts.append(line)
                total += len(line) + 1
                if total > _MAX_EXTRACT_CHARS:
                    break
            if total > _MAX_EXTRACT_CHARS:
                break
        try:
            wb.close()
        except Exception:  # noqa: BLE001
            pass
        text = "\n".join(parts).strip()
    elif low.endswith(".xls"):
        raise ValueError("legacy .xls is not supported — save it as .xlsx and re-upload")
    elif low.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
        text = "\n".join(parts).strip()
    else:
        # Plain-text family: txt, md, markdown, csv, tsv, json, html, xml, yaml, log, etc.
        text = raw.decode("utf-8", errors="replace").strip()
    return text[:_MAX_EXTRACT_CHARS]


# ── S3 staging for large knowledge uploads ───────────────────────────────────
# Big files (sales decks, etc.) can't go through the Vercel proxy (~4.5 MB body
# cap) or Supabase Storage limits, so the browser uploads them straight to S3 via
# a presigned PUT, then the backend pulls the object here and extracts the text.
# Auth is the presigned URL itself (admin-gated when minted at the proxy); the
# bucket is private and the ECS task role has GetObject/PutObject/DeleteObject on
# it only (policy mase-knowledge-s3).
_S3_BUCKET = os.getenv("MASE_KNOWLEDGE_S3_BUCKET", "mase-knowledge-uploads-022187637784")
_S3_REGION = os.getenv("AWS_REGION") or os.getenv("AWS_DEFAULT_REGION") or "ap-south-1"
_PRESIGN_EXPIRY_S = int(os.getenv("MASE_PRESIGN_EXPIRY_S", "900"))  # 15 min
_s3_client = None


def _get_s3():
    global _s3_client
    if _s3_client is None:
        import boto3
        from botocore.config import Config
        # Pin the REGIONAL virtual-hosted endpoint + SigV4. Without this, boto3 signs
        # presigned URLs against the global s3.amazonaws.com host, which 307-redirects
        # for non-us-east-1 buckets and breaks the (host-bound) signature on PUT.
        _s3_client = boto3.client(
            "s3", region_name=_S3_REGION,
            endpoint_url=f"https://s3.{_S3_REGION}.amazonaws.com",
            config=Config(signature_version="s3v4", s3={"addressing_style": "virtual"}),
        )
    return _s3_client


def _s3_presign_put(key: str) -> str:
    return _get_s3().generate_presigned_url(
        "put_object", Params={"Bucket": _S3_BUCKET, "Key": key}, ExpiresIn=_PRESIGN_EXPIRY_S)


def _s3_download(key: str) -> bytes:
    obj = _get_s3().get_object(Bucket=_S3_BUCKET, Key=key)
    raw = obj["Body"].read()
    if len(raw) > _MAX_FILE_BYTES:
        raise ValueError(f"file too large (max ~{_MAX_FILE_BYTES // 1_000_000} MB)")
    return raw


def _s3_delete(key: str) -> None:
    try:
        _get_s3().delete_object(Bucket=_S3_BUCKET, Key=key)
    except Exception as _e:  # noqa: BLE001
        print(f"[KNOWLEDGE] S3 cleanup failed for {key}: {_e}")


async def _documents_has_doc_type() -> bool:
    """True iff the documents table has a doc_type column (cached after first probe).
    Lets upload_document store doc_type natively once the migration lands, without a
    redeploy, and fall back to name-encoding before then."""
    global _DOCS_DOC_TYPE_COL
    if _DOCS_DOC_TYPE_COL is not None:
        return _DOCS_DOC_TYPE_COL
    try:
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(
            None, lambda: supabase.table("documents").select("doc_type").limit(1).execute())
        _DOCS_DOC_TYPE_COL = True
    except Exception as e:  # noqa: BLE001
        # Only CACHE False on a genuine "column does not exist" (Postgres 42703);
        # a transient REST/network blip returns False WITHOUT caching, so the next
        # upload re-probes (and picks up the column once the migration lands, no
        # restart needed).
        msg = str(e).lower()
        if "42703" in msg or "does not exist" in msg or "doc_type" in msg:
            _DOCS_DOC_TYPE_COL = False
        return False
    return _DOCS_DOC_TYPE_COL


@app.post("/api/documents/upload")
async def upload_document(request_body: dict):
    """Upload a document, chunk it, generate embeddings, and store in Supabase.

    Request body:
        - content (str): The text content of the document
        - name (str): Document name/filename
        - project_id (str, optional): Project ID to associate the document with
        - chat_id (str, optional): Chat ID to link the document to (for chat-scoped access)

    At least one of project_id or chat_id must be provided.
    """
    if not supabase:
        raise HTTPException(status_code=500, detail="Supabase not configured")

    content = request_body.get("content", "")
    name = request_body.get("name", "Untitled Document")
    project_id = request_body.get("project_id")
    chat_id = request_body.get("chat_id")
    doc_type = (request_body.get("doc_type") or "").strip() or None
    file_b64 = request_body.get("file_b64")
    # `filename` carries the original extension for format detection; `name` is the
    # admin-chosen display title. Fall back to name when no separate filename given.
    file_name = request_body.get("filename") or name

    # PDF/DOCX support: when no plain text is supplied but a base64 file is, extract
    # the text server-side. Extraction is CPU-bound and can be slow on adversarial
    # input, so run it in a thread (never block the event loop) with a wall-clock
    # timeout. Plain-text uploads still pass `content` directly.
    if not content and file_b64:
        try:
            _loop = asyncio.get_event_loop()
            content = await asyncio.wait_for(
                _loop.run_in_executor(None, _extract_text_from_file, file_b64, file_name),
                timeout=120)
        except asyncio.TimeoutError:
            raise HTTPException(status_code=400, detail=f"Text extraction from '{file_name}' timed out")
        except HTTPException:
            raise
        except Exception as _e:  # noqa: BLE001
            raise HTTPException(status_code=400, detail=f"Could not extract text from '{file_name}': {_e}")

    if not content:
        raise HTTPException(status_code=400, detail="Document content is required (text, or a base64 file in file_b64)")
    if not project_id and not chat_id:
        raise HTTPException(status_code=400, detail="At least one of project_id or chat_id is required")

    try:
        import httpx

        doc_id = str(uuid.uuid4())
        # doc_type: store natively once the column exists; otherwise encode it into
        # the name so it stays visible/filterable in retrieval results pre-migration.
        if doc_type and not (await _documents_has_doc_type()) and not name.startswith("["):
            name = f"[{doc_type}] {name}"
        doc_row = {
            "id": doc_id,
            "name": name,
            "file_path": f"uploads/{doc_id}_{name}",
            "content": content,
        }
        if project_id:
            doc_row["project_id"] = project_id
        if doc_type and await _documents_has_doc_type():
            doc_row["doc_type"] = doc_type

        loop = asyncio.get_event_loop()
        await loop.run_in_executor(
            None,
            lambda: supabase.table("documents").insert(doc_row).execute()
        )

        chunk_size = 1000
        overlap = 200
        chunks = []
        start = 0
        while start < len(content):
            end = start + chunk_size
            chunk_text = content[start:end]
            if chunk_text.strip():
                chunks.append(chunk_text)
            start += chunk_size - overlap

        print(f"[DOC UPLOAD] Document '{name}' split into {len(chunks)} chunks")

        openai_key = config.OPENAI_API_KEY
        if not openai_key:
            raise HTTPException(status_code=500, detail="OpenAI API key not configured (needed for embeddings)")

        batch_size = 50
        all_embeddings = []
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i:i + batch_size]
            async with httpx.AsyncClient(timeout=60.0) as client:
                resp = await client.post(
                    "https://api.openai.com/v1/embeddings",
                    headers={
                        "Authorization": f"Bearer {openai_key}",
                        "Content-Type": "application/json",
                    },
                    json={"model": "text-embedding-ada-002", "input": batch},
                )
                resp.raise_for_status()
                data = resp.json()
                batch_embeddings = [item["embedding"] for item in data["data"]]
                all_embeddings.extend(batch_embeddings)

        chunk_rows = []
        for i, (chunk_text, embedding) in enumerate(zip(chunks, all_embeddings)):
            row = {
                "id": str(uuid.uuid4()),
                "document_id": doc_id,
                "content": chunk_text,
                "embedding": embedding,
            }
            if project_id:
                row["project_id"] = project_id
            chunk_rows.append(row)

        insert_batch_size = 25
        for i in range(0, len(chunk_rows), insert_batch_size):
            batch = chunk_rows[i:i + insert_batch_size]
            await loop.run_in_executor(
                None,
                lambda b=batch: supabase.table("document_chunks").insert(b).execute()
            )

        if chat_id:
            try:
                await loop.run_in_executor(
                    None,
                    lambda: supabase.table("chat_documents").insert({
                        "chat_id": chat_id,
                        "document_id": doc_id,
                    }).execute()
                )
            except Exception as e:
                print(f"[DOC UPLOAD] Warning: Could not link to chat_documents: {e}")

        print(f"[DOC UPLOAD] Successfully uploaded '{name}': {len(chunks)} chunks with embeddings")

        return JSONResponse(content={
            "status": "success",
            "document_id": doc_id,
            "name": name,
            "chunks_count": len(chunks),
            "project_id": project_id,
            "chat_id": chat_id,
        })

    except HTTPException:
        raise
    except Exception as e:
        print(f"[DOC UPLOAD] Error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Document upload failed: {str(e)}")


@app.get("/api/documents")
async def list_documents(project_id: str = None, chat_id: str = None):
    """List documents accessible for a given project and/or chat."""
    if not supabase:
        raise HTTPException(status_code=500, detail="Supabase not configured")
    if not project_id and not chat_id:
        raise HTTPException(status_code=400, detail="Provide project_id or chat_id")

    try:
        loop = asyncio.get_event_loop()
        all_docs = []

        if project_id:
            result = await loop.run_in_executor(
                None,
                lambda: supabase.table("documents")
                    .select("id,name,file_path,created_at")
                    .eq("project_id", project_id)
                    .execute()
            )
            all_docs.extend(result.data or [])

        if chat_id:
            try:
                chat_doc_result = await loop.run_in_executor(
                    None,
                    lambda: supabase.table("chat_documents")
                        .select("document_id")
                        .eq("chat_id", chat_id)
                        .execute()
                )
                chat_doc_ids = [r["document_id"] for r in (chat_doc_result.data or [])]
                if chat_doc_ids:
                    existing_ids = {d["id"] for d in all_docs}
                    new_ids = [did for did in chat_doc_ids if did not in existing_ids]
                    if new_ids:
                        docs_result = await loop.run_in_executor(
                            None,
                            lambda: supabase.table("documents")
                                .select("id,name,file_path,created_at")
                                .in_("id", new_ids)
                                .execute()
                        )
                        all_docs.extend(docs_result.data or [])
            except Exception:
                pass

        return JSONResponse(content={
            "documents": all_docs,
            "count": len(all_docs),
        })
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.websocket("/ws/chat")
async def websocket_chat(websocket: WebSocket):
    """WebSocket endpoint for real-time chat"""
    await websocket.accept()
    try:
        agent = await agent_manager.get_agent()
        while True:
            data = await websocket.receive_json()
            messages = data.get("messages", [])
            # Parity fix (cost task #40, 2026-05-22): the REST + async +
            # structured endpoints all run incoming history through
            # context_manager.summarize_conversation_history before astream;
            # this WS endpoint did not, so a long restored chat fed every
            # turn of historical messages to the LLM uncompressed. The
            # summarizer is internally short-circuited when under threshold,
            # so we just delegate the whole decision to it.
            try:
                messages = await context_manager.summarize_conversation_history(messages)
            except Exception as _e:
                print(f"[WS-CHAT] summarize_conversation_history skipped: {_e}")
            # Watchdog-wrapped iteration (same pattern as
            # _agent_astream_autocontinue) — prevents a hung LLM call from
            # leaving the WebSocket worker pinned indefinitely.
            _ws_stream = agent.astream(
                {"messages": messages},
                stream_mode="values",
                config={"recursion_limit": _RECURSION_LIMIT},
            )
            while True:
                try:
                    chunk = await asyncio.wait_for(
                        _ws_stream.__anext__(),
                        timeout=config.WATCHDOG_STALL_SECONDS,
                    )
                except StopAsyncIteration:
                    break
                except asyncio.TimeoutError:
                    try:
                        await _ws_stream.aclose()
                    except Exception:
                        pass
                    raise RuntimeError(
                        f"WATCHDOG: no agent chunk for "
                        f"{config.WATCHDOG_STALL_SECONDS}s — websocket run stalled."
                    )
                if "messages" in chunk and chunk["messages"]:
                    last_message = chunk["messages"][-1]
                    if hasattr(last_message, 'content'):
                        await websocket.send_json({
                            "type":
                            "message",
                            "content":
                            str(last_message.content),
                            "done":
                            False
                        })
            await websocket.send_json({
                "type": "message",
                "content": "",
                "done": True
            })
    except WebSocketDisconnect:
        print("WebSocket disconnected")
    except Exception as e:
        await websocket.send_json({"type": "error", "content": str(e)})
        await websocket.close()


@app.get("/api/config")
async def get_config():
    return {
        "model": config.MODEL,
        "mcp_config": agent_manager.mcp_config_manager.config,
        "custom_tools_dir": config.CUSTOM_TOOLS_DIR,
        "context_management": {
            "summarizer_model":
            config.SUMMARIZER_MODEL,
            "tool_response_summarize_threshold":
            config.TOOL_RESPONSE_SUMMARIZE_THRESHOLD,
            "conversation_summarize_token_threshold":
            config.CONVERSATION_SUMMARIZE_TOKEN_THRESHOLD,
            "conversation_keep_recent_messages":
            config.CONVERSATION_KEEP_RECENT_MESSAGES,
        },
        "mcp_truncation_limits": {
            "max_response_size": config.MCP_MAX_RESPONSE_SIZE,
            "max_string_length": config.MCP_MAX_STRING_LENGTH,
            "max_list_items": config.MCP_MAX_LIST_ITEMS
        }
    }


@app.post("/api/config")
async def update_config(request: ConfigRequest):
    try:
        await agent_manager.reinitialize_agent(
            instructions=request.instructions, headless=request.headless)
        return {"status": "success", "message": "Agent reinitialized"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ---------------------------------------------------------------------------
# Jarvis — a single cross-analysis agent scoped to a GLOBAL "enabled analyses"
# toggle list (jarvis_settings). It searches the enabled analyses first via the
# jarvis_* tools and only falls back to the native read-only tool catalog when the
# answer isn't in those analyses. Salesforce writes stay blocked (shared denylist).
# ---------------------------------------------------------------------------

# The editable persona/instructions shown in the Jarvis settings tab. When the
# user hasn't customised the prompt (jarvis_settings.system_prompt == ""), this is
# used as the base. It is the ONLY user-editable part — the live scope listing and
# the operating rules below are always appended so the search-analyses-first +
# read-only contract holds no matter what the user types.
DEFAULT_JARVIS_INSTRUCTIONS = (
    "You are Jarvis, a cross-analysis research assistant. Unlike the per-analysis "
    "tools, you are NOT tied to a single analysis — you can read across MANY "
    "analyses at once, but ONLY the ones enabled in Jarvis settings. Be concise and "
    "cite which analysis (by title) each fact came from."
)


def _build_jarvis_system_prompt(enabled: list, custom_instructions: Optional[str] = None) -> str:
    base = (custom_instructions or "").strip() or DEFAULT_JARVIS_INSTRUCTIONS
    if enabled:
        listing = "\n".join(
            f"  - {a.get('title') or '(untitled)'} "
            f"(analysis_id={a['id']}, status={a.get('status')})"
            for a in enabled
        )
        scope = f"You currently have access to these {len(enabled)} analyses:\n{listing}"
    else:
        scope = ("No analyses are currently enabled in Jarvis settings. Tell the "
                 "user to enable analyses in the Jarvis settings tab before you can "
                 "answer from analysis data.")
    # Non-negotiable operating rules — appended after the (possibly custom) persona
    # so tool routing and read-only safety can't be edited away from the settings tab.
    rules = (
        "OPERATING RULES (always apply, regardless of the instructions above):\n"
        "1. ALWAYS look in the enabled analyses FIRST. Use `jarvis_list_analyses` to "
        "confirm your scope, `jarvis_search` to find where something is mentioned, "
        "`jarvis_filter_rows` to find opportunities matching a condition, and "
        "`jarvis_get_cells` to read a cell's full value. These tools already span "
        "every enabled analysis — you do not need to pass an analysis_id.\n"
        "2. ONLY IF the answer is not present in the enabled analyses should you fall "
        "back to the native tools (Salesforce/CRM reads, web search, other MCP "
        "tools).\n"
        "3. You are strictly READ-ONLY. Never write to Salesforce or any external "
        "system.\n"
        "4. Cite which analysis (by title) each fact came from. If the data isn't in "
        "the enabled analyses and the native tools can't find it, say so plainly "
        "rather than guessing."
    )
    return f"{base}\n\n{scope}\n\n{rules}"


@app.get("/api/jarvis/settings")
async def jarvis_get_settings():
    """Return the global enabled-analyses list plus every analysis flagged with its
    enabled state, so the frontend settings tab can render the toggles directly."""
    import jarvis_store
    import analysis_store as store
    try:
        settings = await _aw(jarvis_store.get_settings)
        enabled_ids = settings["enabled_analysis_ids"]
        all_analyses = await _aw(store.list_analyses, limit=200)
        enabled_set = set(enabled_ids)
        items = [{
            "id": a["id"], "title": a.get("title"), "status": a.get("status"),
            "project_id": a.get("project_id"), "updated_at": a.get("updated_at"),
            "enabled": a["id"] in enabled_set,
        } for a in all_analyses]
        return {
            "enabled_analysis_ids": enabled_ids,
            "system_prompt": settings["system_prompt"],
            "default_system_prompt": DEFAULT_JARVIS_INSTRUCTIONS,
            "count": len(items),
            "analyses": items,
        }
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.put("/api/jarvis/settings")
async def jarvis_put_settings(request: JarvisSettingsRequest):
    """Update the global Jarvis settings — the enabled-analyses toggles and/or the
    editable system prompt. Only the fields present in the body are changed; send
    system_prompt="" to reset it to the backend default."""
    import jarvis_store
    try:
        saved = await _aw(jarvis_store.set_settings,
                          enabled_analysis_ids=request.enabled_analysis_ids,
                          system_prompt=request.system_prompt)
        return {
            "enabled_analysis_ids": saved["enabled_analysis_ids"],
            "system_prompt": saved["system_prompt"],
            "count": len(saved["enabled_analysis_ids"]),
        }
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/jarvis/chat/async")
async def jarvis_chat_async(request: ChatRequest):
    """Jarvis chat — mirrors /api/chat/async but injects the Jarvis system prompt
    (scoped to the enabled analyses). Runs the agent inside a keepalive stream and
    persists all events to Supabase; the client reads results via Supabase realtime."""
    import jarvis_store
    chat_id = request.chat_id or str(uuid.uuid4())
    _reserve_run_slot(chat_id)
    try:
        if not supabase:
            raise HTTPException(
                status_code=500,
                detail="Supabase not configured - async mode requires Supabase")

        active_count = sum(1 for t in _running_tasks.values() if t and not t.done())
        if active_count >= config.MAX_CONCURRENT_SESSIONS:
            global _sessions_rejected
            _sessions_rejected += 1
            raise HTTPException(
                status_code=503,
                detail=f"Server at capacity ({config.MAX_CONCURRENT_SESSIONS} concurrent sessions). Please try again shortly.")

        print(f"\n{'='*60}")
        print(f"[API REQUEST] /api/jarvis/chat/async chat_id={chat_id}, model={request.model}")

        settings = await asyncio.to_thread(jarvis_store.get_settings)
        enabled = await asyncio.to_thread(jarvis_store.get_enabled_analyses)
        jarvis_prompt = _build_jarvis_system_prompt(enabled, settings.get("system_prompt"))
        if request.system_prompt:
            jarvis_prompt = jarvis_prompt + "\n\nAdditional user instructions:\n" + request.system_prompt

        _current_chat_id.set(chat_id)
        await agent_manager.reinitialize_agent(
            instructions=jarvis_prompt, model=request.model, headless=request.headless)
        agent = await agent_manager.get_agent()

        messages = [_build_message_content(msg) for msg in request.messages]
        total_chars = sum(len(str(m.get("content", ""))) for m in messages)
        if total_chars // 4 > config.CONVERSATION_SUMMARIZE_TOKEN_THRESHOLD:
            messages = await context_manager.summarize_conversation_history(messages)
        elif len(messages) > 10:
            messages = messages[-10:]

        chunk_queue = asyncio.Queue()
        stream_done = asyncio.Event()

        async def consume_stream():
            try:
                await run_agent_and_save(chat_id, messages, agent, request.model,
                                         project_id=request.project_id)
            except Exception as e:  # noqa: BLE001
                await chunk_queue.put(e)
            finally:
                stream_done.set()

        consumer_task = asyncio.create_task(consume_stream())
        _running_tasks[chat_id] = consumer_task
        _session_start_times[chat_id] = asyncio.get_event_loop().time()

        def _cleanup_session(t, cid=chat_id):
            _running_tasks.pop(cid, None)
            _session_start_times.pop(cid, None)
            _supabase_seq_counters.pop(cid, None)
            _dedupe_completed.pop(cid, None)
            _dedupe_inflight.pop(cid, None)
            _approved_campaigns_cache.pop(cid, None)
            _starting_chats.discard(cid)

        consumer_task.add_done_callback(_cleanup_session)

        async def run_and_keepalive():
            KEEPALIVE_INTERVAL = 5
            yield f"data: {json.dumps({'type': 'chat_id', 'chat_id': chat_id})}\n\n"
            while not stream_done.is_set():
                try:
                    msg = await asyncio.wait_for(chunk_queue.get(),
                                                 timeout=KEEPALIVE_INTERVAL)
                    if isinstance(msg, Exception):
                        yield f"data: {json.dumps({'type': 'error', 'content': str(msg)})}\n\n"
                        break
                except asyncio.TimeoutError:
                    yield f"data: {json.dumps({'type': 'ping'})}\n\n"
            await consumer_task
            yield f"data: {json.dumps({'type': 'done', 'chat_id': chat_id})}\n\n"

        return StreamingResponse(run_and_keepalive(),
                                 media_type="text/event-stream",
                                 headers={
                                     "X-Accel-Buffering": "no",
                                     "Cache-Control": "no-cache",
                                     "Connection": "keep-alive",
                                 })

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"[JARVIS ERROR] {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        _release_run_slot(chat_id)


@app.get("/api/mcp/servers")
async def get_mcp_servers():
    return agent_manager.mcp_config_manager.config.get("mcp_servers", {})


@app.post("/api/mcp/servers/{server_name}")
async def add_mcp_server(server_name: str, server_config: MCPServerConfig):
    try:
        current_config = agent_manager.mcp_config_manager.config
        if "mcp_servers" not in current_config:
            current_config["mcp_servers"] = {}
        current_config["mcp_servers"][server_name] = {
            "command": server_config.command,
            "args": server_config.args,
            "env": server_config.env or {},
            "enabled": server_config.enabled
        }
        agent_manager.mcp_config_manager.save_config(current_config)
        if server_config.enabled:
            await agent_manager.reinitialize_agent()
        return {
            "status": "success",
            "message": f"MCP server '{server_name}' configured"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/api/mcp/servers/{server_name}")
async def delete_mcp_server(server_name: str):
    try:
        current_config = agent_manager.mcp_config_manager.config
        if "mcp_servers" in current_config and server_name in current_config[
                "mcp_servers"]:
            del current_config["mcp_servers"][server_name]
            agent_manager.mcp_config_manager.save_config(current_config)
            await agent_manager.reinitialize_agent()
            return {
                "status": "success",
                "message": f"MCP server '{server_name}' deleted"
            }
        else:
            raise HTTPException(status_code=404, detail="Server not found")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/tools")
async def list_tools():
    agent = await agent_manager.get_agent()
    tools_info = []
    try:
        nodes = agent.get_graph().nodes
        for node_name, node in nodes.items():
            if hasattr(node, 'tools'):
                for t in node.tools:
                    tools_info.append({
                        "name": t.name,
                        "description": t.description,
                    })
    except Exception:
        pass
    if not tools_info:
        enabled = agent_manager.mcp_config_manager.get_enabled_servers()
        tools_info = [
            {
                "name": "web_search",
                "description": "Search the web using DuckDuckGo"
            },
            {
                "name": "get_current_time",
                "description": "Get the current date and time"
            },
        ]
        for server_name in enabled:
            tools_info.append({
                "name": f"[{server_name}]",
                "description": f"MCP server with multiple tools"
            })
    return {"tools": tools_info}


@app.get("/api/health")
async def health_check():
    return {
        "status": "healthy",
        "agent_initialized": agent_manager.agent is not None,
        "model": config.MODEL,
        "mcp_tools_loaded": agent_manager.mcp_tools_loaded,
        "mcp_servers": agent_manager.mcp_loading_status,
        "concurrency": {
            "active_sessions": sum(1 for t in _running_tasks.values() if t and not t.done()),
            "max_sessions": config.MAX_CONCURRENT_SESSIONS,
            "sessions_rejected": _sessions_rejected,
        },
        "context_management": {
            "summarizer_model":
            config.SUMMARIZER_MODEL,
            "tool_response_summarize_threshold":
            config.TOOL_RESPONSE_SUMMARIZE_THRESHOLD,
            "conversation_summarize_token_threshold":
            config.CONVERSATION_SUMMARIZE_TOKEN_THRESHOLD,
        },
        "mcp_truncation_limits": {
            "max_response_size": config.MCP_MAX_RESPONSE_SIZE,
            "max_string_length": config.MCP_MAX_STRING_LENGTH,
            "max_list_items": config.MCP_MAX_LIST_ITEMS
        }
    }


async def _run_sf_pull_and_cache(
    mid: str,
    muid: str,
    opp_id,
    acct_id,
    contacts,
    meeting_subject=None,
    meeting_start_at=None,
):
    """Deterministic SOQL pull + cache refresh for one Avoma meeting.

    Re-introduces the deterministic Salesforce refresh removed from the webhook
    on 2026-05-28 (Salesforce CDC is not enabled for Opportunity, so the cache
    tables would otherwise go stale). For the resolved opp/account/contact IDs:
      1. Runs `sf_meeting_enricher.pull_structured_sf_data` (no gpt-4o-mini
         summaries — structured data only).
      2. Persists the three `*_data` columns + `sf_pull_status` /
         `sf_pull_duration_ms` onto the avoma_event_reports row.
      3. Calls `cache_sync.update_cache_from_report` to refresh
         opportunity_cache / meeting_cache / field_history_cache.

    Marks `sf_pull_status=running` first, then completed/failed. Never raises —
    failures are caught and recorded so the rest of the pipeline survives.
    """
    if supabase is None:
        return
    loop = asyncio.get_running_loop()
    import sf_meeting_enricher as enricher
    import cache_sync
    try:
        await loop.run_in_executor(
            None,
            lambda: supabase.table("avoma_event_reports").update(
                {"sf_pull_status": "running"}
            ).eq("message_id", mid).execute(),
        )
        print(f"[SF-PULL] ▶ msg={mid} opp={opp_id} acct={acct_id} contacts={len(contacts or [])}", flush=True)
        pull = await loop.run_in_executor(
            None,
            lambda: enricher.pull_structured_sf_data(opp_id, acct_id, contacts or []),
        )
        await loop.run_in_executor(
            None,
            lambda: supabase.table("avoma_event_reports").update({
                "deal_health_data":      pull["deal_health_data"],
                "account_briefing_data": pull["account_briefing_data"],
                "full_snapshot_data":    pull["full_snapshot_data"],
                "sf_pull_status":        pull["status"],
                "sf_pull_duration_ms":   pull["pull_duration_ms"],
            }).eq("message_id", mid).execute(),
        )
        print(
            f"[SF-PULL] ✅ msg={mid} status={pull['status']} dur={pull['pull_duration_ms']}ms "
            f"err={pull.get('error')}",
            flush=True,
        )

        # Look up the report row id for latest_report_id linkage in the cache.
        rid = None
        try:
            r = await loop.run_in_executor(
                None,
                lambda: supabase.table("avoma_event_reports")
                    .select("id").eq("message_id", mid).limit(1).execute(),
            )
            rid = (r.data or [{}])[0].get("id")
        except Exception:
            pass

        report = {
            "meeting_uuid":          muid,
            "meeting_subject":       meeting_subject,
            "meeting_start_at":      meeting_start_at,
            "sf_opportunity_id":     opp_id,
            "sf_account_id":         acct_id,
            "deal_health_data":      pull["deal_health_data"],
            "account_briefing_data": pull["account_briefing_data"],
            "full_snapshot_data":    pull["full_snapshot_data"],
        }
        cache_out = await loop.run_in_executor(
            None,
            lambda: cache_sync.update_cache_from_report(supabase, report, report_id=rid),
        )
        print(
            f"[CACHE-SYNC] msg={mid} opp_upserts={cache_out['opp_upserts']} "
            f"meeting_inserts={cache_out['meeting_inserts']} "
            f"field_history_inserts={cache_out['field_history_inserts']} "
            f"errors={cache_out['errors']}",
            flush=True,
        )
    except Exception as e:
        print(f"[SF-PULL] ❌ msg={mid} deterministic pull failed: {e}", flush=True)
        try:
            await loop.run_in_executor(
                None,
                lambda: supabase.table("avoma_event_reports").update(
                    {"sf_pull_status": f"failed: {str(e)[:200]}"}
                ).eq("message_id", mid).execute(),
            )
        except Exception:
            pass


# ---------------------------------------------------------------------------
# SNS webhook hardening (added 2026-05-29)
# ---------------------------------------------------------------------------
import base64 as _sns_b64
from urllib.parse import urlparse as _sns_urlparse

# Topic allowlist — only messages from these TopicArns are processed. Override
# via SNS_ALLOWED_TOPIC_ARNS (comma-separated) e.g. to add a prod topic.
_SNS_ALLOWED_TOPIC_ARNS = {
    a.strip()
    for a in os.environ.get(
        "SNS_ALLOWED_TOPIC_ARNS",
        "arn:aws:sns:us-east-1:385817851343:avoma-meeting-events-dev",
    ).split(",")
    if a.strip()
}
_SNS_MAX_MESSAGE_AGE_MINUTES = int(os.environ.get("SNS_MAX_MESSAGE_AGE_MINUTES", "15"))
# Defense-in-depth config validators (NOT secrets). Cross-check the TopicArn's
# region + account, and pin the signing-cert host to the allowed region(s).
_SNS_ALLOWED_REGIONS = {
    r.strip()
    for r in os.environ.get("SNS_ALLOWED_REGIONS", "us-east-1").split(",")
    if r.strip()
}
_SNS_ALLOWED_ACCOUNT_IDS = {
    a.strip()
    for a in os.environ.get("SNS_ALLOWED_ACCOUNT_IDS", "385817851343").split(",")
    if a.strip()
}
# Signature verification is ON by default (fail-closed). Set SNS_VERIFY_SIGNATURE=false
# to disable (NOT recommended — TopicArn alone is spoofable in the request body).
_SNS_VERIFY_SIGNATURE = os.environ.get("SNS_VERIFY_SIGNATURE", "true").strip().lower() not in (
    "0", "false", "no", "off",
)
_sns_cert_cache: Dict[str, Any] = {}


# SNS publishes its signing cert and confirmation endpoint only on
# sns.<region>.amazonaws.com. Pinning to that exact pattern (not any
# *.amazonaws.com host) prevents both SSRF and a signature-bypass where an
# attacker hosts a self-signed cert on some other amazonaws subdomain.
_SNS_HOST_RE = re.compile(r"^sns\.([a-z0-9-]+)\.amazonaws\.com$")


def _sns_host_allowed(url: str) -> bool:
    """True only for https URLs whose host is sns.<region>.amazonaws.com where
    <region> is in _SNS_ALLOWED_REGIONS.

    Blocks SSRF (internal addresses like 169.254.169.254, arbitrary hosts),
    a forged-cert bypass via non-SNS amazonaws hosts, AND cert fetches from
    unexpected regions.
    """
    try:
        p = _sns_urlparse(url or "")
        host = (p.hostname or "").lower()
        m = _SNS_HOST_RE.match(host)
        return p.scheme == "https" and bool(m) and m.group(1) in _SNS_ALLOWED_REGIONS
    except Exception:
        return False


def _sns_arn_region_account(arn: str):
    """Parse (region, account_id) from an SNS TopicArn.

    Format: arn:aws:sns:<region>:<account-id>:<topic-name>. Returns
    (None, None) on anything that isn't a well-formed SNS ARN.
    """
    parts = (arn or "").split(":")
    if len(parts) >= 6 and parts[0] == "arn" and parts[2] == "sns":
        return parts[3], parts[4]
    return None, None


def _sns_is_message_fresh(timestamp_str: str) -> bool:
    """Reject replays of captured messages older than the configured window."""
    from datetime import timezone as _tz, timedelta as _td
    try:
        msg_time = datetime.fromisoformat((timestamp_str or "").replace("Z", "+00:00"))
        if msg_time.tzinfo is None:
            msg_time = msg_time.replace(tzinfo=_tz.utc)
        age = datetime.now(_tz.utc) - msg_time
        return _td(minutes=-5) <= age < _td(minutes=_SNS_MAX_MESSAGE_AGE_MINUTES)
    except Exception:
        return False


def _sns_canonical_string(body: dict) -> Optional[str]:
    """Build the exact string AWS signed, per SNS message type.

    Field set + order is defined by AWS:
    https://docs.aws.amazon.com/sns/latest/dg/sns-verify-signature-of-message.html
    """
    t = body.get("Type")
    if t == "Notification":
        keys = ["Message", "MessageId"]
        if "Subject" in body:
            keys.append("Subject")
        keys += ["Timestamp", "TopicArn", "Type"]
    elif t in ("SubscriptionConfirmation", "UnsubscribeConfirmation"):
        keys = ["Message", "MessageId", "SubscribeURL", "Timestamp", "Token", "TopicArn", "Type"]
    else:
        return None
    parts = []
    for k in keys:
        if k not in body:
            return None
        parts.append(k)
        parts.append(str(body[k]))
    return "\n".join(parts) + "\n"


async def _sns_verify_signature(body: dict) -> bool:
    """Verify the AWS SNS message signature (SignatureVersion 1=SHA1, 2=SHA256).

    Fetches the signing certificate from SigningCertURL (host-validated to
    *.amazonaws.com, cached per-URL) and checks the RSA/PKCS1v15 signature over
    the canonical string. Returns False on any failure (fail-closed).
    """
    from cryptography.hazmat.primitives import hashes as _hashes
    from cryptography.hazmat.primitives.asymmetric import padding as _padding
    from cryptography.x509 import load_pem_x509_certificate as _load_cert
    try:
        cert_url = body.get("SigningCertURL") or body.get("SigningCertUrl")
        if not _sns_host_allowed(cert_url):
            print(f"[SNS-WEBHOOK] ⛔ rejected SigningCertURL host: {cert_url}", flush=True)
            return False
        canonical = _sns_canonical_string(body)
        if canonical is None:
            print("[SNS-WEBHOOK] ⛔ could not build canonical string", flush=True)
            return False
        public_key = _sns_cert_cache.get(cert_url)
        if public_key is None:
            async with httpx.AsyncClient(timeout=5.0) as client:
                r = await client.get(cert_url)
                r.raise_for_status()
            public_key = _load_cert(r.content).public_key()
            _sns_cert_cache[cert_url] = public_key
        sigver = str(body.get("SignatureVersion", "1"))
        if sigver not in ("1", "2"):
            print(f"[SNS-WEBHOOK] ⛔ unsupported SignatureVersion: {sigver!r}", flush=True)
            return False
        sig = _sns_b64.b64decode(body["Signature"])
        algo = _hashes.SHA256() if sigver == "2" else _hashes.SHA1()
        public_key.verify(sig, canonical.encode("utf-8"), _padding.PKCS1v15(), algo)
        return True
    except Exception as e:
        print(f"[SNS-WEBHOOK] ⛔ signature verification failed: {e}", flush=True)
        return False


@app.post("/webhook")
async def sns_webhook(request: Request):
    """
    AWS SNS webhook endpoint.

    Handles three SNS message types per
    https://docs.aws.amazon.com/sns/latest/dg/sns-message-and-json-formats.html:

      - SubscriptionConfirmation: auto-confirm by GETting the SubscribeURL.
        Required once per subscription, otherwise the topic stays in
        PendingConfirmation and never delivers.
      - Notification: log + persist raw envelope to mcp_output/sns_*.json for
        forensics. No downstream dispatch yet — wire that up once the
        intended behavior per topic is decided.
      - UnsubscribeConfirmation: log only.

    Hardened (2026-05-29): every message is gated by (1) a TopicArn allowlist,
    (2) a freshness window, and (3) AWS SNS signature verification (fail-closed)
    before ANY side effect (disk write, SubscribeURL GET, enrichment). This
    closes the prior SSRF + spoofing exposure on this public endpoint.
    """
    import json as _json
    from datetime import datetime as _dt
    try:
        raw = await request.body()
        body = _json.loads(raw.decode("utf-8") or "{}")
    except Exception as e:
        print(f"[SNS-WEBHOOK] ❌ bad body: {e}", flush=True)
        return JSONResponse({"error": "invalid json"}, status_code=400)

    msg_type = body.get("Type") or request.headers.get("x-amz-sns-message-type", "")
    topic_arn = body.get("TopicArn", "")
    msg_id = body.get("MessageId", "")
    print(f"[SNS-WEBHOOK] ⬅ Type={msg_type} TopicArn={topic_arn} MessageId={msg_id}", flush=True)

    # --- Hardening gates (run BEFORE any side effect) ---
    # 1. TopicArn allowlist + region/account cross-check (defense-in-depth).
    if topic_arn not in _SNS_ALLOWED_TOPIC_ARNS:
        print(f"[SNS-WEBHOOK] ⛔ rejected unauthorized topic: {topic_arn!r}", flush=True)
        return JSONResponse({"error": "unauthorized topic"}, status_code=403)
    _arn_region, _arn_account = _sns_arn_region_account(topic_arn)
    if _arn_region not in _SNS_ALLOWED_REGIONS:
        print(f"[SNS-WEBHOOK] ⛔ rejected unauthorized region: {_arn_region!r}", flush=True)
        return JSONResponse({"error": "unauthorized region"}, status_code=403)
    if _arn_account not in _SNS_ALLOWED_ACCOUNT_IDS:
        print(f"[SNS-WEBHOOK] ⛔ rejected unauthorized account: {_arn_account!r}", flush=True)
        return JSONResponse({"error": "unauthorized account"}, status_code=403)
    # 2. Freshness (anti-replay).
    if not _sns_is_message_fresh(body.get("Timestamp", "")):
        print(f"[SNS-WEBHOOK] ⛔ rejected stale message: id={msg_id} ts={body.get('Timestamp')!r}", flush=True)
        return JSONResponse({"error": "message too old"}, status_code=403)
    # 3. Signature verification (fail-closed unless explicitly disabled).
    if _SNS_VERIFY_SIGNATURE:
        if not await _sns_verify_signature(body):
            print(f"[SNS-WEBHOOK] ⛔ rejected invalid signature: id={msg_id}", flush=True)
            return JSONResponse({"error": "invalid signature"}, status_code=403)
    else:
        print("[SNS-WEBHOOK] ⚠ signature verification DISABLED (SNS_VERIFY_SIGNATURE=false)", flush=True)

    try:
        os.makedirs("mcp_output", exist_ok=True)
        ts = _dt.utcnow().strftime("%Y%m%d_%H%M%S")
        fname = f"mcp_output/sns_{msg_type or 'unknown'}_{ts}_{msg_id[:8] or 'noid'}.json"
        with open(fname, "w") as f:
            _json.dump(body, f, indent=2)
        print(f"[SNS-WEBHOOK] 💾 raw saved to {fname}", flush=True)
    except Exception as e:
        print(f"[SNS-WEBHOOK] ⚠ failed to persist raw: {e}", flush=True)

    if msg_type == "SubscriptionConfirmation":
        subscribe_url = body.get("SubscribeURL", "")
        if not subscribe_url:
            print("[SNS-WEBHOOK] ❌ SubscriptionConfirmation missing SubscribeURL", flush=True)
            return JSONResponse({"error": "missing SubscribeURL"}, status_code=400)
        # SSRF guard: only GET https://*.amazonaws.com (defense-in-depth even if
        # signature verification is disabled).
        if not _sns_host_allowed(subscribe_url):
            print(f"[SNS-WEBHOOK] ⛔ rejected SubscribeURL host: {subscribe_url}", flush=True)
            return JSONResponse({"error": "invalid SubscribeURL"}, status_code=400)
        try:
            async with httpx.AsyncClient(timeout=15.0) as client:
                resp = await client.get(subscribe_url)
            print(f"[SNS-WEBHOOK] ✅ confirmed subscription | http={resp.status_code} | topic={topic_arn}", flush=True)
            return {"status": "confirmed", "http": resp.status_code, "topic": topic_arn}
        except Exception as e:
            print(f"[SNS-WEBHOOK] ❌ confirmation GET failed: {e}", flush=True)
            return JSONResponse({"error": f"confirm failed: {e}"}, status_code=502)

    if msg_type == "Notification":
        subject = body.get("Subject", "")
        message_preview = str(body.get("Message", ""))[:300]
        print(f"[SNS-WEBHOOK] 📬 Notification | Subject={subject!r} | Message[:300]={message_preview!r}", flush=True)

        # Parse Avoma payload: Message is a JSON string with meeting_id at top level.
        meeting_uuid = None
        try:
            inner = body.get("Message")
            if isinstance(inner, str) and inner.strip().startswith("{"):
                inner = _json.loads(inner)
            if isinstance(inner, dict):
                meeting_uuid = (
                    inner.get("meeting_id")
                    or inner.get("uuid")
                    or (inner.get("full_data") or {}).get("uuid")
                )
        except Exception as e:
            print(f"[SNS-WEBHOOK] ⚠ failed to parse Message JSON: {e}", flush=True)

        if not meeting_uuid:
            print(f"[SNS-WEBHOOK] ⚠ no meeting_uuid in Notification — skipping enrichment", flush=True)
            return {"status": "received_no_meeting", "message_id": msg_id}

        # Idempotency: check if we've already processed this SNS MessageId.
        # If status is terminal (completed/completed_with_errors/no_sf_links),
        # return 200 without re-enriching — prevents duplicate SF queries and
        # LLM spend on SNS redeliveries. Re-process only failed or stuck-pending rows.
        if supabase is not None:
            try:
                existing = await asyncio.get_running_loop().run_in_executor(
                    None,
                    lambda: supabase.table("avoma_event_reports")
                                    .select("status,processed_at")
                                    .eq("message_id", msg_id).limit(1).execute(),
                )
                rows = existing.data or []
                if rows and rows[0].get("status") in (
                    "completed",
                    "completed_with_errors",
                    "no_sf_links",
                    "completed_no_opportunity",
                ):
                    print(f"[SNS-WEBHOOK] 🔁 duplicate MessageId={msg_id} (status={rows[0]['status']}) — skipping enrichment", flush=True)
                    return {"status": "duplicate", "message_id": msg_id, "prior_status": rows[0]["status"]}
            except Exception as e:
                print(f"[SNS-WEBHOOK] ⚠ idempotency check failed (proceeding): {e}", flush=True)

        # Insert a pending row immediately for audit trail, then run SF
        # enrichment + summarization in the background. The handler returns
        # 200 fast so SNS doesn't retry.
        if supabase is not None:
            try:
                await asyncio.get_running_loop().run_in_executor(
                    None,
                    lambda: supabase.table("avoma_event_reports").upsert(
                        {
                            "message_id": msg_id,
                            "meeting_uuid": meeting_uuid,
                            "status": "pending",
                            "raw_sns_envelope": body,
                        },
                        on_conflict="message_id",
                    ).execute(),
                )
            except Exception as e:
                print(f"[SNS-WEBHOOK] ⚠ pending-row insert failed: {e}", flush=True)

        async def _enrich_bg(mid: str, muid: str):
            """Lean webhook pipeline:
              1. Resolve SF ids from the Avoma meeting (one HTTP call).
              2. Persist meeting metadata + sf links to avoma_event_reports.
              3. If we have an Opportunity id, run the opportunity-analysis
                 agent (Sonnet, scoped to Salesforce + Avoma MCP tools) and
                 store the JSON record on the same row.

            The legacy 3-tier enrichment (deal_health / account_briefing /
            full_snapshot) + per-tier gpt-4o-mini summaries + cache_sync hook
            were removed 2026-05-28 on user request: the new agent is the
            single source of truth for per-meeting analysis. cache_sync still
            runs via the bulk-import + /cron/sync-sf-to-cache endpoints.
            """
            from datetime import datetime as __dt
            loop = asyncio.get_running_loop()
            try:
                import sf_meeting_enricher as enricher
                t0 = time.time()
                print(f"[SNS-ENRICH] ▶ msg={mid} meeting={muid} resolving sf ids", flush=True)
                ids = await loop.run_in_executor(None, enricher.extract_sf_ids_from_meeting, muid)
                meeting = ids.get("meeting") or {}
                opp_id     = ids.get("opportunity_id")
                acct_id    = ids.get("account_id")
                contacts   = ids.get("contact_ids") or []
                dur = int((time.time() - t0) * 1000)
                print(
                    f"[SNS-ENRICH] ✅ msg={mid} opp={opp_id} acct={acct_id} "
                    f"contacts={len(contacts)} dur={dur}ms",
                    flush=True,
                )

                status = "no_sf_links" if not (opp_id or acct_id or contacts) else (
                    "pending_analysis" if opp_id else "completed_no_opportunity"
                )
                if supabase is not None:
                    row = {
                        "meeting_uuid":      muid,
                        "meeting_subject":   meeting.get("subject"),
                        "meeting_start_at":  meeting.get("start_at"),
                        "sf_opportunity_id": opp_id,
                        "sf_account_id":     acct_id,
                        "sf_contact_ids":    contacts,
                        "pull_duration_ms":  dur,
                        "status":            status,
                        "processed_at":      __dt.utcnow().isoformat(),
                    }
                    await loop.run_in_executor(
                        None,
                        lambda: supabase.table("avoma_event_reports").update(row).eq("message_id", mid).execute(),
                    )
                    print(f"[SNS-ENRICH] 💾 msg={mid} metadata persisted (status={status})", flush=True)

                # --- deterministic SF pull + cache refresh ---
                # Runs whenever we resolved any SF link (opp/account/contacts),
                # independent of the analysis-agent path below. CDC is off for
                # Opportunity, so this is what keeps the cache tables fresh.
                if opp_id or acct_id or contacts:
                    await _run_sf_pull_and_cache(
                        mid, muid, opp_id, acct_id, contacts,
                        meeting_subject=meeting.get("subject"),
                        meeting_start_at=meeting.get("start_at"),
                    )

                # --- opportunity-analysis agent (only path now) ---
                if not opp_id:
                    print(f"[OPP-ANALYZER] ⏭ msg={mid} skipped (no sf_opportunity_id)", flush=True)
                    return
                try:
                    import opportunity_analyzer
                    print(f"[OPP-ANALYZER] ▶ msg={mid} opp={opp_id} starting", flush=True)
                    an = await opportunity_analyzer.analyze_opportunity(agent_manager, opp_id)
                    print(
                        f"[OPP-ANALYZER] ✅ msg={mid} status={an['status']} "
                        f"dur={an['duration_ms']}ms err={an.get('error')}",
                        flush=True,
                    )
                    if supabase is not None:
                        await loop.run_in_executor(
                            None,
                            lambda: supabase.table("avoma_event_reports").update({
                                "opportunity_analysis_data":   an["data"],
                                "opportunity_analysis_status": (
                                    an["status"] if an["status"] == "completed"
                                    else f"{an['status']}: {an.get('error') or 'no detail'}"
                                ),
                                "status":                      "completed" if an["status"] == "completed" else "completed_with_errors",
                            }).eq("message_id", mid).execute(),
                        )
                except Exception as ae:
                    print(
                        f"[OPP-ANALYZER] ⚠ msg={mid} opp={opp_id} analyzer failed: {ae}",
                        flush=True,
                    )
                    if supabase is not None:
                        try:
                            await loop.run_in_executor(
                                None,
                                lambda: supabase.table("avoma_event_reports").update({
                                    "opportunity_analysis_status": f"failed: {str(ae)[:200]}",
                                    "status":                      "completed_with_errors",
                                }).eq("message_id", mid).execute(),
                            )
                        except Exception:
                            pass
            except Exception as e:
                print(f"[SNS-ENRICH] ❌ msg={mid} background task failed: {e}", flush=True)
                if supabase is not None:
                    try:
                        await asyncio.get_running_loop().run_in_executor(
                            None,
                            lambda: supabase.table("avoma_event_reports").update(
                                {"status": "failed", "error": str(e)[:1000]}
                            ).eq("message_id", mid).execute(),
                        )
                    except Exception:
                        pass

        asyncio.create_task(_enrich_bg(msg_id, meeting_uuid))
        return {"status": "received", "message_id": msg_id, "meeting_uuid": meeting_uuid, "enrichment": "scheduled"}

    if msg_type == "UnsubscribeConfirmation":
        print(f"[SNS-WEBHOOK] 👋 UnsubscribeConfirmation | topic={topic_arn}", flush=True)
        return {"status": "unsubscribe_acknowledged"}

    print(f"[SNS-WEBHOOK] ⚠ unknown Type={msg_type!r} — ignored", flush=True)
    return {"status": "ignored", "type": msg_type}


@app.post("/api/avoma/reports/{message_id}/reanalyze")
async def reanalyze_avoma_report(message_id: str):
    """Re-run the opportunity-analysis agent for an existing report row.

    Looks up the row by message_id, takes its sf_opportunity_id, runs the
    analyzer, and patches opportunity_analysis_{data,status} on the row.
    Useful for backfilling old reports + smoke-testing the analyzer.
    """
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    loop = asyncio.get_running_loop()
    row = await loop.run_in_executor(
        None,
        lambda: supabase.table("avoma_event_reports")
            .select("message_id,sf_opportunity_id,meeting_subject")
            .eq("message_id", message_id).limit(1).execute(),
    )
    if not row.data:
        return JSONResponse({"error": f"no report for message_id={message_id}"}, status_code=404)
    opp_id = row.data[0].get("sf_opportunity_id")
    if not opp_id:
        return JSONResponse({"error": "report has no sf_opportunity_id"}, status_code=400)
    # Fire-and-forget — analyzer can take minutes; survive client disconnect.
    async def _bg():
        import opportunity_analyzer
        try:
            print(f"[OPP-ANALYZER] ▶ reanalyze msg={message_id} opp={opp_id}", flush=True)
            an = await opportunity_analyzer.analyze_opportunity(agent_manager, opp_id)
            print(
                f"[OPP-ANALYZER] ✅ reanalyze msg={message_id} status={an['status']} "
                f"dur={an['duration_ms']}ms err={an.get('error')}",
                flush=True,
            )
            await asyncio.get_running_loop().run_in_executor(
                None,
                lambda: supabase.table("avoma_event_reports").update({
                    "opportunity_analysis_data":   an["data"],
                    "opportunity_analysis_status": (
                        an["status"] if an["status"] == "completed"
                        else f"{an['status']}: {an.get('error') or 'no detail'}"
                    ),
                }).eq("message_id", message_id).execute(),
            )
        except Exception as e:
            print(f"[OPP-ANALYZER] ❌ reanalyze msg={message_id} crashed: {e}", flush=True)
            try:
                await asyncio.get_running_loop().run_in_executor(
                    None,
                    lambda: supabase.table("avoma_event_reports").update({
                        "opportunity_analysis_status": f"failed: {str(e)[:200]}",
                    }).eq("message_id", message_id).execute(),
                )
            except Exception:
                pass

    # Mark row as in-progress BEFORE kicking the background task — otherwise
    # a fast-completing task can be overwritten back to "running".
    await loop.run_in_executor(
        None,
        lambda: supabase.table("avoma_event_reports").update({
            "opportunity_analysis_status": "running",
            "opportunity_analysis_data":   None,
        }).eq("message_id", message_id).execute(),
    )
    asyncio.create_task(_bg())
    return {
        "status": "started",
        "message_id": message_id,
        "sf_opportunity_id": opp_id,
        "meeting_subject": row.data[0].get("meeting_subject"),
        "poll_url": f"/api/avoma/reports/{message_id}",
        "note": "analysis runs in background (typically 30-180s); poll the report row for opportunity_analysis_{status,data}",
    }


@app.post("/api/avoma/reports/{message_id}/refresh-sf")
async def refresh_sf_avoma_report(message_id: str):
    """Manually re-run the deterministic SF pull + cache refresh for one report.

    Looks up the report row by message_id, takes its resolved SF ids, then runs
    the same deterministic pull + cache refresh as the webhook background
    pipeline (`_run_sf_pull_and_cache`). Fire-and-forget so a curl client
    disconnect doesn't kill the task; poll the report row for `sf_pull_status`.
    """
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    loop = asyncio.get_running_loop()
    row = await loop.run_in_executor(
        None,
        lambda: supabase.table("avoma_event_reports")
            .select("message_id,meeting_uuid,meeting_subject,meeting_start_at,"
                    "sf_opportunity_id,sf_account_id,sf_contact_ids")
            .eq("message_id", message_id).limit(1).execute(),
    )
    if not row.data:
        return JSONResponse({"error": f"no report for message_id={message_id}"}, status_code=404)
    r = row.data[0]
    opp_id   = r.get("sf_opportunity_id")
    acct_id  = r.get("sf_account_id")
    contacts = r.get("sf_contact_ids") or []
    muid     = r.get("meeting_uuid")
    if not (opp_id or acct_id or contacts):
        return JSONResponse({"error": "report has no resolved SF links to pull"}, status_code=400)

    # Mark in-progress BEFORE kicking the background task so a fast-completing
    # task isn't overwritten back to running, mirroring /reanalyze.
    await loop.run_in_executor(
        None,
        lambda: supabase.table("avoma_event_reports").update(
            {"sf_pull_status": "running"}
        ).eq("message_id", message_id).execute(),
    )

    async def _bg():
        await _run_sf_pull_and_cache(
            message_id, muid, opp_id, acct_id, contacts,
            meeting_subject=r.get("meeting_subject"),
            meeting_start_at=r.get("meeting_start_at"),
        )

    asyncio.create_task(_bg())
    return {
        "status": "started",
        "message_id": message_id,
        "sf_opportunity_id": opp_id,
        "sf_account_id": acct_id,
        "sf_contact_ids": contacts,
        "meeting_subject": r.get("meeting_subject"),
        "poll_url": f"/api/avoma/reports/{message_id}",
        "note": "deterministic SF pull + cache refresh runs in background; poll the report row for sf_pull_status + the *_data columns",
    }


@app.get("/api/avoma/reports")
async def list_avoma_reports(limit: int = 50, status: Optional[str] = None):
    """Browse recent Avoma-triggered SF enrichment reports.

    Returns lightweight rows (no jsonb bodies) so the list stays cheap.
    Use /api/avoma/reports/{message_id} for the full payload.
    """
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    try:
        cols = ("id,message_id,meeting_uuid,meeting_subject,meeting_start_at,"
                "sf_opportunity_id,sf_account_id,sf_contact_ids,pull_duration_ms,"
                "sf_pull_status,sf_pull_duration_ms,"
                "status,error,created_at,processed_at")
        q = supabase.table("avoma_event_reports").select(cols).order("created_at", desc=True).limit(min(limit, 200))
        if status:
            q = q.eq("status", status)
        res = await asyncio.get_running_loop().run_in_executor(None, q.execute)
        return {"count": len(res.data or []), "reports": res.data or []}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/avoma/reports/{message_id}")
async def get_avoma_report(message_id: str):
    """Full enrichment report for one SNS MessageId, including all 3 tiers + summaries."""
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    try:
        res = await asyncio.get_running_loop().run_in_executor(
            None,
            lambda: supabase.table("avoma_event_reports").select("*").eq("message_id", message_id).limit(1).execute(),
        )
        rows = res.data or []
        if not rows:
            return JSONResponse({"error": "not found", "message_id": message_id}, status_code=404)
        return rows[0]
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


# ============================================================================
# Cache layer: opportunity_cache / meeting_cache / field_history_cache
# ============================================================================

@app.post("/api/cache/bulk_import")
async def cache_bulk_import(request: Request):
    """One-time import of deal_summaries_bulk_parallel.json into opportunity_cache.

    Body: { "path": "attached_assets/deal_summaries_bulk_parallel_1779965448034.json" }
    or empty body -> uses default path.
    """
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    try:
        body = {}
        try:
            body = await request.json()
        except Exception:
            pass
        path = body.get("path") or "attached_assets/deal_summaries_bulk_parallel_1779965448034.json"
        import cache_sync, os as _os
        if not _os.path.exists(path):
            return JSONResponse({"error": f"file not found: {path}"}, status_code=404)
        out = await asyncio.get_running_loop().run_in_executor(
            None, lambda: cache_sync.bulk_import_opportunities(supabase, path)
        )
        return out
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/cache/bulk_import_field_history")
async def cache_bulk_import_field_history(request: Request):
    """One-time bulk import of an OpportunityFieldHistory JSON dump.

    Body: { "path": "attached_assets/field_history_cache_*.json" }
    Idempotent — dedups against UNIQUE(opp,field,changed_date,old,new).
    """
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    try:
        body = {}
        try:
            body = await request.json()
        except Exception:
            pass
        path = body.get("path")
        if not path:
            import glob as _glob
            cands = sorted(_glob.glob("attached_assets/field_history_cache_*.json"))
            if not cands:
                return JSONResponse({"error": "no field_history_cache_*.json in attached_assets/"}, status_code=404)
            path = cands[-1]
        import cache_sync, os as _os
        if not _os.path.exists(path):
            return JSONResponse({"error": f"file not found: {path}"}, status_code=404)
        out = await asyncio.get_running_loop().run_in_executor(
            None, lambda: cache_sync.bulk_import_field_history(supabase, path)
        )
        return {"path": path, **out}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/cache/tenant/{opp_id}")
async def cache_tenant_state(opp_id: str, meetings_limit: int = 10):
    """Full current state for one opportunity in one parallel-fanout call.

    Returns: { opportunity, meetings[], history[], history_summary }
    All 3 caches read concurrently via asyncio.gather (latency = max, not sum).
    """
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    try:
        import cache_sync, time as _t
        t0 = _t.perf_counter()
        out = await cache_sync.get_tenant_state(supabase, opp_id, meetings_limit=meetings_limit)
        out["_latency_ms"] = round((_t.perf_counter() - t0) * 1000, 1)
        if not out.get("opportunity"):
            return JSONResponse({**out, "error": "opportunity not found in cache"}, status_code=404)
        return out
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/cron/sync-sf-to-cache")
async def cron_sync_sf(lookback_minutes: int = 20):
    """Poll OpportunityFieldHistory for recent changes; refresh affected cache rows."""
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    try:
        import cache_sync, sf_meeting_enricher as enricher
        sf = enricher._sf()
        out = await asyncio.get_running_loop().run_in_executor(
            None, lambda: cache_sync.sync_sf_changes_to_cache(supabase, sf, lookback_minutes=lookback_minutes)
        )
        return out
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


def _env_bool(name: str, default: bool) -> bool:
    """Parse a boolean env var. Truthy/falsey strings recognised; unset → default."""
    raw = os.getenv(name)
    if raw is None:
        return default
    return raw.strip().lower() not in ("false", "0", "no", "off", "")


_SF_PULL_CRON_RUNNING = False


async def _select_sf_pull_targets(lookback_hours: int, limit: int, opp_only: bool):
    """Scan avoma_event_reports within the lookback window, dedupe to the most
    recent report per opportunity, and return the picked rows (capped by limit)."""
    from datetime import datetime, timezone, timedelta
    loop = asyncio.get_running_loop()
    since = (datetime.now(timezone.utc) - timedelta(hours=max(lookback_hours, 1))).isoformat()
    res = await loop.run_in_executor(
        None,
        lambda: supabase.table("avoma_event_reports")
            .select("message_id,meeting_uuid,meeting_subject,meeting_start_at,"
                    "sf_opportunity_id,sf_account_id,sf_contact_ids,created_at")
            .gte("created_at", since)
            .order("created_at", desc=True)
            .limit(2000).execute(),
    )
    rows = res.data or []
    seen = set()
    picked = []
    for r in rows:
        opp = r.get("sf_opportunity_id")
        if opp_only and not opp:
            continue
        if not (opp or r.get("sf_account_id") or r.get("sf_contact_ids")):
            continue
        key = opp or f"msg:{r['message_id']}"
        if key in seen:
            continue
        seen.add(key)
        picked.append(r)
        if len(picked) >= max(limit, 1):
            break
    return picked


async def _run_sf_pull_batch(items, tag="CRON-SF-PULL", clear_flag=True):
    """Run the deterministic SF pull + cache refresh for each item SEQUENTIALLY.
    Clears the overlap flag in finally unless clear_flag=False (the nightly run
    holds the lock across several sub-jobs and releases it itself). Never raises
    (helper is safe)."""
    global _SF_PULL_CRON_RUNNING
    try:
        print(f"[{tag}] ▶ refreshing {len(items)} opportunity(ies)", flush=True)
        for r in items:
            # _run_sf_pull_and_cache never raises (catches internally), so one
            # failed opp won't abort the rest of the batch.
            await _run_sf_pull_and_cache(
                r["message_id"], r.get("meeting_uuid"),
                r.get("sf_opportunity_id"), r.get("sf_account_id"),
                r.get("sf_contact_ids") or [],
                meeting_subject=r.get("meeting_subject"),
                meeting_start_at=r.get("meeting_start_at"),
            )
        print(f"[{tag}] ✅ batch of {len(items)} complete", flush=True)
    finally:
        if clear_flag:
            _SF_PULL_CRON_RUNNING = False


async def _run_nightly_sf_pull():
    """One nightly run, three sub-jobs under a single overlap guard:
      (A) meeting-linked deterministic SF pull + cache refresh for opps seen in
          the avoma_event_reports lookback window;
      (B) new-opportunity discovery — scans Salesforce directly for opps created
          (and optionally modified) within the lookback, upserting them into
          opportunity_cache with NO dependency on an Avoma meeting;
      (C) global field-delta sync — OpportunityFieldHistory → opportunity_cache +
          field_history_cache.

    (B) and (C) run regardless of whether (A) found any meeting-linked opps, and
    each is independently toggleable via env vars. Awaited (not fire-and-forget)
    so the scheduler knows when it's done. Respects the overlap guard and returns
    a summary dict (also used by the on-demand /cron/nightly-sf-pull endpoint)."""
    global _SF_PULL_CRON_RUNNING
    summary = {
        "status": "completed",
        "meeting_linked": None,
        "discovery": None,
        "delta_sync": None,
        "deal_engine_discovery": None,
        "total_duration_ms": 0,
    }
    if supabase is None:
        print("[NIGHTLY-SF-PULL] supabase not configured; skip", flush=True)
        summary["status"] = "skipped_no_supabase"
        return summary
    # Atomic check-and-set: asyncio is single-threaded, so no other coroutine can
    # run between the `if` and the assignment (no await in between). This closes
    # the race the endpoint + scheduler would otherwise share.
    if _SF_PULL_CRON_RUNNING:
        print("[NIGHTLY-SF-PULL] a batch is already running; skip this tick", flush=True)
        summary["status"] = "skipped_already_running"
        return summary
    _SF_PULL_CRON_RUNNING = True
    import time as _t
    run_t0 = _t.perf_counter()
    try:
        # ---------- Sub-job A: meeting-linked deterministic refresh ----------
        try:
            lookback = int(os.getenv("SF_PULL_CRON_LOOKBACK_HOURS", "26"))
        except (TypeError, ValueError):
            lookback = 26
        try:
            limit = int(os.getenv("SF_PULL_CRON_LIMIT", "500"))
        except (TypeError, ValueError):
            limit = 500
        a_t0 = _t.perf_counter()
        try:
            picked = await _select_sf_pull_targets(lookback, limit, opp_only=True)
            print(f"[NIGHTLY-SF-PULL] (A) meeting-linked: selected {len(picked)} opp(s) "
                  f"(lookback={lookback}h, limit={limit})", flush=True)
            if picked:
                await _run_sf_pull_batch(picked, tag="NIGHTLY-SF-PULL", clear_flag=False)
            summary["meeting_linked"] = {
                "selected": len(picked),
                "lookback_hours": lookback,
                "limit": limit,
                "duration_ms": round((_t.perf_counter() - a_t0) * 1000, 1),
            }
        except Exception as e:
            print(f"[NIGHTLY-SF-PULL] (A) meeting-linked failed: {e}", flush=True)
            summary["meeting_linked"] = {"error": str(e),
                                         "duration_ms": round((_t.perf_counter() - a_t0) * 1000, 1)}

        # ---------- Sub-job B: new-opportunity discovery ----------
        if _env_bool("SF_PULL_CRON_DISCOVERY_ENABLED", True):
            try:
                disc_lookback = int(os.getenv("SF_PULL_CRON_DISCOVERY_LOOKBACK_HOURS", "26"))
            except (TypeError, ValueError):
                disc_lookback = 26
            include_modified = _env_bool("SF_PULL_CRON_DISCOVERY_INCLUDE_MODIFIED", False)
            b_t0 = _t.perf_counter()
            try:
                import cache_sync, sf_meeting_enricher as enricher
                sf = enricher._sf()
                disc = await asyncio.get_running_loop().run_in_executor(
                    None,
                    lambda: cache_sync.discover_new_opportunities(
                        supabase, sf,
                        lookback_hours=disc_lookback,
                        include_modified=include_modified,
                    ),
                )
                disc["duration_ms"] = round((_t.perf_counter() - b_t0) * 1000, 1)
                disc["lookback_hours"] = disc_lookback
                summary["discovery"] = disc
                print(f"[NIGHTLY-SF-PULL] (B) discovery: window={disc['window_opps_found']} "
                      f"new={disc['new_opps_found']} upserted={disc['opps_upserted']} "
                      f"(lookback={disc_lookback}h, include_modified={include_modified}) "
                      f"in {disc['duration_ms']}ms", flush=True)
            except Exception as e:
                print(f"[NIGHTLY-SF-PULL] (B) discovery failed: {e}", flush=True)
                summary["discovery"] = {"error": str(e),
                                        "duration_ms": round((_t.perf_counter() - b_t0) * 1000, 1)}
        else:
            print("[NIGHTLY-SF-PULL] (B) discovery disabled via SF_PULL_CRON_DISCOVERY_ENABLED", flush=True)
            summary["discovery"] = {"skipped": "disabled"}

        # ---------- Sub-job C: global field-delta sync ----------
        if _env_bool("SF_PULL_CRON_DELTA_ENABLED", True):
            try:
                delta_minutes = int(os.getenv("SF_PULL_CRON_DELTA_LOOKBACK_MINUTES", "1560"))
            except (TypeError, ValueError):
                delta_minutes = 1560
            c_t0 = _t.perf_counter()
            try:
                import cache_sync, sf_meeting_enricher as enricher
                sf = enricher._sf()
                delta = await asyncio.get_running_loop().run_in_executor(
                    None,
                    lambda: cache_sync.sync_sf_changes_to_cache(
                        supabase, sf, lookback_minutes=delta_minutes
                    ),
                )
                delta["duration_ms"] = round((_t.perf_counter() - c_t0) * 1000, 1)
                delta["lookback_minutes"] = delta_minutes
                summary["delta_sync"] = delta
                print(f"[NIGHTLY-SF-PULL] (C) delta-sync: changes={delta['changes_seen']} "
                      f"opps_refreshed={delta['opps_refreshed']} "
                      f"field_history_inserts={delta['field_history_inserts']} "
                      f"(lookback={delta_minutes}m) in {delta['duration_ms']}ms", flush=True)
            except Exception as e:
                print(f"[NIGHTLY-SF-PULL] (C) delta-sync failed: {e}", flush=True)
                summary["delta_sync"] = {"error": str(e),
                                         "duration_ms": round((_t.perf_counter() - c_t0) * 1000, 1)}
        else:
            print("[NIGHTLY-SF-PULL] (C) delta-sync disabled via SF_PULL_CRON_DELTA_ENABLED", flush=True)
            summary["delta_sync"] = {"skipped": "disabled"}

        # ---------- Sub-job D: Deal Engine new-opp discovery + sweep ----------
        # Keeps the Deal Engine (deal_records -> MASE) self-refreshing: open team
        # opps with no canonical record yet get swept here, so brand-new deals
        # appear within one scheduled cycle. Distinct from sub-job B, which only
        # feeds opportunity_cache. Gated + capped via DEAL_DISCOVERY_* env.
        if _env_bool("DEAL_ENGINE_DISCOVERY_ENABLED", True):
            d_t0 = _t.perf_counter()
            try:
                import deal_engine_sweep as _dsweep
                # First reconcile book membership against the MASE report (single
                # source of truth): deactivate exits, reactivate re-entrants, and
                # sweep brand-new entrants — all in this cycle.
                rec = await _dsweep.reconcile_membership(
                    agent_manager, source="scheduled_reconcile")
                summary["deal_engine_reconcile"] = rec
                print(f"[NIGHTLY-SF-PULL] (D) report reconcile: "
                      f"report={rec.get('report_count')} "
                      f"added={len(rec.get('added') or [])} "
                      f"removed={len(rec.get('removed') or [])} "
                      f"reactivated={len(rec.get('reactivated') or [])} "
                      f"removal_ran={rec.get('removal_ran')}", flush=True)
                # Then the watermark change re-sweep over the (now correct) book.
                de = await _dsweep.discover_and_sweep_new(
                    agent_manager, source="scheduled_discovery")
                de["duration_ms"] = round((_t.perf_counter() - d_t0) * 1000, 1)
                summary["deal_engine_discovery"] = de
                print(f"[NIGHTLY-SF-PULL] (D) deal-engine discovery: "
                      f"discovered={de.get('discovered')} new={de.get('new')} "
                      f"completed={de.get('completed')} failed={de.get('failed')} "
                      f"in {de['duration_ms']}ms", flush=True)
            except Exception as e:
                print(f"[NIGHTLY-SF-PULL] (D) deal-engine discovery failed: {e}", flush=True)
                summary["deal_engine_discovery"] = {
                    "error": str(e),
                    "duration_ms": round((_t.perf_counter() - d_t0) * 1000, 1)}
        else:
            print("[NIGHTLY-SF-PULL] (D) deal-engine discovery disabled via DEAL_ENGINE_DISCOVERY_ENABLED", flush=True)
            summary["deal_engine_discovery"] = {"skipped": "disabled"}
    except Exception as e:
        summary["status"] = "error"
        summary["error"] = str(e)
        print(f"[NIGHTLY-SF-PULL] run failed: {e}", flush=True)
    finally:
        _SF_PULL_CRON_RUNNING = False
    summary["total_duration_ms"] = round((_t.perf_counter() - run_t0) * 1000, 1)
    print(f"[NIGHTLY-SF-PULL] ✅ run complete in {summary['total_duration_ms']}ms", flush=True)
    return summary


async def _nightly_sf_pull_scheduler():
    """Background loop that fires _run_nightly_sf_pull once per day at the
    configured local hour. Defaults to 00:00 (midnight) UTC.

    Env overrides:
      - SF_PULL_CRON_ENABLED   (default "true") — set false to disable.
      - SF_PULL_CRON_HOUR      (default 0)       — hour of day, 0-23.
      - SF_PULL_CRON_TZ        (default "UTC")   — IANA tz, e.g. "America/New_York".
      - SF_PULL_CRON_LOOKBACK_HOURS / SF_PULL_CRON_LIMIT — scope of sub-job (A),
        the meeting-linked refresh.
      - SF_PULL_CRON_DISCOVERY_ENABLED          (default "true") — toggle sub-job
        (B) new-opportunity discovery independently.
      - SF_PULL_CRON_DISCOVERY_LOOKBACK_HOURS   (default 26) — discovery window.
      - SF_PULL_CRON_DISCOVERY_INCLUDE_MODIFIED (default "false") — also pull opps
        whose SystemModstamp falls in the window, not just newly created ones.
      - SF_PULL_CRON_DELTA_ENABLED              (default "true") — toggle sub-job
        (C) global OpportunityFieldHistory delta sync independently.
      - SF_PULL_CRON_DELTA_LOOKBACK_MINUTES     (default 1560 = 26h) — delta window.
    """
    if os.getenv("SF_PULL_CRON_ENABLED", "true").strip().lower() in ("false", "0", "no", "off"):
        print("[NIGHTLY-SF-PULL] scheduler disabled via SF_PULL_CRON_ENABLED", flush=True)
        return
    from datetime import datetime, timedelta, timezone
    tz_name = os.getenv("SF_PULL_CRON_TZ", "UTC")
    try:
        from zoneinfo import ZoneInfo
        tz = ZoneInfo(tz_name)
    except Exception as e:
        print(f"[NIGHTLY-SF-PULL] bad SF_PULL_CRON_TZ={tz_name!r} ({e}); falling back to UTC", flush=True)
        tz = timezone.utc  # explicit UTC, not naive/local
        tz_name = "UTC"
    try:
        hour = max(0, min(23, int(os.getenv("SF_PULL_CRON_HOUR", "0"))))
    except (TypeError, ValueError):
        hour = 0
    while True:
        # Each iteration is fully guarded so a transient error (clock, runtime)
        # can never permanently kill the scheduler task.
        try:
            now = datetime.now(tz)
            nxt = now.replace(hour=hour, minute=0, second=0, microsecond=0)
            if nxt <= now:
                nxt = nxt + timedelta(days=1)
            sleep_s = max(1.0, (nxt - now).total_seconds())
            print(f"[NIGHTLY-SF-PULL] next run at {nxt.isoformat()} "
                  f"({tz_name}) — sleeping {int(sleep_s)}s", flush=True)
            await asyncio.sleep(sleep_s)
            await _run_nightly_sf_pull()
            await asyncio.sleep(60)  # avoid a double-fire within the same minute
        except asyncio.CancelledError:
            raise
        except Exception as e:
            print(f"[NIGHTLY-SF-PULL] scheduler loop error: {e}; retrying in 300s", flush=True)
            await asyncio.sleep(300)


async def _nightly_hard_refresh_scheduler():
    """Background loop that fires the token-free hard-fact reconciliation once per
    day at the configured local hour. This re-reads the live Salesforce hard facts
    (stage / amount / dates / owner+manager / competitor / next step) for every
    persisted deal record and overwrites them with no AI cost, so reps never see
    stale facts between paid AI sweeps.

    sweep.hard_refresh_all() owns ALL the mutual exclusion: it serializes against
    itself, refuses while the in-process AI sweep is running, and refuses while the
    durable sweep_queue has waiting/working rows (the out-of-process worker), so it
    can never clobber — or be clobbered by — the paid sweep / queue worker.

    Runs an hour after the SF-pull cron by default so cache/discovery has settled.
    Discovery hygiene (deleting deals slipped back to 'Initial Interest') is OFF by
    default — this job is facts-only; report reconciliation owns membership.

    Env overrides:
      - DEAL_HARD_REFRESH_CRON_ENABLED   (default "true") — set false to disable.
      - DEAL_HARD_REFRESH_CRON_HOUR      (default 1)       — hour of day, 0-23.
      - DEAL_HARD_REFRESH_CRON_TZ        (default "UTC")   — IANA tz.
      - DEAL_HARD_REFRESH_CRON_DELETE_II (default "false") — also delete deals now
        back at 'Initial Interest' (discovery hygiene; off for a pure fact sync).
    """
    import deal_engine_sweep as sweep
    if os.getenv("DEAL_HARD_REFRESH_CRON_ENABLED", "true").strip().lower() in ("false", "0", "no", "off"):
        print("[NIGHTLY-HARD-REFRESH] scheduler disabled via DEAL_HARD_REFRESH_CRON_ENABLED", flush=True)
        return
    from datetime import datetime, timedelta, timezone
    tz_name = os.getenv("DEAL_HARD_REFRESH_CRON_TZ", "UTC")
    try:
        from zoneinfo import ZoneInfo
        tz = ZoneInfo(tz_name)
    except Exception as e:
        print(f"[NIGHTLY-HARD-REFRESH] bad DEAL_HARD_REFRESH_CRON_TZ={tz_name!r} ({e}); falling back to UTC", flush=True)
        tz = timezone.utc
        tz_name = "UTC"
    try:
        hour = max(0, min(23, int(os.getenv("DEAL_HARD_REFRESH_CRON_HOUR", "1"))))
    except (TypeError, ValueError):
        hour = 1
    delete_ii = os.getenv("DEAL_HARD_REFRESH_CRON_DELETE_II", "false").strip().lower() in ("true", "1", "yes", "on")
    while True:
        # Each iteration is fully guarded so a transient error can never
        # permanently kill the scheduler task.
        try:
            now = datetime.now(tz)
            nxt = now.replace(hour=hour, minute=0, second=0, microsecond=0)
            if nxt <= now:
                nxt = nxt + timedelta(days=1)
            sleep_s = max(1.0, (nxt - now).total_seconds())
            print(f"[NIGHTLY-HARD-REFRESH] next run at {nxt.isoformat()} "
                  f"({tz_name}) — sleeping {int(sleep_s)}s", flush=True)
            await asyncio.sleep(sleep_s)
            summary = await sweep.hard_refresh_all(
                agent_manager,
                delete_initial_interest=delete_ii,
                source="nightly_cron",
            )
            if summary.get("skipped"):
                print(f"[NIGHTLY-HARD-REFRESH] skipped: {summary['skipped']} "
                      "(AI sweep / queue / another refresh active)", flush=True)
            else:
                print(f"[NIGHTLY-HARD-REFRESH] done: records={summary.get('records')} "
                      f"matched={summary.get('matched')} updated={summary.get('updated')} "
                      f"removed={summary.get('removed')} unmatched={summary.get('unmatched')} "
                      f"failed={summary.get('failed')}", flush=True)
            await asyncio.sleep(60)  # avoid a double-fire within the same minute
        except asyncio.CancelledError:
            raise
        except Exception as e:
            print(f"[NIGHTLY-HARD-REFRESH] scheduler loop error: {e}; retrying in 300s", flush=True)
            await asyncio.sleep(300)


@app.get("/cron/sf-pull-refresh")
async def cron_sf_pull_refresh(lookback_hours: int = 72, limit: int = 25, opp_only: bool = True):
    """Re-run the deterministic 3-tier SF pull + cache refresh for recently-seen
    opportunities, to keep cached deal data fresh between Avoma meetings.

    This is the heavyweight companion to /cron/sync-sf-to-cache:
      - /cron/sync-sf-to-cache is a LIGHT poll of OpportunityFieldHistory that
        only touches opp_cache rows whose SF fields changed in the last N minutes.
      - /cron/sf-pull-refresh re-runs the FULL structured pull
        (`_run_sf_pull_and_cache` → deal_health / account_briefing / full_snapshot
        + opportunity_cache / meeting_cache / field_history_cache) for each
        avoma_event_reports row that has a resolved SF link within the lookback
        window. Because Salesforce CDC is off for Opportunity, this is what keeps
        the richer cache data fresh for deals that haven't had a new meeting.

    Dedupes to the MOST RECENT report per opportunity so each opp is pulled once.
    Fire-and-forget: pulls run sequentially in a background task (to avoid
    hammering Salesforce with parallel SOQL); the endpoint returns the queued
    list immediately so a scheduler's HTTP call doesn't block or time out.
    Poll each report's `sf_pull_status` for completion.

    Query params:
      - lookback_hours: how far back to scan reports by created_at (default 72).
      - limit: max distinct opportunities to refresh this run (default 25).
      - opp_only: when true (default) only refresh reports with a resolved
        sf_opportunity_id; when false also include account/contact-only reports.
    """
    global _SF_PULL_CRON_RUNNING
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    # Overlap guard: a batch is heavyweight (sequential SOQL per opp). If one is
    # already in flight, don't start a second — concurrent runs would re-pull the
    # same opps and race sf_pull_status / cache writes.
    if _SF_PULL_CRON_RUNNING:
        return JSONResponse(
            {"status": "already_running",
             "note": "an sf-pull-refresh batch is already in progress; try again after it finishes"},
            status_code=409,
        )
    # Atomic check-and-set (no await between the check above and here), so the
    # endpoint and the nightly scheduler can never both pass the guard.
    _SF_PULL_CRON_RUNNING = True
    try:
        picked = await _select_sf_pull_targets(lookback_hours, limit, opp_only)

        if picked:
            task = asyncio.create_task(_run_sf_pull_batch(picked))

            def _on_done(t):
                try:
                    exc = t.exception()
                    if exc:
                        print(f"[CRON-SF-PULL] ❌ batch task crashed: {exc}", flush=True)
                except asyncio.CancelledError:
                    pass
            task.add_done_callback(_on_done)
        else:
            # Nothing to do — release the guard we just took.
            _SF_PULL_CRON_RUNNING = False

        return {
            "status": "started" if picked else "noop",
            "queued": len(picked),
            "lookback_hours": lookback_hours,
            "opp_only": opp_only,
            "opportunities": [r.get("sf_opportunity_id") or f"msg:{r['message_id']}" for r in picked],
            "note": "deterministic SF pull + cache refresh running sequentially in background; poll each report's sf_pull_status",
        }
    except Exception as e:
        # Selection failed before a batch task could take ownership of the guard;
        # release it so the next call isn't permanently blocked.
        _SF_PULL_CRON_RUNNING = False
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/cron/nightly-sf-pull")
async def cron_nightly_sf_pull():
    """Run the full combined nightly Salesforce sync on demand and return a
    JSON summary of what each sub-job did.

    Runs the SAME logic the in-process nightly scheduler fires once a day:
      (A) meeting-linked deterministic SF pull + cache refresh,
      (B) new-opportunity discovery (opps created/modified since the lookback,
          with no dependency on an Avoma meeting),
      (C) global field-delta sync (OpportunityFieldHistory → opportunity_cache +
          field_history_cache).

    All three share the single `_SF_PULL_CRON_RUNNING` overlap guard, so calling
    this while a nightly run (or another invocation) is in flight returns
    status="skipped_already_running" rather than starting a second run.

    Sub-jobs (B) and (C) are individually toggleable via env vars
    (SF_PULL_CRON_DISCOVERY_ENABLED / SF_PULL_CRON_DELTA_ENABLED).
    """
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    try:
        summary = await _run_nightly_sf_pull()
        if summary.get("status") == "skipped_already_running":
            return JSONResponse(summary, status_code=409)
        return summary
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/ceo-query")
async def ceo_query(request: Request):
    """Thin keyword router over the cache tables (50-100ms response).

    For arbitrary questions, use the agent (POST /api/chat) which has the
    supabase_query MCP tool and can write SQL against opportunity_cache /
    meeting_cache / field_history_cache directly.
    """
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    try:
        data = await request.json()
        q = (data.get("question") or "").strip()
        if not q:
            return JSONResponse({"error": "question required"}, status_code=400)
        ql = q.lower()
        loop = asyncio.get_running_loop()

        # "stalled deals"
        if "stalled" in ql or "at risk" in ql:
            res = await loop.run_in_executor(None, lambda:
                supabase.table("opportunity_cache").select("*").eq("momentum", "Stalled").order("health_score").limit(50).execute())
            return {"pattern": "stalled_deals", "count": len(res.data or []), "deals": res.data}

        # "deals with no meetings" / "zero meetings"
        if "no meeting" in ql or "zero meeting" in ql or "0 meetings" in ql:
            res = await loop.run_in_executor(None, lambda:
                supabase.table("opportunity_cache").select("*").eq("meetings_count", 0).gte("amount", 100000).order("amount", desc=True).limit(50).execute())
            return {"pattern": "no_meetings_high_value", "count": len(res.data or []), "deals": res.data}

        # "amount history" / "amount before"
        if "amount" in ql and ("history" in ql or "before" in ql or "changed" in ql):
            # extract first 2-3 word capitalized phrase as account hint
            import re as _re
            m = _re.search(r"([A-Z][a-zA-Z0-9&'.\-]+(?:\s+[A-Z][a-zA-Z0-9&'.\-]+){0,3})", q)
            hint = (m.group(1) if m else q.split()[0]).strip()
            opp = await loop.run_in_executor(None, lambda:
                supabase.table("opportunity_cache").select("opportunity_id,opportunity_name,amount")
                .or_(f"opportunity_name.ilike.%{hint}%,account_name.ilike.%{hint}%").limit(1).execute())
            if not opp.data:
                return {"pattern": "amount_history", "error": f"no opp matched '{hint}'"}
            opp_id = opp.data[0]["opportunity_id"]
            hist = await loop.run_in_executor(None, lambda:
                supabase.table("field_history_cache").select("*")
                .eq("opportunity_id", opp_id).eq("field_name", "Amount")
                .order("changed_date", desc=True).limit(20).execute())
            return {"pattern": "amount_history", "opportunity": opp.data[0], "history": hist.data}

        # Fallback: name/account search → return matching opps
        import re as _re
        _STOP = {"show","me","tell","about","find","get","what","is","the","a","an","for","on","of","in","to","please","status","deal","deals","opportunity","account","report"}
        toks = [t for t in _re.findall(r"[A-Za-z0-9&'.\-]+", q) if t.lower() not in _STOP]
        hint = " ".join(toks[:3]) if toks else q
        hint = hint.strip().strip("?.,!")
        res = await loop.run_in_executor(None, lambda:
            supabase.table("opportunity_cache").select("*")
            .or_(f"opportunity_name.ilike.%{hint}%,account_name.ilike.%{hint}%")
            .order("amount", desc=True).limit(10).execute())
        return {"pattern": "name_search", "hint": hint, "count": len(res.data or []), "deals": res.data}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/ask")
async def ask_cache(request: Request):
    """Natural-language Q&A answered ONLY from the cached tables + reports.

    A small read-only LLM (gpt-4o-mini) uses tools hard-scoped to
    opportunity_cache / meeting_cache / field_history_cache / avoma_event_reports
    / opportunity_observatory. It never touches live Salesforce/Avoma and has no
    write path. Powers the GET /ask UI.
    """
    if supabase is None:
        return JSONResponse({"error": "supabase not configured"}, status_code=503)
    try:
        data = await request.json()
        question = (data.get("question") or "").strip()
        if not question:
            return JSONResponse({"error": "question required"}, status_code=400)
        import cache_qa
        result = await cache_qa.answer_question(supabase, question)
        return result
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/ask", response_class=HTMLResponse)
async def ask_cache_ui():
    """Dummy 'ask anything' UI over the cached data + reports."""
    return HTMLResponse(content=r"""<!doctype html>
<html><head><meta charset="utf-8"><meta name="viewport" content="width=device-width, initial-scale=1">
<title>Ask the Deal Brain</title>
<style>
  :root { --bg:#0f1419; --panel:#1a1f2e; --border:#2a3142; --border2:#3a4258; --text:#d4d4d4; --muted:#8a93a6; --accent:#4a9eff; }
  * { box-sizing: border-box; }
  body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", system-ui, sans-serif; margin:0; background:var(--bg); color:var(--text); }
  header { background:var(--panel); padding:16px 24px; border-bottom:1px solid var(--border); }
  header h1 { margin:0; font-size:18px; font-weight:600; }
  header p { margin:4px 0 0; font-size:13px; color:var(--muted); }
  main { max-width:860px; margin:0 auto; padding:24px 20px 60px; }
  .box { display:flex; gap:10px; margin-bottom:8px; }
  .box input { flex:1; background:var(--panel); color:var(--text); border:1px solid var(--border2); border-radius:8px; padding:12px 14px; font-size:15px; }
  .box input:focus { outline:none; border-color:var(--accent); }
  .box button { background:var(--accent); color:#06121f; border:none; border-radius:8px; padding:0 20px; font-size:15px; font-weight:600; cursor:pointer; }
  .box button:disabled { opacity:.5; cursor:default; }
  .hints { font-size:12px; color:var(--muted); margin-bottom:24px; }
  .hints b { color:var(--text); font-weight:500; }
  .chip { display:inline-block; background:var(--panel); border:1px solid var(--border2); border-radius:14px; padding:4px 10px; margin:4px 6px 0 0; font-size:12px; cursor:pointer; color:var(--text); }
  .chip:hover { border-color:var(--accent); }
  .answer { background:var(--panel); border:1px solid var(--border); border-radius:10px; padding:18px 20px; white-space:pre-wrap; line-height:1.55; font-size:14.5px; margin-top:10px; }
  .answer.empty { color:var(--muted); }
  .steps { margin-top:14px; font-size:12px; color:var(--muted); }
  .steps summary { cursor:pointer; }
  .steps code { background:#11161f; padding:2px 6px; border-radius:4px; color:#9bd1ff; }
  .step { margin:6px 0; }
  .spinner { display:inline-block; width:14px; height:14px; border:2px solid var(--border2); border-top-color:var(--accent); border-radius:50%; animation:spin .7s linear infinite; vertical-align:middle; }
  @keyframes spin { to { transform:rotate(360deg); } }
  .meta { font-size:11px; color:var(--muted); margin-top:8px; }
</style></head>
<body>
<header>
  <h1>Ask the Deal Brain</h1>
  <p>Answers come only from your cached data + reports (no live Salesforce/Avoma). If it isn't cached, it'll say so.</p>
</header>
<main>
  <div class="box">
    <input id="q" type="text" placeholder="e.g. Which deals are stalled? What's the health of the Sabic opportunity?" autofocus />
    <button id="go">Ask</button>
  </div>
  <div class="hints">
    Try:
    <span class="chip">Show me all stalled deals</span>
    <span class="chip">High-value deals with no meetings</span>
    <span class="chip">What changed on the biggest opportunity?</span>
    <span class="chip">Summarize the latest meeting analysis for Sabic</span>
  </div>
  <div id="answer" class="answer empty">Your answer will appear here.</div>
  <div id="meta" class="meta"></div>
  <div id="steps"></div>
</main>
<script>
  const q = document.getElementById('q');
  const go = document.getElementById('go');
  const answerEl = document.getElementById('answer');
  const stepsEl = document.getElementById('steps');
  const metaEl = document.getElementById('meta');

  document.querySelectorAll('.chip').forEach(c => c.addEventListener('click', () => { q.value = c.textContent; ask(); }));
  q.addEventListener('keydown', e => { if (e.key === 'Enter') ask(); });
  go.addEventListener('click', ask);

  async function ask() {
    const question = q.value.trim();
    if (!question) return;
    go.disabled = true;
    answerEl.className = 'answer';
    answerEl.innerHTML = '<span class="spinner"></span> Thinking…';
    stepsEl.innerHTML = '';
    metaEl.textContent = '';
    try {
      const r = await fetch('/api/ask', {
        method: 'POST', headers: { 'Content-Type': 'application/json' },
        body: JSON.stringify({ question })
      });
      const data = await r.json();
      if (!r.ok || data.error) {
        answerEl.className = 'answer empty';
        answerEl.textContent = 'Error: ' + (data.error || r.status);
      } else {
        answerEl.className = 'answer';
        answerEl.textContent = data.answer || '(no answer)';
        if (data.model) metaEl.textContent = 'model: ' + data.model;
        if (data.steps && data.steps.length) {
          let h = '<details class="steps"><summary>' + data.steps.length + ' data lookup(s)</summary>';
          data.steps.forEach(s => {
            h += '<div class="step"><code>' + s.tool + '</code> ' +
                 (s.args ? JSON.stringify(s.args) : '') + '</div>';
          });
          h += '</details>';
          stepsEl.innerHTML = h;
        }
      }
    } catch (e) {
      answerEl.className = 'answer empty';
      answerEl.textContent = 'Request failed: ' + e;
    } finally {
      go.disabled = false;
    }
  }
</script>
</body></html>""")


# ============================================================================
# Analysis feature (Task #52) — spreadsheet-style analyses over opportunities.
# Writes go through the service-role store (Bearer-gated endpoints); the browser
# reads live updates via Supabase realtime on the analysis_* tables.
# ============================================================================

async def _aw(fn, *args, **kwargs):
    """Run a sync analysis_store function in a thread (keeps the loop free)."""
    import functools
    return await asyncio.get_running_loop().run_in_executor(
        None, functools.partial(fn, *args, **kwargs))


@app.get("/api/analysis/models")
async def analysis_models():
    """Curated model suggestions for AI columns (any supported provider:model is accepted)."""
    import analysis_engine as engine
    return {"default": engine._DEFAULT_MODEL, "models": engine.MODEL_SUGGESTIONS,
            "providers": sorted(engine.SUPPORTED_PROVIDERS)}


@app.post("/api/analysis")
async def analysis_create_ep(request: Request):
    import analysis_store as store
    try:
        d = await request.json()
        if not (d.get("title") or "").strip():
            return JSONResponse({"error": "title required"}, status_code=400)
        a = await _aw(store.create_analysis, d["title"],
                      description=d.get("description"), project_id=d.get("project_id"),
                      chat_id=d.get("chat_id"), created_by=d.get("created_by"),
                      source_config=d.get("source_config"))
        return a
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/analysis")
async def analysis_list_ep(project_id: str = "", chat_id: str = "", limit: int = 50):
    import analysis_store as store
    try:
        rows = await _aw(store.list_analyses, project_id=project_id or None,
                         chat_id=chat_id or None, limit=limit)
        return {"count": len(rows), "analyses": rows}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/analysis/{analysis_id}")
async def analysis_get_ep(analysis_id: str):
    import analysis_store as store
    try:
        full = await _aw(store.get_full_analysis, analysis_id)
        if not full:
            return JSONResponse({"error": "analysis not found"}, status_code=404)
        return full
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.patch("/api/analysis/{analysis_id}")
async def analysis_update_ep(analysis_id: str, request: Request):
    import analysis_store as store
    try:
        d = await request.json()
        allowed = {k: d[k] for k in ("title", "description", "status", "source_config") if k in d}
        if not allowed:
            return JSONResponse({"error": "nothing to update"}, status_code=400)
        a = await _aw(store.update_analysis, analysis_id, allowed)
        return a or JSONResponse({"error": "analysis not found"}, status_code=404)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.delete("/api/analysis/{analysis_id}")
async def analysis_delete_ep(analysis_id: str):
    import analysis_store as store
    try:
        await _aw(store.delete_analysis, analysis_id)
        return {"deleted": analysis_id}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/analysis/{analysis_id}/columns")
async def analysis_add_column_ep(analysis_id: str, request: Request):
    import analysis_store as store, analysis_engine as engine
    try:
        d = await request.json()
        col_type = d.get("type", "data")
        config = d.get("config") or {}
        if col_type == "ai" and config.get("model"):
            config["model"] = engine.validate_model(config["model"])
        col = await _aw(store.add_column, analysis_id, d.get("name", ""), col_type,
                        config=config, position=d.get("position"))
        if col_type == "data":
            await _aw(store.populate_data_cells, analysis_id)
        return col
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.patch("/api/analysis/columns/{column_id}")
async def analysis_update_column_ep(column_id: str, request: Request):
    import analysis_store as store
    try:
        d = await request.json()
        patch = {k: d[k] for k in ("name", "position", "config") if k in d}
        if not patch:
            return JSONResponse({"error": "nothing to update"}, status_code=400)
        col = await _aw(store.update_column, column_id, patch)
        return col or JSONResponse({"error": "column not found"}, status_code=404)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.delete("/api/analysis/columns/{column_id}")
async def analysis_delete_column_ep(column_id: str):
    import analysis_store as store
    try:
        await _aw(store.delete_column, column_id)
        return {"deleted_column": column_id}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/analysis/{analysis_id}/rows")
async def analysis_add_rows_ep(analysis_id: str, request: Request):
    """Add rows either from a cache source ({"source": "opportunity_cache", ...filters})
    or explicitly ({"rows": [{entity_ref, label, source}]})."""
    import analysis_store as store
    try:
        d = await request.json()
        if d.get("source"):
            out = await _aw(store.add_rows_from_source, analysis_id, d["source"],
                            stage=d.get("stage"), momentum=d.get("momentum"),
                            min_amount=d.get("min_amount"), max_amount=d.get("max_amount"),
                            account_contains=d.get("account_contains"),
                            name_contains=d.get("name_contains"),
                            is_closed=d.get("is_closed"), limit=d.get("limit", 25))
            return out
        rows = d.get("rows") or []
        added = await _aw(store.add_rows, analysis_id, rows)
        await _aw(store.populate_data_cells, analysis_id)
        return {"added": len(added), "rows": added}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.delete("/api/analysis/rows/{row_id}")
async def analysis_delete_row_ep(row_id: str):
    import analysis_store as store
    try:
        await _aw(store.delete_row, row_id)
        return {"deleted_row": row_id}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/analysis/{analysis_id}/run")
async def analysis_run_ep(analysis_id: str):
    import analysis_store as store, analysis_engine as engine
    try:
        if not await _aw(store.get_analysis, analysis_id):
            return JSONResponse({"error": "analysis not found"}, status_code=404)
        return engine.start_run(analysis_id)
    except engine.AnalysisRunError as e:
        return JSONResponse({"error": str(e)}, status_code=409)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/analysis/{analysis_id}/stop")
async def analysis_stop_ep(analysis_id: str):
    import analysis_engine as engine
    try:
        return engine.stop_run(analysis_id)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/analysis/{analysis_id}/runs")
async def analysis_runs_ep(analysis_id: str, limit: int = 20):
    import analysis_store as store, analysis_engine as engine
    try:
        runs = await _aw(store.list_runs, analysis_id, limit=limit)
        return {"is_running": engine.is_running(analysis_id), "count": len(runs), "runs": runs}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/analysis/{analysis_id}/query")
async def analysis_query_ep(analysis_id: str, request: Request):
    import analysis_engine as engine
    try:
        d = await request.json()
        q = (d.get("question") or "").strip()
        if not q:
            return JSONResponse({"error": "question required"}, status_code=400)
        return await engine.query_analysis(analysis_id, q, model=d.get("model"))
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.patch("/api/analysis/cells/{cell_id}")
async def analysis_edit_cell_ep(cell_id: str, request: Request):
    """Manual cell edit — sets the value and marks the cell done."""
    import analysis_store as store
    try:
        d = await request.json()
        if "value" not in d:
            return JSONResponse({"error": "value required"}, status_code=400)
        cell = await _aw(store.edit_cell, cell_id, d["value"])
        return cell or JSONResponse({"error": "cell not found"}, status_code=404)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/analysis/{analysis_id}/resume")
async def analysis_resume_ep(analysis_id: str):
    """Resume an interrupted run: recompute ONLY the non-done AI cells (already
    done cells keep their values) and reset the analysis status when finished."""
    import analysis_engine as engine
    try:
        return engine.start_resume(analysis_id)
    except engine.AnalysisRunError as e:
        return JSONResponse({"error": str(e)}, status_code=409)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/analysis/{analysis_id}/cells/rerun")
async def analysis_rerun_cell_ep(analysis_id: str, request: Request):
    """Re-run a single AI cell. Body: {cell_id} or {row_id, column_id}."""
    import analysis_engine as engine
    try:
        d = await request.json()
        return await engine.rerun_cell(analysis_id, cell_id=d.get("cell_id"),
                                       row_id=d.get("row_id"), column_id=d.get("column_id"))
    except engine.AnalysisNotFound as e:
        return JSONResponse({"error": str(e)}, status_code=404)
    except engine.AnalysisRunError as e:
        return JSONResponse({"error": str(e)}, status_code=409)
    except ValueError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


# ----- Dashboards (Task #53): spec-driven chart views over an analysis -------
# Spec is validated against an allowlist of widget/aggregation types and every
# referenced column_id is checked against the analysis's columns before persist
# (dashboard_store.validate_spec). The model never supplies SQL/table names.

@app.post("/api/analysis/{analysis_id}/dashboards")
async def dashboard_create_ep(analysis_id: str, request: Request):
    import dashboard_store as dstore
    try:
        d = await request.json()
        spec = d.get("spec")
        if spec is None and "widgets" in d:
            spec = {"widgets": d.get("widgets") or []}
        dash = await _aw(dstore.create_dashboard, analysis_id,
                         d.get("title", ""), spec or {"widgets": []},
                         description=d.get("description"),
                         created_by=d.get("created_by"))
        return dash
    except dstore.DashboardNotFound as e:
        return JSONResponse({"error": str(e)}, status_code=404)
    except dstore.DashboardError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/analysis/{analysis_id}/dashboards/suggest")
async def dashboard_suggest_ep(analysis_id: str, request: Request):
    """Auto-suggest a starter dashboard spec from an analysis's columns.

    Returns a validated draft spec (KPIs per numeric column, a bar chart of the
    first numeric grouped by the first categorical, and a table of primary
    columns). By default nothing is persisted; pass {"persist": true} to save it
    as a new dashboard immediately.
    """
    import dashboard_store as dstore
    try:
        body = {}
        try:
            body = await request.json()
        except Exception:  # noqa: BLE001 — empty/no body is fine
            body = {}
        max_widgets = body.get("max_widgets")
        spec = await _aw(dstore.suggest_spec, analysis_id,
                         int(max_widgets) if max_widgets is not None else None)
        if body.get("persist"):
            dash = await _aw(dstore.create_dashboard, analysis_id,
                             body.get("title") or spec.get("title") or "Starter dashboard",
                             spec, description=body.get("description"),
                             created_by=body.get("created_by"))
            return {"persisted": True, "dashboard": dash}
        return {"persisted": False, "analysis_id": analysis_id, "spec": spec}
    except dstore.DashboardNotFound as e:
        return JSONResponse({"error": str(e)}, status_code=404)
    except dstore.DashboardError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/analysis/{analysis_id}/dashboards")
async def dashboard_list_ep(analysis_id: str, limit: int = 100):
    import dashboard_store as dstore
    try:
        rows = await _aw(dstore.list_dashboards, analysis_id, limit=limit)
        return {"count": len(rows), "dashboards": rows}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/analysis/dashboards/{dashboard_id}")
async def dashboard_get_ep(dashboard_id: str):
    import dashboard_store as dstore
    try:
        dash = await _aw(dstore.get_dashboard, dashboard_id)
        return dash or JSONResponse({"error": "dashboard not found"}, status_code=404)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.patch("/api/analysis/dashboards/{dashboard_id}")
async def dashboard_update_ep(dashboard_id: str, request: Request):
    import dashboard_store as dstore
    try:
        d = await request.json()
        patch = {k: d[k] for k in ("title", "description", "spec") if k in d}
        if "spec" not in patch and "widgets" in d:
            patch["spec"] = {"widgets": d.get("widgets") or []}
        if not patch:
            return JSONResponse({"error": "nothing to update"}, status_code=400)
        dash = await _aw(dstore.update_dashboard, dashboard_id, patch)
        return dash or JSONResponse({"error": "dashboard not found"}, status_code=404)
    except dstore.DashboardNotFound as e:
        return JSONResponse({"error": str(e)}, status_code=404)
    except dstore.DashboardError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.delete("/api/analysis/dashboards/{dashboard_id}")
async def dashboard_delete_ep(dashboard_id: str):
    import dashboard_store as dstore
    try:
        await _aw(dstore.delete_dashboard, dashboard_id)
        return {"deleted_dashboard": dashboard_id}
    except dstore.DashboardNotFound as e:
        return JSONResponse({"error": str(e)}, status_code=404)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


# ============================================================================
# Deal Intelligence Engine — Deals / Espresso (to-do) / Matcha / Chat.
# Read-only, deterministic derivations over the stored evidence-anchored deal
# records (one canonical JSON per opportunity in deal_records). Namespaced under
# /api/deal-engine/* so they never collide with the agent's /api/chat etc.
# Writes (the sweep) go through the service-role store; the browser reads live.
# ============================================================================

# Strategist persona for the Deal chat (RevOps strategist over the whole book).
_DEAL_ENGINE_CHAT_SYSTEM = (
    "You are a RevOps strategist for the Zycus sales team, reasoning over a book "
    "of evidence-anchored deal records (Salesforce facts plus dated, cited AI "
    "analysis). Operating rules:\n"
    "- Test rep-stated probability and forecast labels against a 7-point "
    "qualification drill: engagement, access to power, champion, competition, "
    "product fit, risk, value. Call out claims the evidence does not support.\n"
    "- Weight recent evidence over stale evidence. Always name dates; never say "
    "'recently' or 'lately'.\n"
    "- Use fiscal quarters running April to March.\n"
    "- Describe AI appetite strictly as 'AI Hungry', 'AI Curious', or 'AI "
    "Resistant'.\n"
    "- No fabrication. Every claim must trace to a record field, a dated "
    "activity, or a cited quote in the provided book. If the book does not "
    "support an answer, say so plainly.\n"
    "- Plain English. No em dashes. Be specific and prescriptive: name the move "
    "and who should be in the room."
)

# Appended (by code, not admin-editable) to the chat system prompt so the chat agent
# always knows its tools and EXACTLY what the Todo Runner can / cannot do.
_CHAT_CAPABILITIES = (
    "\nTOOLS YOU CAN USE:\n"
    "- search_knowledge(query): retrieve from the MASE knowledge base — the shared, "
    "isolated store the sweep and Todo Runner also use (uploaded playbooks, Showpad index "
    "cards, competitive battlecards, capability decks). Use it to ground answers in real "
    "collateral and cite what you use. Never invent facts that aren't in the book or the "
    "knowledge base.\n"
    "- run_todo(task, account?, contact?, opportunity_id?): delegate ONE tactical, "
    "prospect-facing to-do to the Todo Runner, which DRAFTS a single outbound email to "
    "complete it. Use this whenever the user asks you to draft / write / follow up with an "
    "email for a specific to-do — do NOT write that email yourself; delegate it.\n"
    "  The Todo Runner CAN: draft one outbound email for a to-do that needs no internal "
    "collaboration; pull real facts from Showpad, Salesforce (real closed-won references) "
    "and the knowledge base; attach relevant Showpad collateral as shareable links; and it "
    "never invents customers, prices, or claims.\n"
    "  The Todo Runner CANNOT / WILL NOT: send the email (a human reviews and sends); handle "
    "anything needing a manager or exec, legal, security/infosec, the pricing desk, a sales "
    "engineer, product, or a partner — for those it returns one line 'NEEDS HUMAN: <who and "
    "why>'. When you call run_todo, present its draft to the user and surface any 'NEEDS "
    "HUMAN' verbatim with a short note on what's needed.\n"
)

# Persona overlays for the Deal chat. The UI's persona tabs send `persona`; the
# matching directive is appended to the system prompt for that run, switching the
# agent's lens. "strategist" is the default RevOps voice (no overlay).
_CHAT_PERSONAS = {
    "strategist": "",
    "vp": ("PERSONA — VP OF SALES: Answer as a VP of Sales reviewing the whole book. Lead with "
           "forecast risk and accuracy, pipeline coverage, where reps should spend time, what to "
           "inspect/escalate, and the few moves that most change the quarter. Be decisive and "
           "prioritised; think in roll-ups, not single-deal minutiae."),
    "coach": ("PERSONA — DEAL COACH: Answer as a hands-on deal coach for the specific deal(s) in "
              "question. Give concrete next moves, who to involve, how to handle the live "
              "objection/competition, and the exact next step to advance the deal. Practical and "
              "per-deal, not strategic generalities."),
    "qual": ("PERSONA — QUALIFICATION EXPERT: Answer as a qualification expert. Run the 7-point "
             "drill (engagement, access to power, champion, competition, product fit, risk, value) "
             "and score MEDDPICC/BANT. Expose what is unproven, what evidence is missing, and "
             "exactly what to verify next."),
}


def _deal_engine_model() -> str:
    # Spec default is Claude, but Anthropic is unavailable in this environment,
    # so the Deal chat runs on OpenAI. Overridable via DEAL_ENGINE_MODEL.
    return os.environ.get("DEAL_ENGINE_MODEL", "gpt-4o")


@app.get("/api/deal-engine")
async def deal_engine_descriptor():
    """Descriptor: endpoints, auth mode, record count."""
    import deal_engine_store as dstore
    try:
        count = await _aw(dstore.count_records)
    except Exception:  # noqa: BLE001
        count = None
    auth_on = bool(_api_auth_token()) if "_api_auth_token" in globals() else None
    return {
        "name": "Deal Intelligence Engine API",
        "record_count": count,
        "auth": "bearer" if auth_on else "open",
        "endpoints": [
            "GET /api/deal-engine", "GET /api/deal-engine/health",
            "GET /api/deal-engine/team",
            "GET /api/deal-engine/opportunities?owner=",
            "GET /api/deal-engine/opportunities/{opp_id}",
            "GET /api/deal-engine/todo?owner=",
            "POST /api/deal-engine/todo/push",
            "GET /api/deal-engine/todo/dashboard",
            "GET /api/deal-engine/matcha?owner=",
            "GET /api/deal-engine/deltas?owner=&limit=&group_by=owner",
            "GET /api/deal-engine/deltas/{opp_id}",
            "POST /api/deal-engine/records",
            "POST /api/deal-engine/backfill-packets",
            "POST /api/deal-engine/sweep",
            "POST /api/deal-engine/sweep/trigger",
            "POST /api/deal-engine/sweep/discover-new",
            "GET /api/deal-engine/sweep/discover?owner=",
            "GET /api/deal-engine/sweep/status",
            "GET /api/deal-engine/sweep/dashboard",
            "POST /api/deal-engine/sweep/{opp_id}",
            "GET /api/deal-engine/trigger-logs",
            "GET /api/deal-engine/trigger-logs/{opp_id}",
            "POST /api/deal-engine/chat",
        ],
    }


@app.get("/api/deal-engine/health")
async def deal_engine_health():
    return {"ok": True, "hasKey": bool(os.environ.get("OPENAI_API_KEY")),
            "model": _deal_engine_model()}


@app.get("/api/deal-engine/team")
async def deal_engine_team():
    import deal_engine_store as dstore
    try:
        return await _aw(dstore.get_team)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/chat/prompt")
async def get_chat_prompt():
    """Return the admin instruction override for the chat/completion agent plus the
    built-in default, for the Admin -> Agent Control 'Instructions' editor. Admin
    enforcement lives in the frontend proxy (this path is admin-gated there)."""
    import agent_prompt_store as aps
    try:
        override = await _aw(aps.get_prompt)
    except Exception:  # noqa: BLE001
        override = ""
    return {
        "prompt": override,
        # The RevOps chat's built-in base prompt. The book of deals + a fixed
        # tools/capabilities block (search_knowledge + run_todo delegation) are
        # appended automatically at runtime, so they are NOT part of this editable text.
        "default": _DEAL_ENGINE_CHAT_SYSTEM,
        "is_override": bool((override or "").strip()),
        "note": ("Base 'personality / strategy' prompt for the RevOps chat agent. The "
                 "book of deals and the tools block (search_knowledge over the shared "
                 "knowledge base + run_todo delegation to the Todo Runner) are always "
                 "appended automatically. Leave empty + save to use the built-in default."),
    }


@app.post("/api/deal-engine/chat/prompt")
async def set_chat_prompt(request: Request):
    """Persist the admin instruction override for the chat/completion agent. Send
    {"prompt": "..."}; empty clears it. Applies to the next message of every run."""
    import agent_prompt_store as aps
    try:
        body = await request.json()
    except Exception:  # noqa: BLE001
        body = {}
    prompt = (body.get("prompt") or "").strip()
    try:
        await _aw(aps.set_prompt, prompt)
        return {"ok": True, "is_override": bool(prompt)}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


from pathlib import Path as _PathTR
_TODO_RUNNER_PROMPT_PATH = _PathTR(__file__).parent / "prompts" / "todo_runner_system_prompt.md"


def _todo_runner_seed_prompt() -> str:
    """The version-controlled cold-start SEED for the todo-runner (Tactical
    Fulfillment / 'Run with AI') agent.

    DEPRECATED as the source of truth: Supabase (agent_prompt_store ID_TODO_RUNNER)
    is authoritative — edit the prompt from Admin -> Agent Control -> Todo Runner,
    NOT this file. Used only when Supabase has no row; mirrors the fallback constant
    in the frontend AgentRun.tsx. Its leading DEPRECATION banner is stripped so it
    never enters the prompt."""
    try:
        import agent_prompt_store as _aps
        return _aps.strip_leading_banner(
            _TODO_RUNNER_PROMPT_PATH.read_text(encoding="utf-8")).strip()
    except Exception:  # noqa: BLE001
        return ""


@app.get("/api/deal-engine/todo-runner/prompt")
async def get_todo_runner_prompt():
    """Return the admin override for the TODO RUNNER agent's system prompt (the
    Tactical Fulfillment / 'Run with AI' drafting agent), stored in Supabase
    (agent_prompt_store, key ID_TODO_RUNNER), plus the on-disk seed/default.

    This GET is read by BOTH the Admin editor AND the frontend agent-run panel
    (so a rep's 'Run with AI' picks up admin edits), so it is NOT admin-gated; the
    write (POST) is admin-gated at the proxy. Prompts are fetched from Supabase at
    runtime — the seed below is the version that ships on disk."""
    import agent_prompt_store as aps
    try:
        override = await _aw(aps.get_prompt, aps.ID_TODO_RUNNER)
    except Exception:  # noqa: BLE001
        override = ""
    default = await _aw(_todo_runner_seed_prompt)
    return {
        "prompt": override,
        "default": default,
        "is_override": bool((override or "").strip()),
        "note": ("Tactical Fulfillment ('Run with AI') agent system prompt. Stored "
                 "in Supabase (the runtime source of truth); leave empty + save to "
                 "fall back to the shipped default shown here. Applies to the next "
                 "'Run with AI' run, no redeploy needed."),
    }


@app.post("/api/deal-engine/todo-runner/prompt")
async def set_todo_runner_prompt(request: Request):
    """Persist the admin override for the TODO RUNNER agent's system prompt to
    Supabase (key ID_TODO_RUNNER). Send {"prompt": "..."}; empty clears it (falls
    back to the on-disk seed). Each 'Run with AI' fetches the effective prompt at
    run time, so there is no agent cache to reset."""
    import agent_prompt_store as aps
    try:
        body = await request.json()
    except Exception:  # noqa: BLE001
        body = {}
    prompt = (body.get("prompt") or "").strip()
    try:
        await _aw(aps.set_prompt, prompt, aps.ID_TODO_RUNNER)
        return {"ok": True, "is_override": bool(prompt)}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


# The seed user-message every TODO RUNNER ('Run with AI') run opens with (see the
# frontend AgentRun.buildSeedPrompt). We identify todo-runner runs by this prefix —
# the chat agent's conversations never start with it — so we can list them without a
# schema change to the shared chats/chat_messages tables.
_TODO_RUNNER_SEED_PREFIX = "Complete this sales to-do by drafting"


def _list_todo_runner_runs(limit: int) -> list:
    """Recent todo-runner runs, classified from the shared chat tables via the
    service-role client (bypasses RLS). Sync; called via _aw."""
    if not supabase:
        return []
    import re as _re
    limit = max(1, min(int(limit or 25), 50))
    # Recent chats, then their opening user message (the seed) — opening messages
    # only (sequence<=3) to keep the scan small.
    chats = supabase.table("chats").select("id,updated_at").order(
        "updated_at", desc=True).limit(200).execute()
    ids = [c["id"] for c in (chats.data or []) if c.get("id")]
    if not ids:
        return []
    opens = supabase.table("chat_messages").select(
        "chat_id,content,created_at").eq("role", "user").lte("sequence", 3).in_(
        "chat_id", ids).execute()
    runs = []
    for m in (opens.data or []):
        c = m.get("content") or ""
        if not c.startswith(_TODO_RUNNER_SEED_PREFIX):
            continue
        def _field(label: str) -> str:
            mm = _re.search(rf"^{_re.escape(label)}:\s*(.+)$", c, _re.M)
            return mm.group(1).strip() if mm else ""
        todo_m = _re.search(r"TO-DO\s*\(([^)]*)\):\s*(.+)", c)
        runs.append({
            "chat_id": m.get("chat_id"),
            "account": _field("Account"),
            "opp": _field("Opportunity"),
            "owner": _field("Deal owner (the rep you draft for)"),
            "category": (todo_m.group(1).strip() if todo_m else ""),
            "todo": (todo_m.group(2).strip() if todo_m else c[:120]),
            "status": "running",
            "created_at": m.get("created_at"),
        })
    runs.sort(key=lambda r: r.get("created_at") or "", reverse=True)
    runs = runs[:limit]
    # Terminal messages -> status (draft_ready / needs_human / error).
    run_ids = [r["chat_id"] for r in runs if r.get("chat_id")]
    if run_ids:
        try:
            terms = supabase.table("chat_messages").select("chat_id,type,content").in_(
                "chat_id", run_ids).in_("type", ["final", "error"]).execute()
            st = {}
            for t in (terms.data or []):
                cid, ty, tc = t.get("chat_id"), t.get("type"), (t.get("content") or "")
                if ty == "error":
                    st[cid] = "error"
                elif ty == "final" and st.get(cid) != "error":
                    st[cid] = "needs_human" if _re.match(r"\s*NEEDS HUMAN", tc, _re.I) else "draft_ready"
            for r in runs:
                r["status"] = st.get(r["chat_id"], "running")
        except Exception:  # noqa: BLE001 — status is best-effort
            pass
    return runs


@app.get("/api/deal-engine/todo-runner/runs")
async def list_todo_runner_runs(limit: int = 25):
    """Recent TODO RUNNER ('Run with AI') agent runs for the Admin -> Execution
    view, kept separate from the deal-sweep runs. Identified by the run's seed
    user-message; status derived from the run's terminal message. Admin-gated at
    the proxy."""
    try:
        runs = await _aw(_list_todo_runner_runs, limit)
        return {"runs": runs, "count": len(runs)}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e), "runs": []}, status_code=500)


@app.get("/api/deal-engine/sweep/prompt")
async def get_sweep_prompt():
    """Return the admin override for the DEAL SWEEP system prompt (stored in
    Supabase) plus the on-disk seed/default, for the Admin -> Agent Control
    'Deal Sweep Prompt' editor. This governs the Deal Intelligence Engine SWEEP
    agent (deal_engine_sweep.py) — NOT the chat / todo-runner agent. Prompts are
    fetched from the Supabase jarvis_settings table at runtime (see
    agent_prompt_store, key ID_DEAL_SWEEP); the seed below is the version that
    ships on disk. Admin enforcement lives in the frontend proxy (admin-gated)."""
    import agent_prompt_store as aps
    try:
        override = await _aw(aps.get_prompt, aps.ID_DEAL_SWEEP)
    except Exception:  # noqa: BLE001
        override = ""
    default = ""
    try:
        import deal_engine_sweep as des
        default = await _aw(des._disk_prompt)
    except Exception:  # noqa: BLE001 — seed unreadable; editor still works on the override
        default = ""
    return {
        "prompt": override,
        "default": default,
        "is_override": bool((override or "").strip()),
        "note": ("Deal Intelligence Engine sweep system prompt. Stored in Supabase "
                 "(the runtime source of truth); leave empty + save to fall back to "
                 "the shipped default shown here. Takes effect on the next "
                 "opportunity swept, no redeploy needed."),
    }


@app.post("/api/deal-engine/sweep/prompt")
async def set_sweep_prompt(request: Request):
    """Persist the admin override for the DEAL SWEEP system prompt to Supabase
    (key ID_DEAL_SWEEP). Send {"prompt": "..."}; empty clears it (falls back to the
    on-disk seed). The cached sweep agent is reset so the new prompt rebuilds on the
    next opportunity in this process; other processes (the worker) pick it up via
    the TTL re-check in deal_engine_sweep._get_agent."""
    import agent_prompt_store as aps
    try:
        body = await request.json()
    except Exception:  # noqa: BLE001
        body = {}
    prompt = (body.get("prompt") or "").strip()
    try:
        await _aw(aps.set_prompt, prompt, aps.ID_DEAL_SWEEP)
        try:
            import deal_engine_sweep as des
            des.reset()
        except Exception:  # noqa: BLE001 — save succeeded; cache clears via TTL anyway
            pass
        return {"ok": True, "is_override": bool(prompt)}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/knowledge/presign")
async def mase_knowledge_presign(request_body: dict):
    """Mint a short-lived presigned S3 PUT URL so the browser can upload a large file
    directly to S3 (bypassing the ~4.5 MB Vercel proxy body cap and Supabase Storage
    limits). The browser then calls POST /knowledge with the returned `key` as `s3_key`.
    Admin-gated at the proxy (path starts with `knowledge`)."""
    filename = (request_body.get("filename") or "").strip()
    if not filename:
        raise HTTPException(status_code=400, detail="filename is required")
    safe = re.sub(r"[^A-Za-z0-9._-]", "_", filename)[-180:] or "file"
    key = f"uploads/{uuid.uuid4().hex}/{safe}"
    try:
        _loop = asyncio.get_event_loop()
        url = await _loop.run_in_executor(None, _s3_presign_put, key)
    except Exception as _e:  # noqa: BLE001
        raise HTTPException(status_code=500, detail=f"Could not create upload URL: {_e}")
    return {"url": url, "key": key, "bucket": _S3_BUCKET, "expires_in": _PRESIGN_EXPIRY_S}


@app.post("/api/deal-engine/knowledge")
async def mase_knowledge_upload(request_body: dict):
    """Upload a doc into MASE's ISOLATED knowledge store (mase_documents/mase_document_chunks)
    — completely separate from VIBE's documents/projects. Accepts {name|title, doc_type,
    content}, a base64 file (file_b64 + filename) for small inline uploads, or an
    `s3_key` (+ filename) for a file the browser already PUT to S3 via /knowledge/presign.
    Admin-gated at the proxy."""
    import mase_knowledge as mk
    name = (request_body.get("name") or request_body.get("title") or "").strip()
    doc_type = request_body.get("doc_type")
    content = request_body.get("content") or ""
    file_b64 = request_body.get("file_b64")
    s3_key = request_body.get("s3_key")
    file_name = request_body.get("filename") or name
    if not content and s3_key:
        # Large-file path: pull the object from S3, extract text, then delete the temp object.
        _loop = asyncio.get_event_loop()
        try:
            raw = await asyncio.wait_for(_loop.run_in_executor(None, _s3_download, s3_key), timeout=300)
        except asyncio.TimeoutError:
            raise HTTPException(status_code=400, detail=f"Download of '{file_name}' from storage timed out")
        except Exception as _e:  # noqa: BLE001
            raise HTTPException(status_code=400, detail=f"Could not read '{file_name}' from storage: {_e}")
        try:
            content = await asyncio.wait_for(
                _loop.run_in_executor(None, _extract_text_from_bytes, raw, file_name), timeout=180)
        except asyncio.TimeoutError:
            _loop.run_in_executor(None, _s3_delete, s3_key)
            raise HTTPException(status_code=400, detail=f"Text extraction from '{file_name}' timed out")
        except Exception as _e:  # noqa: BLE001
            _loop.run_in_executor(None, _s3_delete, s3_key)
            raise HTTPException(status_code=400, detail=f"Could not extract text from '{file_name}': {_e}")
        _loop.run_in_executor(None, _s3_delete, s3_key)
    elif not content and file_b64:
        try:
            _loop = asyncio.get_event_loop()
            content = await asyncio.wait_for(
                _loop.run_in_executor(None, _extract_text_from_file, file_b64, file_name), timeout=120)
        except asyncio.TimeoutError:
            raise HTTPException(status_code=400, detail=f"Text extraction from '{file_name}' timed out")
        except HTTPException:
            raise
        except Exception as _e:  # noqa: BLE001
            raise HTTPException(status_code=400, detail=f"Could not extract text from '{file_name}': {_e}")
    if not (content or "").strip():
        if s3_key or file_b64:
            raise HTTPException(
                status_code=400,
                detail=(f"No text could be extracted from '{file_name}'. It may be image-only/"
                        "scanned, empty, or password-protected. Upload a text-based version, or "
                        "paste the text directly."))
        raise HTTPException(status_code=400, detail="Document content is required (text, a base64 file in file_b64, or an s3_key)")
    if not name:
        raise HTTPException(status_code=400, detail="Document name/title is required")
    if not supabase:
        raise HTTPException(status_code=500, detail="Supabase not configured")
    try:
        res = await mk.upload(supabase, name=name, content=content, doc_type=doc_type)
        return {"status": "success", **res}
    except HTTPException:
        raise
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": f"Upload failed: {e}"}, status_code=500)


@app.get("/api/deal-engine/knowledge")
async def mase_knowledge_list():
    """List docs in MASE's isolated knowledge store. Admin-gated at the proxy."""
    import mase_knowledge as mk
    if not supabase:
        return {"documents": []}
    try:
        return {"documents": await mk.list_docs(supabase)}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e), "documents": []}, status_code=500)


@app.get("/api/deal-engine/knowledge/{doc_id}")
async def mase_knowledge_get(doc_id: str):
    """Return one MASE knowledge doc + its full reconstructed text (for the viewer modal).
    Admin-gated at the proxy."""
    import mase_knowledge as mk
    if not supabase:
        raise HTTPException(status_code=500, detail="Supabase not configured")
    try:
        doc = await mk.get_doc(supabase, doc_id)
        if not doc:
            return JSONResponse({"error": "not found"}, status_code=404)
        return doc
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.delete("/api/deal-engine/knowledge/{doc_id}")
async def mase_knowledge_delete(doc_id: str):
    """Delete a doc (+ its chunks) from MASE's isolated knowledge store. Admin-gated."""
    import mase_knowledge as mk
    if not supabase:
        raise HTTPException(status_code=500, detail="Supabase not configured")
    try:
        await mk.delete_doc(supabase, doc_id)
        return {"status": "deleted", "id": doc_id}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/deals-count")
async def deal_engine_deals_count():
    """Total number of tracked (active) deals — for the Admin panel stat."""
    import deal_engine_store as dstore
    try:
        return {"count": await _aw(dstore.count_active_records)}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e), "count": None}, status_code=500)


@app.get("/api/deal-engine/opportunities")
async def deal_engine_opportunities(owner: str = "", slim: bool = False, paged: bool = False,
                                    q: str = "", sort: str = "close_date", dir: str = "asc",
                                    limit: int = 50, offset: int = 0, owners: str = ""):
    import deal_engine_store as dstore
    try:
        if paged:
            # Server-side one-page slice for the Deals table: scope (owners) + search (q)
            # + sort + range all run in Postgres, so the request returns ONE page + the
            # total, instead of the whole book.
            owner_list = [o.strip() for o in owners.split(",") if o.strip()] if owners else []
            if owner and not owner_list:
                owner_list = [owner]
            recs, total = await _aw(
                dstore.list_records_page, owners=owner_list or None, q=q,
                sort=sort, direction=dir, limit=limit, offset=offset)
            return {"records": recs, "total": total, "count": len(recs)}
        records = await _aw(dstore.list_records, owner or None)
        if slim:
            # LIST/aggregate payload: attach pulse (from the full record) then strip the
            # heavy ai narratives — the Deals list, Matcha and filters only need hard +
            # verdict + ai-fit + pulse. The full ai is fetched per-deal on drawer open
            # (GET /opportunities/{opp_id}). ~10-25x smaller -> fast first load.
            records = [dstore.slim_record(dstore.attach_pulse(r)) for r in records]
        else:
            # Frontend contract: every record carries `pulse`; also stamp each
            # recommended_move with its todo_key + apply user edit/delete overrides
            # (read the override table ONCE for the whole list).
            ovr = await _aw(dstore._overrides_index)
            records = [dstore.stamp_move_overrides(dstore.attach_pulse(r), ovr) for r in records]
        return {"count": len(records), "records": records}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/opportunities/{opp_id}")
async def deal_engine_opportunity(opp_id: str):
    import deal_engine_store as dstore
    try:
        rec = await _aw(dstore.get_record, opp_id)
        if not rec:
            return JSONResponse({"error": "opportunity not found"}, status_code=404)
        # Frontend contract: `record["pulse"]` always present; also stamp
        # recommended_moves with todo_key + apply user edit/delete overrides.
        return dstore.stamp_move_overrides(dstore.attach_pulse(rec))
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/todo")
async def deal_engine_todo(owner: str = ""):
    import deal_engine_store as dstore
    try:
        return await _aw(dstore.derive_todo, owner or None)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


def _build_todo_task_fields(category: str, item: dict) -> tuple[str, str]:
    """Build the Salesforce Task Subject + Description for a ticked to-do.

    Subject = the to-do's primary action text for its category (truncated to
    Salesforce's 255-char limit by the writer). Description = supporting context
    (category, trigger/due/date, expected effect or grounding quote, the rep /
    intervention owner). Mirrors the field intent of the create_task MCP tool."""
    cat = (category or "").strip()
    primary = {
        "critical": "action",
        "important": "commitment",
        "explicitRequirements": "requirement",
        "implicit": "inferred_need",
        "bestPractice": "flag",
    }.get(cat)
    # Fall back across the known text fields so a mislabelled category still
    # produces a usable subject.
    subject = ""
    if primary:
        subject = str(item.get(primary) or "").strip()
    if not subject:
        for k in ("action", "commitment", "requirement", "inferred_need",
                  "flag", "subject"):
            v = str(item.get(k) or "").strip()
            if v:
                subject = v
                break

    _CAT_LABELS = {
        "critical": "Critical move",
        "important": "Open commitment",
        "explicitRequirements": "Explicit requirement",
        "implicit": "Implicit need",
        "bestPractice": "Best-practice flag",
    }
    lines: list[str] = []
    label = _CAT_LABELS.get(cat, cat or "To-do")
    lines.append(f"Deal Engine to-do ({label}) pushed from Espresso.")
    acct = str(item.get("account_name") or "").strip()
    opp = str(item.get("opp_name") or "").strip()
    if acct or opp:
        lines.append(f"Deal: {acct}{(' — ' + opp) if opp else ''}".strip(" —"))
    # Dated context (whichever applies to this category).
    for lbl, key in (("Trigger", "trigger"), ("Trigger date", "trigger_date"),
                     ("Act by", "act_by"), ("Due", "due"), ("Date", "date"),
                     ("Raised by", "said_by"), ("Status", "status")):
        v = str(item.get(key) or "").strip()
        if v:
            lines.append(f"{lbl}: {v}")
    eff = str(item.get("expected_effect") or "").strip()
    if eff:
        lines.append(f"Expected effect: {eff}")
    quote = str(item.get("grounding_quote") or "").strip()
    if quote:
        lines.append(f"Grounding quote: {quote}")
    who = str(item.get("intervention_owner") or item.get("owner_name") or "").strip()
    if who:
        lines.append(f"Owner: {who}")
    return subject, "\n".join(lines)


@app.post("/api/deal-engine/todo/push")
async def deal_engine_todo_push(request: Request):
    """Push a ticked Espresso to-do to Salesforce as a COMPLETED Task on its
    Opportunity (human-confirmed). Idempotent by todo_key: a second push of the
    same to-do returns the existing Task (already_pushed=true) without a second
    Salesforce write. The write is a DIRECT server-side simple-salesforce call,
    NOT the agent tool catalog, so the Salesforce write lockdown stays intact.

    Body: {todo_key, opp_id, category, <display fields: action/commitment/
    requirement/inferred_need/flag + dates + context>, pushed_by?}.
    Returns {ok, sf_task_id, already_pushed, pushed_at, subject}."""
    from datetime import date as _date
    import deal_engine_store as dstore
    import salesforce_task_writer as sfw
    try:
        d = await request.json()
    except Exception:  # noqa: BLE001
        return JSONResponse({"error": "invalid JSON body"}, status_code=400)
    if not isinstance(d, dict):
        return JSONResponse({"error": "body must be a JSON object"}, status_code=400)
    key = str(d.get("todo_key") or "").strip()
    opp_id = str(d.get("opp_id") or "").strip()
    category = str(d.get("category") or "").strip()
    pushed_by = str(d.get("pushed_by") or "").strip() or None
    if not key:
        return JSONResponse({"error": "todo_key required"}, status_code=400)
    if not opp_id:
        return JSONResponse({"error": "opp_id required"}, status_code=400)
    try:
        # Idempotency: if already pushed, return the existing Task (no 2nd write).
        existing = await _aw(dstore.get_push, key)
        if existing:
            return {"ok": True, "already_pushed": True,
                    "sf_task_id": existing.get("sf_task_id"),
                    "todo_key": key, "pushed_at": existing.get("pushed_at"),
                    "subject": existing.get("subject")}
        # Resolve the Opportunity record (handles the 15- vs 18-char id join).
        rec = await _aw(dstore.get_record, opp_id)
        if not rec:
            return JSONResponse({"error": f"opportunity not found: {opp_id}"},
                                status_code=404)
        what_id = (rec.get("opp_id") or opp_id)
        subject, description = _build_todo_task_fields(category, d)
        if not subject:
            return JSONResponse(
                {"error": "could not build a Task subject from the to-do "
                          "(no action/commitment/requirement/need/flag text)"},
                status_code=400)
        who_id = str(d.get("who_id") or "").strip() or None
        # Per-user push: the proxy injects the rep's OAuth token + instance_url
        # when they've connected their Salesforce, so the Task is created AS the
        # rep (CreatedBy + Owner = them). If the rep is connected but their token
        # fails, we DO NOT silently fall back to the shared user (that's exactly
        # the "shows as Aleen" problem) — we surface a reconnect error. Only an
        # UNconnected caller falls back to the shared integration user.
        sf_token = str(d.get("sf_access_token") or "").strip() or None
        sf_instance = str(d.get("sf_instance_url") or "").strip() or None
        try:
            if sf_token and sf_instance:
                result = await _aw(
                    sfw.create_completed_task_oauth,
                    access_token=sf_token, instance_url=sf_instance,
                    subject=subject, what_id=what_id,
                    activity_date=_date.today().isoformat(),
                    description=description, who_id=who_id,
                )
            else:
                result = await _aw(
                    sfw.create_completed_task,
                    subject=subject, what_id=what_id,
                    activity_date=_date.today().isoformat(),
                    description=description, who_id=who_id,
                )
        except sfw.SalesforceWriteError as e:
            msg = (f"Salesforce write failed — reconnect your Salesforce account: {e}"
                   if sf_token else f"Salesforce write failed: {e}")
            return JSONResponse({"error": msg}, status_code=502)
        if not (isinstance(result, dict) and result.get("success")
                and result.get("id")):
            # Do NOT persist a push record on failure, so the rep can retry.
            return JSONResponse(
                {"error": "Salesforce did not confirm Task creation",
                 "salesforce": result}, status_code=502)
        sf_task_id = result.get("id")
        rec_subject = sfw.truncate_subject(subject)
        push = await _aw(dstore.insert_push, todo_key=key, opp_id=what_id,
                         category=category or None, subject=rec_subject,
                         sf_task_id=sf_task_id, pushed_by=pushed_by, payload=d)
        return {"ok": True, "already_pushed": False, "sf_task_id": sf_task_id,
                "todo_key": key, "pushed_at": push.get("pushed_at"),
                "subject": rec_subject}
    except dstore.DealEngineError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/todo/override")
async def deal_engine_todo_override(request: Request):
    """Persist a user EDIT or DELETE of an AI-derived to-do, keyed by todo_key so it
    survives the daily re-sweep (the regenerated item with the same key picks the
    override back up). Body: {opp_id, todo_key, action:'edit'|'delete', text?, due?,
    by?}. For 'delete', text/due are ignored."""
    import deal_engine_store as dstore
    try:
        d = await request.json()
    except Exception:  # noqa: BLE001
        return JSONResponse({"error": "invalid JSON body"}, status_code=400)
    if not isinstance(d, dict):
        return JSONResponse({"error": "body must be a JSON object"}, status_code=400)
    key = str(d.get("todo_key") or "").strip()
    opp_id = str(d.get("opp_id") or "").strip()
    action = str(d.get("action") or "").strip()
    if not key:
        return JSONResponse({"error": "todo_key required"}, status_code=400)
    if not opp_id:
        return JSONResponse({"error": "opp_id required"}, status_code=400)
    if action not in ("edit", "delete"):
        return JSONResponse({"error": "action must be 'edit' or 'delete'"}, status_code=400)
    text = (str(d.get("text") or "").strip() or None) if action == "edit" else None
    due = (str(d.get("due") or "").strip() or None) if action == "edit" else None
    if action == "edit" and not text:
        return JSONResponse({"error": "text required for an edit"}, status_code=400)
    by = str(d.get("by") or d.get("pushed_by") or "").strip() or None
    # capture what/where so the Learning Observatory can mine deletes & edits
    cat = str(d.get("category") or "").strip() or None
    orig = str(d.get("orig_text") or d.get("text") or "").strip() or None
    stage = str(d.get("stage") or "").strip() or None
    try:
        await _aw(dstore.upsert_override, todo_key=key, opp_id=opp_id, action=action,
                  edited_text=text, edited_due=due, created_by=by,
                  category=cat, orig_text=orig, stage=stage)
        return {"ok": True, "todo_key": key, "action": action}
    except dstore.DealEngineError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/todo/override/clear")
async def deal_engine_todo_override_clear(request: Request):
    """Undo an edit/delete override for one todo_key. Body: {todo_key}."""
    import deal_engine_store as dstore
    try:
        d = await request.json()
    except Exception:  # noqa: BLE001
        return JSONResponse({"error": "invalid JSON body"}, status_code=400)
    key = str((d or {}).get("todo_key") or "").strip()
    if not key:
        return JSONResponse({"error": "todo_key required"}, status_code=400)
    try:
        await _aw(dstore.clear_override, key)
        return {"ok": True, "todo_key": key}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/todo/update")
async def deal_engine_todo_update(request: Request):
    """Add a manual update on a deal, BRANCHING to one of three destinations. All
    writes are direct, human-initiated simple-salesforce calls (NOT the agent tool
    catalog), so the Salesforce write lockdown stays intact. The in-app row is
    always saved so the update is never lost.

    Body: {opp_id, note, destination?, due_date?, done_date?, by?,
           sf_access_token?, sf_instance_url?}.
      destination = 'completed' (default) -> Completed Salesforce Task on the opp.
                  = 'todo'                -> MASE to-do + OPEN Salesforce Task
                                             (Status='Planned', ActivityDate=due_date).
                  = 'next_step'           -> APPEND to Opportunity.Next_Step__c,
                                             newest on top, full prior trail kept.
    Every branch carries a due date (completed defaults to today). Per-user OAuth
    token (sf_access_token/sf_instance_url) makes the rep the author when present.
    Returns {ok, destination, update, sf_task_id?, next_step_updated?, sf_error?}."""
    from datetime import date as _date
    import html as _html
    import deal_engine_store as dstore
    import salesforce_task_writer as sfw
    try:
        d = await request.json()
    except Exception:  # noqa: BLE001
        return JSONResponse({"error": "invalid JSON body"}, status_code=400)
    if not isinstance(d, dict):
        return JSONResponse({"error": "body must be a JSON object"}, status_code=400)
    opp_id = str(d.get("opp_id") or "").strip()
    note = str(d.get("note") or "").strip()
    destination = str(d.get("destination") or "completed").strip().lower()
    if destination not in ("completed", "todo", "next_step"):
        return JSONResponse(
            {"error": "destination must be completed | todo | next_step"},
            status_code=400)
    due_date = str(d.get("due_date") or "").strip()
    done_date = str(d.get("done_date") or "").strip() or _date.today().isoformat()
    by = str(d.get("by") or d.get("pushed_by") or "").strip() or None
    sf_token = str(d.get("sf_access_token") or "").strip() or None
    sf_instance = str(d.get("sf_instance_url") or "").strip() or None
    if not opp_id:
        return JSONResponse({"error": "opp_id required"}, status_code=400)
    if not note:
        return JSONResponse({"error": "note required"}, status_code=400)
    if destination in ("todo", "next_step") and not due_date:
        return JSONResponse(
            {"error": "due_date required for the todo / next_step destinations"},
            status_code=400)

    def _fmt(iso: str) -> str:
        try:
            return _date.fromisoformat(iso).strftime("%b %d, %Y")
        except Exception:  # noqa: BLE001
            return iso or ""

    def _sf_err(e) -> str:
        return (f"Salesforce write failed — reconnect your Salesforce account: {e}"
                if sf_token else f"Salesforce write failed: {e}")

    try:
        rec = await _aw(dstore.get_record, opp_id)
        what_id = (rec.get("opp_id") if rec else None) or opp_id
        sf_task_id = None
        sf_error = None
        next_step_updated = False

        if destination == "next_step":
            # Newest-on-top append to the HTML rich-text Next_Step__c, full prior
            # trail preserved. Entry is an HTML <p> block matching the field's
            # existing convention; note text is escaped.
            stamp = f"{_fmt(done_date)} (due {_fmt(due_date)})"
            entry = f"<p>{stamp}: {_html.escape(note)}</p>"
            try:
                if sf_token and sf_instance:
                    await _aw(sfw.append_next_step_oauth, access_token=sf_token,
                              instance_url=sf_instance, opp_id=what_id, entry=entry)
                else:
                    await _aw(sfw.append_next_step, opp_id=what_id, entry=entry)
                next_step_updated = True
            except sfw.SalesforceWriteError as e:
                sf_error = _sf_err(e)
        elif destination == "todo":
            # MASE to-do (the in-app row below) + an OPEN Salesforce activity.
            subject = sfw.truncate_subject(note)
            description = ("Deal Engine to-do (open) logged from MASE.\n"
                           f"To-do: {note}\nDue: {_fmt(due_date)}"
                           + (f"\nBy: {by}" if by else ""))
            try:
                if sf_token and sf_instance:
                    result = await _aw(sfw.create_open_task_oauth,
                                       access_token=sf_token, instance_url=sf_instance,
                                       subject=subject, what_id=what_id,
                                       activity_date=due_date, description=description,
                                       who_id=None)
                else:
                    result = await _aw(sfw.create_open_task, subject=subject,
                                       what_id=what_id, activity_date=due_date,
                                       description=description, who_id=None)
                if isinstance(result, dict) and result.get("success") and result.get("id"):
                    sf_task_id = result.get("id")
                else:
                    sf_error = "Salesforce did not confirm Task creation"
            except sfw.SalesforceWriteError as e:
                sf_error = _sf_err(e)
        else:  # 'completed' — unchanged behaviour (a completed Task on the opp)
            subject = sfw.truncate_subject(note)
            description = ("Deal Engine manual update (completed) logged from MASE.\n"
                           f"Update: {note}\nDone: {done_date}"
                           + (f"\nBy: {by}" if by else ""))
            act = due_date or done_date
            try:
                if sf_token and sf_instance:
                    result = await _aw(sfw.create_completed_task_oauth,
                                       access_token=sf_token, instance_url=sf_instance,
                                       subject=subject, what_id=what_id,
                                       activity_date=act, description=description,
                                       who_id=None)
                else:
                    result = await _aw(sfw.create_completed_task, subject=subject,
                                       what_id=what_id, activity_date=act,
                                       description=description, who_id=None)
                if isinstance(result, dict) and result.get("success") and result.get("id"):
                    sf_task_id = result.get("id")
                else:
                    sf_error = "Salesforce did not confirm Task creation"
            except sfw.SalesforceWriteError as e:
                sf_error = _sf_err(e)

        # Always persist the in-app row so the update is never lost — this is also
        # the MASE-side record for the 'todo' and 'next_step' destinations.
        note_for_log = note if destination == "completed" else f"[{destination}] {note}"
        row = await _aw(dstore.insert_manual_update, opp_id=what_id, note=note_for_log,
                        done_date=(due_date or done_date), sf_task_id=sf_task_id,
                        created_by=by)
        return {"ok": True, "destination": destination, "update": row,
                "sf_task_id": sf_task_id, "next_step_updated": next_step_updated,
                "sf_error": sf_error}
    except dstore.DealEngineError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/learnings")
async def deal_engine_learnings(status: str = ""):
    """List Learning Observatory entries (optionally by status). Each carries
    category, stage_scope, scope, status, source, evidence, weight."""
    import deal_engine_store as dstore
    try:
        return {"learnings": await _aw(dstore.list_learnings, status or None)}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/learnings/signals")
async def deal_engine_learning_signals():
    """The raw operator-behaviour signals the observatory learns from — deleted /
    edited / completed to-dos and manual updates, grouped by stage + category. Feeds
    both the UI 'Signals' panel and the daily miner."""
    import deal_engine_store as dstore
    try:
        return await _aw(dstore.mine_signals)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/learnings")
async def deal_engine_learning_create(request: Request):
    """Add a learning. Admin 'doc entry point' (source=manual) and the daily miner
    (source=mined) both use this. Defaults to status='candidate' — promotion to
    'active' is an explicit switch via the PATCH route. Body: {title, body, category?,
    stage_scope?, scope?, scope_selector?, status?, source?, evidence?, weight?, by?}."""
    import deal_engine_store as dstore
    try:
        d = await request.json()
    except Exception:  # noqa: BLE001
        return JSONResponse({"error": "invalid JSON body"}, status_code=400)
    if not isinstance(d, dict):
        return JSONResponse({"error": "body must be a JSON object"}, status_code=400)
    title = str(d.get("title") or "").strip()
    body = str(d.get("body") or "").strip()
    if not title or not body:
        return JSONResponse({"error": "title and body are required"}, status_code=400)
    allowed_cat = {"risk", "equity", "assurance", "differentiation", "general"}
    category = str(d.get("category") or "general").strip().lower()
    if category not in allowed_cat:
        category = "general"
    allowed_status = {"candidate", "active", "paused", "retired"}
    status = str(d.get("status") or "candidate").strip().lower()
    if status not in allowed_status:
        status = "candidate"
    source = "mined" if str(d.get("source") or "").strip().lower() == "mined" else "manual"
    try:
        row = await _aw(dstore.insert_learning, title=title, body=body, category=category,
                        stage_scope=str(d.get("stage_scope") or "any").strip(),
                        scope=str(d.get("scope") or "global").strip(),
                        scope_selector=d.get("scope_selector") if isinstance(d.get("scope_selector"), dict) else {},
                        status=status, source=source,
                        evidence=d.get("evidence") if isinstance(d.get("evidence"), list) else [],
                        weight=int(d.get("weight") or 0),
                        created_by=str(d.get("by") or "").strip() or None)
        return {"ok": True, "learning": row}
    except dstore.DealEngineError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/learnings/{learning_id}")
async def deal_engine_learning_update(learning_id: str, request: Request):
    """Patch one learning — primarily to flip the switch (status: candidate -> active
    -> paused -> retired) or tweak body/category/stage_scope/weight. Body: any subset
    of {status, title, body, category, stage_scope, scope, weight}."""
    import deal_engine_store as dstore
    try:
        d = await request.json()
    except Exception:  # noqa: BLE001
        return JSONResponse({"error": "invalid JSON body"}, status_code=400)
    if not isinstance(d, dict) or not d:
        return JSONResponse({"error": "non-empty JSON body required"}, status_code=400)
    patch = {k: d[k] for k in ("status", "title", "body", "category", "stage_scope", "scope", "weight") if k in d}
    if "status" in patch and patch["status"] not in ("candidate", "active", "paused", "retired"):
        return JSONResponse({"error": "invalid status"}, status_code=400)
    if not patch:
        return JSONResponse({"error": "no updatable fields in body"}, status_code=400)
    try:
        row = await _aw(dstore.update_learning, learning_id, patch)
        if not row:
            return JSONResponse({"error": "learning not found"}, status_code=404)
        return {"ok": True, "learning": row}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/matcha")
async def deal_engine_matcha(owner: str = ""):
    import deal_engine_store as dstore
    try:
        return await _aw(dstore.derive_matcha, owner or None)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/deltas")
async def deal_engine_deltas(owner: str = "", limit: int = 200, group_by: str = ""):
    """The book-wide 'what changed' feed (newest first). Each entry carries its
    deal context (opp_id/account/opp/owner/stage for the back-link) plus a human
    `label` and a rep-facing `group` (added/changed/resolved/dormant), with a
    per-owner roll-up. Optional ?owner= scopes to one RSD; ?limit= caps the feed;
    ?group_by=owner adds a `groups` array bucketed by deal owner."""
    import deal_engine_store as dstore
    try:
        return await _aw(dstore.list_deltas, owner or None, limit, group_by or None)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/deltas/{opp_id}")
async def deal_engine_deltas_one(opp_id: str):
    """One deal's 'What changed' panel: deal context + labelled change log
    (newest first) with per-group counts for tab badges."""
    import deal_engine_store as dstore
    try:
        return await _aw(dstore.get_deltas_view, opp_id)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/records")
async def deal_engine_upsert_record(request: Request):
    """Upsert one canonical deal record (used by the sweep agent / ingest)."""
    import deal_engine_store as dstore
    try:
        d = await request.json()
        if not (d.get("opp_id") or "").strip():
            return JSONResponse({"error": "opp_id required"}, status_code=400)
        return await _aw(dstore.upsert_record, d)
    except dstore.DealEngineError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/backfill-packets")
async def deal_engine_backfill_packets(request: Request):
    """One-time, idempotent migration: seed a living-memory packets baseline onto
    every deal record that predates living memory, so the 'What changed' feed is
    populated across the whole book immediately (instead of only after each deal's
    natural next sweep). Seeding does NOT emit `added` deltas — pre-existing facts
    are not treated as changes. Body (optional): {dry_run?: bool}. Safe to re-run:
    records already carrying a packet baseline (schema_version >= 2) are skipped."""
    import deal_engine_store as dstore
    try:
        d = {}
        try:
            d = await request.json()
        except Exception:  # noqa: BLE001
            pass
        dry_run = bool(d.get("dry_run"))
        stats = await _aw(dstore.backfill_packets, dry_run=dry_run)
        return {"status": "ok", "dry_run": dry_run, "stats": stats}
    except dstore.DealEngineError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/sweep")
async def deal_engine_sweep_start(request: Request):
    """Kick off the AI sweep that POPULATES the book. Body (all optional):
    {owner?, opp_ids?: [..], limit?}. Discovers open opps for the team (or one
    owner), runs a Salesforce+Avoma agent per opp, and upserts canonical records.
    Heavy + token-costly; one sweep runs at a time. Returns {run_id, status}."""
    import deal_engine_sweep as sweep
    try:
        d = {}
        try:
            d = await request.json()
        except Exception:  # noqa: BLE001
            pass
        owner = (d.get("owner") or "").strip() or None
        opp_ids = d.get("opp_ids") if isinstance(d.get("opp_ids"), list) else None
        limit = int(d.get("limit") or 500)
        concurrency = int(d["concurrency"]) if d.get("concurrency") is not None else None
        max_retries = int(d["max_retries"]) if d.get("max_retries") is not None else None
        header = await sweep.start_sweep(agent_manager, owner=owner,
                                         opp_ids=opp_ids, limit=limit,
                                         concurrency=concurrency, max_retries=max_retries)
        return header
    except RuntimeError as e:
        return JSONResponse({"error": str(e)}, status_code=409)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/sweep/trigger")
async def deal_engine_sweep_trigger(request: Request):
    """Salesforce-update webhook: re-run the analysis for one (or a few) opps.

    Intended for a Salesforce Flow / outbound message that fires when an
    Opportunity changes and POSTs the record id here. Returns 202 IMMEDIATELY
    and runs the analysis in the background (each opp takes minutes — far longer
    than a Salesforce callout will wait), so the Flow doesn't time out. Repeat
    updates for an opp already being analyzed are deduped.

    Auth: same gate as the rest of /api/deal-engine (Authorization: Bearer
    <token>, or ?key=<token>). Body accepts the opp id under any of: opp_id,
    oppId, id, Id, opportunity_id, opportunityId, OpportunityId, recordId — or
    a list under opp_ids.
    """
    import deal_engine_sweep as sweep
    try:
        d = {}
        try:
            d = await request.json()
        except Exception:  # noqa: BLE001
            pass
        opp_id = None
        for k in ("opp_id", "oppId", "id", "Id", "opportunity_id",
                  "opportunityId", "OpportunityId", "recordId"):
            v = d.get(k)
            if isinstance(v, str) and v.strip():
                opp_id = v.strip()
                break
        ids = [str(i).strip() for i in d.get("opp_ids", [])
               if str(i).strip()] if isinstance(d.get("opp_ids"), list) else []
        if opp_id and opp_id not in ids:
            ids.insert(0, opp_id)
        if not ids:
            return JSONResponse(
                {"error": "provide an opportunity id (opp_id / id / recordId) "
                          "or opp_ids[]"}, status_code=400)
        if sweep.queue_enabled():
            # Queue mode: enqueue a durable `waiting` row per opp; the separate
            # worker.py drains it. The web process does NOT run the analysis, so
            # a burst of Salesforce updates can't starve it.
            results = {oid: await sweep.enqueue_trigger(agent_manager, oid)
                       for oid in ids}
        else:
            results = {oid: sweep.trigger_opp_async(agent_manager, oid)
                       for oid in ids}
        return JSONResponse(
            {"status": "accepted", "results": results, "count": len(ids),
             "note": "analysis runs in the background; poll the record or the "
                     "sweep dashboard to see it refresh"},
            status_code=202)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/sweep/discover")
async def deal_engine_sweep_discover(owner: str = "", limit: int = 500):
    """Dry-run discovery: the open opps that a sweep WOULD analyze for the team
    (or one owner). One SOQL query, no AI, no cost. Lets the UI preview/count
    before kicking off the heavy run."""
    import deal_engine_sweep as sweep
    try:
        opps = await sweep.discover_opps(agent_manager, owner or None, limit=limit)
        return {"count": len(opps), "owner": owner or "all-team", "opps": opps}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/sweep/discover-new")
async def deal_engine_discover_new(request: Request):
    """Discover open team opps that have NO canonical record yet and sweep them
    in the background (the same step the scheduled job runs). Returns 202
    immediately — poll GET /api/deal-engine/opportunities or the sweep dashboard
    to watch them appear. Body (all optional): {limit?, concurrency?, max_new?}.
    Skips while a full sweep is already running; capped + bounded for cost."""
    import deal_engine_sweep as sweep
    try:
        d = {}
        try:
            d = await request.json()
        except Exception:  # noqa: BLE001
            pass
        limit = int(d.get("limit") or 500)
        concurrency = int(d["concurrency"]) if d.get("concurrency") is not None else None
        max_new = int(d["max_new"]) if d.get("max_new") is not None else None
        t = asyncio.create_task(sweep.discover_and_sweep_new(
            agent_manager, limit=limit, concurrency=concurrency,
            max_new=max_new, source="manual_discovery"))
        sweep._discovery_tasks.add(t)
        t.add_done_callback(sweep._discovery_tasks.discard)
        return JSONResponse(
            {"status": "accepted",
             "note": "discovery + sweep of new opps runs in the background; poll "
                     "/api/deal-engine/opportunities or the sweep dashboard"},
            status_code=202)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/sweep/reconcile")
async def deal_engine_reconcile(request: Request):
    """Reconcile book membership against the MASE report (single source of truth)
    NOW, synchronously, and return the summary. Deactivates opps that left the
    report, reactivates re-entrants, and sweeps brand-new/re-entered opps (capped).
    Aborts safely if the report read fails/empties or shrinks the book implausibly.
    Body (all optional): {sweep_new?: bool, concurrency?: int, max_new?: int}."""
    import deal_engine_sweep as sweep
    try:
        d = {}
        try:
            d = await request.json()
        except Exception:  # noqa: BLE001
            pass
        sweep_new = d.get("sweep_new")
        sweep_new = True if sweep_new is None else bool(sweep_new)
        concurrency = int(d["concurrency"]) if d.get("concurrency") is not None else None
        max_new = int(d["max_new"]) if d.get("max_new") is not None else None
        summary = await sweep.reconcile_membership(
            agent_manager, sweep_new=sweep_new,
            concurrency=concurrency, max_new=max_new, source="manual_reconcile")
        return summary
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/hard-refresh")
async def deal_engine_hard_refresh(request: Request):
    """Token-free refresh of the HARD Salesforce fields (stage, amount, products,
    close date, owner, next step) across EVERY persisted deal record — no AI, no
    agent, ~zero token cost. Merges the live values onto each record (preserving
    the AI analysis + history) and DELETES any deal now back at 'Initial Interest'.
    Runs synchronously and returns a summary. Body (all optional):
    {delete_initial_interest?: bool, concurrency?: int}. Skips while a full sweep
    is running."""
    import deal_engine_sweep as sweep
    try:
        d = {}
        try:
            d = await request.json()
        except Exception:  # noqa: BLE001
            pass
        delete_ii = d.get("delete_initial_interest")
        delete_ii = True if delete_ii is None else bool(delete_ii)
        concurrency = int(d["concurrency"]) if d.get("concurrency") is not None else None
        summary = await sweep.hard_refresh_all(
            agent_manager,
            delete_initial_interest=delete_ii,
            concurrency=concurrency,
            source="manual",
        )
        return summary
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/hard-refresh/status")
async def deal_engine_hard_refresh_status():
    """Summary of the most recent token-free hard-fact reconciliation (records /
    matched / updated / removed / unmatched / failed / finished_at / source) and
    whether one is running right now. Survives a server restart (persisted to
    disk), so the last nightly run can be checked later."""
    import deal_engine_sweep as sweep
    try:
        return {
            "running": bool(getattr(sweep, "_hard_refresh_running", False)),
            "last": sweep.get_hard_refresh_last(),
        }
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/hard-refresh/history")
async def deal_engine_hard_refresh_history(limit: int = 200):
    """Append-only history of token-free hard-fact reconciliation runs, newest
    first. Each row records timestamp, source (manual / nightly), and the
    records / matched / updated / removed / unmatched / failed counts — so the
    nightly schedule is auditable and an anomalous run (unusually high/low
    updated/removed counts) can be spotted over time. `limit` caps the rows
    returned (default 200)."""
    import deal_hard_refresh_log as hard_refresh_log
    try:
        return {"runs": hard_refresh_log.list_runs(limit=int(limit))}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/sweep/status")
async def deal_engine_sweep_status():
    """Live progress of the current/last sweep (poll this from the UI)."""
    import deal_engine_sweep as sweep
    try:
        return await sweep.get_status()
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/sweep/dashboard", response_class=HTMLResponse)
async def deal_engine_sweep_dashboard():
    """Self-contained real-time dashboard for the sweep. Open in a browser with
    ?key=<API_AUTH_TOKEN> the first time — the auth gate sets a cookie and strips
    the key from the URL, after which the page polls /sweep/status on its own."""
    return HTMLResponse(content=r"""<!doctype html>
<html lang="en"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Deal Sweep — Live</title>
<style>
  :root { --bg:#0f1419; --panel:#1a1f2e; --line:#2a3142; --muted:#8b95a7; --txt:#d4d4d4; }
  * { box-sizing:border-box; }
  body { font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",system-ui,sans-serif; margin:0; background:var(--bg); color:var(--txt); }
  header { background:var(--panel); padding:14px 24px; border-bottom:1px solid var(--line); display:flex; gap:18px; align-items:center; flex-wrap:wrap; }
  header h1 { margin:0; font-size:17px; font-weight:600; }
  .pill { font-size:12px; padding:3px 10px; border-radius:999px; background:var(--line); color:var(--txt); }
  .pill.running { background:#1e3a5f; color:#7cc4ff; }
  .pill.succeeded { background:#1e4620; color:#7ee787; }
  .pill.partial { background:#5c4a16; color:#f0c674; }
  .pill.failed { background:#5a1e1e; color:#ff8585; }
  .pill.idle { background:var(--line); color:var(--muted); }
  main { padding:20px 24px; }
  .cards { display:flex; gap:14px; flex-wrap:wrap; margin-bottom:18px; }
  .card { background:var(--panel); border:1px solid var(--line); border-radius:8px; padding:14px 18px; min-width:120px; }
  .card .n { font-size:26px; font-weight:700; }
  .card .l { font-size:12px; color:var(--muted); margin-top:2px; text-transform:uppercase; letter-spacing:.04em; }
  .card.done .n { color:#7ee787; } .card.failed .n { color:#ff8585; }
  .card.running .n { color:#7cc4ff; } .card.queued .n { color:var(--muted); }
  .bar { height:10px; border-radius:6px; background:var(--line); overflow:hidden; display:flex; margin-bottom:18px; }
  .bar > i { display:block; height:100%; }
  .bar .b-done { background:#3fb950; } .bar .b-failed { background:#f85149; } .bar .b-run { background:#388bfd; }
  table { width:100%; border-collapse:collapse; font-size:13px; }
  th,td { text-align:left; padding:8px 10px; border-bottom:1px solid #1f2533; }
  th { color:var(--muted); font-weight:600; font-size:11px; text-transform:uppercase; letter-spacing:.04em; position:sticky; top:0; background:var(--bg); cursor:pointer; }
  td.acct { font-weight:600; }
  .s { font-size:11px; padding:2px 8px; border-radius:999px; white-space:nowrap; }
  .s-queued { background:var(--line); color:var(--muted); }
  .s-running { background:#1e3a5f; color:#7cc4ff; }
  .s-retrying { background:#5c4a16; color:#f0c674; }
  .s-completed { background:#1e4620; color:#7ee787; }
  .s-failed,.s-parse_error { background:#5a1e1e; color:#ff8585; }
  .err { color:#ff8585; font-size:11px; max-width:360px; overflow:hidden; text-overflow:ellipsis; white-space:nowrap; }
  .muted { color:var(--muted); }
  .controls { margin-left:auto; display:flex; gap:8px; align-items:center; }
  select,button { background:var(--line); color:var(--txt); border:1px solid #3a4258; padding:5px 10px; border-radius:5px; font-size:12px; cursor:pointer; }
  #updated { font-size:11px; color:var(--muted); }
  .hr-panel { background:var(--panel); border:1px solid var(--line); border-radius:8px; padding:12px 18px; margin-bottom:18px; display:flex; gap:22px; align-items:center; flex-wrap:wrap; }
  .hr-panel .hr-title { font-size:12px; font-weight:600; text-transform:uppercase; letter-spacing:.04em; color:var(--muted); }
  .hr-panel .hr-src { font-size:11px; padding:3px 10px; border-radius:999px; background:var(--line); color:var(--txt); }
  .hr-panel .hr-src.nightly_cron { background:#1e3a5f; color:#7cc4ff; }
  .hr-panel .hr-src.manual { background:#5c4a16; color:#f0c674; }
  .hr-panel .hr-when { font-size:12px; color:var(--muted); }
  .hr-panel .hr-stats { display:flex; gap:18px; margin-left:auto; flex-wrap:wrap; }
  .hr-panel .hr-stat { text-align:center; }
  .hr-panel .hr-stat .v { font-size:18px; font-weight:700; }
  .hr-panel .hr-stat .k { font-size:10px; color:var(--muted); text-transform:uppercase; letter-spacing:.04em; margin-top:1px; }
  .hr-panel .hr-stat.updated .v { color:#7ee787; }
  .hr-panel .hr-stat.removed .v { color:#f0c674; }
  .hr-panel .hr-stat.failed .v { color:#ff8585; }
  .hr-panel .hr-run { font-size:11px; padding:3px 10px; border-radius:999px; background:#1e3a5f; color:#7cc4ff; }
</style></head>
<body>
<header>
  <h1>Deal Sweep</h1>
  <span id="status" class="pill idle">idle</span>
  <span id="runid" class="muted" style="font-size:12px"></span>
  <span id="eta" class="muted" style="font-size:12px"></span>
  <div class="controls">
    <label class="muted" style="font-size:12px">Filter</label>
    <select id="filter">
      <option value="">All</option>
      <option value="running">Running</option>
      <option value="retrying">Retrying</option>
      <option value="queued">Queued</option>
      <option value="completed">Completed</option>
      <option value="failed">Failed</option>
    </select>
    <span id="updated"></span>
  </div>
</header>
<main>
  <div class="cards">
    <div class="card"><div class="n" id="c-total">0</div><div class="l">Total</div></div>
    <div class="card done"><div class="n" id="c-done">0</div><div class="l">Done</div></div>
    <div class="card running"><div class="n" id="c-run">0</div><div class="l">In progress</div></div>
    <div class="card queued"><div class="n" id="c-queued">0</div><div class="l">Queued</div></div>
    <div class="card failed"><div class="n" id="c-failed">0</div><div class="l">Failed</div></div>
  </div>
  <div class="bar"><i class="b-done" id="bar-done"></i><i class="b-run" id="bar-run"></i><i class="b-failed" id="bar-failed"></i></div>
  <div class="hr-panel" id="hr-panel" style="display:none">
    <span class="hr-title">Hard-fact sync</span>
    <span class="hr-src" id="hr-src">—</span>
    <span class="hr-when" id="hr-when"></span>
    <span class="hr-run" id="hr-run" style="display:none">running…</span>
    <div class="hr-stats">
      <div class="hr-stat"><div class="v" id="hr-records">0</div><div class="k">Records</div></div>
      <div class="hr-stat"><div class="v" id="hr-matched">0</div><div class="k">Matched</div></div>
      <div class="hr-stat updated"><div class="v" id="hr-updated">0</div><div class="k">Updated</div></div>
      <div class="hr-stat removed"><div class="v" id="hr-removed">0</div><div class="k">Removed</div></div>
      <div class="hr-stat failed"><div class="v" id="hr-failed">0</div><div class="k">Failed</div></div>
    </div>
  </div>
  <table>
    <thead><tr><th>#</th><th>Account</th><th>Owner</th><th>Opp ID</th><th>Status</th><th>Tries</th><th>Duration</th><th>Error</th></tr></thead>
    <tbody id="rows"></tbody>
  </table>
</main>
<script>
const $ = id => document.getElementById(id);
const esc = s => (s==null?'':String(s)).replace(/[&<>"']/g, c => ({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[c]));
function dur(ms){ if(!ms) return ''; const s=Math.round(ms/1000); if(s<60) return s+'s'; const m=Math.floor(s/60); return m+'m'+(s%60)+'s'; }
async function tick(){
  let d;
  try { const r = await fetch('/api/deal-engine/sweep/status', {credentials:'same-origin'});
        if(r.status===401){ $('status').textContent='unauthorized — add ?key=TOKEN'; return; }
        d = await r.json(); }
  catch(e){ $('updated').textContent='connection error'; return; }
  const st = d.status||'idle';
  const sp = $('status'); sp.textContent = st; sp.className = 'pill ' + st;
  $('runid').textContent = d.run_id ? ('run '+d.run_id) : '';
  const opps = d.opps||[];
  const total = d.total||opps.length||0;
  const done = d.done||0, failed = d.failed||0, run = d.in_progress||0;
  const queued = Math.max(0, total - done - failed - run);
  $('c-total').textContent=total; $('c-done').textContent=done;
  $('c-run').textContent=run; $('c-queued').textContent=queued; $('c-failed').textContent=failed;
  $('bar-done').style.width=(total?100*done/total:0)+'%';
  $('bar-run').style.width=(total?100*run/total:0)+'%';
  $('bar-failed').style.width=(total?100*failed/total:0)+'%';
  // simple ETA from average completed duration
  const compl = opps.filter(o=>o.status==='completed' && o.duration_ms);
  if(compl.length && (queued+run)>0 && d.concurrency){
    const avg = compl.reduce((a,o)=>a+o.duration_ms,0)/compl.length;
    const remainMs = avg * (queued+run) / Math.max(1,d.concurrency);
    $('eta').textContent = '~'+dur(remainMs)+' left · '+(d.concurrency)+' at a time';
  } else { $('eta').textContent = d.concurrency ? (d.concurrency+' at a time') : ''; }
  const f = $('filter').value;
  const rows = opps.map((o,i)=>({o,i})).filter(({o})=>!f|| (f==='retrying'?o.status==='retrying':o.status===f));
  $('rows').innerHTML = rows.map(({o,i})=>`<tr>
    <td class="muted">${i+1}</td>
    <td class="acct">${esc(o.account)||'<span class="muted">—</span>'}</td>
    <td>${esc(o.owner_name)||'<span class="muted">—</span>'}</td>
    <td class="muted">${esc(o.opp_id)}</td>
    <td><span class="s s-${esc(o.status)}">${esc(o.status)}</span></td>
    <td class="muted">${o.attempts||0}</td>
    <td class="muted">${dur(o.duration_ms)}</td>
    <td class="err" title="${esc(o.error)}">${esc(o.error)||''}</td>
  </tr>`).join('');
  $('updated').textContent = 'updated ' + new Date().toLocaleTimeString();
}
async function tickHardRefresh(){
  let d;
  try { const r = await fetch('/api/deal-engine/hard-refresh/status', {credentials:'same-origin'});
        if(!r.ok) return;
        d = await r.json(); }
  catch(e){ return; }
  const last = d && d.last;
  const running = !!(d && d.running);
  $('hr-run').style.display = running ? '' : 'none';
  if(!last || !Object.keys(last).length){
    if(running){ $('hr-panel').style.display=''; $('hr-src').style.display='none'; $('hr-when').textContent='no completed sync yet'; }
    return;
  }
  $('hr-panel').style.display='';
  $('hr-src').style.display='';
  const src = last.source || 'manual';
  $('hr-src').textContent = src;
  $('hr-src').className = 'hr-src ' + esc(src);
  const fin = last.finished_at;
  let when = '';
  if(fin){ const dt = new Date(fin); when = isNaN(dt) ? esc(fin) : ('last synced ' + dt.toLocaleString()); }
  $('hr-when').textContent = when;
  $('hr-records').textContent = last.records!=null ? last.records : 0;
  $('hr-matched').textContent = last.matched!=null ? last.matched : 0;
  $('hr-updated').textContent = last.updated!=null ? last.updated : 0;
  $('hr-removed').textContent = last.removed!=null ? last.removed : 0;
  $('hr-failed').textContent = last.failed!=null ? last.failed : 0;
}
$('filter').addEventListener('change', tick);
tick(); setInterval(tick, 2500);
tickHardRefresh(); setInterval(tickHardRefresh, 10000);
</script>
</body></html>""")


@app.get("/api/deal-engine/todo/dashboard", response_class=HTMLResponse)
async def deal_engine_todo_dashboard():
    """Self-contained Espresso (to-do) dashboard. A rep can tick a to-do and
    confirm to log a COMPLETED Salesforce Task on its Opportunity via
    POST /api/deal-engine/todo/push. Pushed state is reload-safe (the /todo feed
    carries pushed / sf_task_id from the deal_todo_pushes ledger).

    Open in a browser with ?key=<API_AUTH_TOKEN> the first time — the auth gate
    sets a cookie and strips the key from the URL, after which the page's own
    fetch calls authenticate via the cookie."""
    return HTMLResponse(content=r"""<!doctype html>
<html lang="en"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>Espresso — To-do</title>
<style>
  :root { --bg:#0f1419; --panel:#1a1f2e; --line:#2a3142; --muted:#8b95a7; --txt:#d4d4d4; }
  * { box-sizing:border-box; }
  body { font-family:-apple-system,BlinkMacSystemFont,"Segoe UI",system-ui,sans-serif; margin:0; background:var(--bg); color:var(--txt); }
  header { background:var(--panel); padding:14px 24px; border-bottom:1px solid var(--line); display:flex; gap:18px; align-items:center; flex-wrap:wrap; }
  header h1 { margin:0; font-size:17px; font-weight:600; }
  .sub { color:var(--muted); font-size:12px; }
  .controls { margin-left:auto; display:flex; gap:8px; align-items:center; }
  label.lbl { color:var(--muted); font-size:12px; }
  select,button { background:var(--line); color:var(--txt); border:1px solid #3a4258; padding:6px 11px; border-radius:6px; font-size:12px; cursor:pointer; }
  button:disabled { opacity:.5; cursor:default; }
  main { padding:20px 24px; max-width:1100px; margin:0 auto; }
  #updated { font-size:11px; color:var(--muted); }
  .empty { color:var(--muted); padding:40px 0; text-align:center; }
  section.cat { margin-bottom:26px; }
  section.cat > h2 { font-size:13px; font-weight:700; text-transform:uppercase; letter-spacing:.05em; margin:0 0 10px; display:flex; align-items:center; gap:9px; }
  section.cat > h2 .dot { width:9px; height:9px; border-radius:50%; }
  section.cat > h2 .count { font-size:11px; color:var(--muted); font-weight:600; }
  .item { background:var(--panel); border:1px solid var(--line); border-left:3px solid var(--line); border-radius:8px; padding:13px 15px; margin-bottom:9px; display:flex; gap:14px; align-items:flex-start; }
  .item .body { flex:1; min-width:0; }
  .item .txt { font-size:14px; font-weight:600; line-height:1.4; }
  .item .deal { font-size:12px; color:var(--muted); margin-top:3px; }
  .item .meta { margin-top:7px; display:flex; gap:6px; flex-wrap:wrap; }
  .chip { font-size:10.5px; padding:2px 8px; border-radius:999px; background:var(--line); color:var(--muted); white-space:nowrap; }
  .chip.u-overdue { background:#5a1e1e; color:#ff8585; }
  .chip.u-next_14_days { background:#5c4a16; color:#f0c674; }
  .chip.u-next_30_days { background:#1e3a5f; color:#7cc4ff; }
  .act { flex-shrink:0; display:flex; flex-direction:column; gap:6px; align-items:flex-end; }
  .pushbtn { background:#1e4620; color:#7ee787; border-color:#2c5e30; font-weight:600; }
  .pushbtn:hover:not(:disabled) { background:#235426; }
  .pushed { font-size:11px; color:#7ee787; display:flex; align-items:center; gap:5px; white-space:nowrap; }
  .pushed a { color:#7ee787; }
  .taskid { font-size:10px; color:var(--muted); font-family:ui-monospace,Menlo,monospace; }
  /* modal */
  .ov { position:fixed; inset:0; background:rgba(0,0,0,.6); display:none; align-items:center; justify-content:center; padding:20px; z-index:20; }
  .ov.on { display:flex; }
  .modal { background:var(--panel); border:1px solid var(--line); border-radius:10px; max-width:460px; width:100%; padding:20px 22px; }
  .modal h3 { margin:0 0 6px; font-size:15px; }
  .modal .m-sub { font-size:12px; color:var(--muted); margin-bottom:14px; }
  .modal .field { font-size:12px; margin:8px 0; }
  .modal .field b { color:var(--muted); font-weight:600; display:block; font-size:10.5px; text-transform:uppercase; letter-spacing:.04em; margin-bottom:2px; }
  .modal .row { display:flex; gap:9px; justify-content:flex-end; margin-top:18px; }
  .modal .row .ghost { background:transparent; }
  .modal .err { color:#ff8585; font-size:12px; margin-top:10px; min-height:1em; }
  .toast { position:fixed; bottom:22px; left:50%; transform:translateX(-50%); background:#1e4620; color:#7ee787; border:1px solid #2c5e30; padding:10px 18px; border-radius:8px; font-size:13px; opacity:0; transition:opacity .2s; pointer-events:none; z-index:30; }
  .toast.on { opacity:1; }
  .toast.bad { background:#5a1e1e; color:#ff8585; border-color:#7a2a2a; }
</style></head>
<body>
<header>
  <h1>Espresso</h1>
  <span class="sub">ticked to-dos log a completed task on the deal in Salesforce</span>
  <div class="controls">
    <label class="lbl">RSD</label>
    <select id="owner"><option value="">All</option></select>
    <button id="refresh">Refresh</button>
    <span id="updated"></span>
  </div>
</header>
<main><div id="list"><div class="empty">Loading…</div></div></main>

<div class="ov" id="ov">
  <div class="modal">
    <h3>Log this to-do in Salesforce?</h3>
    <div class="m-sub">A completed activity will be added to the opportunity today. This can't be undone here.</div>
    <div class="field"><b>Task</b><span id="m-subject"></span></div>
    <div class="field"><b>Deal</b><span id="m-deal"></span></div>
    <div class="err" id="m-err"></div>
    <div class="row">
      <button class="ghost" id="m-cancel">Cancel</button>
      <button class="pushbtn" id="m-confirm">Confirm &amp; log in Salesforce</button>
    </div>
  </div>
</div>
<div class="toast" id="toast"></div>

<script>
const $ = id => document.getElementById(id);
const esc = s => (s==null?'':String(s)).replace(/[&<>"']/g, c => ({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[c]));
const CATS = [
  {key:'critical',             label:'Critical moves',        text:'action',        accent:'#ff8585'},
  {key:'important',            label:'Open commitments',      text:'commitment',    accent:'#f0c674'},
  {key:'explicitRequirements', label:'Explicit requirements', text:'requirement',   accent:'#7cc4ff'},
  {key:'implicit',             label:'Implicit needs',        text:'inferred_need', accent:'#b392f0'},
  {key:'bestPractice',         label:'Best-practice flags',   text:'flag',          accent:'#7ee787'},
];
const U_LABEL = {overdue:'Overdue', next_14_days:'Next 14 days', next_30_days:'Next 30 days', later:'Later', undated:'Undated'};
let DATA = {};        // category -> [items]
let pending = null;   // {item, category, btn}

function toast(msg, bad){ const t=$('toast'); t.textContent=msg; t.className='toast on'+(bad?' bad':''); setTimeout(()=>{t.className='toast'+(bad?' bad':'');}, 2600); }

function chips(it){
  const out=[];
  if(it.urgency && it.urgency!=='undated') out.push(`<span class="chip u-${esc(it.urgency)}">${esc(U_LABEL[it.urgency]||it.urgency)}</span>`);
  for(const [lbl,k] of [['act by','act_by'],['due','due'],['date','date'],['trigger','trigger_date']]){
    if(it[k]) { out.push(`<span class="chip">${lbl} ${esc(it[k])}</span>`); break; }
  }
  if(it.status) out.push(`<span class="chip">${esc(it.status)}</span>`);
  const who = it.intervention_owner || it.who || it.said_by || it.owner_name;
  if(who) out.push(`<span class="chip">${esc(who)}</span>`);
  return out.join('');
}

function itemHTML(it, cat){
  const txt = esc(it[cat.text] || it.action || it.commitment || it.requirement || it.inferred_need || it.flag || '—');
  const deal = [it.account_name, it.opp_name].filter(Boolean).map(esc).join(' — ') || '<span class="chip">no deal label</span>';
  let act;
  if(it.pushed){
    const tid = it.sf_task_id ? `<span class="taskid">${esc(it.sf_task_id)}</span>` : '';
    act = `<span class="pushed">&#10003; Logged in Salesforce</span>${tid}`;
  } else {
    act = `<button class="pushbtn" data-key="${esc(it.todo_key)}">Push to Salesforce</button>`;
  }
  return `<div class="item" style="border-left-color:${cat.accent}" data-key="${esc(it.todo_key)}">
    <div class="body">
      <div class="txt">${txt}</div>
      <div class="deal">${deal}</div>
      <div class="meta">${chips(it)}</div>
    </div>
    <div class="act">${act}</div>
  </div>`;
}

function render(){
  const sel = $('list'); const parts=[]; let total=0;
  for(const cat of CATS){
    const items = DATA[cat.key] || [];
    if(!items.length) continue;
    total += items.length;
    parts.push(`<section class="cat"><h2><span class="dot" style="background:${cat.accent}"></span>${esc(cat.label)}<span class="count">${items.length}</span></h2>`
      + items.map(it=>itemHTML(it, cat)).join('') + `</section>`);
  }
  sel.innerHTML = total ? parts.join('') : '<div class="empty">No open to-dos for this view.</div>';
}

async function load(){
  const owner = $('owner').value;
  $('list').innerHTML='<div class="empty">Loading…</div>';
  try {
    const r = await fetch('/api/deal-engine/todo?owner='+encodeURIComponent(owner), {credentials:'same-origin'});
    if(r.status===401){ $('list').innerHTML='<div class="empty">Unauthorized — open this page with ?key=YOUR_TOKEN once.</div>'; return; }
    DATA = await r.json();
  } catch(e){ $('list').innerHTML='<div class="empty">Connection error.</div>'; return; }
  render();
  $('updated').textContent = 'updated '+new Date().toLocaleTimeString();
}

async function loadOwners(){
  try {
    const r = await fetch('/api/deal-engine/team', {credentials:'same-origin'});
    if(!r.ok) return;
    const t = await r.json();
    const names = [];
    if(Array.isArray(t.rsds)) for(const x of t.rsds) names.push(typeof x==='string'?x:(x.name||x.owner_name));
    else if(Array.isArray(t.members)) for(const x of t.members) names.push(typeof x==='string'?x:(x.name||x.owner_name));
    const seen=new Set();
    for(const n of names.filter(Boolean)){ if(seen.has(n))continue; seen.add(n);
      const o=document.createElement('option'); o.value=n; o.textContent=n; $('owner').appendChild(o); }
  } catch(e){ /* owner filter is optional */ }
}

// find an item across all categories by todo_key
function findItem(key){
  for(const cat of CATS){ for(const it of (DATA[cat.key]||[])){ if(it.todo_key===key) return {it, cat}; } }
  return null;
}

function openModal(key){
  const f = findItem(key); if(!f) return;
  pending = {item:f.it, category:f.cat.key};
  $('m-subject').textContent = f.it[f.cat.text] || f.it.action || f.it.commitment || f.it.requirement || f.it.inferred_need || f.it.flag || '—';
  $('m-deal').textContent = [f.it.account_name, f.it.opp_name].filter(Boolean).join(' — ') || '(no label)';
  $('m-err').textContent='';
  $('m-confirm').disabled=false;
  $('ov').classList.add('on');
}
function closeModal(){ $('ov').classList.remove('on'); pending=null; }

async function confirmPush(){
  if(!pending) return;
  const {item, category} = pending;
  $('m-confirm').disabled=true; $('m-err').textContent='';
  const payload = Object.assign({}, item, {category});
  try {
    const r = await fetch('/api/deal-engine/todo/push', {
      method:'POST', credentials:'same-origin',
      headers:{'Content-Type':'application/json'}, body:JSON.stringify(payload),
    });
    const j = await r.json().catch(()=>({}));
    if(!r.ok || j.ok===false){ $('m-err').textContent = j.error || ('Failed ('+r.status+')'); $('m-confirm').disabled=false; return; }
    // mark pushed locally so the row updates without a full reload
    const f = findItem(item.todo_key);
    if(f){ f.it.pushed=true; f.it.sf_task_id=j.sf_task_id; }
    render();
    closeModal();
    toast(j.already_pushed ? 'Already logged in Salesforce' : 'Logged in Salesforce');
  } catch(e){ $('m-err').textContent='Connection error.'; $('m-confirm').disabled=false; }
}

$('list').addEventListener('click', e=>{ const b=e.target.closest('.pushbtn'); if(b && b.dataset.key) openModal(b.dataset.key); });
$('owner').addEventListener('change', load);
$('refresh').addEventListener('click', load);
$('m-cancel').addEventListener('click', closeModal);
$('m-confirm').addEventListener('click', confirmPush);
$('ov').addEventListener('click', e=>{ if(e.target===$('ov')) closeModal(); });
document.addEventListener('keydown', e=>{ if(e.key==='Escape') closeModal(); });

loadOwners(); load();
</script>
</body></html>""")


@app.post("/api/deal-engine/sweep/{opp_id}")
async def deal_engine_sweep_one(opp_id: str, request: Request):
    """Re-run the sweep for a single opportunity (synchronous). Refreshes that
    one canonical record in place. Returns {opp_id, status, duration_ms, error}."""
    import deal_engine_sweep as sweep
    try:
        d = {}
        try:
            d = await request.json()
        except Exception:  # noqa: BLE001
            pass
        # Enrich from live Salesforce first so this path also re-syncs StageName
        # (and account/owner labels) instead of trusting the request body. Fall
        # back to the body-provided labels for any field the lookup can't resolve.
        enriched = await sweep._enrich_opp_ids(agent_manager, [opp_id])
        opp = enriched[0] if enriched else {"id": opp_id}
        opp.setdefault("name", None)
        opp.setdefault("account", None)
        opp.setdefault("owner_name", None)
        opp.setdefault("owner_id", None)
        opp["name"] = opp.get("name") or d.get("name")
        opp["account"] = opp.get("account") or d.get("account")
        opp["owner_name"] = opp.get("owner_name") or d.get("owner")
        res = await sweep.analyze_one(agent_manager, opp, source="manual")
        code = 200 if res.get("status") == "completed" else 502
        return JSONResponse(res, status_code=code)
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/trigger-logs")
async def deal_engine_trigger_logs(limit: int = 500):
    """Dashboard list: the latest analysis run per opportunity, newest first.
    Each row carries the opp labels plus last_run_at, status, source,
    duration_ms, model, token counts, cost_usd and (if it failed) error."""
    import deal_trigger_log as tlog
    try:
        rows = await _aw(tlog.list_latest, int(limit))
        return {"count": len(rows), "rows": rows}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/api/deal-engine/trigger-logs/{opp_id}")
async def deal_engine_trigger_logs_one(opp_id: str, limit: int = 200):
    """Drill-in: the full run history for one opportunity (every analysis, newest
    first) — time taken, cost, model, token usage, status and failures."""
    import deal_trigger_log as tlog
    try:
        rows = await _aw(tlog.list_runs_for_opp, opp_id, int(limit))
        return {"opp_id": opp_id, "count": len(rows), "rows": rows}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/chat")
async def deal_engine_chat(request: Request):
    """RevOps strategist over the whole book. Body: {messages:[{role,content}],
    owner?, model?}. Returns {answer, usage}."""
    import deal_engine_store as dstore
    try:
        d = await request.json()
        messages = d.get("messages") or []
        if not isinstance(messages, list) or not messages:
            return JSONResponse({"error": "messages required"}, status_code=400)
        owner = (d.get("owner") or "").strip() or None
        opp_ids = [str(i).strip() for i in d.get("opp_ids", [])
                   if str(i).strip()] if isinstance(d.get("opp_ids"), list) else []
        owners = [str(o).strip() for o in d.get("owners", [])
                  if str(o).strip()] if isinstance(d.get("owners"), list) else []
        model_name = (d.get("model") or "").strip() or _deal_engine_model()
        if not os.environ.get("OPENAI_API_KEY"):
            return JSONResponse({"error": "OPENAI_API_KEY is not configured"}, status_code=400)

        book = await _aw(dstore.chat_book_context, owner,
                         owners or None, opp_ids or None)

        # Scope label mirrors the precedence used by the store: opp_ids > owners > owner.
        if opp_ids:
            scope = f" (filtered to {len(book)} selected opps)"
        elif owners:
            scope = f" (filtered to {', '.join(owners)})"
        elif owner:
            scope = f" (filtered to {owner})"
        else:
            scope = ""

        # Base prompt: admin override (Supabase ID_CHAT) wins, else the built-in default.
        # Editable from Admin -> Agent Control -> Chat Agent. The capabilities block + book
        # are appended by code (not editable) so tool-awareness is always present.
        _base = ""
        try:
            import agent_prompt_store as _aps
            _base = (await _aw(_aps.get_prompt)) or ""
        except Exception:  # noqa: BLE001
            _base = ""
        if not _base.strip():
            _base = _DEAL_ENGINE_CHAT_SYSTEM
        sys_text = (
            f"{_base}\n\n"
            f"THE BOOK{scope} — {len(book)} opportunities (compact view; ask for a "
            f"specific opp for full detail):\n{json.dumps(book, default=str)}"
        )
        # FAST one-shot completion (kept synchronous so it returns well within the
        # Vercel proxy timeout). The tool-using version (search_knowledge over the
        # shared MASE KB + run_todo delegation to the Todo Runner) is being rebuilt on
        # the streaming/realtime path — a synchronous tool loop + nested sub-agent runs
        # for tens of seconds to minutes and times out at the proxy, so it can't live
        # behind this blocking endpoint. See deal_engine_chat_agent.py (kept for that
        # build) and CHANGELOG 2026-06-19.
        from langchain_openai import ChatOpenAI
        llm = ChatOpenAI(model=model_name, temperature=0.2)
        lc_messages = [{"role": "system", "content": sys_text}]
        for m in messages:
            role = m.get("role")
            if role in ("user", "assistant") and m.get("content"):
                lc_messages.append({"role": role, "content": m["content"]})
        resp = await llm.ainvoke(lc_messages)
        usage = getattr(resp, "response_metadata", {}).get("token_usage", {}) or \
            getattr(resp, "usage_metadata", {}) or {}
        return {"answer": resp.content, "usage": usage}
    except Exception as e:  # noqa: BLE001
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/api/deal-engine/chat/async")
async def deal_engine_chat_async(request: Request):
    """Streaming/realtime RevOps chat. Builds the SAME book + system prompt as the
    sync /api/deal-engine/chat, but instead of a blocking one-shot completion it
    spawns the tool-using chat agent (search_knowledge over the MASE KB + run_todo
    delegation) as a BACKGROUND task that streams its thinking / tool calls / final
    answer into the shared `chat_messages` table via save_to_supabase. The browser
    subscribes to that table over Supabase realtime.

    Body: {chat_id?, messages:[{role,content}], owner?, owners?, opp_ids?, model?}.
    Returns FAST JSON {chat_id} immediately (NOT a long-lived stream) so the
    buffering proxy returns at once; everything else flows over realtime.
    """
    import deal_engine_store as dstore
    d = await request.json()
    chat_id = (d.get("chat_id") or "").strip() or str(uuid.uuid4())

    # One run per chat_id (backend backstop against double-submit). Reserved
    # before the try so the 409 propagates without the finally releasing a slot
    # owned by another request.
    _reserve_run_slot(chat_id)
    try:
        if not supabase:
            raise HTTPException(
                status_code=500,
                detail="Supabase not configured - realtime chat requires Supabase")

        messages = d.get("messages") or []
        if not isinstance(messages, list) or not messages:
            return JSONResponse({"error": "messages required"}, status_code=400)
        owner = (d.get("owner") or "").strip() or None
        opp_ids = [str(i).strip() for i in d.get("opp_ids", [])
                   if str(i).strip()] if isinstance(d.get("opp_ids"), list) else []
        owners = [str(o).strip() for o in d.get("owners", [])
                  if str(o).strip()] if isinstance(d.get("owners"), list) else []
        model_name = (d.get("model") or "").strip() or _deal_engine_model()

        active_count = sum(1 for t in _running_tasks.values() if t and not t.done())
        if active_count >= config.MAX_CONCURRENT_SESSIONS:
            global _sessions_rejected
            _sessions_rejected += 1
            raise HTTPException(
                status_code=503,
                detail=f"Server at capacity ({config.MAX_CONCURRENT_SESSIONS} concurrent sessions). Please try again shortly.")

        # Build the book exactly like the sync endpoint (opp_ids > owners > owner).
        book = await _aw(dstore.chat_book_context, owner,
                         owners or None, opp_ids or None)
        if opp_ids:
            scope = f" (filtered to {len(book)} selected opps)"
        elif owners:
            scope = f" (filtered to {', '.join(owners)})"
        elif owner:
            scope = f" (filtered to {owner})"
        else:
            scope = ""

        # Editable base prompt (Supabase ID_CHAT) + code-appended capabilities +
        # the book JSON. The capabilities block tells the agent about its tools so
        # it knows when to call search_knowledge / delegate run_todo.
        _base = ""
        try:
            import agent_prompt_store as _aps
            _base = (await _aw(_aps.get_prompt)) or ""
        except Exception:  # noqa: BLE001
            _base = ""
        if not _base.strip():
            _base = _DEAL_ENGINE_CHAT_SYSTEM
        # Persona overlay (from the UI persona tabs) — switches the agent's lens by
        # appending its directive to the system prompt for this run.
        _persona = (d.get("persona") or "").strip().lower()
        _pdir = _CHAT_PERSONAS.get(_persona, "")
        sys_text = (
            f"{_base}\n{_CHAT_CAPABILITIES}\n"
            + (f"\n{_pdir}\n" if _pdir else "")
            + f"\nTHE BOOK{scope} — {len(book)} opportunities (compact view; ask for a "
            f"specific opp for full detail):\n{json.dumps(book, default=str)}"
        )

        # Build the tool-using chat agent. If it can't be built (tools not loaded),
        # write a terminal row so the UI stops spinning, and still return {chat_id}.
        try:
            import deal_engine_chat_agent

            # Live nested trace for the Todo Runner: this callback writes the
            # Todo Runner's own steps as sub-rows to the SAME chat_id, tagged
            # group="todo". They share _supabase_seq_counters[chat_id] with the
            # parent run, so they're sequenced BETWEEN the parent's run_todo
            # tool_call and tool_result — ordering stays correct. The frontend
            # renders group="todo" rows as a nested "Todo Runner working…"
            # sub-accordion. (deal_engine_chat_agent must not import server.py.)
            async def _emit(t, c, meta):
                await save_to_supabase(
                    chat_id, t, c, {**(meta or {}), "group": "todo"})

            agent = deal_engine_chat_agent.build_chat_agent(
                agent_manager, sys_text, emit=_emit)
        except Exception as e:  # noqa: BLE001
            import traceback
            print(f"[DEAL CHAT ASYNC] build_chat_agent failed: {traceback.format_exc()}")
            await save_to_supabase(
                chat_id, "error",
                f"Could not start the strategist agent: {e}",
                {"status": "failed", "kind": "agent_build_failed"})
            return {"chat_id": chat_id}

        # Only the real conversation turns go to the agent.
        conv = [{"role": m.get("role"), "content": m.get("content")}
                for m in messages
                if m.get("role") in ("user", "assistant") and m.get("content")]

        # Spawn the run as a tracked background task — mirror /api/chat/async.
        # MASE_KNOWLEDGE_PROJECT_ID routes the agent's search_knowledge to the
        # isolated MASE namespace (it's also the realtime routing marker).
        consumer_task = asyncio.create_task(
            run_agent_and_save(chat_id, conv, agent, model_name,
                               deal_engine_chat_agent.MASE_KNOWLEDGE_PROJECT_ID))
        _running_tasks[chat_id] = consumer_task
        _session_start_times[chat_id] = asyncio.get_event_loop().time()

        def _cleanup_session(t, cid=chat_id):
            _running_tasks.pop(cid, None)
            _session_start_times.pop(cid, None)
            _supabase_seq_counters.pop(cid, None)
            _dedupe_completed.pop(cid, None)
            _dedupe_inflight.pop(cid, None)
            _approved_campaigns_cache.pop(cid, None)
            _starting_chats.discard(cid)

        consumer_task.add_done_callback(_cleanup_session)

        # Fast JSON — the proxy buffers, so do NOT return a long-lived stream.
        return {"chat_id": chat_id}

    except HTTPException:
        raise
    except Exception as e:  # noqa: BLE001
        import traceback
        print(f"[DEAL CHAT ASYNC ERROR] {traceback.format_exc()}")
        return JSONResponse({"error": str(e)}, status_code=500)
    finally:
        # Idempotent release on every exit path. Once the task is registered it
        # owns the slot via _running_tasks, so this discard is harmless; if setup
        # failed before registration it frees the chat so it isn't 409-blocked.
        _release_run_slot(chat_id)


@app.post("/api/deal-engine/chat/stop")
async def deal_engine_chat_stop(chat_id: str = None):
    """Stop a running MASE strategist chat by chat_id. Mirrors /api/chat/stop but
    lives under the deal-engine prefix so the frontend proxy can reach it. Marks
    the chat cancelled and cancels its asyncio task; writes a `cancelled` status
    row so the UI (realtime + polling) stops the spinner immediately."""
    if not chat_id:
        raise HTTPException(status_code=400, detail="chat_id is required")
    print(f"[DEAL CHAT STOP] stop request chat_id={chat_id}")
    _cancelled_chats.add(chat_id)
    task = _running_tasks.get(chat_id)
    if task and not task.done():
        try:
            task.cancel()
        except Exception as e:  # noqa: BLE001
            print(f"[DEAL CHAT STOP] cancel failed (non-fatal): {e}")
        try:
            await save_to_supabase(
                chat_id, "status", "Agent stopped by user.",
                {"status": "cancelled", "source": "deal_stop_endpoint"})
        except Exception as e:  # noqa: BLE001
            print(f"[DEAL CHAT STOP] supabase save failed (non-fatal): {e}")
        return {"chat_id": chat_id, "status": "stopped", "message": "Agent stopped."}
    _running_tasks.pop(chat_id, None)
    return {"chat_id": chat_id, "status": "not_found",
            "message": "No running agent task for this chat_id."}


@app.get("/avoma/reports", response_class=HTMLResponse)
async def avoma_reports_ui():
    """Lightweight UI for browsing Avoma-triggered SF enrichment reports."""
    return HTMLResponse(content=r"""<!doctype html>
<html><head><meta charset="utf-8"><title>Avoma → SF Enrichment Reports</title>
<style>
  body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", system-ui, sans-serif; margin: 0; background: #0f1419; color: #d4d4d4; }
  header { background: #1a1f2e; padding: 14px 24px; border-bottom: 1px solid #2a3142; display: flex; gap: 16px; align-items: center; }
  header h1 { margin: 0; font-size: 18px; font-weight: 600; }
  header .refresh, header select { background: #2a3142; color: #d4d4d4; border: 1px solid #3a4258; padding: 6px 12px; border-radius: 4px; cursor: pointer; font-size: 13px; }
  header .refresh:hover { background: #3a4258; }
  main { display: flex; height: calc(100vh - 53px); }
  #list { width: 42%; overflow-y: auto; border-right: 1px solid #2a3142; }
  #detail { flex: 1; overflow-y: auto; padding: 20px 28px; }
  .row { padding: 12px 20px; border-bottom: 1px solid #1f2533; cursor: pointer; }
  .row:hover { background: #1a1f2e; }
  .row.active { background: #1f2533; border-left: 3px solid #4a9eff; padding-left: 17px; }
  .row .subj { font-weight: 500; color: #f0f0f0; margin-bottom: 4px; }
  .row .meta { font-size: 11px; color: #888; display: flex; gap: 10px; flex-wrap: wrap; }
  .badge { display: inline-block; padding: 2px 8px; border-radius: 10px; font-size: 11px; font-weight: 500; }
  .badge.completed { background: #1e3a2a; color: #6fcf97; }
  .badge.completed_with_errors { background: #3a2f1e; color: #f2c94c; }
  .badge.completed_no_opportunity { background: #2a2a3a; color: #888; }
  .badge.no_sf_links { background: #2a2a3a; color: #888; }
  .badge.pending { background: #1e2e3a; color: #4a9eff; }
  .badge.pending_analysis { background: #1e2e3a; color: #4a9eff; }
  .badge.running { background: #1e2e3a; color: #4a9eff; }
  .badge.parse_error { background: #3a2f1e; color: #f2c94c; }
  .badge.failed { background: #3a1e1e; color: #eb5757; }
  .rerun { background: #2a3142; color: #d4d4d4; border: 1px solid #3a4258; padding: 5px 12px; border-radius: 4px; cursor: pointer; font-size: 12px; margin-left: 8px; }
  .rerun:hover { background: #3a4258; }
  .rerun:disabled { opacity: 0.5; cursor: not-allowed; }
  .kv { display: grid; grid-template-columns: 180px 1fr; gap: 6px 14px; font-size: 13px; margin: 8px 0 14px; }
  .kv .k { color: #888; }
  .kv .v { color: #f0f0f0; }
  .pill { display: inline-block; padding: 2px 8px; border-radius: 10px; font-size: 11px; background: #2a3142; color: #d4d4d4; margin-right: 6px; }
  .pill.likely { background: #1e3a2a; color: #6fcf97; }
  .pill.possible { background: #1e2e3a; color: #4a9eff; }
  .pill.cautiously_possible { background: #3a2f1e; color: #f2c94c; }
  .pill.unlikely { background: #3a1e1e; color: #eb5757; }
  .conflict { background: #2a1f1f; border-left: 3px solid #eb5757; padding: 8px 12px; margin: 6px 0; font-size: 12px; }
  .gap { background: #2a2519; border-left: 3px solid #f2c94c; padding: 8px 12px; margin: 6px 0; font-size: 12px; }
  #detail h2 { margin: 0 0 4px; font-size: 20px; color: #f0f0f0; }
  #detail .ids { font-size: 12px; color: #888; margin-bottom: 20px; font-family: monospace; }
  .tier { background: #1a1f2e; border-radius: 6px; margin-bottom: 16px; }
  .tier h3 { margin: 0; padding: 12px 16px; font-size: 14px; border-bottom: 1px solid #2a3142; cursor: pointer; display: flex; justify-content: space-between; align-items: center; }
  .tier .body { padding: 16px; }
  .tier .summary { white-space: pre-wrap; line-height: 1.6; color: #d4d4d4; font-size: 13px; }
  .tier details { margin-top: 12px; }
  .tier details summary { cursor: pointer; color: #888; font-size: 12px; }
  .tier pre { background: #0f1419; padding: 12px; border-radius: 4px; overflow-x: auto; font-size: 11px; color: #aaa; max-height: 400px; }
  .empty { color: #666; padding: 40px; text-align: center; }
  a { color: #4a9eff; text-decoration: none; }
</style></head>
<body>
<header>
  <h1>Avoma → SF Enrichment</h1>
  <select id="filter">
    <option value="">all</option>
    <option value="completed">completed</option>
    <option value="completed_no_opportunity">completed_no_opportunity</option>
    <option value="no_sf_links">no_sf_links</option>
    <option value="pending_analysis">pending_analysis</option>
    <option value="failed">failed</option>
  </select>
  <button class="refresh" onclick="loadList()">↻ refresh</button>
  <span id="count" style="color:#888;font-size:12px;"></span>
</header>
<main>
  <div id="list"></div>
  <div id="detail"><div class="empty">← pick a report</div></div>
</main>
<script>
let _active = null;
async function loadList() {
  const f = document.getElementById('filter').value;
  const r = await fetch('/api/avoma/reports?limit=100' + (f ? '&status=' + f : ''));
  const d = await r.json();
  document.getElementById('count').textContent = d.count + ' reports';
  document.getElementById('list').innerHTML = (d.reports || []).map(row => {
    const created = row.created_at ? new Date(row.created_at).toLocaleString() : '';
    return `<div class="row${row.message_id===_active?' active':''}" onclick="loadDetail('${row.message_id}')">
      <div class="subj">${esc(row.meeting_subject || '(no subject)')}</div>
      <div class="meta">
        <span class="badge ${row.status}">${row.status}</span>
        <span>${created}</span>
        ${row.sf_opportunity_id ? `<span>opp ${row.sf_opportunity_id.slice(0,12)}</span>` : ''}
        ${row.sf_account_id ? `<span>acct ${row.sf_account_id.slice(0,12)}</span>` : ''}
        ${row.pull_duration_ms ? `<span>${row.pull_duration_ms}ms</span>` : ''}
      </div>
    </div>`;
  }).join('') || '<div class="empty">no reports yet</div>';
}
async function loadDetail(mid) {
  _active = mid;
  document.querySelectorAll('.row').forEach(r => r.classList.remove('active'));
  event && event.currentTarget && event.currentTarget.classList.add('active');
  const r = await fetch('/api/avoma/reports/' + mid);
  const d = await r.json();
  if (d.error) { document.getElementById('detail').innerHTML = `<div class="empty">${esc(d.error)}</div>`; return; }
  renderDetail(d);
}
function renderDetail(d) {
  const an = d.opportunity_analysis_data || null;
  const ast = d.opportunity_analysis_status || '—';
  const hasOpp = !!d.sf_opportunity_id;
  const isRunning = ast === 'running';
  const rerunBtn = hasOpp
    ? `<button class="rerun" id="rerunBtn" onclick="rerun('${d.message_id}')" ${isRunning ? 'disabled' : ''}>${isRunning ? '… running' : '↻ rerun analysis'}</button>`
    : '';
  document.getElementById('detail').innerHTML = `
    <h2>${esc(d.meeting_subject || '(no subject)')}${rerunBtn}</h2>
    <div class="ids">
      meeting ${d.meeting_uuid || '—'}<br>
      ${d.sf_opportunity_id ? 'opp ' + d.sf_opportunity_id + ' • ' : ''}
      ${d.sf_account_id ? 'acct ' + d.sf_account_id + ' • ' : ''}
      ${d.sf_contact_ids?.length ? d.sf_contact_ids.length + ' contacts • ' : ''}
      <span class="badge ${d.status}">${d.status}</span>
      ${d.pull_duration_ms ? d.pull_duration_ms + 'ms' : ''}
      ${d.error ? '<br><span style="color:#eb5757">' + esc(d.error) + '</span>' : ''}
    </div>
    <div class="tier">
      <h3>opportunity analysis
        <span style="color:#888;font-size:12px;font-weight:normal">
          <span class="badge ${ast.split(':')[0]}">${esc(ast)}</span>
          ${an ? JSON.stringify(an).length.toLocaleString() + ' chars' : ''}
        </span>
      </h3>
      <div class="body">${renderAnalysis(an, ast)}</div>
    </div>`;
}
function renderAnalysis(an, ast) {
  if (!an) {
    if (ast === 'running') return '<div class="summary" style="color:#4a9eff">⏳ analysis in progress — refresh in a few minutes</div>';
    if (ast && ast.startsWith('failed')) return `<div class="summary" style="color:#eb5757">${esc(ast)}</div>`;
    return '<div class="summary" style="color:#666">(no analysis — opportunity not linked, or analyzer not run)</div>';
  }
  if (an._error) {
    return `<div class="summary" style="color:#eb5757">parse error: ${esc(an._error)}</div>
      <details style="margin-top:12px"><summary>raw model output (${(an._raw_len||0).toLocaleString()} chars)</summary>
      <pre>${esc(an._raw || '')}</pre></details>`;
  }
  const meta = an.meta || {};
  const ident = an.identity || {};
  const health = an.health || {};
  const conflicts = (meta.conflicts || []);
  const gaps = (meta.unresolved_gaps || []);
  const wl = health.win_likelihood || '';
  const wlPill = wl ? `<span class="pill ${wl}">${esc(wl)}</span>` : '';
  return `
    <div class="kv">
      ${ident.account_name ? `<div class="k">account</div><div class="v">${esc(ident.account_name)}</div>` : ''}
      ${ident.opportunity_name ? `<div class="k">opportunity</div><div class="v">${esc(ident.opportunity_name)}</div>` : ''}
      ${ident.stage ? `<div class="k">stage</div><div class="v">${esc(ident.stage)} ${ident.probability != null ? '(' + ident.probability + '%)' : ''}</div>` : ''}
      ${ident.amount != null ? `<div class="k">amount</div><div class="v">${esc(String(ident.amount))} ${ident.currency || ''}</div>` : ''}
      ${ident.close_date ? `<div class="k">close date</div><div class="v">${esc(ident.close_date)}</div>` : ''}
      ${ident.owner ? `<div class="k">owner</div><div class="v">${esc(ident.owner)}</div>` : ''}
      ${health.overall ? `<div class="k">overall health</div><div class="v">${esc(health.overall)}</div>` : ''}
      ${wl ? `<div class="k">win likelihood</div><div class="v">${wlPill}</div>` : ''}
      ${health.forecast_confidence ? `<div class="k">forecast confidence</div><div class="v">${esc(health.forecast_confidence)}</div>` : ''}
      ${meta.run_status ? `<div class="k">run_status</div><div class="v">${esc(meta.run_status)}</div>` : ''}
      ${meta.calls_analyzed ? `<div class="k">calls analyzed</div><div class="v">${(meta.calls_analyzed.analyzed||[]).length} of ${(meta.calls_analyzed.all||[]).length}</div>` : ''}
    </div>
    ${conflicts.length ? `<div style="margin-top:6px"><div style="color:#888;font-size:12px;margin-bottom:4px">conflicts (${conflicts.length})</div>
      ${conflicts.map(c => `<div class="conflict">${esc(JSON.stringify(c))}</div>`).join('')}</div>` : ''}
    ${gaps.length ? `<div style="margin-top:6px"><div style="color:#888;font-size:12px;margin-bottom:4px">unresolved gaps (${gaps.length})</div>
      ${gaps.map(g => `<div class="gap">${esc(typeof g === 'string' ? g : JSON.stringify(g))}</div>`).join('')}</div>` : ''}
    <details style="margin-top:14px"><summary>full record (jsonb)</summary><pre>${esc(JSON.stringify(an, null, 2))}</pre></details>
  `;
}
async function rerun(mid) {
  const btn = document.getElementById('rerunBtn');
  if (btn) { btn.disabled = true; btn.textContent = '… kicking off'; }
  const r = await fetch('/api/avoma/reports/' + mid + '/reanalyze', { method: 'POST' });
  const d = await r.json();
  if (d.error) { alert('rerun failed: ' + d.error); if (btn) { btn.disabled = false; btn.textContent = '↻ rerun analysis'; } return; }
  const poll = async () => {
    const rr = await fetch('/api/avoma/reports/' + mid);
    const dd = await rr.json();
    renderDetail(dd);
    if ((dd.opportunity_analysis_status || '') === 'running') setTimeout(poll, 5000);
    else loadList();
  };
  setTimeout(poll, 3000);
}
function esc(s) { return String(s).replace(/[&<>"]/g, c => ({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;'}[c])); }
document.getElementById('filter').onchange = loadList;
loadList();
setInterval(() => { if (!_active) loadList(); }, 10000);
</script>
</body></html>""")


@app.get("/api/mcp/status")
async def mcp_status():
    per_server_tools = {k: len(v) for k, v in agent_manager._cached_mcp_tools_by_server.items()}
    enabled_servers = list(agent_manager.mcp_config_manager.get_enabled_servers().keys())
    missing_servers = [s for s in enabled_servers if s not in agent_manager._cached_mcp_tools_by_server or not agent_manager._cached_mcp_tools_by_server[s]]
    return {
        "mcp_tools_loaded": agent_manager.mcp_tools_loaded,
        "mcp_tools_count": len(agent_manager._cached_mcp_tools),
        "per_server_tool_counts": per_server_tools,
        "expected_tool_counts": agent_manager._expected_tool_counts,
        "missing_servers": missing_servers,
        "servers": agent_manager.mcp_loading_status,
        "background_task_running":
        agent_manager._mcp_load_task is not None
        and not agent_manager._mcp_load_task.done()
        if agent_manager._mcp_load_task else False,
        "health_check_running":
        agent_manager._mcp_health_task is not None
        and not agent_manager._mcp_health_task.done()
        if agent_manager._mcp_health_task else False,
    }


if GOOGLE_SHEETS_ENABLED and sheets_auth:

    @app.get("/oauth2callback")
    async def oauth2callback(code: str):
        try:
            sheets_auth.handle_callback(code)
            return RedirectResponse(url="/")
        except Exception as e:
            raise HTTPException(status_code=400, detail=str(e))

    @app.get("/api/sheets/status")
    async def sheets_auth_status():
        return {
            "enabled":
            GOOGLE_SHEETS_ENABLED,
            "authenticated":
            sheets_auth.is_authenticated() if sheets_auth else False
        }

    @app.get("/api/sheets/auth-url")
    async def get_sheets_auth_url():
        if not sheets_auth:
            raise HTTPException(status_code=503,
                                detail="Google Sheets not configured")
        auth_url = sheets_auth.get_auth_url()
        if not auth_url:
            raise HTTPException(status_code=503,
                                detail="Cannot generate auth URL")
        return {"auth_url": auth_url}


# ============================================================================
# STARTUP
# ============================================================================


@app.on_event("startup")
async def startup_event():
    print("Initializing DeepAgent server...")
    print(f"Context Window Management:")
    print(f"  - Summarizer model: {config.SUMMARIZER_MODEL}")
    print(
        f"  - Tool response summarize threshold: {config.TOOL_RESPONSE_SUMMARIZE_THRESHOLD:,} chars"
    )
    print(
        f"  - Conversation summarize threshold: {config.CONVERSATION_SUMMARIZE_TOKEN_THRESHOLD:,} tokens"
    )
    print(
        f"  - Keep recent messages: {config.CONVERSATION_KEEP_RECENT_MESSAGES}"
    )
    print(f"MCP Truncation Limits:")
    print(f"  - MAX_RESPONSE_SIZE: {config.MCP_MAX_RESPONSE_SIZE:,} chars")
    print(f"  - MAX_STRING_LENGTH: {config.MCP_MAX_STRING_LENGTH:,} chars")
    print(f"  - MAX_LIST_ITEMS: {config.MCP_MAX_LIST_ITEMS} items")
    await agent_manager.initialize_agent(skip_mcp=True)
    print("Server ready! (MCP servers loading in background...)")
    agent_manager.start_mcp_background_loading()
    # Capture the server event loop so cancel_running_chat() can route
    # cross-thread cancellations safely (e.g. from sync LangChain tools
    # running in run_in_executor). See 2026-05-22 fixes for chat 8359d7a6.
    _register_server_event_loop()
    global _server_start_time
    _server_start_time = asyncio.get_event_loop().time()
    asyncio.create_task(_cleanup_expired_sessions())
    print(f"Concurrency: max {config.MAX_CONCURRENT_SESSIONS} sessions, timeout {config.SESSION_TIMEOUT_MINUTES} min")
    _qa_timeout = int(os.getenv("QUERY_AGENT_TIMEOUT", "300"))
    print(f"query_agent timeout: {_qa_timeout}s ({_qa_timeout // 60} min) — override with QUERY_AGENT_TIMEOUT env var")
    if _api_auth_token():
        _src = "API_AUTH_TOKEN" if os.environ.get("API_AUTH_TOKEN") else "DISPATCH_SECRET"
        print(f"[API-AUTH] ✅ ENABLED for all non-/mcp routes (token via {_src}). "
              f"Public allowlist: {sorted(_API_AUTH_PUBLIC_EXACT)} + /.well-known/*")
    else:
        print("[API-AUTH] ⚠️  DISABLED — no API_AUTH_TOKEN or DISPATCH_SECRET set; "
              "all data/admin endpoints are PUBLIC. Set API_AUTH_TOKEN to enable.")

    async def _run_mcp_lifespan():
        async with _mcp.session_manager.run():
            await _mcp_lifespan_exit.wait()

    global _mcp_lifespan_task
    _mcp_lifespan_task = asyncio.create_task(_run_mcp_lifespan())
    print("[MCP-HTTP] Streamable HTTP endpoint active at /mcp")

    # Nightly deterministic SF pull + cache refresh (default midnight UTC).
    asyncio.create_task(_nightly_sf_pull_scheduler())

    # Nightly token-free hard-fact reconciliation (default 01:00 UTC) — re-reads
    # live Salesforce hard facts onto every deal record with no AI cost, keeping
    # the book accurate between paid AI sweeps.
    asyncio.create_task(_nightly_hard_refresh_scheduler())


@app.on_event("shutdown")
async def shutdown_event():
    # Graceful drain on ECS stop (deploy / scale-in): uvicorn runs this on SIGTERM,
    # before the SIGKILL that follows the task's stopTimeout (120s). Give in-flight agent
    # runs a short grace to finish naturally; then CANCEL any stragglers and await them so
    # each run's OWN finally / cancel handler writes its single terminal row. We do NOT
    # inject a terminal row from here: the run paths already write exactly one terminal on
    # cancel (chat/agent finally safety-net, structured-async CancelledError handler), so
    # injecting one would double-write and violate the one-terminal-row contract, and would
    # also fabricate a row for pure HTTP-streaming runs that have no chat_messages terminal.
    # The whole drain is bounded (grace + cancel-wait) to stay well under stopTimeout.
    try:
        inflight = [t for _, t in list(_running_tasks.items()) if t and not t.done()]
        if inflight:
            grace = float(os.getenv("SHUTDOWN_DRAIN_GRACE_S", "15"))
            print(f"[SHUTDOWN] {len(inflight)} agent run(s) in flight; {grace}s grace to finish")
            await asyncio.wait(inflight, timeout=grace)
            still = [t for t in inflight if not t.done()]
            if still:
                print(f"[SHUTDOWN] cancelling {len(still)} unfinished run(s) so each writes its own terminal row")
                for t in still:
                    t.cancel()
                # Let each cancelled run's finally / CancelledError path write its single
                # terminal row (this is what unblocks the chat UI on deploy).
                await asyncio.wait(still, timeout=float(os.getenv("SHUTDOWN_CANCEL_WAIT_S", "10")))
    except Exception as _e:  # noqa: BLE001
        print(f"[SHUTDOWN] drain error: {_e}")

    _mcp_lifespan_exit.set()
    if _mcp_lifespan_task and not _mcp_lifespan_task.done():
        try:
            await asyncio.wait_for(_mcp_lifespan_task, timeout=5.0)
        except (asyncio.TimeoutError, asyncio.CancelledError):
            pass


@app.post("/api/admin/backfill-lake")
async def backfill_lake(request: Request):
    """Trigger the historical lake backfill for OD projects.

    Auth: Authorization: Bearer <LAKE_BACKFILL_TOKEN>
    Streams newline-delimited progress logs; final line is a summary.
    """
    # --- Auth check ---
    auth_header = request.headers.get("Authorization", "")
    provided = auth_header.removeprefix("Bearer ").strip()
    expected = os.environ.get("LAKE_BACKFILL_TOKEN", "")
    if not expected:
        raise HTTPException(status_code=503, detail="LAKE_BACKFILL_TOKEN not configured")
    if not provided or provided != expected:
        raise HTTPException(status_code=401, detail="Unauthorized")

    if not supabase:
        raise HTTPException(status_code=503, detail="Supabase not configured")

    openai_key = config.OPENAI_API_KEY or os.getenv("OPENAI_API_KEY", "")

    progress_queue: asyncio.Queue = asyncio.Queue()
    done_event = asyncio.Event()

    async def _run_backfill():
        """Background task: processes all OD chats and pushes progress lines to queue."""
        counters = {"written": 0, "exists": 0, "skip": 0, "fail": 0}
        failures = []
        loop = asyncio.get_event_loop()

        try:
            await progress_queue.put("[BACKFILL] Starting lake backfill for OD projects\n")

            # Step 1: discover chats by project_id from the chats table
            try:
                resp = await loop.run_in_executor(
                    None,
                    lambda: (
                        supabase.table("chats")
                        .select("id,project_id")
                        .in_("project_id", list(_lake.OD_PROJECT_IDS))
                        .execute()
                    )
                )
                project_chats: dict = {}
                for row in resp.data or []:
                    pid = row.get("project_id")
                    cid = row.get("id")
                    if pid and cid:
                        project_chats.setdefault(pid, []).append(cid)
            except Exception as e:
                await progress_queue.put(f"[BACKFILL] ERROR discovering chats: {e}\n")
                return

            total = sum(len(v) for v in project_chats.values())
            await progress_queue.put(
                f"[BACKFILL] Found {total} OD chats across {len(project_chats)} project(s)\n"
            )
            for pid, cids in project_chats.items():
                await progress_queue.put(f"  project={pid}  chats={len(cids)}\n")

            if total == 0:
                await progress_queue.put(
                    "[BACKFILL] Nothing to process — check that chats.project_id is set\n"
                )
                return

            # Step 2: process each chat
            for project_id, chat_ids in project_chats.items():
                await progress_queue.put(
                    f"[BACKFILL] Processing project {project_id}\n"
                )
                for chat_id in chat_ids:
                    try:
                        # Fetch last type='final' message (run_at = raw created_at)
                        fm_resp = await loop.run_in_executor(
                            None,
                            lambda cid=chat_id: (
                                supabase.table("chat_messages")
                                .select("content,created_at")
                                .eq("chat_id", cid)
                                .eq("type", "final")
                                .order("created_at", desc=True)
                                .limit(1)
                                .execute()
                            )
                        )
                        rows = fm_resp.data or []
                        if not rows:
                            counters["skip"] += 1
                            await progress_queue.put(
                                f"  SKIP    {chat_id}  (no type=final message)\n"
                            )
                            continue

                        run_at: str = rows[0]["created_at"]
                        final_response: str = rows[0].get("content") or ""

                        # Use backfill writer — raises on failure, signals existing rows
                        status, detail = await _lake.write_lake_diagnosis_backfill(
                            chat_id=chat_id,
                            project_id=project_id,
                            final_response=final_response,
                            supabase_client=supabase,
                            openai_api_key=openai_key,
                            run_at=run_at,
                        )
                        if status == "exists":
                            counters["exists"] += 1
                            await progress_queue.put(
                                f"  EXISTS  {chat_id}  (already in lake)\n"
                            )
                        else:
                            counters["written"] += 1
                            await progress_queue.put(
                                f"  WRITTEN {chat_id}  run_at={run_at}"
                                f"  account={detail}\n"
                            )

                    except Exception as e:
                        first_line = str(e).split("\n")[0][:200]
                        counters["fail"] += 1
                        failures.append(f"{chat_id}: {first_line}")
                        await progress_queue.put(
                            f"  FAIL    {chat_id}  {first_line}\n"
                        )

            # Step 3: summary
            await progress_queue.put(
                f"[BACKFILL] DONE — "
                f"written={counters['written']}  "
                f"already_existed={counters['exists']}  "
                f"skipped_no_final={counters['skip']}  "
                f"failed={counters['fail']}\n"
            )
            if failures:
                await progress_queue.put(f"[BACKFILL] Failures ({len(failures)}):\n")
                for f in failures:
                    await progress_queue.put(f"  {f}\n")

        except Exception as e:
            await progress_queue.put(f"[BACKFILL] FATAL ERROR: {e}\n")
        finally:
            done_event.set()

    asyncio.create_task(_run_backfill())

    async def _stream_progress():
        KEEPALIVE_INTERVAL = 10
        try:
            while not done_event.is_set() or not progress_queue.empty():
                try:
                    line = await asyncio.wait_for(
                        progress_queue.get(), timeout=KEEPALIVE_INTERVAL
                    )
                    yield line
                except asyncio.TimeoutError:
                    yield "[BACKFILL] ping\n"
        except Exception:
            pass

    return StreamingResponse(
        _stream_progress(),
        media_type="text/plain",
        headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"},
    )


@app.get("/api/lake/diagnoses/{account_id}")
async def get_account_diagnoses(account_id: str, request: Request, limit: int = 20):
    """Return the most recent Opportunity Diagnosis runs for an account.

    Auth: Authorization: Bearer <LAKE_BACKFILL_TOKEN> (same token as
    /api/admin/backfill-lake). Returns a JSON array (possibly empty) of
    diagnosis rows ordered most-recent-first.
    """
    auth_header = request.headers.get("Authorization", "")
    provided = auth_header.removeprefix("Bearer ").strip()
    expected = os.environ.get("LAKE_BACKFILL_TOKEN", "")
    if not expected:
        raise HTTPException(status_code=503, detail="LAKE_BACKFILL_TOKEN not configured")
    if not provided or provided != expected:
        raise HTTPException(status_code=401, detail="Unauthorized")

    if not supabase:
        raise HTTPException(status_code=503, detail="Supabase not configured")

    # Clamp limit to a sane range
    try:
        limit = max(1, min(int(limit), 100))
    except (TypeError, ValueError):
        limit = 20

    def _coerce_json_list(value):
        """top_risks / recommendations may be stored as JSON string or list."""
        if value is None:
            return []
        if isinstance(value, list):
            return value
        if isinstance(value, str):
            try:
                parsed = json.loads(value)
                return parsed if isinstance(parsed, list) else []
            except Exception:
                return []
        return []

    try:
        loop = asyncio.get_event_loop()
        resp = await loop.run_in_executor(
            None,
            lambda: (
                supabase
                .schema("lake")
                .table("opportunity_diagnoses")
                .select(
                    "run_at,opportunity_name,stage,amount,"
                    "momentum_verdict,health_rating,top_risks,recommendations"
                )
                .eq("account_id", account_id)
                .order("run_at", desc=True)
                .limit(limit)
                .execute()
            ),
        )
        rows = resp.data or []
    except Exception as e:
        # Lake table may not exist yet, or transient DB error — log and return empty
        print(f"[LAKE] /api/lake/diagnoses query failed: {e}", flush=True)
        return []

    result = []
    for row in rows:
        result.append({
            "run_at": row.get("run_at"),
            "opportunity_name": row.get("opportunity_name"),
            "stage": row.get("stage"),
            "amount": row.get("amount"),
            "momentum_verdict": row.get("momentum_verdict"),
            "health_rating": row.get("health_rating"),
            "top_risks": _coerce_json_list(row.get("top_risks")),
            "recommendations": _coerce_json_list(row.get("recommendations")),
        })
    return result


@app.get("/api/usage/{chat_id}")
async def get_chat_usage(chat_id: str):
    """Return LLM token usage and estimated cost for a specific chat."""
    if not supabase:
        raise HTTPException(status_code=503, detail="Supabase not configured")
    loop = asyncio.get_event_loop()
    try:
        result = await loop.run_in_executor(
            None,
            lambda: supabase.table("chat_usage").select("*").eq("chat_id", chat_id).execute(),
        )
        if not result.data:
            return {"chat_id": chat_id, "usage": None, "message": "No usage data found"}
        row = result.data[0]
        row["cost_usd"] = float(row.get("cost_usd", 0))
        return {"chat_id": chat_id, "usage": row}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/usage")
async def list_usage(limit: int = 100, offset: int = 0):
    """List LLM token usage across all chats ordered by most recent, with aggregate totals."""
    if not supabase:
        raise HTTPException(status_code=503, detail="Supabase not configured")
    loop = asyncio.get_event_loop()
    try:
        page_result, count_result, agg_result = await asyncio.gather(
            loop.run_in_executor(
                None,
                lambda: supabase.table("chat_usage")
                    .select("*")
                    .order("updated_at", desc=True)
                    .range(offset, offset + limit - 1)
                    .execute(),
            ),
            loop.run_in_executor(
                None,
                lambda: supabase.table("chat_usage")
                    .select("*", count="exact")
                    .limit(0)
                    .execute(),
            ),
            loop.run_in_executor(
                None,
                lambda: supabase.table("chat_usage")
                    .select("input_tokens,output_tokens,total_tokens,cost_usd")
                    .execute(),
            ),
        )
        rows = page_result.data or []
        for row in rows:
            row["cost_usd"] = float(row.get("cost_usd", 0))
        total_count = count_result.count or 0
        all_rows = agg_result.data or []
        total_input = sum(r.get("input_tokens", 0) or 0 for r in all_rows)
        total_output = sum(r.get("output_tokens", 0) or 0 for r in all_rows)
        total_tokens = sum(r.get("total_tokens", 0) or 0 for r in all_rows)
        total_cost = sum(float(r.get("cost_usd", 0) or 0) for r in all_rows)
        return {
            "usage": rows,
            "count": total_count,
            "totals": {
                "input_tokens": total_input,
                "output_tokens": total_output,
                "total_tokens": total_tokens,
                "cost_usd": round(total_cost, 6),
            },
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ---------------------------------------------------------------------------
# MCP Streamable HTTP endpoint — Claude Desktop / Claude.ai integration
# ---------------------------------------------------------------------------
from mcp.server.fastmcp import FastMCP as _FastMCP
from mcp.server.transport_security import TransportSecuritySettings as _TransportSecuritySettings

# Disable DNS rebinding protection — we enforce auth via Bearer token in _MCPGateway,
# so the built-in host/origin header check is not needed and would reject Claude.ai's
# external IPs (which send the Replit dev domain as the Host header).
_mcp = _FastMCP(
    "DeepAgent",
    instructions=(
        "This server exposes these tools:\n\n"
        "1. get_apps — returns all connected integrations (apps) with their live status "
        "and tool count. Currently 15 healthy apps: Salesforce, Avoma, Apollo, ZoomInfo, "
        "Seamless.ai, Lemlist, ZeroBounce, Showpad, Gojiberry, Wiza, Lusha, LinkedIn, "
        "Clearout, Eloqua, and Mailchimp.\n\n"
        "2. get_tools — returns the full list of tools for a given app (pass app_name), "
        "or all tools across every app (leave app_name empty). Each tool includes its "
        "name, app, description, and FULL parameter schema (type, required, description).\n\n"
        "3. call_tool — invoke any tool by name with arguments. Supports fuzzy name "
        "matching so approximate tool names are auto-corrected.\n\n"
        "4. smart_call — PREFERRED for natural language requests. Pass app_name + "
        "intent (what you want to do) + context (dates, names, etc. from the user's "
        "message). The server automatically selects the right tool, infers the correct "
        "arguments, and executes — all in one step. Use this whenever the user describes "
        "what they want in plain English.\n\n"
        "Opportunity Observatory (pre-computed deal dossiers) — call these THREE by "
        "exact name; do NOT route them through call_tool/smart_call:\n"
        "5. list_opportunity_dossiers — lightweight header rows (find the right opp).\n"
        "6. get_opportunity_dossier — one full dossier by opportunity_id (optionally a "
        "subset of sections).\n"
        "7. search_opportunity_dossiers — fuzzy search by opportunity/account name.\n\n"
        "Structured opportunity + Avoma meeting data (cache + diagnosis history) — also "
        "call these EIGHT by exact name; all are read-only and return structured JSON with "
        "limit/offset pagination (has_more + next_offset):\n"
        "8. list_cached_opportunities — list/filter the opportunity_cache (by momentum, "
        "stage, amount, is_closed, etc.).\n"
        "9. search_cached_opportunities — substring search opportunity_cache by "
        "opportunity/account name.\n"
        "10. get_cached_opportunity — one opportunity's full cached state + its linked "
        "meetings + recent field-history, by opportunity_id.\n"
        "11. get_opportunity_field_history — field-change history for one opportunity "
        "(optionally filtered to one field_name).\n"
        "12. list_opportunity_diagnoses — prior diagnosis runs (full text + momentum "
        "verdict, health rating, top risks, timestamp) by account_id and/or opportunity_id.\n"
        "13. list_opportunity_meetings — Avoma meetings linked to one opportunity "
        "(meeting_cache: title, date, transcript summary), by opportunity_id.\n"
        "14. get_meeting_analysis — per-meeting Avoma AI analysis reports for one "
        "opportunity (conflicts, win likelihood, evidence), by opportunity_id.\n"
        "15. find_meetings_by_name — find a deal's meetings (and optionally analysis) "
        "directly from a company/deal name, no opportunity_id needed."
    ),
    transport_security=_TransportSecuritySettings(enable_dns_rebinding_protection=False),
)


# ---------------------------------------------------------------------------
# Helper: extract full parameter schema from a LangChain/MCP tool
# ---------------------------------------------------------------------------

def _extract_tool_param_schema(t) -> list:
    """Return a rich parameter list: name, type, required, description, default, enum.

    Tries multiple approaches to handle both Pydantic v1 (.schema()) and
    v2 (.model_json_schema()), as well as LangChain's .args_schema and .args fallbacks.
    """
    def _parse_schema_dict(schema: dict) -> list:
        props = schema.get("properties", {})
        required_set = set(schema.get("required", []))
        out = []
        for pname, pinfo in props.items():
            entry = {
                "name":        pname,
                "type":        pinfo.get("type", "any"),
                "required":    pname in required_set,
                "description": pinfo.get("description", ""),
            }
            if "default" in pinfo:
                entry["default"] = pinfo["default"]
            if "enum" in pinfo:
                entry["enum"] = pinfo["enum"]
            out.append(entry)
        return out

    # Strategy 1: args_schema.model_json_schema() — Pydantic v2
    try:
        schema = t.args_schema.model_json_schema()
        result = _parse_schema_dict(schema)
        if result:
            return result
    except Exception:
        pass

    # Strategy 2: args_schema.schema() — Pydantic v1
    try:
        schema = t.args_schema.schema()
        result = _parse_schema_dict(schema)
        if result:
            return result
    except Exception:
        pass

    # Strategy 3: get_input_schema().model_json_schema() — Pydantic v2 via LangChain
    try:
        schema = t.get_input_schema().model_json_schema()
        result = _parse_schema_dict(schema)
        if result:
            return result
    except Exception:
        pass

    # Strategy 4: get_input_schema().schema() — Pydantic v1 via LangChain
    try:
        schema = t.get_input_schema().schema()
        result = _parse_schema_dict(schema)
        if result:
            return result
    except Exception:
        pass

    # Strategy 5: t.args — LangChain dict of {name: type_annotation}
    try:
        return [
            {"name": k, "type": str(v), "required": False, "description": ""}
            for k, v in (t.args or {}).items()
        ]
    except Exception:
        pass

    return []


# ---------------------------------------------------------------------------
# MCP Gateway tools: get_apps  /  get_tools
# ---------------------------------------------------------------------------

@_mcp.tool()
async def get_apps() -> dict:
    """Return every app (integration) registered in the server's MCP config,
    together with its live status and tool count.

    Returns a JSON object with:
      - total_apps: total number of apps in mcp_config.json
      - enabled_count: number of apps currently enabled
      - apps: list of app objects, each with:
          name        – app identifier (e.g. "mailchimp", "salesforce")
          enabled     – whether the app is turned on in config
          status      – "healthy" | "loading" | "failed" | "disabled"
          tool_count  – number of tools loaded for this app (0 if not loaded yet)
    """
    all_servers = agent_manager.mcp_config_manager.config.get("mcp_servers", {})
    cached      = agent_manager._cached_mcp_tools_by_server
    loading_st  = getattr(agent_manager, "mcp_loading_status", {})

    apps = []
    for name, cfg in all_servers.items():
        enabled = cfg.get("enabled", False)
        tool_count = len(cached.get(name, []))
        raw_status = loading_st.get(name, "")
        if not enabled:
            status = "disabled"
        elif tool_count > 0:
            status = "healthy"
        elif "failed" in raw_status.lower():
            status = "failed"
        else:
            status = "loading"

        apps.append({
            "name":       name,
            "enabled":    enabled,
            "status":     status,
            "tool_count": tool_count,
        })

    # sort: enabled+healthy first, then by name
    apps.sort(key=lambda a: (0 if a["status"] == "healthy" else 1, a["name"]))

    return {
        "total_apps":    len(apps),
        "enabled_count": sum(1 for a in apps if a["enabled"]),
        "healthy_count": sum(1 for a in apps if a["status"] == "healthy"),
        "apps":          apps,
    }


@_mcp.tool()
async def get_tools(app_name: str = "") -> dict:
    """Return the tools available on this server, optionally filtered to a single app.

    Args:
        app_name: Name of the app to filter by (e.g. "mailchimp", "salesforce").
                  Leave empty to retrieve tools for ALL apps.

    Returns a JSON object with:
      - total_tools: number of tools returned
      - app_filter:  the app_name argument that was used (empty = all)
      - tools: list of tool objects, each with:
          name        – tool name (e.g. "mailchimp_full_campaign_report")
          app         – which app this tool belongs to
          description – full tool description
          parameters  – list of parameter names accepted by the tool
    """
    cached = agent_manager._cached_mcp_tools_by_server
    filter_name = (app_name or "").strip().lower()

    tools_out = []

    servers_to_scan = (
        {filter_name: cached.get(filter_name, [])}
        if filter_name
        else cached
    )

    _is_admin = _is_admin_ctx.get()
    for srv_name, tool_list in servers_to_scan.items():
        for t in tool_list:
            if not _is_admin and t.name in _ADMIN_ONLY_TOOLS:
                continue
            tools_out.append({
                "name":        t.name,
                "app":         srv_name,
                "description": (t.description or "").strip(),
                "parameters":  _extract_tool_param_schema(t),
            })

    # stable sort: by app then tool name
    tools_out.sort(key=lambda x: (x["app"], x["name"]))

    if filter_name and filter_name not in cached:
        known = sorted(cached.keys())
        return {
            "error":       f"App '{filter_name}' not found or not loaded yet.",
            "known_apps":  known,
            "total_tools": 0,
            "app_filter":  filter_name,
            "tools":       [],
        }

    return {
        "total_tools": len(tools_out),
        "app_filter":  filter_name,
        "tools":       tools_out,
    }


@_mcp.tool()
async def call_tool(tool_name: str, arguments: dict = {}) -> str:
    """Invoke any tool connected to DeepAgent by name.

    Use get_apps to see which apps are connected, then get_tools(app_name=...) to
    browse available tools and their parameters. Then call this tool to run them.

    Supports fuzzy name matching — if the exact tool name is not found, the closest
    match is used automatically and noted in the response.

    Args:
        tool_name:  Tool name (e.g. "mailchimp_full_campaign_report", "soql").
                    Approximate names are auto-corrected via fuzzy matching.
        arguments:  Dict of parameter names → values for the tool.
                    Pass {} or omit for tools that take no arguments.

    Returns:
        The tool's result as a string. On error, returns a descriptive message.
    """
    import difflib as _difflib

    cached = agent_manager._cached_mcp_tools_by_server

    # build a flat name → tool lookup across all servers
    tool_map: dict = {}
    for tool_list in cached.values():
        for t in tool_list:
            tool_map[t.name] = t

    autocorrect_note = ""
    if tool_name not in tool_map:
        close = _difflib.get_close_matches(tool_name, list(tool_map.keys()), n=3, cutoff=0.45)
        if close:
            matched = close[0]
            autocorrect_note = f"[Auto-corrected tool name: '{tool_name}' → '{matched}']\n\n"
            tool_name = matched
        else:
            known_count = len(tool_map)
            return (
                f"Tool '{tool_name}' not found and no close match. "
                f"Use get_tools() to browse all {known_count} available tools."
            )

    if tool_name in _ADMIN_ONLY_TOOLS and not _is_admin_ctx.get():
        return (
            f"Tool '{tool_name}' is admin-only. "
            f"Reconnect using ADMIN_DISPATCH_SECRET as your Bearer token to use it."
        )

    tool = tool_map[tool_name]
    try:
        result = await tool.ainvoke(arguments or {})
        result_str = result if isinstance(result, str) else json.dumps(result, default=str)
        return autocorrect_note + result_str
    except Exception as exc:
        return f"{autocorrect_note}Tool '{tool_name}' raised {type(exc).__name__}: {exc}"


@_mcp.tool()
async def smart_call(app_name: str, intent: str, context: str = "") -> str:
    """Discover the right tool and execute it in ONE step from a natural language request.

    PREFER this over calling get_tools + call_tool separately. Pass what the user
    wants in plain English — the server selects the correct tool, infers argument
    values (including date parsing, boolean flags, etc.), and runs it automatically.

    Args:
        app_name: Integration to use, e.g. "mailchimp", "salesforce", "lemlist".
                  Use get_apps() if unsure which apps are available.
        intent:   What the user wants to do in plain English.
                  e.g. "get campaign performance for April 20 to 24"
        context:  Any extra values from the user's message: specific dates, names,
                  IDs, email addresses, filters, etc.
                  e.g. "April 20-24 2026, exclude bot activity"

    Returns:
        JSON with tool_selected, arguments_used, reasoning, and the tool result.
    """
    import difflib as _difflib
    import httpx as _httpx
    from datetime import datetime as _datetime

    OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "")
    if not OPENAI_API_KEY:
        return json.dumps({"error": "OPENAI_API_KEY not configured — smart_call requires it."})

    cached = agent_manager._cached_mcp_tools_by_server
    filter_name = (app_name or "").strip().lower()

    if filter_name not in cached:
        known = sorted(cached.keys())
        return json.dumps({
            "error": f"App '{filter_name}' not found or not loaded yet.",
            "known_apps": known,
        })

    _is_admin = _is_admin_ctx.get()
    tool_list = [t for t in cached[filter_name] if _is_admin or t.name not in _ADMIN_ONLY_TOOLS]
    # flat map for execution (also gated)
    tool_map: dict = {}
    for srv_tools in cached.values():
        for t in srv_tools:
            if not _is_admin and t.name in _ADMIN_ONLY_TOOLS:
                continue
            tool_map[t.name] = t

    # Build compact tool catalog for the LLM (cap description at 300 chars to save tokens)
    catalog = []
    for t in tool_list:
        params = _extract_tool_param_schema(t)
        catalog.append({
            "name":        t.name,
            "description": (t.description or "").strip()[:300],
            "parameters":  params,
        })

    today = _datetime.utcnow().strftime("%Y-%m-%d")

    system_prompt = (
        f"You are a tool-selection assistant. Given the user's intent and a list of "
        f"available tools for the \"{app_name}\" integration, select the single best "
        f"tool and extract the correct argument values from the context.\n\n"
        f"Today's date (UTC): {today}\n\n"
        "Rules:\n"
        "- Select EXACTLY ONE tool that best matches the intent.\n"
        "- Extract argument values from the context. Infer from today's date when "
        "  relevant (e.g. 'April 20-24' → ISO 8601 datetimes for the current year).\n"
        "- Datetime parameters must be ISO 8601: \"YYYY-MM-DDTHH:MM:SSZ\".\n"
        "- Boolean parameters use true/false (not strings).\n"
        "- Only include parameters that are required or clearly implied by the intent.\n"
        "- Respond ONLY with valid JSON — no markdown, no explanation outside JSON.\n\n"
        "Response format:\n"
        "{\n"
        "  \"tool_name\": \"exact_tool_name\",\n"
        "  \"arguments\": {\"param\": value},\n"
        "  \"reasoning\": \"one sentence\"\n"
        "}"
    )

    user_message = (
        f"Intent: {intent}\n"
        f"Context: {context}\n\n"
        f"Available tools:\n{json.dumps(catalog, indent=2)}"
    )

    try:
        resp = _httpx.post(
            "https://api.openai.com/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENAI_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "model": "gpt-4o-mini",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user",   "content": user_message},
                ],
                "response_format": {"type": "json_object"},
                "temperature": 0,
                "max_tokens": 600,
            },
            timeout=30.0,
        )
        resp.raise_for_status()
    except Exception as exc:
        return json.dumps({"error": f"LLM call failed: {type(exc).__name__}: {exc}"})

    try:
        selection = json.loads(resp.json()["choices"][0]["message"]["content"])
    except Exception as exc:
        return json.dumps({"error": f"Failed to parse LLM response: {exc}"})

    selected_name = selection.get("tool_name", "")
    selected_args = selection.get("arguments", {})
    reasoning     = selection.get("reasoning", "")

    # Fuzzy-correct if LLM returned a slightly wrong name
    if selected_name not in tool_map:
        close = _difflib.get_close_matches(selected_name, list(tool_map.keys()), n=1, cutoff=0.45)
        if close:
            selected_name = close[0]
        else:
            return json.dumps({
                "error":      f"LLM selected '{selected_name}' which doesn't match any loaded tool.",
                "llm_output": selection,
            })

    tool = tool_map[selected_name]
    try:
        result = await tool.ainvoke(selected_args)
        result_str = result if isinstance(result, str) else json.dumps(result, default=str)
        return json.dumps({
            "tool_selected":  selected_name,
            "arguments_used": selected_args,
            "reasoning":      reasoning,
            "result":         result_str,
        }, default=str)
    except Exception as exc:
        return json.dumps({
            "tool_selected":  selected_name,
            "arguments_used": selected_args,
            "reasoning":      reasoning,
            "error":          f"{type(exc).__name__}: {exc}",
        })


# ---------------------------------------------------------------------------
# Opportunity Observatory — first-class MCP tools (exposed by EXACT name).
#
# These wrap the read-only @tool functions in custom_tools/opportunity_observatory.py
# so external MCP clients (Claude Desktop / Claude.ai) see them directly in
# tools/list and call them by name — instead of falling through to call_tool's
# fuzzy matcher (which previously mis-routed `list_opportunity_dossiers` to an
# unrelated `list_users` tool). Each is hard-scoped to the single
# opportunity_observatory table; no write path. Sync httpx reads run in an
# executor so they don't block the event loop.
# ---------------------------------------------------------------------------
from custom_tools import opportunity_observatory as _obs


@_mcp.tool()
async def list_opportunity_dossiers(
    limit: int = 50,
    stage: Optional[str] = None,
    account_name_contains: Optional[str] = None,
    name_contains: Optional[str] = None,
) -> str:
    """List Opportunity Observatory dossiers (lightweight — header fields only).

    The Observatory holds one rich, pre-computed dossier per opportunity (SF
    90-day evidence, Avoma evidence, outbound/campaign intelligence, and four
    diagnostic bundles A-D plus a final diagnosis sheet). This returns ONLY the
    header fields so you can find the right opportunity before pulling its full
    dossier with get_opportunity_dossier.

    Args:
        limit:                 Max rows (default 50, hard cap 200).
        stage:                 Exact stage filter (e.g. 'Qualified', 'Shortlisted').
        account_name_contains: Case-insensitive substring match on account_name.
        name_contains:         Case-insensitive substring match on opportunity name.

    Returns:
        JSON string: {count, dossiers:[{opportunity_id, name, opportunity_owner,
        close_date, amount, stage, account_name, updated_at}]}.
    """
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, lambda: _obs.list_opportunity_dossiers.invoke({
        "limit": limit,
        "stage": stage,
        "account_name_contains": account_name_contains,
        "name_contains": name_contains,
    }))


@_mcp.tool()
async def get_opportunity_dossier(opportunity_id: str, sections: Optional[str] = None) -> str:
    """Fetch one full Opportunity Observatory dossier by opportunity_id.

    Returns the header fields plus long-form analysis sections (multi-paragraph
    markdown). Pull only the sections you need for large dossiers.

    Available sections: sf_90day_evidence, avoma_evidence,
    outbound_campaign_intelligence, bundle_a_deal_progress,
    bundle_b_competition_fit, bundle_c_stakeholder_map, bundle_d_vulnerabilities,
    diagnosis_sheet.

    Args:
        opportunity_id: Salesforce Opportunity Id (from list/search).
        sections:       Optional comma-separated subset of section names above.
                        Omit to return ALL sections (can be large).

    Returns:
        JSON string with the header fields + requested section(s), or an error.
    """
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, lambda: _obs.get_opportunity_dossier.invoke({
        "opportunity_id": opportunity_id,
        "sections": sections,
    }))


@_mcp.tool()
async def search_opportunity_dossiers(query: str, limit: int = 20) -> str:
    """Substring search the Observatory across opportunity name + account name.

    Case-insensitive ILIKE match on either `name` or `account_name`. Use this
    when you have a fuzzy company or deal name. Returns lightweight header rows;
    follow up with get_opportunity_dossier for the full content.

    Args:
        query: Search text (e.g. 'Bright Horizons', 'Anora'). Matched as a
               substring against both name and account_name.
        limit: Max rows (default 20, hard cap 200).

    Returns:
        JSON string: {count, dossiers:[...header fields...]}.
    """
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, lambda: _obs.search_opportunity_dossiers.invoke({
        "query": query,
        "limit": limit,
    }))


# ---------------------------------------------------------------------------
# Structured opportunity data exposed first-class on /mcp (cache + diagnoses).
#
# These wrap the read-only query paths in cache_qa.py (opportunity_cache,
# meeting_cache, field_history_cache) and lake.py (lake.opportunity_diagnoses)
# so external clients get precise, named, structured-JSON read tools instead of
# falling through call_tool's fuzzy matcher. Sync DB work runs in an executor
# following the Observatory pattern; every tool is strictly read-only and
# hard-scoped to a single table by the delegate. Pagination via limit + offset
# returns has_more / next_offset so a client can page reliably.
# ---------------------------------------------------------------------------

import cache_qa as _cqa


def _cache_ready() -> Optional[str]:
    """Return a JSON error string if Supabase isn't configured, else None."""
    if not supabase:
        return json.dumps({"error": "Supabase not configured"})
    return None


@_mcp.tool()
async def list_cached_opportunities(
    limit: int = 25,
    offset: int = 0,
    momentum: Optional[str] = None,
    stage: Optional[str] = None,
    min_amount: Optional[float] = None,
    max_amount: Optional[float] = None,
    max_meetings: Optional[int] = None,
    is_closed: Optional[bool] = None,
    order_by: str = "amount",
    descending: bool = True,
) -> str:
    """List/filter opportunities from the fast Salesforce mirror (opportunity_cache).

    Each row is the cached state of one Salesforce Opportunity with computed
    health fields. Use this to find opportunities by momentum or stage before
    pulling one with get_cached_opportunity.

    Args:
        limit:        Max rows per page (default 25, hard cap 100).
        offset:       Row offset for pagination (default 0).
        momentum:     Filter by momentum bucket (Active/Moderate/Slow/Stalled).
        stage:        Case-insensitive substring match on stage name.
        min_amount:   Minimum deal amount.
        max_amount:   Maximum deal amount.
        max_meetings: Max meetings_count (use 0 to find deals with no meetings).
        is_closed:    Filter on the IsClosed flag.
        order_by:     amount | health_score | days_in_stage | probability |
                      days_since_last_meeting (default amount).
        descending:   Sort descending (default true).

    Returns:
        JSON string: {count, opportunities:[...], has_more, next_offset}.
    """
    err = _cache_ready()
    if err:
        return err
    loop = asyncio.get_running_loop()
    try:
        out = await loop.run_in_executor(None, lambda: _cqa.read_filter_opportunities(
            supabase, momentum=momentum, stage=stage, min_amount=min_amount,
            max_amount=max_amount, max_meetings=max_meetings, is_closed=is_closed,
            order_by=order_by, descending=descending, limit=limit, offset=offset))
        return json.dumps(out, default=str)
    except Exception as e:  # noqa: BLE001
        return json.dumps({"error": f"list_cached_opportunities failed: {type(e).__name__}: {e}"})


@_mcp.tool()
async def search_cached_opportunities(query: str, limit: int = 10, offset: int = 0) -> str:
    """Substring search opportunity_cache by opportunity name OR account name.

    Case-insensitive match. Use when you have a fuzzy company/deal name; follow
    up with get_cached_opportunity for the full record.

    Args:
        query:  Search text (matched against opportunity_name and account_name).
        limit:  Max rows per page (default 10, hard cap 50).
        offset: Row offset for pagination (default 0).

    Returns:
        JSON string: {count, opportunities:[...], has_more, next_offset}.
    """
    err = _cache_ready()
    if err:
        return err
    loop = asyncio.get_running_loop()
    try:
        out = await loop.run_in_executor(None, lambda: _cqa.read_search_opportunities(
            supabase, query, limit=limit, offset=offset))
        return json.dumps(out, default=str)
    except Exception as e:  # noqa: BLE001
        return json.dumps({"error": f"search_cached_opportunities failed: {type(e).__name__}: {e}"})


@_mcp.tool()
async def get_cached_opportunity(opportunity_id: str) -> str:
    """Fetch one opportunity's full cached state + its linked meetings + recent
    field-history changes, by Salesforce Opportunity Id.

    Args:
        opportunity_id: Salesforce Opportunity Id (from the list/search tools).

    Returns:
        JSON string: {opportunity:{...all cached fields...},
        meetings:[{meeting_uuid, meeting_title, meeting_date, transcript_summary}],
        field_history:[...recent changes...]}, or {error:...} if unknown.
    """
    err = _cache_ready()
    if err:
        return err
    loop = asyncio.get_running_loop()
    try:
        out = await loop.run_in_executor(None, lambda: _cqa.read_opportunity_detail(
            supabase, opportunity_id))
        return json.dumps(out, default=str)
    except Exception as e:  # noqa: BLE001
        return json.dumps({"error": f"get_cached_opportunity failed: {type(e).__name__}: {e}"})


@_mcp.tool()
async def get_opportunity_field_history(
    opportunity_id: str,
    field_name: Optional[str] = None,
    limit: int = 30,
    offset: int = 0,
) -> str:
    """Field-change history for one opportunity (e.g. how Amount or StageName
    changed over time), newest first, from field_history_cache.

    Args:
        opportunity_id: Salesforce Opportunity Id.
        field_name:     Optional exact field name to filter to (e.g. StageName).
        limit:          Max rows per page (default 30, hard cap 100).
        offset:         Row offset for pagination (default 0).

    Returns:
        JSON string: {count, history:[...changes...], has_more, next_offset}.
    """
    err = _cache_ready()
    if err:
        return err
    loop = asyncio.get_running_loop()
    try:
        out = await loop.run_in_executor(None, lambda: _cqa.read_field_history(
            supabase, opportunity_id, field_name=field_name, limit=limit, offset=offset))
        return json.dumps(out, default=str)
    except Exception as e:  # noqa: BLE001
        return json.dumps({"error": f"get_opportunity_field_history failed: {type(e).__name__}: {e}"})


@_mcp.tool()
async def list_opportunity_diagnoses(
    account_id: Optional[str] = None,
    opportunity_id: Optional[str] = None,
    limit: int = 20,
    offset: int = 0,
) -> str:
    """List prior Opportunity Diagnosis runs from the lake (lake.opportunity_diagnoses),
    newest first. Filter by account_id and/or opportunity_id.

    Each run includes the full diagnosis text plus the structured verdict fields.

    Args:
        account_id:     Salesforce Account Id (001…) to filter by.
        opportunity_id: Salesforce Opportunity Id (006…) to filter by.
        limit:          Max rows per page (default 20, hard cap 50).
        offset:         Row offset for pagination (default 0).

    Returns:
        JSON string: {count, diagnoses:[{run_at, account_id, account_name,
        opportunity_id, opportunity_name, stage, amount, momentum_verdict,
        health_rating, top_risks, recommendations, diagnosis_md}], has_more,
        next_offset}.
    """
    err = _cache_ready()
    if err:
        return err
    loop = asyncio.get_running_loop()
    try:
        out = await loop.run_in_executor(None, lambda: _lake.read_diagnoses(
            supabase, account_id=account_id, opportunity_id=opportunity_id,
            limit=limit, offset=offset))
        return json.dumps(out, default=str)
    except Exception as e:  # noqa: BLE001
        return json.dumps({"error": f"list_opportunity_diagnoses failed: {type(e).__name__}: {e}"})


@_mcp.tool()
async def list_opportunity_meetings(
    opportunity_id: str,
    limit: int = 20,
    offset: int = 0,
) -> str:
    """List the Avoma meetings linked to one opportunity, newest first, from
    meeting_cache.

    Args:
        opportunity_id: Salesforce Opportunity Id.
        limit:          Max rows per page (default 20, hard cap 100).
        offset:         Row offset for pagination (default 0).

    Returns:
        JSON string: {count, meetings:[{meeting_uuid, meeting_title,
        meeting_date, transcript_summary}], has_more, next_offset}.
    """
    err = _cache_ready()
    if err:
        return err
    loop = asyncio.get_running_loop()
    try:
        out = await loop.run_in_executor(None, lambda: _cqa.read_opportunity_meetings(
            supabase, opportunity_id, limit=limit, offset=offset))
        return json.dumps(out, default=str)
    except Exception as e:  # noqa: BLE001
        return json.dumps({"error": f"list_opportunity_meetings failed: {type(e).__name__}: {e}"})


@_mcp.tool()
async def get_meeting_analysis(
    opportunity_id: str,
    limit: int = 5,
    offset: int = 0,
) -> str:
    """Per-meeting Avoma AI analysis reports for one opportunity (conflicts, win
    likelihood, evidence), newest first, from avoma_event_reports.

    Args:
        opportunity_id: Salesforce Opportunity Id.
        limit:          Max rows per page (default 5, hard cap 50).
        offset:         Row offset for pagination (default 0).

    Returns:
        JSON string: {count, reports:[{message_id, meeting_uuid,
        sf_opportunity_id, opportunity_analysis_data, opportunity_analysis_status,
        status, created_at}], has_more, next_offset}.
    """
    err = _cache_ready()
    if err:
        return err
    loop = asyncio.get_running_loop()
    try:
        out = await loop.run_in_executor(None, lambda: _cqa.read_meeting_analysis(
            supabase, opportunity_id, limit=limit, offset=offset))
        return json.dumps(out, default=str)
    except Exception as e:  # noqa: BLE001
        return json.dumps({"error": f"get_meeting_analysis failed: {type(e).__name__}: {e}"})


@_mcp.tool()
async def find_meetings_by_name(
    query: str,
    include_analysis: bool = False,
    opp_limit: int = 5,
    meeting_limit: int = 20,
    analysis_limit: int = 5,
) -> str:
    """Find a deal's Avoma meetings directly from a company OR deal name — no
    opportunity_id required.

    Substring-matches the name against opportunity_name AND account_name, then
    returns each matched opportunity with its linked meetings (and, when
    include_analysis is true, its per-meeting Avoma AI analysis). Use this when
    you only know a fuzzy company/deal name; it folds the
    search_cached_opportunities -> list_opportunity_meetings steps into one call.

    Args:
        query:          Company or deal name to search for (case-insensitive).
        include_analysis: When true, also attach each opportunity's Avoma AI
                        analysis reports (default false).
        opp_limit:      Max matched opportunities to return (default 5, cap 25).
        meeting_limit:  Max meetings per opportunity (default 20, hard cap 100).
        analysis_limit: Max analysis reports per opportunity when
                        include_analysis (default 5, hard cap 50).

    Returns:
        JSON string: {query, matched_opportunities, opportunities:[{opportunity_id,
        opportunity_name, account_name, meetings:[{meeting_uuid, meeting_title,
        meeting_date, transcript_summary}], meetings_has_more, meetings_next_offset,
        analysis?:[...], analysis_has_more?, analysis_next_offset?}],
        opportunities_has_more, opportunities_next_offset}.
    """
    err = _cache_ready()
    if err:
        return err
    loop = asyncio.get_running_loop()
    try:
        out = await loop.run_in_executor(None, lambda: _cqa.read_meetings_by_name(
            supabase, query, include_analysis=include_analysis,
            opp_limit=opp_limit, meeting_limit=meeting_limit,
            analysis_limit=analysis_limit))
        return json.dumps(out, default=str)
    except Exception as e:  # noqa: BLE001
        return json.dumps({"error": f"find_meetings_by_name failed: {type(e).__name__}: {e}"})


# ---------------------------------------------------------------------------
# OAuth 2.0 – Authorization Code flow for Claude.ai custom connectors.
#
# Claude.ai redirects the user to GET /authorize with response_type=code.
# The user enters their DISPATCH_SECRET on a simple HTML form.
# On success we generate a one-time code and redirect back to Claude.
# Claude then exchanges the code for a Bearer token at POST /oauth/token.
# Subsequent MCP calls use  Authorization: Bearer <token>  which our
# _MCPGateway already validates.
#
# To connect in Claude.ai → Settings → Integrations → Add custom connector:
#   URL: https://<domain>/mcp   (no OAuth Client ID/Secret needed)
# ---------------------------------------------------------------------------

import secrets as _secrets_mod
from datetime import datetime as _dt, timedelta as _td

# In-memory store of one-time auth codes: code -> {client_id, redirect_uri, expires_at}
_auth_codes: dict = {}


@app.get("/.well-known/oauth-authorization-server")
async def oauth_server_metadata(request: Request):
    print(f"[OAUTH_DISCOVERY] /.well-known/oauth-authorization-server | ip={request.client.host if request.client else '?'} | ua={request.headers.get('user-agent','')[:80]}", flush=True)
    base = str(request.base_url).rstrip("/")
    return {
        "issuer": base,
        "authorization_endpoint": f"{base}/authorize",
        "token_endpoint": f"{base}/oauth/token",
        "registration_endpoint": f"{base}/oauth/register",
        "grant_types_supported": ["authorization_code", "client_credentials"],
        "response_types_supported": ["code"],
        "code_challenge_methods_supported": ["S256", "plain"],
        "token_endpoint_auth_methods_supported": [
            "none",
            "client_secret_post",
            "client_secret_basic",
        ],
        "scopes_supported": ["mcp"],
    }


@app.get("/authorize")
async def oauth_authorize_get(request: Request):
    """Show a simple login page for Claude.ai's Authorization Code OAuth flow."""
    from fastapi.responses import HTMLResponse as _HTML

    response_type = request.query_params.get("response_type", "code")
    client_id = request.query_params.get("client_id", "")
    redirect_uri = request.query_params.get("redirect_uri", "")
    state = request.query_params.get("state", "")
    code_challenge = request.query_params.get("code_challenge", "")
    code_challenge_method = request.query_params.get("code_challenge_method", "")
    print(f"[OAUTH_AUTHORIZE] GET /authorize | response_type={response_type!r} | client_id={client_id!r} | redirect_uri={redirect_uri!r}", flush=True)

    # Accept any response_type Claude sends — always issue an auth code.
    # (Strict checking blocks Claude.ai which may not always send response_type=code)

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>DeepAgent — Authorize</title>
  <style>
    * {{ box-sizing: border-box; margin: 0; padding: 0; }}
    body {{ font-family: system-ui, -apple-system, sans-serif; background: #0f0f1a;
            color: #e2e8f0; display: flex; align-items: center; justify-content: center;
            min-height: 100vh; padding: 20px; }}
    .card {{ background: #1a1a2e; border: 1px solid #2d2d4e; border-radius: 12px;
             padding: 40px 36px; width: 100%; max-width: 400px; }}
    h2 {{ font-size: 22px; color: #818cf8; margin-bottom: 8px; }}
    .subtitle {{ color: #94a3b8; font-size: 14px; margin-bottom: 28px; line-height: 1.5; }}
    label {{ display: block; font-size: 13px; color: #94a3b8; margin-bottom: 6px; }}
    input[type=password] {{ width: 100%; padding: 10px 14px; background: #0f0f1a;
                            border: 1px solid #3d3d5e; border-radius: 8px; color: #e2e8f0;
                            font-size: 15px; outline: none; }}
    input[type=password]:focus {{ border-color: #818cf8; }}
    button {{ width: 100%; padding: 11px; background: #818cf8; color: white; border: none;
              border-radius: 8px; font-size: 15px; font-weight: 600; cursor: pointer;
              margin-top: 16px; transition: background 0.2s; }}
    button:hover {{ background: #6366f1; }}
    .hint {{ color: #64748b; font-size: 12px; margin-top: 16px; text-align: center; }}
  </style>
</head>
<body>
  <div class="card">
    <h2>🤖 DeepAgent</h2>
    <p class="subtitle">
      <strong style="color:#e2e8f0">Claude</strong> is requesting access to your
      DeepAgent server with 173 AI tools.<br><br>
      Enter your secret key to authorize.
    </p>
    <form method="POST" action="/authorize">
      <input type="hidden" name="response_type" value="{response_type}">
      <input type="hidden" name="client_id" value="{client_id}">
      <input type="hidden" name="redirect_uri" value="{redirect_uri}">
      <input type="hidden" name="state" value="{state}">
      <input type="hidden" name="code_challenge" value="{code_challenge}">
      <input type="hidden" name="code_challenge_method" value="{code_challenge_method}">
      <label for="secret">Secret Key (DISPATCH_SECRET)</label>
      <input type="password" id="secret" name="secret" placeholder="Paste your secret key" required autofocus>
      <button type="submit">Authorize Access</button>
    </form>
    <p class="hint">Find your key in Replit Secrets → DISPATCH_SECRET</p>
  </div>
</body>
</html>"""
    return _HTML(html)


@app.post("/authorize")
async def oauth_authorize_post(request: Request):
    """Handle the authorization form and redirect back to Claude with an auth code."""
    from fastapi.responses import HTMLResponse as _HTML, RedirectResponse as _Redir
    import urllib.parse as _up

    form = await request.form()
    client_id = form.get("client_id", "")
    redirect_uri = form.get("redirect_uri", "")
    state = form.get("state", "")
    code_challenge = form.get("code_challenge", "")
    code_challenge_method = form.get("code_challenge_method", "")
    user_secret = form.get("secret", "")
    print(f"[OAUTH_AUTHORIZE] POST /authorize | client_id={client_id!r} | redirect_uri={redirect_uri!r}", flush=True)

    secret = os.environ.get("DISPATCH_SECRET", "")
    if not secret or user_secret.strip() != secret:
        return _HTML(
            "<html><body style='font-family:sans-serif;padding:40px'>"
            "<h2 style='color:#e53e3e'>Invalid secret</h2>"
            "<p>The key you entered is incorrect. "
            "<a href='javascript:history.back()' style='color:#818cf8'>Go back and try again.</a></p>"
            "</body></html>",
            status_code=401,
        )

    # Issue a one-time authorization code (valid 10 minutes)
    code = _secrets_mod.token_urlsafe(32)
    _auth_codes[code] = {
        "client_id": client_id,
        "redirect_uri": redirect_uri,
        "code_challenge": code_challenge,
        "code_challenge_method": code_challenge_method,
        "expires_at": _dt.utcnow() + _td(minutes=10),
    }

    params: dict = {"code": code}
    if state:
        params["state"] = state
    sep = "&" if "?" in redirect_uri else "?"
    return _Redir(url=redirect_uri + sep + _up.urlencode(params), status_code=302)


@app.post("/oauth/token")
async def oauth_token(request: Request):
    """Token endpoint — handles authorization_code and client_credentials grants."""
    from fastapi.responses import JSONResponse as _JSON
    import hashlib as _hl, base64 as _b64

    secret = os.environ.get("DISPATCH_SECRET", "")
    client_id_expected = os.environ.get("OAUTH_CLIENT_ID", "claude")

    content_type = request.headers.get("content-type", "")
    if "application/x-www-form-urlencoded" in content_type or "multipart/form-data" in content_type:
        form = await request.form()
        params = dict(form)
    else:
        try:
            params = await request.json()
        except Exception:
            params = {}

    grant_type = params.get("grant_type", "")
    client_id = params.get("client_id", "")
    client_secret = params.get("client_secret", "")

    # HTTP Basic auth fallback
    auth_header = request.headers.get("authorization", "")
    if auth_header.startswith("Basic "):
        try:
            decoded = _b64.b64decode(auth_header[6:]).decode("utf-8")
            basic_id, _, basic_secret = decoded.partition(":")
            if not client_id:
                client_id = basic_id
            if not client_secret:
                client_secret = basic_secret
        except Exception:
            pass

    print(f"[OAUTH_TOKEN] POST /oauth/token | grant_type={grant_type!r} | client_id={client_id!r}", flush=True)

    if not secret:
        return _JSON({"error": "server_error", "error_description": "DISPATCH_SECRET not configured"}, status_code=500)

    # ---- Authorization Code grant ----
    if grant_type == "authorization_code":
        code = params.get("code", "")
        redirect_uri = params.get("redirect_uri", "")
        code_verifier = params.get("code_verifier", "")

        entry = _auth_codes.pop(code, None)
        if entry is None:
            return _JSON({"error": "invalid_grant", "error_description": "Unknown or expired code"}, status_code=400)
        if _dt.utcnow() > entry["expires_at"]:
            return _JSON({"error": "invalid_grant", "error_description": "Code expired"}, status_code=400)

        # PKCE verification (if code_challenge was set during /authorize)
        challenge = entry.get("code_challenge", "")
        method = entry.get("code_challenge_method", "")
        if challenge and code_verifier:
            if method == "S256":
                computed = _b64.urlsafe_b64encode(
                    _hl.sha256(code_verifier.encode()).digest()
                ).rstrip(b"=").decode()
            else:
                computed = code_verifier
            # Strip padding from both sides before comparing (some clients send with padding)
            if computed.rstrip("=") != challenge.rstrip("="):
                print(f"[OAUTH_TOKEN] PKCE fail | method={method!r} | computed={computed!r} | challenge={challenge!r}", flush=True)
                return _JSON({"error": "invalid_grant", "error_description": "PKCE verification failed"}, status_code=400)

        scope_granted = params.get("scope", "mcp")
        print(f"[OAUTH_TOKEN] → 200 authorization_code OK | scope={scope_granted!r}", flush=True)
        return _JSON({"access_token": secret, "token_type": "Bearer", "expires_in": 86400, "scope": scope_granted})

    # ---- Client Credentials grant (Claude Desktop / API clients) ----
    if grant_type == "client_credentials":
        if client_id != client_id_expected or client_secret != secret:
            return _JSON({"error": "invalid_client"}, status_code=401)
        print(f"[OAUTH_TOKEN] → 200 client_credentials OK", flush=True)
        return _JSON({"access_token": secret, "token_type": "Bearer", "expires_in": 86400, "scope": "mcp"})

    return _JSON({"error": "unsupported_grant_type"}, status_code=400)


# ---- RFC 9728: OAuth 2.0 Protected Resource Metadata ----
# Claude.ai fetches this to discover which authorization server protects /mcp.
@app.get("/.well-known/oauth-protected-resource")
async def oauth_protected_resource(request: Request):
    print(f"[OAUTH_DISCOVERY] /.well-known/oauth-protected-resource | ip={request.client.host if request.client else '?'} | ua={request.headers.get('user-agent','')[:80]}", flush=True)
    base = str(request.base_url).rstrip("/")
    return {
        "resource": f"{base}/mcp",
        "authorization_servers": [base],
        "bearer_methods_supported": ["header"],
        "resource_documentation": f"{base}/mcp",
    }


# ---- RFC 7591: OAuth 2.0 Dynamic Client Registration ----
# Claude.ai posts here before starting the auth flow.
# We accept any registration and return a generated client_id (no secret
# since we use PKCE + user-entered DISPATCH_SECRET on the /authorize page).
@app.post("/oauth/register")
async def oauth_register(request: Request):
    from fastapi.responses import JSONResponse as _JSON
    print(f"[OAUTH_DCR] POST /oauth/register | ip={request.client.host if request.client else '?'} | ua={request.headers.get('user-agent','')[:80]}", flush=True)
    try:
        body = await request.json()
    except Exception:
        body = {}

    client_id = _secrets_mod.token_urlsafe(16)
    redirect_uris = body.get("redirect_uris", ["https://claude.ai/api/mcp/auth_callback"])

    return _JSON({
        "client_id": client_id,
        "client_id_issued_at": int(_dt.utcnow().timestamp()),
        "redirect_uris": redirect_uris,
        "grant_types": ["authorization_code"],
        "response_types": ["code"],
        "token_endpoint_auth_method": "none",
        "client_name": body.get("client_name", "MCP Client"),
    }, status_code=201)


_mcp_inner_asgi = _mcp.streamable_http_app()


# ---------------------------------------------------------------------------
# Admin tier — tools that require ADMIN_DISPATCH_SECRET on /mcp Bearer auth.
# Non-admin callers (authenticated with the regular DISPATCH_SECRET) cannot
# see these in get_tools() and cannot invoke them via call_tool / smart_call.
# ---------------------------------------------------------------------------
import contextvars as _contextvars

_ADMIN_ONLY_TOOLS: set = {
    "create_record",
    "update_record",
    "delete_record",
    "create_task",
}

_is_admin_ctx: _contextvars.ContextVar = _contextvars.ContextVar("_is_admin_ctx", default=False)


class _MCPGateway:
    """Auth + path normalisation gateway for the FastMCP Streamable HTTP handler.

    Intercepts every request whose path starts with /mcp BEFORE FastAPI routing,
    so Starlette's Mount never has a chance to issue a 307 redirect.

    Auth: if DISPATCH_SECRET is set, every HTTP request must carry
    Authorization: Bearer <DISPATCH_SECRET>.  If empty, auth is skipped.
    On 401, returns WWW-Authenticate with as_uri so Claude.ai can discover
    our OAuth authorization server automatically.

    Path: FastMCP's streamable_http_app() routes at internal path /mcp.
    We normalise every incoming path to /mcp so the handler always matches.

    Claude Desktop config:
        {"mcpServers": {"deep-agent": {
            "type": "http",
            "url": "https://<domain>/mcp",
            "headers": {"Authorization": "Bearer <DISPATCH_SECRET>"}
        }}}
    """

    def __init__(self, asgi_app):
        self._app = asgi_app

    async def __call__(self, scope, receive, send):
        secret = os.environ.get("DISPATCH_SECRET", "")
        admin_secret = os.environ.get("ADMIN_DISPATCH_SECRET", "")
        # Testing escape hatch: when MCP_ALLOW_UNAUTH is truthy, skip the Bearer
        # gate entirely so the /mcp endpoint (incl. the read-only opp/Avoma tools)
        # can be hit without a token. Leave unset/false in production.
        allow_unauth = os.environ.get("MCP_ALLOW_UNAUTH", "").strip().lower() in (
            "1", "true", "yes", "on")
        is_admin = False
        if secret and not allow_unauth and scope.get("type") == "http":
            headers_raw = dict(scope.get("headers", []))
            method = scope.get("method", "?")
            path = scope.get("path", "?")
            client = scope.get("client", ("?", 0))
            client_ip = client[0] if client else "?"
            auth_header = headers_raw.get(b"authorization", b"").decode("utf-8", errors="replace")
            ua = headers_raw.get(b"user-agent", b"").decode("utf-8", errors="replace")
            # Determine tier: admin secret > regular secret
            if admin_secret and auth_header == f"Bearer {admin_secret}":
                is_admin = True
                tier = "ADMIN"
            elif auth_header == f"Bearer {secret}":
                tier = "REGULAR"
            else:
                tier = "NONE"
            print(
                f"[MCP_GATEWAY] {method} {path} | ip={client_ip} "
                f"| auth={'YES' if auth_header.startswith('Bearer ') else 'NO'} "
                f"| tier={tier} | ua={ua[:80]}",
                flush=True
            )
            if tier == "NONE":
                from starlette.responses import Response as _StResp
                # Build the base URL so Claude.ai can discover our OAuth server
                # via the as_uri parameter (MCP spec + RFC 9728).
                _hmap = headers_raw
                _host = _hmap.get(b"host", b"localhost").decode("utf-8", errors="replace")
                _proto = _hmap.get(b"x-forwarded-proto", b"https").decode("utf-8", errors="replace")
                _as_uri = f"{_proto}://{_host}"
                print(f"[MCP_GATEWAY] → 401 | as_uri={_as_uri}", flush=True)
                resp = _StResp(
                    "Unauthorized",
                    status_code=401,
                    headers={
                        "WWW-Authenticate": f'Bearer realm="DeepAgent", as_uri="{_as_uri}"',
                    },
                )
                await resp(scope, receive, send)
                return

        scope = {**scope, "path": "/mcp", "raw_path": b"/mcp"}

        # Propagate admin tier flag to tool handlers via ContextVar.
        _token = _is_admin_ctx.set(is_admin)

        # Intercept the response to log status code
        _resp_status: list = []
        async def _send_logging(message):
            if message.get("type") == "http.response.start":
                _resp_status.append(message.get("status", 0))
            await send(message)

        try:
            await self._app(scope, receive, _send_logging)
        finally:
            _is_admin_ctx.reset(_token)
        print(f"[MCP_GATEWAY] → {_resp_status[0] if _resp_status else '?'}", flush=True)


_mcp_gateway = _MCPGateway(_mcp_inner_asgi)

# Keep a reference to the FastAPI app so the outer wrapper can delegate
# lifespan events and non-/mcp HTTP traffic back to it.
_fastapi_app = app


# ---------------------------------------------------------------------------
# Data/admin API authentication (added 2026-05-29)
# ---------------------------------------------------------------------------
# Mirrors the /mcp _MCPGateway Bearer auth, extended to every other HTTP +
# WebSocket route. Each non-/mcp request must carry a valid token, EXCEPT the
# public allowlist below (health check, AWS SNS webhook, OAuth flow, LinkedIn
# OAuth redirect, CORS preflight). /mcp keeps its own _MCPGateway auth.
#
# Token source: API_AUTH_TOKEN, falling back to the same DISPATCH_SECRET that
# /mcp uses, so the data API and /mcp share one token by default and auth is
# active out of the box. If neither is set, auth is skipped (fail-open) and a
# loud warning is logged at startup.
#
# A valid token may be supplied three ways:
#   1. Authorization: Bearer <token>   — for API clients / curl (same as /mcp)
#   2. ?key=<token> query param        — for first browser visit; on success an
#      HttpOnly `api_auth` cookie is set so the page's same-origin fetch() calls
#      (e.g. /ask -> /api/ask) keep working without re-supplying the key.
#   3. api_auth cookie                 — set automatically by (2).
_API_AUTH_PUBLIC_EXACT = {
    "/api/health",
    "/webhook",
    "/authorize",
    "/oauth/token",
    "/oauth/register",
    "/oauth2callback",
    "/auth/linkedin",
    "/auth/linkedin/callback",
    # Pre-existing endpoints the user explicitly chose to leave PUBLIC
    # (2026-05-29) — treated as externally-managed / not owned by this app.
    # ⚠️ This re-opens cost (chat) + control-plane (config) surfaces. Remove
    # these entries to put them back behind the token.
    "/api/chat",
    "/api/chat/async",
    "/api/chat/structured/async",
    "/api/config",
    # NOTE: /api/documents/upload was REMOVED from the public allowlist
    # (2026-06-18) — it writes to the knowledge base, so it must require the
    # API token. The admin-gated frontend proxy (/api/documents/*) supplies it;
    # direct unauthenticated uploads are now rejected by the global gate.
}
_API_AUTH_PUBLIC_PREFIX = (
    "/.well-known/",
    # /api/mcp/servers/{name} (POST add + DELETE remove) — user-exempted
    # 2026-05-29. Trailing slash keeps the bare GET /api/mcp/servers list gated.
    "/api/mcp/servers/",
)


def _api_auth_token() -> str:
    return os.environ.get("API_AUTH_TOKEN", "") or os.environ.get("DISPATCH_SECRET", "")


def _api_path_is_public(path: str) -> bool:
    if path in _API_AUTH_PUBLIC_EXACT:
        return True
    return any(path.startswith(p) for p in _API_AUTH_PUBLIC_PREFIX)


def _api_auth_source(scope, expected: str):
    """Return 'header' | 'query' | 'cookie' for a valid token, else None.

    Uses hmac.compare_digest for constant-time comparison (no early-exit timing
    leak on the secret)."""
    import hmac
    headers = dict(scope.get("headers", []))
    auth = headers.get(b"authorization", b"").decode("utf-8", "replace")
    if auth.startswith("Bearer ") and hmac.compare_digest(auth[7:].strip(), expected):
        return "header"
    from urllib.parse import parse_qs
    qs = parse_qs(scope.get("query_string", b"").decode("utf-8", "replace"))
    key = (qs.get("key") or [""])[0]
    if key and hmac.compare_digest(key, expected):
        return "query"
    cookie = headers.get(b"cookie", b"").decode("utf-8", "replace")
    for part in cookie.split(";"):
        if "=" in part:
            k, v = part.strip().split("=", 1)
            if k == "api_auth" and hmac.compare_digest(v, expected):
                return "cookie"
    return None


class _AppWithMCP:
    """Outer ASGI wrapper that routes /mcp traffic to FastMCP and everything
    else (including startup/shutdown lifespan) to the inner FastAPI app.

    Using a wrapper rather than app.mount() avoids the Starlette 307 redirect
    that fires when a request path equals the mount prefix exactly (/mcp → /mcp/).

    Also enforces data/admin API auth (see notes above) for all non-/mcp
    HTTP + WebSocket traffic.
    """

    def __init__(self, fastapi_app, mcp_handler):
        self._fastapi = fastapi_app
        self._mcp = mcp_handler

    async def __call__(self, scope, receive, send):
        stype = scope.get("type")
        path = scope.get("path", "")

        if stype in ("http", "websocket") and path.startswith("/mcp"):
            await self._mcp(scope, receive, send)
            return

        # --- Auth gate for everything else (http + websocket) ---
        expected = _api_auth_token()
        if (
            expected
            and stype in ("http", "websocket")
            and not (stype == "http" and scope.get("method", "") == "OPTIONS")
            and not _api_path_is_public(path)
        ):
            source = _api_auth_source(scope, expected)
            if source is None:
                if stype == "websocket":
                    await send({"type": "websocket.close", "code": 1008})
                else:
                    from starlette.responses import JSONResponse as _JR
                    resp = _JR(
                        {
                            "error": "unauthorized",
                            "detail": "Missing or invalid API token. Supply "
                            "Authorization: Bearer <token>, ?key=<token>, or "
                            "the api_auth cookie.",
                        },
                        status_code=401,
                        headers={"WWW-Authenticate": 'Bearer realm="DeepAgent"'},
                    )
                    await resp(scope, receive, send)
                print(f"[API-AUTH] 401 {scope.get('method', stype)} {path}", flush=True)
                return

            # If authed via ?key=, persist an HttpOnly cookie so subsequent
            # same-origin requests (the HTML page's fetch calls) authenticate
            # automatically. For GET we 302-redirect to the same path with the
            # key stripped so the secret does not linger in the URL bar /
            # history / proxy logs. Non-GET (rare) sets the cookie inline.
            if source == "query" and stype == "http":
                hdrs = dict(scope.get("headers", []))
                proto = hdrs.get(b"x-forwarded-proto", b"https").decode("utf-8", "replace")
                is_https = proto == "https"

                if scope.get("method", "GET") == "GET":
                    from urllib.parse import parse_qsl, urlencode
                    from starlette.responses import RedirectResponse as _RR
                    clean = urlencode([
                        (k, v) for k, v in parse_qsl(
                            scope.get("query_string", b"").decode("utf-8", "replace"))
                        if k != "key"
                    ])
                    location = path + (f"?{clean}" if clean else "")
                    resp = _RR(url=location, status_code=302)
                    resp.set_cookie(
                        "api_auth", expected, max_age=2592000, path="/",
                        httponly=True, samesite="strict", secure=is_https,
                    )
                    await resp(scope, receive, send)
                    return

                secure = "; Secure" if is_https else ""
                cookie = (
                    f"api_auth={expected}; Path=/; HttpOnly; SameSite=Strict; "
                    f"Max-Age=2592000{secure}"
                ).encode("latin-1")

                async def _send_with_cookie(message):
                    if message.get("type") == "http.response.start":
                        msg_hdrs = list(message.get("headers", []))
                        msg_hdrs.append((b"set-cookie", cookie))
                        message = {**message, "headers": msg_hdrs}
                    await send(message)

                await self._fastapi(scope, receive, _send_with_cookie)
                return

        await self._fastapi(scope, receive, send)


# Mount the multi-phase ABM pipeline orchestrator: POST /api/run-pipeline (+ its
# DELETE/{chat_id} and GET /_active helpers). pipeline_runner declares these on
# its own APIRouter; without this include the routes are never registered and
# VIBE's pipeline dispatch (lib/dispatch-pipeline.ts -> {ALB}/api/run-pipeline)
# gets HTTP 404 {"detail":"Not Found"} — the ABM "project route" failure. Mount
# on the FastAPI instance BEFORE it's wrapped below. (Re-added: a merge dropped
# this include, which is what broke the project route in production.)
try:
    import pipeline_runner as _pipeline_runner
    _fastapi_app.include_router(_pipeline_runner.router)
    print("[STARTUP] pipeline_runner router included (POST /api/run-pipeline)")
except Exception as _e:  # noqa: BLE001
    print(f"[STARTUP] WARNING: pipeline_runner router NOT included: {_e}")

# Reassign module-level `app` so uvicorn.run("server:app", …) picks up the wrapper.
app = _AppWithMCP(_fastapi_app, _mcp_gateway)

_mcp_lifespan_exit: asyncio.Event = asyncio.Event()
_mcp_lifespan_task: asyncio.Task | None = None


if __name__ == "__main__":
    import uvicorn

    print(f"""
DeepAgent Server v2.0 (CONTEXT WINDOW MANAGEMENT)
--------------------------------------------------
Starting on: http://{config.HOST}:{config.PORT}
""")

    uvicorn.run(
        "server:app",
        host=config.HOST,
        port=config.PORT,
        reload=False,
        access_log=False,
        log_level="warning",
        # Trust X-Forwarded-Proto / X-Forwarded-For from Replit's TLS proxy so
        # that request.base_url uses "https://" in production (and in the dev
        # preview).  Without this, OAuth discovery endpoints return "http://"
        # URLs which causes Claude.ai to reject the connector on the production
        # .replit.app domain (the "SSL cert mismatch" symptom).
        proxy_headers=True,
        forwarded_allow_ips="*",
    )
