import io
import os
import httpx
from typing import Optional
from fastmcp import FastMCP

API_KEY = os.environ.get("SHOWPAD_API_KEY", "")
V4_BASE  = "https://zycus.api.showpad.com/v4"
V3_BASE  = "https://zycus.showpad.biz/api/v3"
HEADERS  = {"Authorization": f"Bearer {API_KEY}"}

# Max characters of extracted text returned in a single get_asset_content
# call. Protects the agent's token/context budget against very large decks.
# Callers page past this with the offset parameter.
CONTENT_MAX_CHARS = int(os.environ.get("SHOWPAD_CONTENT_MAX_CHARS", "20000"))
# Hard ceiling on bytes downloaded for extraction so a giant binary cannot
# blow up memory in the stdio subprocess.
CONTENT_MAX_BYTES = int(os.environ.get("SHOWPAD_CONTENT_MAX_BYTES", str(60 * 1024 * 1024)))

mcp = FastMCP("Showpad MCP Server")

def v4_get(path: str, params: dict = None):
    url = f"{V4_BASE}/{path}"
    r = httpx.get(url, headers=HEADERS, params=params, timeout=30)
    r.raise_for_status()
    return r.json()

def v4_post(path: str, json_body: dict = None, params: dict = None):
    url = f"{V4_BASE}/{path}"
    r = httpx.post(url, headers={**HEADERS, "Content-Type": "application/json"}, json=json_body, params=params, timeout=30)
    r.raise_for_status()
    return r.json()

def v3_get(path: str, params: dict = None):
    url = f"{V3_BASE}/{path}"
    r = httpx.get(url, headers=HEADERS, params=params, timeout=30)
    r.raise_for_status()
    return r.json()

@mcp.tool
def list_users(limit: int = 20, offset: int = 0) -> dict:
    """List all Showpad users with pagination."""
    return v4_get("users", {"limit": limit, "offset": offset})

@mcp.tool
def get_user(user_id: str) -> dict:
    """Get details of a specific Showpad user by ID."""
    return v4_get(f"users/{user_id}")

@mcp.tool
def list_assets(limit: int = 20, offset: int = 0, status: Optional[str] = None) -> dict:
    """
    List all Showpad assets with pagination. Does NOT filter by keyword.
    Use search_assets for keyword search or query_assets for ShowQL filtering.
    status: active | inactive | archived
    """
    params = {"limit": limit, "offset": offset}
    if status:
        params["status"] = status
    return v4_get("assets", params)

@mcp.tool
def get_asset(asset_id: str) -> dict:
    """Get full details of a specific Showpad asset by ID."""
    return v4_get(f"assets/{asset_id}")

@mcp.tool
def search_assets(keyword: str, asset_type: Optional[str] = None, shareable_only: bool = False, limit: int = 20, offset: int = 0) -> dict:
    """
    Search Showpad assets by keyword/name using ShowQL name~ operator.
    This is the PRIMARY tool for finding specific content (case studies, presentations, demos, etc.).

    keyword: search term (e.g. 'BFSI', 'case study', 'product demo', 'procurement')
    asset_type: optional filter - 'document', 'video', 'image', etc.
    shareable_only: if True, only return shareable assets
    limit: max results (default 20)
    offset: pagination offset
    """
    parts = [f'name ~ "{keyword}"', 'archivedAt IS EMPTY']
    if asset_type:
        parts.append(f'type = "{asset_type}"')
    if shareable_only:
        parts.append('isShareable = true')
    showql = " AND ".join(parts)
    return v4_post("assets/query", json_body={"query": showql}, params={"limit": limit, "offset": offset})

@mcp.tool
def query_assets(showql_query: str, limit: int = 20, offset: int = 0) -> dict:
    """
    Query Showpad assets using raw ShowQL (Showpad Query Language) via POST /v4/assets/query.

    ShowQL operators:
    - = exact match: name = "file.pdf"
    - ~ contains/like: name ~ "BFSI"
    - IS EMPTY / IS NOT EMPTY: archivedAt IS EMPTY
    - AND / OR / NOT: combine conditions
    - Parentheses for grouping: (tags ~ "id1") OR (tags ~ "id2")

    Filterable fields: name, type, division, tags, archivedAt, isShareable, isDivisionShared, languages, countries, expiresAt

    Examples:
    - 'name ~ "BFSI" AND archivedAt IS EMPTY'
    - 'type = "document" AND archivedAt IS EMPTY'
    - 'tags ~ "TAG_ID" AND archivedAt IS EMPTY'
    - 'isShareable = true AND archivedAt IS EMPTY'
    - '(name ~ "case study") AND type = "document" AND isShareable = true'
    """
    return v4_post("assets/query", json_body={"query": showql_query}, params={"limit": limit, "offset": offset})

@mcp.tool
def search_tags(name: Optional[str] = None, limit: int = 50, offset: int = 0) -> dict:
    """
    List or search Showpad tags. Use to find tag IDs for ShowQL filtering with query_assets.
    name: optional keyword to filter tags by name (v3 API text search)
    """
    if name:
        return v3_get("tags.json", {"name": name, "limit": limit, "offset": offset})
    return v4_get("tags", {"limit": limit, "offset": offset})

@mcp.tool
def get_tag(tag_id: str) -> dict:
    """Get details of a specific tag by ID."""
    return v4_get(f"tags/{tag_id}")

@mcp.tool
def list_divisions() -> dict:
    """List all Showpad divisions."""
    return v4_get("divisions")

@mcp.tool
def get_division(division_id: str) -> dict:
    """Get details of a specific Showpad division by ID."""
    return v4_get(f"divisions/{division_id}")

@mcp.tool
def list_shared_spaces(limit: int = 20, offset: int = 0) -> dict:
    """List all Showpad shared spaces."""
    return v4_get("shared-spaces", {"limit": limit, "offset": offset})

@mcp.tool
def get_shared_space(space_id: str) -> dict:
    """Get details of a specific shared space by ID."""
    return v4_get(f"shared-spaces/{space_id}")

@mcp.tool
def list_channels(limit: int = 20, offset: int = 0) -> dict:
    """List all Showpad channels (uses v3 API)."""
    return v3_get("channels.json", {"limit": limit, "offset": offset})

# --- Content download & text extraction -------------------------------------

# Extensions we can extract text from, grouped by extractor.
_PLAINTEXT_EXTS = {"txt", "md", "markdown", "csv", "tsv", "log", "json", "xml", "rtf"}
_HTML_EXTS = {"html", "htm"}
_BINARY_EXTS = {"pdf", "pptx", "docx", "xlsx"}
_SUPPORTED_EXTS = _PLAINTEXT_EXTS | _HTML_EXTS | _BINARY_EXTS


def _v3_get_json(path: str, params: dict = None):
    url = f"{V3_BASE}/{path}"
    r = httpx.get(url, headers=HEADERS, params=params, timeout=30)
    r.raise_for_status()
    return r.json()


def _download_bytes(url: str) -> bytes:
    """Stream a download link to bytes with the Bearer header, enforcing the
    byte ceiling so a runaway file can't exhaust memory."""
    chunks = []
    total = 0
    with httpx.stream("GET", url, headers=HEADERS, timeout=120,
                      follow_redirects=True) as r:
        r.raise_for_status()
        for chunk in r.iter_bytes():
            chunks.append(chunk)
            total += len(chunk)
            if total > CONTENT_MAX_BYTES:
                raise ValueError(
                    f"File exceeds the {CONTENT_MAX_BYTES} byte download cap "
                    f"for text extraction.")
    return b"".join(chunks)


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {idx} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells]
                    parts.append("\t".join(cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append(f"[Notes] {notes}")
    return "\n".join(parts)


def _extract_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text for c in row.cells]
            parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"--- Sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(cells):
                parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts)


def _extract_html(data: bytes) -> str:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(data.decode("utf-8", errors="replace"), "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _extract_plaintext(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _extract_text(data: bytes, ext: str) -> str:
    if ext in _BINARY_EXTS:
        return {
            "pdf": _extract_pdf,
            "pptx": _extract_pptx,
            "docx": _extract_docx,
            "xlsx": _extract_xlsx,
        }[ext](data)
    if ext in _HTML_EXTS:
        return _extract_html(data)
    return _extract_plaintext(data)


@mcp.tool
def get_asset_content(asset_id: str, offset: int = 0,
                      limit: int = CONTENT_MAX_CHARS) -> dict:
    """Download a Showpad asset and return its EXTRACTED TEXT content.

    Unlike get_asset (which returns only metadata: name, type, tags, IDs),
    this tool downloads the underlying file and extracts its readable text so
    you can answer questions grounded in the document's ACTUAL contents.

    Supported file types: PDF, PPTX (PowerPoint), DOCX (Word), XLSX (Excel),
    and plain-text formats (TXT, MD, CSV, TSV, JSON, XML, HTML/HTM).
    Unsupported types (e.g. legacy .ppt/.doc/.xls, images, video) return an
    explicit error — never a silent empty result. Scanned/image-only PDFs
    yield little or no text (no OCR is performed).

    Large documents are paginated by character: at most `limit` characters are
    returned per call starting at `offset`. When `truncated` is true, call
    again with `offset = next_offset` to fetch the next chunk.

    Args:
        asset_id: The Showpad asset ID (from search_assets / query_assets / get_asset).
        offset: Character offset to start from (default 0). Use next_offset to page.
        limit: Max characters to return this call (default {default}, capped at {cap}).

    Returns:
        A dict with: asset_id, name, extension, file_type, total_chars,
        returned_chars, offset, next_offset, truncated, and content (the text).
        On failure: a dict with an `error` key explaining what went wrong.
    """.format(default=CONTENT_MAX_CHARS, cap=CONTENT_MAX_CHARS)
    if not API_KEY:
        return {"error": "SHOWPAD_API_KEY is not configured."}

    # 1. Resolve the downloadable file link + extension from v3 metadata.
    try:
        meta = _v3_get_json(f"assets/{asset_id}.json")
    except httpx.HTTPStatusError as exc:
        code = exc.response.status_code
        if code == 404:
            return {"error": f"Asset '{asset_id}' not found.", "asset_id": asset_id}
        if code in (401, 403):
            return {"error": f"Not authorized to read asset '{asset_id}' "
                             f"(HTTP {code}).", "asset_id": asset_id}
        return {"error": f"Failed to fetch asset metadata (HTTP {code}).",
                "asset_id": asset_id}
    except Exception as exc:  # noqa: BLE001
        return {"error": f"Failed to fetch asset metadata: {exc}",
                "asset_id": asset_id}

    resp = meta.get("response", meta) if isinstance(meta, dict) else {}
    name = resp.get("name") or resp.get("originalName") or ""
    file_type = resp.get("filetype")
    download_link = resp.get("downloadLink") or resp.get("shortLivedDownloadLink")

    # Determine the file extension from the API field, falling back to the name.
    ext = (resp.get("extension") or "").lower().lstrip(".")
    if not ext:
        base = (resp.get("originalName") or name or "")
        if "." in base:
            ext = base.rsplit(".", 1)[-1].lower()

    if not download_link:
        return {"error": f"Asset '{asset_id}' has no downloadable file link "
                         f"(it may be a link/web asset or not downloadable).",
                "asset_id": asset_id, "name": name, "file_type": file_type}

    if ext not in _SUPPORTED_EXTS:
        return {"error": f"Unsupported file type '{ext or 'unknown'}' for text "
                         f"extraction. Supported: PDF, PPTX, DOCX, XLSX, and "
                         f"plain-text (TXT/MD/CSV/TSV/JSON/XML/HTML).",
                "asset_id": asset_id, "name": name, "extension": ext,
                "file_type": file_type}

    # 2. Download the raw bytes.
    try:
        data = _download_bytes(download_link)
    except httpx.HTTPStatusError as exc:
        return {"error": f"Download failed (HTTP {exc.response.status_code}). "
                         f"The link may have expired.",
                "asset_id": asset_id, "name": name}
    except Exception as exc:  # noqa: BLE001
        return {"error": f"Download failed: {exc}", "asset_id": asset_id,
                "name": name}

    # 3. Extract text via the type-appropriate extractor.
    try:
        text = _extract_text(data, ext)
    except Exception as exc:  # noqa: BLE001
        return {"error": f"Failed to extract text from '{ext}' file: {exc}",
                "asset_id": asset_id, "name": name, "extension": ext}

    if not text or not text.strip():
        return {"error": "No extractable text found in this asset. It may be "
                         "image-only/scanned (no OCR is performed) or empty.",
                "asset_id": asset_id, "name": name, "extension": ext,
                "file_type": file_type}

    # 4. Paginate by character to protect the agent's context budget.
    total_chars = len(text)
    offset = max(0, int(offset))
    limit = max(1, min(int(limit), CONTENT_MAX_CHARS))
    chunk = text[offset:offset + limit]
    end = offset + len(chunk)
    truncated = end < total_chars

    return {
        "asset_id": asset_id,
        "name": name,
        "extension": ext,
        "file_type": file_type,
        "total_chars": total_chars,
        "returned_chars": len(chunk),
        "offset": offset,
        "next_offset": end if truncated else None,
        "truncated": truncated,
        "content": chunk,
    }


if __name__ == "__main__":
    mcp.run()
